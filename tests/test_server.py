"""
Integration tests for MCP server tool functions.

These tests call the tool functions directly (they are regular Python functions
wrapped by the MCP decorator). All file-based operations use tmp_path.

v0.3.0 response contract (issues #88, #98, #99):
    - Success: json.dumps({"ok": true, "result": <dict>})
      — ``result`` is ALWAYS a dict. Human-readable status messages are under
        the ``message`` key (``result["message"]``).
    - Error:   json.dumps({"ok": false, "error": {"code", "message", ...}})

v0.3.0 structured-params contract (issue #97):
    - ``*_json: str`` parameters were removed. Tests pass native Python
      structures (list / dict) to MCP tools.

Tests unwrap via ``_unwrap_ok`` / ``_unwrap_err`` helpers below.
"""

from __future__ import annotations

import json
import os

import pytest
from pptx import Presentation

from pptx_mcp_server.engine.pptx_io import save_pptx, open_pptx
from pptx_mcp_server.server import (
    mcp,
    pptx_create,
    pptx_get_info,
    pptx_read_slide,
    pptx_list_shapes,
    pptx_add_slide,
    pptx_delete_slide,
    pptx_duplicate_slide,
    pptx_set_slide_background,
    pptx_set_dimensions,
    pptx_add_textbox,
    pptx_edit_text,
    pptx_add_paragraph,
    pptx_add_shape,
    pptx_delete_shape,
    pptx_format_shape,
    pptx_add_table,
    pptx_edit_table_cell,
    pptx_edit_table_cells,
    pptx_format_table,
    pptx_add_content_slide,
    pptx_add_section_divider,
    pptx_add_kpi_row_legacy,
    pptx_add_bullet_block,
    pptx_add_image,
    pptx_render_slide,
    pptx_check_layout,
    pptx_add_flex_container,
    pptx_add_responsive_card_row,
    pptx_build_slide,
    pptx_build_deck,
    pptx_add_chart,
    _format_check_layout_summary,
    INSTRUCTIONS,
)


# ── Response-shape helpers (v0.3.0, issues #88, #98) ───────────────────

def _unwrap_ok(payload: str):
    """Parse a tool response and assert it is a success; return the result dict.

    v0.3.0 (#98): ``result`` は **常に dict**。caller は ``result["message"]``
    等のキーに明示的にアクセスする。
    """
    parsed = json.loads(payload)
    assert isinstance(parsed, dict), f"response must be a JSON object, got {type(parsed)}"
    assert parsed.get("ok") is True, f"expected ok=true, got {parsed!r}"
    assert "result" in parsed, f"success payload must include 'result': {parsed!r}"
    result = parsed["result"]
    assert isinstance(result, dict), (
        f"v0.3.0 contract: result must be a dict, got {type(result).__name__}: {result!r}"
    )
    return result


def _unwrap_err(payload: str) -> dict:
    """Parse a tool response and assert it is an error; return the error dict."""
    parsed = json.loads(payload)
    assert isinstance(parsed, dict)
    assert parsed.get("ok") is False, f"expected ok=false, got {parsed!r}"
    err = parsed.get("error")
    assert isinstance(err, dict), f"error must be an object, got {err!r}"
    assert "code" in err and "message" in err, f"error missing required keys: {err!r}"
    return err


@pytest.fixture
def deck(tmp_path):
    """Create a blank deck with one slide, return file path."""
    path = str(tmp_path / "deck.pptx")
    pptx_create(path)
    pptx_add_slide(path, layout_index=6)
    return path


class TestToolRegistration:
    """All registered MCP tools must expose the expected core signatures."""

    def test_all_tools_registered(self):
        # FastMCP stores tools internally; list them via the _tool_manager.
        # v0.3.1 (#108): dropped hardcoded "25 tools" drift.
        # v0.6.0 (#137): tool tiering — advanced tools are hidden from
        # the default registry (env-gate ``PPTX_MCP_INCLUDE_ADVANCED``).
        # The expected list below therefore only contains **default-tier**
        # tools so the assertion holds without setting the env var.
        tool_names = list(mcp._tool_manager._tools.keys())
        expected = [
            "pptx_create",
            "pptx_get_info",
            "pptx_read_slide",
            "pptx_add_slide",
            "pptx_add_image",
            "pptx_render_slide",
        ]
        for name in expected:
            assert name in tool_names, f"Tool '{name}' not registered"
        # dynamic count — matches the live FastMCP registration.
        assert len(tool_names) >= len(expected)


class TestCreatePptx:
    """pptx_create tool creates a valid file."""

    def test_creates_file(self, tmp_path):
        path = str(tmp_path / "new.pptx")
        result = pptx_create(path)
        assert os.path.exists(path)
        payload = _unwrap_ok(result)
        assert "Created" in payload["message"]


class TestFileBased:
    """File-based tool calls modify the underlying PPTX correctly."""

    def test_add_slide(self, deck):
        result = pptx_add_slide(deck)
        assert "Added slide" in _unwrap_ok(result)["message"]
        prs = open_pptx(deck)
        assert len(prs.slides) == 2

    def test_add_textbox(self, deck):
        result = pptx_add_textbox(deck, 0, 1, 1, 4, 0.5, text="Hello")
        assert "Added textbox" in _unwrap_ok(result)["message"]
        prs = open_pptx(deck)
        slide = prs.slides[0]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Hello" in texts

    def test_add_table(self, deck):
        rows = [["Name", "Score"], ["Alice", "95"]]
        result = pptx_add_table(deck, 0, rows, 1, 1, 5)
        assert "Added table" in _unwrap_ok(result)["message"]
        prs = open_pptx(deck)
        slide = prs.slides[0]
        table_shapes = [s for s in slide.shapes if s.has_table]
        assert len(table_shapes) == 1


class TestCompositeTools:
    """Composite tool calls produce correct output."""

    def test_add_content_slide(self, deck):
        result = pptx_add_content_slide(deck, "Revenue Analysis")
        payload = _unwrap_ok(result)
        assert "Added content slide" in payload["message"]
        prs = open_pptx(deck)
        assert len(prs.slides) == 2  # original + content

    def test_add_section_divider(self, deck):
        result = pptx_add_section_divider(deck, "Q1 Results", subtitle="FY2024")
        payload = _unwrap_ok(result)
        assert "Added section divider" in payload["message"]
        prs = open_pptx(deck)
        assert len(prs.slides) == 2


class TestErrorCases:
    """Error cases must return structured JSON error payloads (issue #88)."""

    def test_open_nonexistent_returns_structured_error(self, tmp_path):
        result = pptx_get_info(str(tmp_path / "nope.pptx"))
        err = _unwrap_err(result)
        assert err["code"] == "FILE_NOT_FOUND"
        assert "message" in err

    def test_invalid_slide_returns_structured_error(self, deck):
        result = pptx_read_slide(deck, 99)
        err = _unwrap_err(result)
        assert err["code"] == "SLIDE_NOT_FOUND"

    def test_invalid_shape_type_returns_structured_error(self, deck):
        result = pptx_add_shape(deck, 0, "nonexistent", 1, 1, 2, 2)
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"


# ── Issue #97: structured params (native list/dict replace *_json: str) ──

class TestStructuredParams:
    """Structured-param tools accept native Python values and reject bad types."""

    def test_add_table_with_native_rows(self, deck):
        rows = [["X", "Y"], ["1", "2"]]
        result = pptx_add_table(deck, 0, rows, 1, 1, 5)
        assert "Added table" in _unwrap_ok(result)["message"]

    def test_add_table_rows_none_rejected(self, deck):
        result = pptx_add_table(deck, 0, None, 1, 1, 5)  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "rows"

    def test_add_table_rows_wrong_type_rejected(self, deck):
        result = pptx_add_table(deck, 0, "not a list", 1, 1, 5)  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "rows"

    def test_add_table_col_widths_wrong_type_rejected(self, deck):
        result = pptx_add_table(deck, 0, [["a", "b"]], 1, 1, 5, col_widths="bad")  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "col_widths"

    def test_add_kpi_row_with_native_list(self, deck):
        kpis = [{"value": "99", "label": "Score"}]
        result = pptx_add_kpi_row_legacy(deck, 0, kpis, 2.0)
        payload = _unwrap_ok(result)
        assert "Added 1 KPI" in payload["message"]

    def test_add_kpi_row_wrong_type_rejected(self, deck):
        result = pptx_add_kpi_row_legacy(deck, 0, "not a list", 2.0)  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "kpis"

    def test_add_bullet_block_with_native_list(self, deck):
        bullets = ["Point A", "Point B"]
        result = pptx_add_bullet_block(deck, 0, bullets, 1, 2, 5, 3)
        payload = _unwrap_ok(result)
        assert "Added bullet block" in payload["message"]

    def test_add_bullet_block_wrong_type_rejected(self, deck):
        result = pptx_add_bullet_block(deck, 0, "not a list", 1, 2, 5, 3)  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "bullets"

    def test_edit_table_cells_with_native_list(self, deck):
        # First add a table so shape_index=0 is a table.
        pptx_add_table(deck, 0, [["H1", "H2"], ["v1", "v2"]], 1, 1, 5)
        prs = open_pptx(deck)
        slide = prs.slides[0]
        shape_idx = next(i for i, s in enumerate(slide.shapes) if s.has_table)
        edits = [{"row": 0, "col": 1, "text": "new"}]
        result = pptx_edit_table_cells(deck, 0, shape_idx, edits)
        _unwrap_ok(result)

    def test_edit_table_cells_wrong_type_rejected(self, deck):
        result = pptx_edit_table_cells(deck, 0, 0, "bad")  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "edits"

    def test_build_slide_with_native_dict(self, deck):
        spec = {"layout": "content", "title": "Hello", "elements": []}
        result = pptx_build_slide(deck, spec)
        payload = _unwrap_ok(result)
        assert "Built slide" in payload["message"]

    def test_build_slide_wrong_type_rejected(self, deck):
        result = pptx_build_slide(deck, "not a dict")  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "spec"

    def test_build_deck_with_native_list(self, deck):
        slides = [
            {"layout": "content", "title": "S1", "elements": []},
            {"layout": "content", "title": "S2", "elements": []},
        ]
        result = pptx_build_deck(deck, slides)
        payload = _unwrap_ok(result)
        assert "Built 2 slides" in payload["message"]

    def test_build_deck_wrong_type_rejected(self, deck):
        result = pptx_build_deck(deck, "not a list")  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "slides"

    def test_add_chart_with_native_dict(self, deck):
        chart = {
            "chart_type": "column",
            "categories": ["A", "B"],
            "series": [{"name": "S", "values": [1, 2]}],
            "left": 1.0, "top": 1.0, "width": 5.0, "height": 3.0,
        }
        result = pptx_add_chart(deck, 0, chart)
        _unwrap_ok(result)

    def test_add_chart_wrong_type_rejected(self, deck):
        result = pptx_add_chart(deck, 0, "not a dict")  # type: ignore[arg-type]
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err["parameter"] == "chart"


# ── Issue #33 / #98: pptx_check_layout back-compat (wrapped in message) ───

class TestCheckLayoutBackCompat:
    """``pptx_check_layout`` wraps the legacy string in result.message.

    Issue #33 の legacy 文字列は ``result["message"]`` に入る (#98)。
    ``detailed=True`` は ``result["slides"]`` / ``result["summary"]`` を
    そのまま返し, 以前の double-encoded JSON 文字列 (#99) は撤廃された。
    """

    def test_clean_deck_returns_legacy_string_in_message(self, deck):
        result = pptx_check_layout(deck)
        payload = _unwrap_ok(result)
        msg = payload["message"]
        assert isinstance(msg, str)
        assert msg.startswith("All slides clean") or msg.startswith("Found")

    def test_clean_deck_wording_exact(self, deck):
        payload = _unwrap_ok(pptx_check_layout(deck))
        msg = payload["message"]
        if msg.startswith("All slides clean"):
            assert msg == (
                "All slides clean — no overlaps, out-of-bounds, text overflow, "
                "or readability issues detected."
            )

    def test_deck_with_overlap_returns_found_string(self, deck):
        pptx_add_shape(deck, 0, "rectangle", 1.0, 1.0, 3.0, 2.0, fill_color="2251FF")
        pptx_add_shape(deck, 0, "rectangle", 2.0, 1.5, 3.0, 2.0, fill_color="FF0000")
        payload = _unwrap_ok(pptx_check_layout(deck))
        msg = payload["message"]
        assert isinstance(msg, str)
        assert msg.startswith("Found")
        assert "overlap" in msg.lower()

    def test_detailed_returns_flat_dict_no_inner_string(self, deck):
        """#99: detailed=True must return the dict inline — single json.loads."""
        payload = _unwrap_ok(pptx_check_layout(deck, detailed=True))
        # payload itself is a dict (not a JSON-encoded string needing re-decode).
        assert "slides" in payload
        assert "summary" in payload
        assert isinstance(payload["slides"], list)
        assert set(payload["summary"].keys()) >= {"errors", "warnings", "infos"}

    def test_detailed_single_json_loads_suffices(self, deck):
        """#99 regression: one json.loads on the envelope fully decodes the result."""
        raw = pptx_check_layout(deck, detailed=True)
        parsed = json.loads(raw)
        # No second json.loads needed on parsed["result"].
        assert isinstance(parsed["result"], dict)
        assert "slides" in parsed["result"]

    def test_format_helper_clean(self):
        empty = {"slides": [], "summary": {"errors": 0, "warnings": 0, "infos": 0}}
        assert _format_check_layout_summary(empty).startswith("All slides clean")

    def test_format_helper_problems(self):
        result = {
            "slides": [
                {
                    "index": 1,
                    "overlaps": ["shape A overlaps shape B"],
                    "out_of_bounds": [],
                    "text_overflow": [
                        {
                            "severity": "error",
                            "slide_index": 1,
                            "shape_name": "tb_1",
                            "category": "text_overflow",
                            "message": "text overflows",
                        }
                    ],
                    "unreadable_text": [],
                    "divider_collision": [],
                    "inconsistent_gaps": [],
                }
            ],
            "summary": {"errors": 2, "warnings": 0, "infos": 0},
        }
        s = _format_check_layout_summary(result)
        assert s.startswith("Found 2 layout issues:")
        assert "Slide 1 [error] overlap:" in s
        assert "Slide 1 [error] text_overflow:" in s


# ── Issue #34: card_row save order ─────────────────────────────────

class TestResponsiveCardRowSaveOrder:
    """``pptx_add_responsive_card_row`` が save を最後に行うこと."""

    def test_save_not_called_when_post_processing_raises(self, deck, monkeypatch):
        """placement serialize 段階で例外が起きたら save は呼ばれないこと."""
        from pptx_mcp_server import server as server_mod

        save_called = {"count": 0}
        original_save = server_mod.save_pptx

        def tracking_save(prs, path):
            save_called["count"] += 1
            return original_save(prs, path)

        monkeypatch.setattr(server_mod, "save_pptx", tracking_save)

        class Boom:
            @property
            def left(self):
                raise RuntimeError("boom")
            top = 0.0
            width = 0.0
            height = 0.0

        def fake_add_row(slide, cards, **kwargs):
            return [Boom()], 0.0

        monkeypatch.setattr(server_mod, "add_responsive_card_row", fake_add_row)

        cards = [{"title": "A", "body": "a"}]
        result = server_mod.pptx_add_responsive_card_row(
            file_path=deck,
            slide_index=0,
            cards=cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
        )
        err = _unwrap_err(result)
        # error code should be INTERNAL_ERROR (boom is a RuntimeError).
        assert err["code"] == "INTERNAL_ERROR"
        assert "boom" in err["message"]
        assert save_called["count"] == 0, (
            f"save_pptx should not have been called; was called {save_called['count']} times"
        )

    def test_happy_path_saves_once_and_returns_placements(self, deck):
        from pptx_mcp_server.server import pptx_add_responsive_card_row

        cards = [{"title": "A", "body": "aa"}, {"title": "B", "body": "bb"}]
        result = pptx_add_responsive_card_row(
            file_path=deck,
            slide_index=0,
            cards=cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
        )
        payload = _unwrap_ok(result)
        # render disabled by default → payload carries the cards dict directly.
        assert "cards" in payload
        assert len(payload["cards"]) == 2
        assert {"left", "top", "width", "height"} <= set(payload["cards"][0].keys())
        assert "consumed_height" in payload
        prs = open_pptx(deck)
        assert len(prs.slides) >= 1


# ── Issue #97: flex_container native items (was items_json) ─────────

class TestFlexContainerItems:
    """``pptx_add_flex_container`` が ``items: list[dict]`` を受け取る (#97)."""

    def test_native_list_works(self, deck):
        items = [
            {"sizing": "fixed", "size": 2.0, "type": "rectangle", "fill_color": "2251FF"}
        ]
        result = pptx_add_flex_container(
            file_path=deck,
            slide_index=0,
            items=items,
            left=0.5, top=0.5, width=10.0, height=1.0,
        )
        payload = _unwrap_ok(result)
        assert "allocations" in payload

    def test_non_list_rejected(self, deck):
        result = pptx_add_flex_container(
            file_path=deck,
            slide_index=0,
            items="not a list",  # type: ignore[arg-type]
            left=0.5, top=0.5, width=10.0, height=1.0,
        )
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert err.get("parameter") == "items"

    def test_dict_rejected(self, deck):
        """Passing a dict where a list is expected → INVALID_PARAMETER."""
        result = pptx_add_flex_container(
            file_path=deck,
            slide_index=0,
            items={"not": "array"},  # type: ignore[arg-type]
            left=0.5, top=0.5, width=10.0, height=1.0,
        )
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"


# ── #43: CardSpec JSON strict unknown-key validation ───────────────

class TestResponsiveCardRowStrictKeys:
    """``pptx_add_responsive_card_row`` が未知キーを弾く (#43)."""

    def test_unknown_key_rejected_with_index_and_name(self, deck):
        cards = [{"title": "t", "body": "b", "typo": "x"}]
        result = pptx_add_responsive_card_row(
            file_path=deck,
            slide_index=0,
            cards=cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
        )
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert "typo" in err["message"]
        assert "card[0]" in err["message"]

    def test_unknown_key_index_points_to_bad_card(self, deck):
        cards = [
            {"title": "ok", "body": "ok"},
            {"title": "bad", "body": "b", "oops": 1},
        ]
        result = pptx_add_responsive_card_row(
            file_path=deck,
            slide_index=0,
            cards=cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
        )
        err = _unwrap_err(result)
        assert err["code"] == "INVALID_PARAMETER"
        assert "card[1]" in err["message"]
        assert "oops" in err["message"]

    def test_valid_cards_ok_regression(self, deck):
        cards = [
            {"title": "A", "body": "aa", "accent_color": "2251FF"},
            {"title": "B", "body": "bb", "padding": 0.25},
        ]
        result = pptx_add_responsive_card_row(
            file_path=deck,
            slide_index=0,
            cards=cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
        )
        payload = _unwrap_ok(result)
        assert "cards" in payload
        assert len(payload["cards"]) == 2


# ── Issues #88, #98: structured response contract ──────────────────

class TestStructuredResponseContract:
    """All tool returns follow the v0.3.0 ``{ok, result|error}`` contract.

    v0.3.0 (#98): ``result`` is ALWAYS a dict — tools wrap legacy string
    payloads under ``result["message"]``.
    """

    def test_success_shape_has_ok_and_result_dict(self, tmp_path):
        path = str(tmp_path / "s.pptx")
        response = pptx_create(path)
        parsed = json.loads(response)
        assert parsed["ok"] is True
        assert isinstance(parsed["result"], dict)
        assert "message" in parsed["result"]
        assert isinstance(parsed["result"]["message"], str)

    def test_all_primitive_tools_return_dict_result(self, deck):
        """Every success path returns result: dict (never raw string)."""
        responses = [
            pptx_create(deck.replace(".pptx", "-new.pptx")),
            pptx_get_info(deck),
            pptx_add_slide(deck),
            pptx_add_textbox(deck, 0, 1, 1, 3, 0.5, text="hi"),
            pptx_add_shape(deck, 0, "rectangle", 1, 1, 2, 2),
            pptx_add_content_slide(deck, "Title"),
            pptx_add_section_divider(deck, "Section"),
        ]
        for r in responses:
            parsed = json.loads(r)
            assert parsed["ok"] is True, f"expected ok=true, got {parsed!r}"
            assert isinstance(parsed["result"], dict), (
                f"v0.3.0 contract: result must be a dict, got {parsed!r}"
            )

    def test_caller_can_discriminate_on_ok(self, deck, tmp_path):
        ok_response = pptx_get_info(deck)
        err_response = pptx_get_info(str(tmp_path / "missing.pptx"))
        assert json.loads(ok_response)["ok"] is True
        assert json.loads(err_response)["ok"] is False

    def test_error_payload_has_code_and_message(self, tmp_path):
        err = _unwrap_err(pptx_get_info(str(tmp_path / "nope.pptx")))
        assert err["code"] == "FILE_NOT_FOUND"
        assert isinstance(err["message"], str) and err["message"]

    def test_error_engine_code_preserved(self, deck):
        err = _unwrap_err(pptx_read_slide(deck, 99))
        assert err["code"] == "SLIDE_NOT_FOUND"

    def test_error_can_include_parameter_and_hint(self, deck):
        """flex_container non-list rejection sets parameter + hint + issue."""
        err = _unwrap_err(pptx_add_flex_container(
            file_path=deck,
            slide_index=0,
            items="bad",  # type: ignore[arg-type]
            left=0.5, top=0.5, width=10.0, height=1.0,
        ))
        assert err["parameter"] == "items"
        assert "hint" in err
        assert err.get("issue") == 97

    def test_all_error_paths_return_valid_json(self, tmp_path):
        """Every error path must round-trip through json.loads cleanly."""
        responses = [
            pptx_get_info(str(tmp_path / "x.pptx")),
            pptx_read_slide(str(tmp_path / "x.pptx"), 0),
            pptx_list_shapes(str(tmp_path / "x.pptx"), 0),
        ]
        for r in responses:
            parsed = json.loads(r)
            assert parsed["ok"] is False
            assert parsed["error"]["code"]
            assert parsed["error"]["message"]


# ── Issue #86 / #98: auto_render opt-in + timeout, dict-only result ────

class TestAutoRenderOptIn:
    """Auto-render is OFF by default; opt-in via PPTX_MCP_AUTO_RENDER.

    See issue #86. v0.3.0 (#98): render info is merged into the ``result`` dict
    alongside ``message`` — no more ``{"value": ...}`` wrapper.
    All tests install monkey-patched ``render_slide`` so no real
    LibreOffice subprocess ever fires.
    """

    def test_default_does_not_invoke_renderer(self, deck, monkeypatch):
        """With PPTX_MCP_AUTO_RENDER unset, renderer is NOT called."""
        from pptx_mcp_server import server as server_mod

        monkeypatch.delenv("PPTX_MCP_AUTO_RENDER", raising=False)
        called = {"count": 0}

        def fake_render(*a, **kw):
            called["count"] += 1
            return "/tmp/fake.png"

        monkeypatch.setattr(server_mod, "render_slide", fake_render)

        result = server_mod.pptx_add_content_slide(deck, "Title")
        payload = _unwrap_ok(result)
        assert "preview_path" not in payload
        assert "render_warning" not in payload
        assert called["count"] == 0, "renderer should not fire when auto-render disabled"

    def test_env_var_enables_renderer(self, deck, monkeypatch):
        """With PPTX_MCP_AUTO_RENDER=1, renderer IS called and preview_path lives in result."""
        from pptx_mcp_server import server as server_mod

        monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", "1")
        called = {"count": 0}

        def fake_render(file_path, slide_index=-1, dpi=100):
            called["count"] += 1
            return "/tmp/fake.png"

        monkeypatch.setattr(server_mod, "render_slide", fake_render)

        result = server_mod.pptx_add_content_slide(deck, "Title")
        payload = _unwrap_ok(result)
        assert called["count"] == 1, "renderer should fire when auto-render enabled"
        # v0.3.0 shape: {"message": ..., "preview_path": ...}
        assert payload["preview_path"] == "/tmp/fake.png"
        assert "Added content slide" in payload["message"]

    def test_env_var_various_truthy_values(self, monkeypatch):
        from pptx_mcp_server import server as server_mod
        for v in ("1", "true", "TRUE", "yes", "on"):
            monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", v)
            assert server_mod._auto_render_enabled(), f"{v!r} should enable"
        for v in ("0", "false", "no", "off", ""):
            monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", v)
            assert not server_mod._auto_render_enabled(), f"{v!r} should disable"

    def test_render_timeout_returns_warning_quickly(self, deck, monkeypatch):
        """Slow renderer + tight timeout → tool returns within ~timeout+buffer
        with a render_warning (not a failure)."""
        import time

        from pptx_mcp_server import server as server_mod

        monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", "1")
        monkeypatch.setenv("PPTX_MCP_RENDER_TIMEOUT", "1")

        def slow_render(*a, **kw):
            time.sleep(20)
            return "/tmp/never.png"

        monkeypatch.setattr(server_mod, "render_slide", slow_render)

        t0 = time.time()
        result = server_mod.pptx_add_content_slide(deck, "Title")
        elapsed = time.time() - t0

        # Must return well before the 20s the renderer would take.
        assert elapsed < 5.0, f"timeout did not kick in: {elapsed:.1f}s"
        payload = _unwrap_ok(result)
        assert "render_warning" in payload
        assert payload["render_warning"]["reason"] == "timeout"
        # primary action still succeeded.
        assert "Added content slide" in payload["message"]

    def test_render_failure_does_not_fail_tool(self, deck, monkeypatch):
        """Renderer raising an exception → tool still returns ok=true."""
        from pptx_mcp_server import server as server_mod

        monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", "1")

        def broken_render(*a, **kw):
            raise RuntimeError("libreoffice exploded")

        monkeypatch.setattr(server_mod, "render_slide", broken_render)

        result = server_mod.pptx_add_content_slide(deck, "Title")
        payload = _unwrap_ok(result)  # primary action still succeeds
        assert "render_warning" in payload
        assert payload["render_warning"]["reason"] == "failed"
        assert "libreoffice exploded" in payload["render_warning"]["error"]

    def test_primary_action_error_skips_renderer(self, deck, monkeypatch):
        """If primary action fails, renderer must not be invoked."""
        from pptx_mcp_server import server as server_mod

        monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", "1")
        called = {"count": 0}

        def fake_render(*a, **kw):
            called["count"] += 1
            return "/tmp/x.png"

        monkeypatch.setattr(server_mod, "render_slide", fake_render)

        # invalid slide_index → add_kpi_row raises before _auto_render runs.
        result = server_mod.pptx_add_kpi_row_legacy(
            deck, 999, [{"value": "1", "label": "x"}], 2.0
        )
        _unwrap_err(result)
        assert called["count"] == 0, "renderer fired despite primary-action error"

    def test_render_timeout_env_var_parsing(self, monkeypatch):
        """Invalid / missing PPTX_MCP_RENDER_TIMEOUT falls back to default."""
        from pptx_mcp_server import server as server_mod

        monkeypatch.delenv("PPTX_MCP_RENDER_TIMEOUT", raising=False)
        assert server_mod._auto_render_timeout() == 10.0

        monkeypatch.setenv("PPTX_MCP_RENDER_TIMEOUT", "garbage")
        assert server_mod._auto_render_timeout() == 10.0

        monkeypatch.setenv("PPTX_MCP_RENDER_TIMEOUT", "0")
        assert server_mod._auto_render_timeout() == 10.0

        monkeypatch.setenv("PPTX_MCP_RENDER_TIMEOUT", "2.5")
        assert server_mod._auto_render_timeout() == 2.5

    def test_responsive_card_row_merges_render_info_into_result(self, deck, monkeypatch):
        """#98: responsive_card_row's rich-dict result merges preview_path in-place."""
        from pptx_mcp_server import server as server_mod

        monkeypatch.setenv("PPTX_MCP_AUTO_RENDER", "1")
        monkeypatch.setattr(server_mod, "render_slide", lambda *a, **kw: "/tmp/cr.png")

        cards = [{"title": "A", "body": "a"}]
        result = server_mod.pptx_add_responsive_card_row(
            file_path=deck,
            slide_index=0,
            cards=cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
        )
        payload = _unwrap_ok(result)
        assert "cards" in payload
        assert payload["preview_path"] == "/tmp/cr.png"


# ── Issue #90: INSTRUCTIONS trimmed of product policy ──────────────

class TestInstructionsContent:
    """INSTRUCTIONS prose should describe capability, not prescribe UX."""

    def test_no_ask_user_prompt(self):
        lower = INSTRUCTIONS.lower()
        # The banned phrases from the pre-#90 INSTRUCTIONS.
        assert "always ask" not in lower, (
            "INSTRUCTIONS must not tell agents to 'always ask' the user "
            "about product decisions (issue #90)."
        )
        assert "ask the user which color palette" not in lower
        assert "ask user for color palette" not in lower
        assert "ask for color palette" not in lower

    def test_themes_enumeration_present(self):
        # Neutral themes enumeration replaces the ask-first directive.
        assert "mckinsey" in INSTRUCTIONS.lower()
        assert "deloitte" in INSTRUCTIONS.lower()
        assert "neutral" in INSTRUCTIONS.lower()
        assert "Available Themes" in INSTRUCTIONS or "Available themes" in INSTRUCTIONS

    def test_response_shape_documented(self):
        # Caller-facing response shape must be discoverable in the prompt.
        assert '"ok": true' in INSTRUCTIONS
        assert '"ok": false' in INSTRUCTIONS

    def test_auto_render_opt_in_documented(self):
        assert "PPTX_MCP_AUTO_RENDER" in INSTRUCTIONS


# ── v0.3.1 (#105/#106/#107/#108): OpenAI/Codex gap closures ────────────

class TestStrictNestedValidation:
    """#105: unknown keys in nested specs return structured INVALID_PARAMETER.

    Adversarial shapes that used to leak ``AttributeError`` / ``TypeError``
    are now rejected cleanly at the engine boundary.
    """

    def test_chart_spec_rejects_unknown_series_key(self, deck):
        err = _unwrap_err(
            pptx_add_chart(
                deck,
                slide_index=0,
                chart={
                    "chart_type": "column",
                    "categories": ["A"],
                    "series": [{"name": "S", "values": [1], "unknown_key": 1}],
                },
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "unknown_key" in err["message"]

    def test_chart_spec_rejects_unknown_top_level_key(self, deck):
        err = _unwrap_err(
            pptx_add_chart(
                deck,
                slide_index=0,
                chart={
                    "chart_type": "column",
                    "categories": ["A"],
                    "series": [{"name": "S", "values": [1]}],
                    "bogus_field": 42,
                },
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "bogus_field" in err["message"]

    def test_chart_spec_series_string_not_dict(self, deck):
        # The original adversarial case: series=['oops'] used to raise
        # AttributeError: 'str' object has no attribute 'get'.
        err = _unwrap_err(
            pptx_add_chart(
                deck,
                slide_index=0,
                chart={
                    "chart_type": "column",
                    "categories": ["A"],
                    "series": ["oops"],
                },
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "series[0]" in err["message"]

    def test_chart_spec_valid_still_works(self, deck):
        payload = _unwrap_ok(
            pptx_add_chart(
                deck,
                slide_index=0,
                chart={
                    "chart_type": "column",
                    "categories": ["Q1", "Q2"],
                    "series": [{"name": "Rev", "values": [10, 20]}],
                },
            )
        )
        assert "message" in payload

    def test_build_slide_rejects_unknown_top_level_key(self, deck):
        err = _unwrap_err(
            pptx_build_slide(
                deck,
                spec={"unknown_field": 1, "title": "x"},
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "unknown_field" in err["message"]

    def test_build_slide_valid_still_works(self, deck):
        payload = _unwrap_ok(
            pptx_build_slide(deck, spec={"layout": "content", "title": "x"})
        )
        assert "message" in payload

    def test_build_deck_rejects_unknown_key_in_slide(self, deck):
        err = _unwrap_err(
            pptx_build_deck(
                deck,
                slides=[{"layout": "content", "title": "x", "bogus": True}],
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "bogus" in err["message"]

    def test_kpi_row_rejects_unknown_key(self, deck):
        err = _unwrap_err(
            pptx_add_kpi_row_legacy(
                deck,
                slide_index=0,
                kpis=[{"label": "x", "value": "y", "typo": 1}],
                y=1.2,
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "typo" in err["message"]

    def test_kpi_row_valid_still_works(self, deck):
        payload = _unwrap_ok(
            pptx_add_kpi_row_legacy(
                deck,
                slide_index=0,
                kpis=[{"label": "Rev", "value": "100"}],
                y=1.2,
            )
        )
        assert "message" in payload

    def test_edit_cells_rejects_unknown_key(self, deck):
        # Build a table first.
        pptx_add_table(
            deck, slide_index=0,
            rows=[["H1", "H2"], ["a", "b"]],
            left=1.0, top=1.0, width=6.0,
        )
        err = _unwrap_err(
            pptx_edit_table_cells(
                deck, slide_index=0, shape_index=0,
                edits=[{"row": 0, "col": 0, "text": "X", "typo": 1}],
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "typo" in err["message"]

    def test_edit_cells_rejects_non_int_row(self, deck):
        pptx_add_table(
            deck, slide_index=0,
            rows=[["H1", "H2"], ["a", "b"]],
            left=1.0, top=1.0, width=6.0,
        )
        err = _unwrap_err(
            pptx_edit_table_cells(
                deck, slide_index=0, shape_index=0,
                edits=[{"row": "zero", "col": 0, "text": "X"}],
            )
        )
        assert err["code"] == "INVALID_PARAMETER"
        assert "row" in err["message"]


class TestEngineWrappersStringJsonRemoved:
    """#105 Part B: engine wrappers no longer accept legacy string-JSON input."""

    def test_build_slide_string_input_rejected(self, tmp_path):
        from pptx_mcp_server.engine.composites import build_slide
        from pptx_mcp_server.engine import create_presentation
        from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode

        path = str(tmp_path / "deck.pptx")
        create_presentation(path)
        with pytest.raises(EngineError) as exc:
            build_slide(path, '{"title":"x"}')
        assert exc.value.code == ErrorCode.INVALID_PARAMETER

    def test_build_deck_string_input_rejected(self, tmp_path):
        from pptx_mcp_server.engine.composites import build_deck
        from pptx_mcp_server.engine import create_presentation
        from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode

        path = str(tmp_path / "deck.pptx")
        create_presentation(path)
        with pytest.raises(EngineError) as exc:
            build_deck(path, '[{"title":"x"}]')
        assert exc.value.code == ErrorCode.INVALID_PARAMETER

    def test_add_bullet_block_string_input_rejected(self, tmp_path):
        from pptx_mcp_server.engine.composites import add_bullet_block
        from pptx_mcp_server.engine import create_presentation, add_slide
        from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode

        path = str(tmp_path / "deck.pptx")
        create_presentation(path)
        add_slide(path, layout_index=6)
        with pytest.raises(EngineError) as exc:
            add_bullet_block(
                path, slide_index=0,
                items='["a","b"]',
                left=1.0, top=1.0, width=5.0, height=2.0,
            )
        assert exc.value.code == ErrorCode.INVALID_PARAMETER


class TestCheckLayoutDetailedMessage:
    """#106: ``detailed=True`` result includes a ``message`` field."""

    def test_detailed_clean_has_message(self, deck):
        payload = _unwrap_ok(pptx_check_layout(deck, detailed=True))
        assert "message" in payload
        assert "clean" in payload["message"].lower() or "finding" in payload["message"].lower()

    def test_detailed_with_findings_has_message(self, deck):
        # Introduce an overlap.
        pptx_add_shape(deck, 0, "rectangle", 1.0, 1.0, 3.0, 2.0, fill_color="2251FF")
        pptx_add_shape(deck, 0, "rectangle", 2.0, 1.5, 3.0, 2.0, fill_color="FF0000")
        payload = _unwrap_ok(pptx_check_layout(deck, detailed=True))
        assert "message" in payload
        msg = payload["message"]
        # Must mention "findings" or "clean" so generic agents can parse it.
        assert "findings" in msg or "clean" in msg


class TestCheckLayoutFileNotFound:
    """#107: missing file → FILE_NOT_FOUND (not INTERNAL_ERROR)."""

    def test_missing_file_returns_file_not_found(self, tmp_path):
        nonexistent = str(tmp_path / "nope.pptx")
        err = _unwrap_err(pptx_check_layout(nonexistent))
        assert err["code"] == "FILE_NOT_FOUND"

    def test_missing_file_detailed_returns_file_not_found(self, tmp_path):
        nonexistent = str(tmp_path / "nope.pptx")
        err = _unwrap_err(pptx_check_layout(nonexistent, detailed=True))
        assert err["code"] == "FILE_NOT_FOUND"


class TestNoStale25Tools:
    """#108: no hardcoded '25 tools' assertion left in tests.

    Historical comments referencing the cleanup are allowed; a
    ``len(...) == 25`` hardcoded count assertion is not.
    """

    def test_no_hardcoded_25_tools_assertion(self):
        # v0.3.1 (#108): the test asserts that the live FastMCP
        # registration has at least the core tools, without a brittle
        # hardcoded count.
        # v0.6.0 (#137): with tool tiering, the default surface is
        # ~20 tools (45 - 25 advanced). The historical > 25 bound no
        # longer holds with tiering enabled, so the assertion now
        # guards against accidental collapse of the default surface.
        tool_names = list(mcp._tool_manager._tools.keys())
        assert len(tool_names) >= 15, (
            "Default MCP surface should register at least the block-"
            f"component + batch tier — got {len(tool_names)}"
        )
