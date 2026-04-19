"""validation 拡張 (Issue #5) のテスト.

text_overflow / unreadable_text / divider_collision / inconsistent_gap を
それぞれ独立に検証し、最後に check_deck_extended の集計動作を確認する。
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx.enum.text import MSO_ANCHOR

from pptx_mcp_server.engine.validation import (
    ValidationFinding,
    _axis_groups,
    _effective_paragraph_size,
    _effective_run_size,
    _projected_text_range,
    check_deck_extended,
    check_divider_collision,
    check_inconsistent_gaps,
    check_text_overflow,
    check_unreadable_text,
)


# ---------------------------------------------------------------------------
# ヘルパ
# ---------------------------------------------------------------------------


def _add_textbox(
    slide,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    text: str,
    font_pt: float | None = None,
    name: str | None = None,
):
    """指定位置に text box を追加し、必要ならフォントサイズと name を設定する."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    if font_pt is not None:
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(font_pt)
    if name is not None:
        tb.name = name
    return tb


def _add_rect(
    slide,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    name: str | None = None,
):
    """rectangle auto shape を追加する."""
    from pptx.enum.shapes import MSO_SHAPE

    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    if name is not None:
        sh.name = name
    return sh


# ---------------------------------------------------------------------------
# text_overflow
# ---------------------------------------------------------------------------


class TestTextOverflow:
    """check_text_overflow の挙動を検証する."""

    def test_short_text_in_large_box_no_finding(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1.0, 1.0, 6.0, 3.0, "短いテキスト", font_pt=14)
        findings = check_text_overflow(one_slide_prs)
        assert findings == []

    def test_long_japanese_in_small_box_overflows(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        long_jp = (
            "これは非常に長い日本語のテキストであり、"
            "指定された小さなテキストボックスには到底収まらないほどの分量を持っている。"
            "さらに追加の説明文も含めて計測する。"
        )
        _add_textbox(slide, 1.0, 1.0, 2.0, 0.5, long_jp, font_pt=18, name="SmallBox")
        findings = check_text_overflow(one_slide_prs)
        assert len(findings) >= 1
        f = findings[0]
        assert f.category == "text_overflow"
        assert f.severity == "error"
        assert f.shape_name == "SmallBox"
        assert f.suggested_fix != ""


# ---------------------------------------------------------------------------
# unreadable_text
# ---------------------------------------------------------------------------


class TestUnreadableText:
    """check_unreadable_text の挙動を検証する."""

    def test_6pt_text_warns(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1.0, 1.0, 3.0, 0.5, "tiny", font_pt=6, name="TinyBox")
        findings = check_unreadable_text(one_slide_prs)
        assert len(findings) == 1
        f = findings[0]
        assert f.category == "unreadable_text"
        assert f.severity == "warning"
        assert f.shape_name == "TinyBox"

    def test_9pt_text_no_finding(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1.0, 1.0, 3.0, 0.5, "ok", font_pt=9, name="OkBox")
        findings = check_unreadable_text(one_slide_prs)
        assert findings == []

    def test_footer_named_shape_whitelisted(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        _add_textbox(
            slide, 1.0, 7.0, 3.0, 0.3, "Source: example",
            font_pt=7, name="Footer_Source",
        )
        findings = check_unreadable_text(one_slide_prs)
        assert findings == []


# ---------------------------------------------------------------------------
# divider_collision
# ---------------------------------------------------------------------------


class TestDividerCollision:
    """check_divider_collision の挙動を検証する."""

    def test_title_wraps_into_divider(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        # タイトル: width 3.0" に 14pt 長タイトルを入れ 2 行以上に折り返す想定
        long_title = (
            "これはスライドのアクションタイトルで必ず折り返される日本語テキスト"
        )
        _add_textbox(slide, 0.9, 0.45, 3.0, 0.5, long_title,
                     font_pt=14, name="Title")
        # divider: y=0.95" に薄い水平線 (高さ 0.01")
        _add_rect(slide, 0.9, 0.95, 11.5, 0.02, name="Divider")

        findings = check_divider_collision(one_slide_prs)
        assert len(findings) >= 1
        f = findings[0]
        assert f.category == "divider_collision"
        assert f.severity == "error"
        assert f.shape_name == "Title"

    def test_short_title_no_collision(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        # 十分に収まる短いタイトル (12pt 1 行 ≈ 0.20" → top 0.65 + 0.20 = 0.85 < 0.93)
        _add_textbox(slide, 0.9, 0.65, 11.5, 0.25, "短題",
                     font_pt=12, name="Title")
        _add_rect(slide, 0.9, 0.95, 11.5, 0.02, name="Divider")
        findings = check_divider_collision(one_slide_prs)
        assert findings == []


# ---------------------------------------------------------------------------
# inconsistent_gap
# ---------------------------------------------------------------------------


class TestInconsistentGap:
    """check_inconsistent_gaps の挙動を検証する."""

    def test_uneven_row_gaps(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        # 同じ top (=2.0") に 3 shape、幅 1.0"、x gap を [0.15, 0.30] に設定
        _add_rect(slide, 1.00, 2.0, 1.0, 1.0, name="A")
        _add_rect(slide, 2.15, 2.0, 1.0, 1.0, name="B")  # gap A→B = 0.15
        _add_rect(slide, 3.45, 2.0, 1.0, 1.0, name="C")  # gap B→C = 0.30
        # もう 1 shape で 3 つ目の gap を確保: gaps [0.15, 0.30] 差 = 0.15 > 0.05
        findings = check_inconsistent_gaps(one_slide_prs)
        assert len(findings) >= 1
        assert any(f.category == "inconsistent_gap" for f in findings)
        assert all(f.severity == "info" for f in findings)

    def test_consistent_row_gaps_no_finding(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        # gaps 0.20, 0.20, 0.20 → 差 0
        _add_rect(slide, 1.00, 2.0, 1.0, 1.0, name="A")
        _add_rect(slide, 2.20, 2.0, 1.0, 1.0, name="B")
        _add_rect(slide, 3.40, 2.0, 1.0, 1.0, name="C")
        _add_rect(slide, 4.60, 2.0, 1.0, 1.0, name="D")
        findings = check_inconsistent_gaps(one_slide_prs)
        # 行方向は一定、列方向も揃っていないので誤検出しない
        row_findings = [f for f in findings if "row" in f.message]
        assert row_findings == []


# ---------------------------------------------------------------------------
# check_deck_extended 集計
# ---------------------------------------------------------------------------


class TestCheckDeckExtended:
    """check_deck_extended が summary を正しく集計することを確認する."""

    def test_aggregated_summary_counts(self, one_slide_prs: Presentation) -> None:
        slide = one_slide_prs.slides[0]
        # 1) text_overflow (error)
        long_jp = (
            "非常に長い日本語のテキストでありこの幅の箱に収まらない量を持つ。"
            "追加文章もあるので必ずオーバーフローする。"
        )
        _add_textbox(slide, 1.0, 1.0, 2.0, 0.5, long_jp,
                     font_pt=18, name="Overflow")
        # 2) unreadable (warning)
        _add_textbox(slide, 1.0, 5.0, 3.0, 0.3, "tiny",
                     font_pt=6, name="TinyBox")

        result = check_deck_extended(one_slide_prs)

        assert "slides" in result
        assert "summary" in result
        assert len(result["slides"]) == 1
        slide0 = result["slides"][0]
        # 後方互換キーの保持を確認
        assert "overlaps" in slide0
        assert "out_of_bounds" in slide0
        # 新キー
        assert "text_overflow" in slide0
        assert "unreadable_text" in slide0
        assert "divider_collision" in slide0
        assert "inconsistent_gaps" in slide0

        assert len(slide0["text_overflow"]) >= 1
        assert len(slide0["unreadable_text"]) >= 1

        summary = result["summary"]
        assert summary["errors"] >= 1
        assert summary["warnings"] >= 1
        assert summary["infos"] >= 0


# ---------------------------------------------------------------------------
# ValidationFinding dataclass の簡易確認
# ---------------------------------------------------------------------------


def test_validation_finding_to_dict_shape() -> None:
    f = ValidationFinding(
        severity="error",
        slide_index=0,
        shape_name="X",
        category="text_overflow",
        message="m",
        suggested_fix="fix",
    )
    d = f.to_dict()
    assert d["severity"] == "error"
    assert d["slide_index"] == 0
    assert d["shape_name"] == "X"
    assert d["category"] == "text_overflow"
    assert d["message"] == "m"
    assert d["suggested_fix"] == "fix"


# ---------------------------------------------------------------------------
# _axis_groups drift (Issue #21)
# ---------------------------------------------------------------------------


class TestAxisGroupsDrift:
    """``_axis_groups`` が running-mean drift で誤吸収しないことを確認する."""

    def test_linear_drift_not_one_group(self) -> None:
        # tops = [0.0, 0.05, 0.10, 0.15, 0.20] は anchor 基準だと
        # 0.0 基準で tolerance 0.1 を超える要素 (0.15, 0.20) が別 group になる
        values = [
            (0, 0.0),
            (1, 0.05),
            (2, 0.10),
            (3, 0.15),
            (4, 0.20),
        ]
        groups = _axis_groups(values, tolerance=0.1)
        # 5 要素が全部 1 group にまとめられる (= running-mean drift) ことは
        # あってはならない
        assert not (len(groups) == 1 and len(groups[0]) == 5)

    def test_tight_cluster_is_one_group(self) -> None:
        # anchor 1.0 で全要素が tolerance 0.1 以内 → 1 group
        values = [(0, 1.0), (1, 1.0), (2, 1.09)]
        groups = _axis_groups(values, tolerance=0.1)
        assert groups == [[0, 1, 2]]

    def test_boundary_split_into_two_groups(self) -> None:
        # 1.0 と 1.18 は tolerance 0.1 を超えるので別 group
        values = [(0, 1.0), (1, 1.18)]
        groups = _axis_groups(values, tolerance=0.1)
        # どちらも 3 要素未満なので返却されず、空リストになる
        assert groups == []

    def test_regression_3kpi_uneven_gaps_still_flagged(
        self, one_slide_prs: Presentation
    ) -> None:
        # gaps [0.15, 0.15, 0.30] を持つ 4-KPI row が引き続き inconsistent
        # と判定されることを担保する (drift 修正でも既存動作が崩れない)
        slide = one_slide_prs.slides[0]
        _add_rect(slide, 1.00, 2.0, 1.0, 1.0, name="A")
        _add_rect(slide, 2.15, 2.0, 1.0, 1.0, name="B")  # gap 0.15
        _add_rect(slide, 3.30, 2.0, 1.0, 1.0, name="C")  # gap 0.15
        _add_rect(slide, 4.60, 2.0, 1.0, 1.0, name="D")  # gap 0.30
        findings = check_inconsistent_gaps(one_slide_prs)
        assert any(f.category == "inconsistent_gap" for f in findings)


# ---------------------------------------------------------------------------
# check_deck_extended summary aggregation (Issue #22)
# ---------------------------------------------------------------------------


class TestSummaryAggregation:
    """``summary`` が legacy overlaps / out_of_bounds も含めて集計することを確認する."""

    def test_clean_deck_all_zeros(self, one_slide_prs: Presentation) -> None:
        # 何も置かないスライドは summary 全 0
        result = check_deck_extended(one_slide_prs)
        assert result["summary"] == {"errors": 0, "warnings": 0, "infos": 0}

    def test_single_overlap_counts_as_error(
        self, one_slide_prs: Presentation
    ) -> None:
        slide = one_slide_prs.slides[0]
        # 大きく重なる 2 つの矩形 (overlap area >= 0.15 sq in)
        _add_rect(slide, 1.0, 1.0, 3.0, 2.0, name="A")
        _add_rect(slide, 2.0, 1.5, 3.0, 2.0, name="B")
        result = check_deck_extended(one_slide_prs)
        slide0 = result["slides"][0]
        assert len(slide0["overlaps"]) == 1
        assert result["summary"]["errors"] == 1

    def test_two_overlaps_and_one_unreadable(
        self, one_slide_prs: Presentation
    ) -> None:
        slide = one_slide_prs.slides[0]
        # 2 つの overlap 警告が生じる 3-way 重なり配置 (A-B, A-C, B-C の 3 組みで
        # すべて大きな area を持つように調整)
        _add_rect(slide, 1.0, 1.0, 3.0, 2.0, name="A")
        _add_rect(slide, 2.0, 1.5, 3.0, 2.0, name="B")
        _add_rect(slide, 6.0, 1.0, 3.0, 2.0, name="C")
        _add_rect(slide, 7.0, 1.5, 3.0, 2.0, name="D")
        # unreadable 警告 (6pt)
        _add_textbox(slide, 0.5, 6.0, 3.0, 0.3, "tiny", font_pt=6, name="TinyBox")

        result = check_deck_extended(one_slide_prs)
        slide0 = result["slides"][0]
        assert len(slide0["overlaps"]) == 2
        assert result["summary"]["errors"] == 2
        assert result["summary"]["warnings"] == 1


# ---------------------------------------------------------------------------
# paragraph-level font sizing (Issue #23)
# ---------------------------------------------------------------------------


class TestParagraphSizeResolution:
    """段落単位の font サイズ解決と overflow 判定."""

    def test_mixed_size_paragraphs_no_false_overflow(
        self, one_slide_prs: Presentation
    ) -> None:
        # 24pt 1行 + 10pt 1行 が収まる十分な高さの box
        # 24pt → ≒ 0.40" 、10pt → ≒ 0.17"、合計 ≒ 0.57"
        # 旧実装は 24pt × 2 行で ≒ 0.80" と過大評価し false overflow になる。
        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        # 最初の paragraph を 24pt タイトル
        tf.text = "Title"
        tf.paragraphs[0].runs[0].font.size = Pt(24)
        # 2 段落目を 10pt body
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "body"
        r2.font.size = Pt(10)
        tb.name = "MixedBox"

        findings = check_text_overflow(one_slide_prs)
        # MixedBox に text_overflow finding は出てはならない
        matched = [f for f in findings if f.shape_name == "MixedBox"]
        assert matched == []

    def test_defrpr_only_size_detected(self, one_slide_prs: Presentation) -> None:
        from lxml import etree

        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tf = tb.text_frame
        tf.text = "hello"
        para = tf.paragraphs[0]
        # run レベル size を除去
        for run in para.runs:
            run.font.size = None
        # <a:pPr><a:defRPr sz="2400"/></a:pPr> を手動で付与 (24pt)
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap = {"a": a_ns}
        p_elem = para._p
        # 既存 pPr を削除してから新たに差し込む
        for existing in p_elem.findall(f"{{{a_ns}}}pPr"):
            p_elem.remove(existing)
        p_pr = etree.SubElement(p_elem, f"{{{a_ns}}}pPr")
        p_elem.insert(0, p_pr)
        etree.SubElement(p_pr, f"{{{a_ns}}}defRPr", {"sz": "2400"})

        # helper が defRPr 経由で 24pt を返すこと
        eff = _effective_paragraph_size(para, default_pt=18.0)
        assert eff == pytest.approx(24.0)

    def test_single_size_frame_still_detects_overflow(
        self, one_slide_prs: Presentation
    ) -> None:
        # 既存挙動の回帰: 18pt 長文が小さな box からあふれる
        slide = one_slide_prs.slides[0]
        long_jp = (
            "これは非常に長い日本語のテキストであり、"
            "指定された小さなテキストボックスには到底収まらない分量である。"
        )
        _add_textbox(slide, 1.0, 1.0, 2.0, 0.5, long_jp, font_pt=18, name="OverflowBox")
        findings = check_text_overflow(one_slide_prs)
        assert any(f.shape_name == "OverflowBox" for f in findings)


# ---------------------------------------------------------------------------
# Issue #44: unreadable whitelist by heuristic and explicit names
# ---------------------------------------------------------------------------


class TestUnreadableWhitelistExtended:
    """``check_unreadable_text`` の whitelist_names / heuristic 除外の検証."""

    def test_japanese_footer_heuristic_excluded(
        self, one_slide_prs: Presentation
    ) -> None:
        # 脚注: スライド底辺から 0.6" 以内 (y=7.1, slide_h=7.5 → top は 7.1, bbox 完全に下 0.4")
        # 6pt (= min_readable_pt - 2) は閾値 min_readable_pt - 1 を下回るため
        # ヒューリスティックで除外されない。フォント 7pt (= 8 - 1) に調整して
        # フッタ ヒューリスティック該当にする。
        slide = one_slide_prs.slides[0]
        _add_textbox(
            slide, 0.9, 7.1, 3.0, 0.3, "出典: 社内資料",
            font_pt=7, name="脚注",
        )
        findings = check_unreadable_text(one_slide_prs)
        assert findings == []

    def test_japanese_footer_6pt_still_flagged_by_heuristic(
        self, one_slide_prs: Presentation
    ) -> None:
        # 6pt は min_readable_pt - 1 を下回るので heuristic は救わない。
        # ただし明示 whitelist_names=["脚注"] を渡せば除外できることを次のテストで担保。
        slide = one_slide_prs.slides[0]
        _add_textbox(
            slide, 0.9, 7.1, 3.0, 0.3, "出典: 社内資料",
            font_pt=6, name="脚注",
        )
        findings = check_unreadable_text(one_slide_prs)
        # 6pt の「脚注」は heuristic の font しきい値を外れるため警告が出る
        assert len(findings) == 1
        assert findings[0].shape_name == "脚注"

    def test_japanese_footer_6pt_whitelist_excluded(
        self, one_slide_prs: Presentation
    ) -> None:
        # whitelist_names に "脚注" を指定すれば 6pt でも skip される
        slide = one_slide_prs.slides[0]
        _add_textbox(
            slide, 0.9, 7.1, 3.0, 0.3, "出典: 社内資料",
            font_pt=6, name="脚注",
        )
        findings = check_unreadable_text(
            one_slide_prs, whitelist_names=["脚注", "フッター"]
        )
        assert findings == []

    def test_midslide_small_text_still_flagged(
        self, one_slide_prs: Presentation
    ) -> None:
        # mid-slide (y=3.0) にある 6pt text は heuristic で救われない
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1.0, 3.0, 3.0, 0.5, "fact", font_pt=6, name="fact")
        findings = check_unreadable_text(one_slide_prs)
        assert len(findings) == 1
        assert findings[0].shape_name == "fact"

    def test_whitelist_matching_vs_non_matching(
        self, one_slide_prs: Presentation
    ) -> None:
        # whitelist_names に合致するものは skip、合致しないものは flag される
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1.0, 3.0, 3.0, 0.3, "tiny", font_pt=6, name="脚注_1")
        _add_textbox(slide, 1.0, 4.0, 3.0, 0.3, "tiny", font_pt=6, name="Body_Small")
        findings = check_unreadable_text(
            one_slide_prs, whitelist_names=["脚注"]
        )
        names = [f.shape_name for f in findings]
        assert "脚注_1" not in names
        assert "Body_Small" in names

    def test_english_regex_whitelist_regression(
        self, one_slide_prs: Presentation
    ) -> None:
        # 既存英語正規表現 (footer / page_number / source / footnote) が引き続き有効
        slide = one_slide_prs.slides[0]
        _add_textbox(
            slide, 1.0, 3.0, 3.0, 0.3, "1",
            font_pt=6, name="page_number_1",
        )
        findings = check_unreadable_text(one_slide_prs)
        assert findings == []

    def test_whitelist_names_case_insensitive(
        self, one_slide_prs: Presentation
    ) -> None:
        # 部分一致・大文字小文字を区別しないことを確認
        slide = one_slide_prs.slides[0]
        _add_textbox(
            slide, 1.0, 3.0, 3.0, 0.3, "note",
            font_pt=6, name="My_FOOTNOTE_Box",
        )
        findings = check_unreadable_text(
            one_slide_prs, whitelist_names=["footnote"]
        )
        assert findings == []

    def test_deck_extended_propagates_whitelist_names(
        self, one_slide_prs: Presentation
    ) -> None:
        # check_deck_extended も whitelist_names を受け取り unreadable に伝搬する
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1.0, 3.0, 3.0, 0.3, "tiny", font_pt=6, name="脚注")
        result = check_deck_extended(one_slide_prs, whitelist_names=["脚注"])
        assert result["slides"][0]["unreadable_text"] == []
        assert result["summary"]["warnings"] == 0


# ---------------------------------------------------------------------------
# Issue #60: check_unreadable_text per-run inherited-size resolution
# ---------------------------------------------------------------------------


def _install_defrpr(paragraph, sz_hundredths: int) -> None:
    """段落に ``<a:pPr><a:defRPr sz="NN"/></a:pPr>`` を手動で付与する."""
    from lxml import etree

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_elem = paragraph._p
    for existing in p_elem.findall(f"{{{a_ns}}}pPr"):
        p_elem.remove(existing)
    p_pr = etree.SubElement(p_elem, f"{{{a_ns}}}pPr")
    # pPr は段落の先頭子要素にある必要がある
    p_elem.insert(0, p_pr)
    etree.SubElement(p_pr, f"{{{a_ns}}}defRPr", {"sz": str(sz_hundredths)})


class TestUnreadableRunInheritance:
    """check_unreadable_text が run ごとに継承サイズを解決することを検証する (Issue #60)."""

    def test_mixed_run_with_inherited_small_run_flagged(
        self, one_slide_prs: Presentation
    ) -> None:
        # 段落に 11pt (明示) run と size=None の run (defRPr sz=600 → 6pt) を同居させる。
        # 旧実装は「explicit run がある段落は run レベルのみ」で早期確定するため
        # 継承 6pt を取りこぼす。新実装は run ごとに実効サイズを解決し警告する。
        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tb.name = "MixedInheritBox"
        tf = tb.text_frame
        tf.text = "A"
        para = tf.paragraphs[0]
        # run A: 11pt 明示
        run_a = para.runs[0]
        run_a.font.size = Pt(11)
        # run B: size=None (継承)
        run_b = para.add_run()
        run_b.text = "B"
        assert run_b.font.size is None
        # 段落 defRPr = 600 (= 6pt)
        _install_defrpr(para, 600)

        findings = check_unreadable_text(one_slide_prs)
        matched = [f for f in findings if f.shape_name == "MixedInheritBox"]
        assert len(matched) == 1
        assert "6.0pt" in matched[0].message

    def test_mixed_run_all_explicit_readable_no_finding(
        self, one_slide_prs: Presentation
    ) -> None:
        # 同じ段落内に 11pt と 14pt の明示 run → どちらも閾値以上なので警告なし。
        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tb.name = "TwoExplicitBox"
        tf = tb.text_frame
        tf.text = "A"
        para = tf.paragraphs[0]
        para.runs[0].font.size = Pt(11)
        run_b = para.add_run()
        run_b.text = "B"
        run_b.font.size = Pt(14)

        findings = check_unreadable_text(one_slide_prs)
        matched = [f for f in findings if f.shape_name == "TwoExplicitBox"]
        assert matched == []

    def test_all_runs_inherit_readable_defrpr_no_finding(
        self, one_slide_prs: Presentation
    ) -> None:
        # 全 run が size=None、defRPr sz=900 (9pt) → 閾値 8pt を上回るため警告なし。
        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tb.name = "InheritReadableBox"
        tf = tb.text_frame
        tf.text = "X"
        para = tf.paragraphs[0]
        for run in para.runs:
            run.font.size = None
        _install_defrpr(para, 900)

        findings = check_unreadable_text(one_slide_prs)
        matched = [f for f in findings if f.shape_name == "InheritReadableBox"]
        assert matched == []

    def test_all_runs_inherit_small_defrpr_flagged_regression(
        self, one_slide_prs: Presentation
    ) -> None:
        # Issue #23 の回帰: 全 run が size=None、defRPr sz=600 (6pt) → 警告 1 件。
        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tb.name = "InheritSmallBox"
        tf = tb.text_frame
        tf.text = "X"
        para = tf.paragraphs[0]
        for run in para.runs:
            run.font.size = None
        _install_defrpr(para, 600)

        findings = check_unreadable_text(one_slide_prs)
        matched = [f for f in findings if f.shape_name == "InheritSmallBox"]
        assert len(matched) == 1
        assert "6.0pt" in matched[0].message

    def test_effective_run_size_helper_resolves_defrpr(
        self, one_slide_prs: Presentation
    ) -> None:
        # helper 単体: size=None の run が defRPr 経由で 6pt を解決できること。
        slide = one_slide_prs.slides[0]
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0)
        )
        tf = tb.text_frame
        tf.text = "A"
        para = tf.paragraphs[0]
        para.runs[0].font.size = Pt(11)
        run_b = para.add_run()
        run_b.text = "B"
        _install_defrpr(para, 600)

        assert _effective_run_size(para.runs[0], para, default_pt=18.0) == pytest.approx(11.0)
        assert _effective_run_size(run_b, para, default_pt=18.0) == pytest.approx(6.0)


# ---------------------------------------------------------------------------
# Issue #45: divider_collision anchor-aware projection
# ---------------------------------------------------------------------------


def _set_vertical_anchor(textbox, anchor) -> None:
    """textbox の vertical_anchor を設定するヘルパ."""
    textbox.text_frame.vertical_anchor = anchor


class TestDividerCollisionAnchor:
    """vertical_anchor に応じた collision 判定の検証 (Issue #45)."""

    def test_bottom_anchored_no_collision(self, one_slide_prs: Presentation) -> None:
        # McKinsey 風: title at y=0.45 h=0.5 (bottom at 0.95), anchor=bottom
        # 長いタイトルが 2 行に折り返されるが bottom が divider top (0.95") に揃うため
        # 実テキストは上方向に伸びて divider を超えない → no collision
        slide = one_slide_prs.slides[0]
        long_title = (
            "これはスライドのアクションタイトルで必ず折り返される日本語テキスト"
        )
        tb = _add_textbox(
            slide, 0.9, 0.45, 3.0, 0.5, long_title,
            font_pt=14, name="Title",
        )
        _set_vertical_anchor(tb, MSO_ANCHOR.BOTTOM)
        _add_rect(slide, 0.9, 0.95, 11.5, 0.02, name="Divider")

        findings = check_divider_collision(one_slide_prs)
        assert findings == [], f"unexpected findings: {[f.message for f in findings]}"

    def test_top_anchored_collision_detected(
        self, one_slide_prs: Presentation
    ) -> None:
        # 同じ寸法で top-anchor は collision として検出される (回帰)
        slide = one_slide_prs.slides[0]
        long_title = (
            "これはスライドのアクションタイトルで必ず折り返される日本語テキスト"
        )
        tb = _add_textbox(
            slide, 0.9, 0.45, 3.0, 0.5, long_title,
            font_pt=14, name="Title",
        )
        _set_vertical_anchor(tb, MSO_ANCHOR.TOP)
        _add_rect(slide, 0.9, 0.95, 11.5, 0.02, name="Divider")

        findings = check_divider_collision(one_slide_prs)
        assert len(findings) >= 1
        assert findings[0].shape_name == "Title"

    def test_middle_anchored_collision_edge_case(
        self, one_slide_prs: Presentation
    ) -> None:
        # middle-anchor で center=0.70", needed_height≈0.6 → range [0.40, 1.00]
        # 1.00 が divider top 0.95 を越えるため collision
        slide = one_slide_prs.slides[0]
        long_title = (
            "これはスライドのアクションタイトルで必ず折り返される日本語テキスト"
        )
        tb = _add_textbox(
            slide, 0.9, 0.45, 3.0, 0.5, long_title,
            font_pt=14, name="Title",
        )
        _set_vertical_anchor(tb, MSO_ANCHOR.MIDDLE)
        _add_rect(slide, 0.9, 0.95, 11.5, 0.02, name="Divider")

        findings = check_divider_collision(one_slide_prs)
        assert len(findings) >= 1
        assert findings[0].shape_name == "Title"

    def test_bottom_anchored_short_text_no_collision(
        self, one_slide_prs: Presentation
    ) -> None:
        # bottom-anchor + 短文は当然 collision なし
        slide = one_slide_prs.slides[0]
        tb = _add_textbox(
            slide, 0.9, 0.45, 11.5, 0.5, "短題",
            font_pt=14, name="Title",
        )
        _set_vertical_anchor(tb, MSO_ANCHOR.BOTTOM)
        _add_rect(slide, 0.9, 0.95, 11.5, 0.02, name="Divider")

        findings = check_divider_collision(one_slide_prs)
        assert findings == []

    def test_projected_text_range_helper(self, one_slide_prs: Presentation) -> None:
        # _projected_text_range ヘルパの単体挙動を検証
        slide = one_slide_prs.slides[0]
        tb = _add_textbox(
            slide, 0.9, 0.45, 3.0, 0.5, "x",
            font_pt=14, name="Title",
        )
        tf = tb.text_frame

        # TOP
        tf.vertical_anchor = MSO_ANCHOR.TOP
        top, bot = _projected_text_range(tb, tf, needed_height=0.6)
        assert top == pytest.approx(0.45, abs=1e-3)
        assert bot == pytest.approx(1.05, abs=1e-3)

        # BOTTOM (s_bottom = 0.45 + 0.5 = 0.95)
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        top, bot = _projected_text_range(tb, tf, needed_height=0.6)
        assert top == pytest.approx(0.35, abs=1e-3)
        assert bot == pytest.approx(0.95, abs=1e-3)

        # MIDDLE (center = 0.70)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        top, bot = _projected_text_range(tb, tf, needed_height=0.6)
        assert top == pytest.approx(0.40, abs=1e-3)
        assert bot == pytest.approx(1.00, abs=1e-3)
