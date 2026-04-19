"""auto-fit textbox ヘルパーのテスト.

``add_auto_fit_textbox`` が指定ボックスに収まる font size を正しく
選択し、min size でも収まらない場合は省略記号で切り詰める挙動を検証する。
"""

from __future__ import annotations

import pytest
from pptx.enum.text import MSO_ANCHOR

from pptx_mcp_server.engine.shapes import add_auto_fit_textbox


def test_short_text_big_box_keeps_font_size(slide):
    """短いテキスト + 余裕のある box では font size は縮まない."""
    shape, actual = add_auto_fit_textbox(
        slide,
        "Hello",
        left=1.0, top=1.0, width=5.0, height=2.0,
        font_size_pt=11, min_size_pt=7,
    )
    assert actual == 11.0
    assert shape.text_frame.text == "Hello"


def test_long_japanese_in_narrow_box_shrinks(slide):
    """長い日本語 + 狭い box では font size が縮小されるが min 以上に留まる."""
    long_jp = "これは日本語のテキストサンプルです。" * 6
    shape, actual = add_auto_fit_textbox(
        slide,
        long_jp,
        left=1.0, top=1.0, width=3.0, height=0.8,
        font_size_pt=14, min_size_pt=7,
    )
    assert actual < 14.0
    assert actual >= 7.0


def test_overflow_truncates_with_ellipsis(slide):
    """min_size でも入らない量を truncate=True で流し込むと省略記号で終わる."""
    huge = "これは非常に長いテキストです。" * 40
    shape, actual = add_auto_fit_textbox(
        slide,
        huge,
        left=1.0, top=1.0, width=1.0, height=0.3,
        font_size_pt=11, min_size_pt=7,
        truncate_with_ellipsis=True,
    )
    assert actual == 7.0
    assert shape.text_frame.text.endswith("\u2026")


def test_overflow_without_truncate_keeps_full_text(slide):
    """truncate=False の場合は full text をそのまま描画しオーバーフローを許容する."""
    huge = "これは非常に長いテキストです。" * 40
    shape, actual = add_auto_fit_textbox(
        slide,
        huge,
        left=1.0, top=1.0, width=1.0, height=0.3,
        font_size_pt=11, min_size_pt=7,
        truncate_with_ellipsis=False,
    )
    assert actual == 7.0
    assert shape.text_frame.text == huge
    assert "\u2026" not in shape.text_frame.text


def test_vertical_anchor_middle(slide):
    """vertical_anchor='middle' が MSO_ANCHOR.MIDDLE として反映される."""
    shape, _ = add_auto_fit_textbox(
        slide,
        "centered",
        left=1.0, top=1.0, width=3.0, height=1.5,
        vertical_anchor="middle",
    )
    assert shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE


# ---------------------------------------------------------------------------
# Issue #79: wrap=False モード
# ---------------------------------------------------------------------------


class TestWrapFalse:
    """``wrap=False`` で単一行 auto-fit と word_wrap=False が機能することを検証する."""

    def test_short_title_wide_frame_no_shrink(self, slide):
        """短いタイトル + 余裕のある幅なら font size は縮まない."""
        shape, actual = add_auto_fit_textbox(
            slide,
            "Short title",
            left=0.9, top=0.45, width=11.533, height=0.55,
            font_size_pt=22, min_size_pt=12,
            wrap=False,
        )
        assert actual == 22.0
        assert shape.text_frame.word_wrap is False

    def test_long_japanese_title_shrinks_until_single_line_fits(self, slide):
        """長い日本語タイトル + 11.533"×0.55" + wrap=False で単一行に収まる size まで縮む.

        Issue #79 の production シナリオ再現。``wrap=False`` では width 基準で
        縮み、``estimate_text_width`` が usable width 以下になった時点で停止する。
        20pt で開始 → 幅超過のため縮小 → 収まる size まで下がる。
        """
        from pptx_mcp_server.engine.text_metrics import estimate_text_width
        from pptx_mcp_server.engine.layout_constants import (
            TEXTBOX_INNER_PADDING_TOTAL,
        )

        title = (
            "Mac乗り換え検討者に対してLenovoはほぼ不可視 — "
            "「Macに匹敵する体験をWindowsで」のナラティブが不在"
        )
        shape, actual = add_auto_fit_textbox(
            slide,
            title,
            left=0.9, top=0.45, width=11.533, height=0.55,
            font_size_pt=20, min_size_pt=12,
            wrap=False,
        )
        # 開始の 20pt では幅超過するため必ず縮む
        assert actual < 20.0
        # min_size_pt 以上には留まる
        assert actual >= 12.0
        # 決定 size で実際に単一行 width に収まっている
        usable = 11.533 - TEXTBOX_INNER_PADDING_TOTAL
        assert estimate_text_width(title, actual) <= usable + 1e-6
        # PowerPoint 側でも wrap しない
        assert shape.text_frame.word_wrap is False

    def test_min_size_still_over_truncates(self, slide):
        """min size でも幅を超える場合、truncate_with_ellipsis=True で末尾省略."""
        huge = "これは非常に長いタイトルです。" * 10
        shape, actual = add_auto_fit_textbox(
            slide,
            huge,
            left=0.9, top=0.45, width=3.0, height=0.55,
            font_size_pt=20, min_size_pt=12,
            wrap=False,
            truncate_with_ellipsis=True,
        )
        assert actual == 12.0
        rendered = shape.text_frame.text
        assert rendered.endswith("\u2026")
        # 切り詰め後の幅は usable width 以下
        from pptx_mcp_server.engine.text_metrics import estimate_text_width
        from pptx_mcp_server.engine.layout_constants import (
            TEXTBOX_INNER_PADDING_TOTAL,
        )

        usable = 3.0 - TEXTBOX_INNER_PADDING_TOTAL
        assert estimate_text_width(rendered, actual) <= usable + 1e-6

    def test_min_size_still_over_no_truncate_preserves_text(self, slide):
        """truncate=False なら min_size で full text を保持する (clip 許容)."""
        huge = "これは非常に長いタイトルです。" * 10
        shape, actual = add_auto_fit_textbox(
            slide,
            huge,
            left=0.9, top=0.45, width=3.0, height=0.55,
            font_size_pt=20, min_size_pt=12,
            wrap=False,
            truncate_with_ellipsis=False,
        )
        assert actual == 12.0
        assert shape.text_frame.text == huge
        assert "\u2026" not in shape.text_frame.text
        assert shape.text_frame.word_wrap is False

    def test_default_wrap_true_unchanged(self, slide):
        """回帰: wrap 指定なし (既定 True) は従来通り高さベースで縮む.

        既存テスト ``test_long_japanese_in_narrow_box_shrinks`` と同等シナリオで
        word_wrap が True のまま残ることを確認する。
        """
        long_jp = "これは日本語のテキストサンプルです。" * 6
        shape, actual = add_auto_fit_textbox(
            slide,
            long_jp,
            left=1.0, top=1.0, width=3.0, height=0.8,
            font_size_pt=14, min_size_pt=7,
        )
        assert actual < 14.0
        assert shape.text_frame.word_wrap is True


def test_production_subtitle_11_5x0_5(slide):
    """production-like: 約 120 文字の日本語を 11.5" x 0.5" の箱に流し込む.

    McKinsey スタイルの 1-line アクションサブタイトルを想定したシナリオ。
    結果の font size は [min_size_pt, 11] の範囲に収まる想定。
    """
    subtitle = (
        "本施策は複数部門の合意形成を前提としつつ、短期では収益性の早期改善、"
        "中期では顧客基盤を軸とした市場シェアの段階的拡大、"
        "長期では持続的な競争優位と価値創出の確立を目指す統合的アプローチとして位置づけるものである"
    )
    assert 100 <= len(subtitle) <= 140  # 仕様 ~120 chars
    shape, actual = add_auto_fit_textbox(
        slide,
        subtitle,
        left=0.9, top=0.5, width=11.5, height=0.5,
        font_size_pt=11, min_size_pt=7,
    )
    assert 7.0 <= actual <= 11.0
    # 切り詰めが起きない場合、本文は保持されているはず (またはオーバーフロー
    # で truncate 発動)。いずれにせよ空ではない。
    assert shape.text_frame.text != ""
