"""
Unit tests for pptx_io -- file operations and color parsing.
"""

from __future__ import annotations

import os
import sys

import pytest
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from pptx_mcp_server.engine import pptx_io
from pptx_mcp_server.engine.pptx_io import (
    EngineError,
    ErrorCode,
    create_presentation,
    open_pptx,
    save_pptx,
    _parse_color,
)


class TestCreatePresentation:
    """create_presentation must produce a valid PPTX with correct dimensions."""

    def test_creates_file(self, tmp_path):
        path = str(tmp_path / "new.pptx")
        result = create_presentation(path)
        assert os.path.exists(path)
        assert "Created presentation" in result

    def test_default_dimensions(self, tmp_path):
        path = str(tmp_path / "wide.pptx")
        create_presentation(path)
        prs = Presentation(path)
        assert abs(prs.slide_width / 914400 - 13.333) < 0.01
        assert abs(prs.slide_height / 914400 - 7.5) < 0.01

    def test_custom_dimensions(self, tmp_path):
        path = str(tmp_path / "custom.pptx")
        create_presentation(path, width_inches=10, height_inches=5.625)
        prs = Presentation(path)
        assert abs(prs.slide_width / 914400 - 10.0) < 0.01
        assert abs(prs.slide_height / 914400 - 5.625) < 0.01


class TestOpenPptx:
    """open_pptx error handling for missing and corrupt files."""

    def test_nonexistent_file_raises_file_not_found(self, tmp_path):
        with pytest.raises(EngineError) as exc_info:
            open_pptx(str(tmp_path / "nope.pptx"))
        assert exc_info.value.code == ErrorCode.FILE_NOT_FOUND

    def test_invalid_file_raises_invalid_pptx(self, tmp_path):
        bad = tmp_path / "bad.pptx"
        bad.write_text("this is not a pptx")
        with pytest.raises(EngineError) as exc_info:
            open_pptx(str(bad))
        assert exc_info.value.code == ErrorCode.INVALID_PPTX


class TestParseColor:
    """_parse_color must handle hex strings with and without '#'."""

    def test_with_hash(self):
        c = _parse_color("#FF0000")
        assert c == RGBColor(0xFF, 0x00, 0x00)

    def test_without_hash(self):
        c = _parse_color("00FF00")
        assert c == RGBColor(0x00, 0xFF, 0x00)

    def test_lowercase_hex(self):
        c = _parse_color("#aabbcc")
        assert c == RGBColor(0xAA, 0xBB, 0xCC)

    def test_invalid_short_hex_raises(self):
        with pytest.raises(EngineError) as exc_info:
            _parse_color("FFF")
        assert exc_info.value.code == ErrorCode.INVALID_PARAMETER

    def test_invalid_long_hex_raises(self):
        with pytest.raises(EngineError) as exc_info:
            _parse_color("#FFAABBCC")
        assert exc_info.value.code == ErrorCode.INVALID_PARAMETER


class TestSavePptx:
    """save_pptx must write a valid file atomically."""

    def test_saves_to_disk(self, tmp_path):
        prs = Presentation()
        path = str(tmp_path / "saved.pptx")
        save_pptx(prs, path)
        assert os.path.exists(path)
        # Verify it's loadable
        loaded = Presentation(path)
        assert loaded is not None

    def test_happy_path_round_trips(self, tmp_path):
        """Happy path: save round-trips through Presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        # Add a blank slide so we can inspect structure after reload.
        layout = prs.slide_layouts[6]
        prs.slides.add_slide(layout)

        path = str(tmp_path / "rt.pptx")
        save_pptx(prs, path)

        assert os.path.exists(path)
        loaded = Presentation(path)
        assert len(loaded.slides) == 1
        assert abs(loaded.slide_width / 914400 - 13.333) < 0.01

        # No lingering temp files.
        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []

    def test_crash_during_save_preserves_original(self, tmp_path, monkeypatch):
        """If Presentation.save raises, the pre-existing target file is untouched."""
        path = tmp_path / "existing.pptx"
        original_bytes = b"ORIGINAL CONTENT -- not a real pptx"
        path.write_bytes(original_bytes)

        prs = Presentation()

        def boom(*_args, **_kwargs):
            raise IOError("simulated mid-write failure")

        # Presentation() is a factory; patch save on the instance's class.
        monkeypatch.setattr(type(prs), "save", boom, raising=True)

        with pytest.raises(IOError, match="simulated"):
            save_pptx(prs, str(path))

        # Original content preserved verbatim.
        assert path.read_bytes() == original_bytes
        # No lingering .tmp.* sibling.
        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []

    def test_crash_with_no_existing_target_leaves_no_file(self, tmp_path, monkeypatch):
        """If target did not exist, a failing save must not create it."""
        path = tmp_path / "fresh.pptx"
        assert not path.exists()

        prs = Presentation()

        def boom(*_args, **_kwargs):
            raise IOError("simulated mid-write failure")

        # Presentation() is a factory; patch save on the instance's class.
        monkeypatch.setattr(type(prs), "save", boom, raising=True)

        with pytest.raises(IOError, match="simulated"):
            save_pptx(prs, str(path))

        assert not path.exists()
        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []

    def test_two_sequential_saves_same_path(self, tmp_path):
        """Back-to-back saves to the same path must both succeed."""
        path = str(tmp_path / "twice.pptx")

        prs1 = Presentation()
        save_pptx(prs1, path)
        assert os.path.exists(path)
        first_mtime = os.path.getmtime(path)

        prs2 = Presentation()
        save_pptx(prs2, path)
        assert os.path.exists(path)

        # File still loadable after the second save.
        loaded = Presentation(path)
        assert loaded is not None

        # No lingering temp files.
        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []
        # Final file is the second save (or at least not removed).
        assert os.path.getmtime(path) >= first_mtime

    def test_relative_path(self, tmp_path, monkeypatch):
        """Relative paths work identically (resolved via os.path.abspath)."""
        monkeypatch.chdir(tmp_path)
        prs = Presentation()
        save_pptx(prs, "relative.pptx")

        target = tmp_path / "relative.pptx"
        assert target.exists()
        loaded = Presentation(str(target))
        assert loaded is not None

        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []


class TestSavePptxFsync:
    """fsync=True opt-in path must call os.fsync; default must not."""

    @staticmethod
    def _install_fsync_counter(monkeypatch):
        """Patch os.fsync via the pptx_io module's reference so the call site
        picks up the mock (the module does `import os` then calls `os.fsync`,
        so patching `os.fsync` on the shared `os` module is what the call
        site sees)."""
        calls = {"count": 0, "fds": []}
        real_fsync = os.fsync

        def mock_fsync(fd):
            calls["count"] += 1
            calls["fds"].append(fd)
            return real_fsync(fd)

        # Patch on the shared os module — pptx_io does `import os` so the
        # attribute lookup goes through this module.
        monkeypatch.setattr(os, "fsync", mock_fsync)
        # Also patch via the engine module's `os` binding explicitly for
        # robustness against any future `from os import fsync`-style import.
        monkeypatch.setattr(pptx_io.os, "fsync", mock_fsync, raising=True)
        return calls

    def test_fsync_true_calls_fsync(self, tmp_path, monkeypatch):
        """fsync=True must call os.fsync at least once (temp file). On
        non-Windows, also calls it for the containing directory (>=2)."""
        calls = self._install_fsync_counter(monkeypatch)

        prs = Presentation()
        path = str(tmp_path / "durable.pptx")
        save_pptx(prs, path, fsync=True)

        assert os.path.exists(path)
        assert calls["count"] >= 1, "fsync=True must fsync the temp file"
        if sys.platform != "win32":
            assert calls["count"] >= 2, (
                "fsync=True on POSIX must fsync both temp file and directory"
            )

    def test_fsync_false_does_not_call_fsync(self, tmp_path, monkeypatch):
        """Default (fsync=False) must not invoke os.fsync."""
        calls = self._install_fsync_counter(monkeypatch)

        prs = Presentation()
        path = str(tmp_path / "fast.pptx")
        save_pptx(prs, path, fsync=False)

        assert os.path.exists(path)
        assert calls["count"] == 0, (
            "default save must not impose an fsync barrier"
        )

    def test_fsync_default_is_false(self, tmp_path, monkeypatch):
        """Omitting the fsync kwarg must behave like fsync=False."""
        calls = self._install_fsync_counter(monkeypatch)

        prs = Presentation()
        path = str(tmp_path / "default.pptx")
        save_pptx(prs, path)  # no kwarg

        assert os.path.exists(path)
        assert calls["count"] == 0

    def test_fsync_true_still_atomic_on_crash(self, tmp_path, monkeypatch):
        """fsync=True must preserve the same atomicity guarantee: if
        Presentation.save raises, the original file is untouched and no
        temp leaks."""
        path = tmp_path / "existing.pptx"
        original_bytes = b"ORIGINAL CONTENT -- not a real pptx"
        path.write_bytes(original_bytes)

        prs = Presentation()

        def boom(*_args, **_kwargs):
            raise IOError("simulated mid-write failure")

        monkeypatch.setattr(type(prs), "save", boom, raising=True)

        with pytest.raises(IOError, match="simulated"):
            save_pptx(prs, str(path), fsync=True)

        # Original bytes preserved, no temp leak.
        assert path.read_bytes() == original_bytes
        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []

    def test_fsync_true_round_trips(self, tmp_path):
        """fsync=True produces a valid, loadable PPTX (no regression)."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6]
        prs.slides.add_slide(layout)

        path = str(tmp_path / "durable_rt.pptx")
        save_pptx(prs, path, fsync=True)

        loaded = Presentation(path)
        assert len(loaded.slides) == 1
        leftovers = [p.name for p in tmp_path.iterdir() if ".tmp." in p.name]
        assert leftovers == []

    def test_fsync_true_skips_dir_fsync_on_windows(self, tmp_path, monkeypatch):
        """On Windows, directory fsync is unsupported and must be skipped:
        only the temp-file fsync runs (1 call total)."""
        calls = self._install_fsync_counter(monkeypatch)
        monkeypatch.setattr(sys, "platform", "win32")
        # pptx_io binds `import sys`; patch that reference too so the
        # module's `sys.platform` check sees "win32".
        monkeypatch.setattr(pptx_io.sys, "platform", "win32", raising=True)

        prs = Presentation()
        path = str(tmp_path / "win.pptx")
        save_pptx(prs, path, fsync=True)

        assert os.path.exists(path)
        assert calls["count"] == 1, (
            "on Windows, only the temp file is fsynced; directory fsync skipped"
        )
