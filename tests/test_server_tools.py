"""``*_file`` wrappers / MCP-adjacent tools の round-trip テスト (#84).

各 file wrapper は「書く → 閉じる → 開き直す → shape を index で拾える」
という契約を満たす必要がある。以前は ``add_auto_fit_textbox_file`` が
``save_pptx`` 後に identity 検索を行っており、``os.replace`` の原子的
rename を挟むことで毎回 ``shape_index=-1`` を返していた (#84)。

本モジュールは ``*_file`` wrappers の専用カバレッジ面として、
- 返却された ``shape_index`` が >= 0 であること
- 保存したファイルを開き直してその index で目的の shape が取れること
- ``shape_index != -1`` を明示する回帰テスト
を提供する。
"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_mcp_server.engine.flex import add_flex_container_file
from pptx_mcp_server.engine.pptx_io import save_pptx
from pptx_mcp_server.engine.shapes import add_auto_fit_textbox_file


@pytest.fixture
def blank_pptx(tmp_path):
    """16:9 blank deck を 1 スライドで作成して path を返す."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    path = str(tmp_path / "deck.pptx")
    save_pptx(prs, path)
    return path


# =============================================================================
# add_auto_fit_textbox_file
# =============================================================================


class TestAddAutoFitTextboxFile:
    """#84 の直接対象: ``add_auto_fit_textbox_file`` が正しい index を返すこと."""

    def test_returns_non_negative_shape_index(self, blank_pptx):
        """shape_index >= 0 (回帰: 以前は常に -1)."""
        result = add_auto_fit_textbox_file(
            blank_pptx,
            slide_index=0,
            text="Hello",
            left=1.0, top=1.0, width=4.0, height=1.0,
        )
        assert result["shape_index"] >= 0

    def test_shape_index_not_minus_one_regression(self, blank_pptx):
        """明示的な回帰テスト: ``shape_index != -1``."""
        result = add_auto_fit_textbox_file(
            blank_pptx,
            slide_index=0,
            text="regression",
            left=1.0, top=1.0, width=4.0, height=1.0,
        )
        assert result["shape_index"] != -1

    def test_round_trip_shape_at_returned_index(self, blank_pptx):
        """保存した deck を開き直し、返却 index で shape が取れる."""
        expected_text = "round-trip check"
        result = add_auto_fit_textbox_file(
            blank_pptx,
            slide_index=0,
            text=expected_text,
            left=1.0, top=1.0, width=4.0, height=1.0,
        )
        idx = result["shape_index"]

        prs = Presentation(blank_pptx)
        slide = prs.slides[0]
        shape = slide.shapes[idx]
        assert shape.has_text_frame
        assert expected_text in shape.text_frame.text

    def test_sequential_calls_return_distinct_indices(self, blank_pptx):
        """連続呼び出しで index が重複しない."""
        r1 = add_auto_fit_textbox_file(
            blank_pptx, slide_index=0, text="first",
            left=0.5, top=0.5, width=3.0, height=0.8,
        )
        r2 = add_auto_fit_textbox_file(
            blank_pptx, slide_index=0, text="second",
            left=0.5, top=2.0, width=3.0, height=0.8,
        )
        r3 = add_auto_fit_textbox_file(
            blank_pptx, slide_index=0, text="third",
            left=0.5, top=3.5, width=3.0, height=0.8,
        )
        assert r1["shape_index"] != r2["shape_index"] != r3["shape_index"]
        assert r2["shape_index"] == r1["shape_index"] + 1
        assert r3["shape_index"] == r2["shape_index"] + 1

        # 各 index で目的の shape が解決できる
        prs = Presentation(blank_pptx)
        slide = prs.slides[0]
        assert "first" in slide.shapes[r1["shape_index"]].text_frame.text
        assert "second" in slide.shapes[r2["shape_index"]].text_frame.text
        assert "third" in slide.shapes[r3["shape_index"]].text_frame.text

    def test_returns_actual_font_size_float(self, blank_pptx):
        """``actual_font_size`` も返却される (回帰対象外だが契約としてチェック)."""
        result = add_auto_fit_textbox_file(
            blank_pptx,
            slide_index=0,
            text="x",
            left=1.0, top=1.0, width=2.0, height=0.5,
            font_size_pt=11, min_size_pt=7,
        )
        assert isinstance(result["actual_font_size"], float)
        assert 7.0 <= result["actual_font_size"] <= 11.0


# =============================================================================
# add_flex_container_file
# =============================================================================


class TestAddFlexContainerFile:
    """``add_flex_container_file`` の round-trip (text 子要素は #84 同類バグが
    隣接パスにあった — 同じく追加直前 index 方式に修正済み)."""

    def test_text_child_shape_index_non_negative(self, blank_pptx):
        """text 子要素の shape_index が >= 0."""
        items = [
            {"type": "text", "text": "A", "sizing": "grow"},
            {"type": "text", "text": "B", "sizing": "grow"},
        ]
        result = add_flex_container_file(
            blank_pptx, 0, items,
            left=0.5, top=0.5, width=10.0, height=1.0,
        )
        shapes = result["shapes"]
        assert len(shapes) == 2
        for s in shapes:
            assert s["shape_index"] >= 0
            assert s["shape_index"] != -1

    def test_text_and_rectangle_round_trip_resolves(self, blank_pptx):
        """mix (text + rectangle) の index が開き直した deck で正しく解決する."""
        items = [
            {"type": "text", "text": "hello", "sizing": "grow"},
            {
                "type": "rectangle",
                "sizing": "fixed",
                "size": 1.0,
                "fill_color": "2251FF",
                "no_line": True,
            },
        ]
        result = add_flex_container_file(
            blank_pptx, 0, items,
            left=0.5, top=0.5, width=10.0, height=1.0,
        )

        prs = Presentation(blank_pptx)
        slide = prs.slides[0]

        shapes = result["shapes"]
        assert shapes[0]["type"] == "text"
        assert shapes[1]["type"] == "rectangle"

        text_shape = slide.shapes[shapes[0]["shape_index"]]
        assert text_shape.has_text_frame
        assert "hello" in text_shape.text_frame.text

        rect_shape = slide.shapes[shapes[1]["shape_index"]]
        # rectangle は has_text_frame でも True になる場合があるので存在だけ確認
        assert rect_shape is not None

    def test_allocations_count_matches_items(self, blank_pptx):
        """allocations の長さが入力 items と一致する."""
        items = [
            {"type": "text", "text": f"item-{i}", "sizing": "grow"}
            for i in range(3)
        ]
        result = add_flex_container_file(
            blank_pptx, 0, items,
            left=0.5, top=0.5, width=10.0, height=1.0,
        )
        assert len(result["allocations"]) == 3
        assert len(result["shapes"]) == 3


# =============================================================================
# responsive_card_row (MCP tool path — placements on reopen)
# =============================================================================


class TestResponsiveCardRowFileRoundTrip:
    """``add_responsive_card_row`` は engine に直接の ``_file`` wrapper を持たず、
    server.py 側の MCP ツールでファイル保存を行う。ここでは engine を直接
    呼びつつ ``save_pptx`` で書き出し、開き直した deck の shape 数が
    placements 数と整合することを確認する。"""

    def test_placements_persist_on_reopen(self, blank_pptx):
        from pptx_mcp_server.engine.cards import CardSpec, add_responsive_card_row
        from pptx_mcp_server.engine.pptx_io import open_pptx, save_pptx

        prs = open_pptx(blank_pptx)
        slide = prs.slides[0]
        shapes_before = len(slide.shapes)

        cards = [
            CardSpec(title="First", body="body1"),
            CardSpec(title="Second", body="body2"),
            CardSpec(title="Third", body="body3"),
        ]
        placements, consumed = add_responsive_card_row(
            slide, cards,
            left=0.5, top=0.5, width=10.0, max_height=2.0,
        )
        save_pptx(prs, blank_pptx)

        assert len(placements) == 3
        assert consumed > 0.0

        # 開き直して placements の bounding box が実 shape と重なっていること
        reopened = Presentation(blank_pptx)
        reopened_slide = reopened.slides[0]
        assert len(reopened_slide.shapes) > shapes_before

        # 各 placement の bbox に 1 個以上の shape が存在する
        for p in placements:
            inside = 0
            for s in reopened_slide.shapes:
                sx = s.left / 914400
                sy = s.top / 914400
                if (
                    p.left - 0.01 <= sx <= p.left + p.width + 0.01
                    and p.top - 0.01 <= sy <= p.top + p.height + 0.01
                ):
                    inside += 1
            assert inside >= 1, f"placement {p} has no shapes inside"


# =============================================================================
# pptx_add_section_header envelope invariant (#120 / #136 regression)
# =============================================================================


class TestAddSectionHeaderMcpEnvelope:
    """``pptx_add_section_header`` must include a ``message`` key in its
    result dict, matching the v0.3.0 envelope invariant (#120)."""

    def test_result_has_message_key(self, blank_pptx):
        """Happy-path MCP call returns ``result["message"]`` (not silently dropped)."""
        from pptx_mcp_server.server import pptx_add_section_header

        out = pptx_add_section_header(
            file_path=blank_pptx,
            slide_index=0,
            title="Section One",
            subtitle="A subtitle line",
        )
        parsed = json.loads(out)
        assert parsed.get("ok") is True
        assert "message" in parsed["result"], (
            "pptx_add_section_header must set result['message'] per envelope invariant (#120)."
        )
        assert isinstance(parsed["result"]["message"], str)
        assert parsed["result"]["message"]  # non-empty
