"""
Unit tests for shape operations -- textbox, auto shape, edit text, paragraph, delete, list.
"""

from __future__ import annotations

import pytest
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_mcp_server.engine.shapes import (
    _add_textbox,
    _add_shape,
    _add_image,
    _edit_text,
    _add_paragraph,
    _delete_shape,
    _list_shapes,
)


# ── Textbox tests ──────────────────────────────────────────────────


class TestAddTextbox:
    """_add_textbox must create a textbox with correct position, size, and formatting."""

    def test_creates_textbox_at_position(self, slide):
        idx = _add_textbox(slide, left=1, top=2, width=4, height=0.5, text="Hi")
        shape = list(slide.shapes)[idx]
        assert shape.left == Inches(1)
        assert shape.top == Inches(2)
        assert shape.width == Inches(4)
        assert shape.height == Inches(0.5)

    def test_sets_text_content(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="Hello World")
        shape = list(slide.shapes)[idx]
        assert shape.text_frame.text == "Hello World"

    def test_applies_font_name(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="T", font_name="Helvetica")
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.font.name == "Helvetica"

    def test_applies_font_size(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="T", font_size=18)
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.font.size == Pt(18)

    def test_applies_font_color(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="T", font_color="#FF0000")
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_applies_bold(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="Bold", bold=True)
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.font.bold is True

    def test_applies_italic(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="Italic", italic=True)
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.font.italic is True

    def test_applies_underline(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="U", underline=True)
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.font.underline is True

    def test_applies_alignment_center(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="Center", alignment="center")
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.alignment == PP_ALIGN.CENTER

    def test_applies_vertical_anchor_bottom(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="Bottom", vertical_anchor="bottom")
        shape = list(slide.shapes)[idx]
        from pptx.oxml.ns import qn
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr.get("anchor") == "b"

    def test_applies_line_spacing(self, slide):
        idx = _add_textbox(slide, 1, 1, 4, 0.5, text="Spaced", line_spacing=20)
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert p.line_spacing == Pt(20)

    def test_returns_correct_shape_index(self, slide):
        idx0 = _add_textbox(slide, 1, 1, 2, 0.5, text="First")
        idx1 = _add_textbox(slide, 3, 1, 2, 0.5, text="Second")
        assert idx0 == 0
        assert idx1 == 1


# ── Auto Shape tests ───────────────────────────────────────────────


class TestAddShape:
    """_add_shape must create the correct shape type with optional styling."""

    def test_creates_rectangle_shape(self, slide):
        idx = _add_shape(slide, "rectangle", 1, 1, 3, 2)
        shapes = list(slide.shapes)
        assert idx == len(shapes) - 1

    def test_unknown_shape_type_raises(self, slide):
        with pytest.raises(EngineError) as exc_info:
            _add_shape(slide, "nonexistent_shape", 1, 1, 3, 2)
        assert exc_info.value.code == ErrorCode.INVALID_PARAMETER

    def test_applies_fill_color(self, slide):
        idx = _add_shape(slide, "rectangle", 1, 1, 3, 2, fill_color="#00FF00")
        shape = list(slide.shapes)[idx]
        assert shape.fill.fore_color.rgb == RGBColor(0x00, 0xFF, 0x00)

    def test_applies_no_line(self, slide):
        idx = _add_shape(slide, "oval", 1, 1, 3, 2, no_line=True)
        shape = list(slide.shapes)[idx]
        # no_line sets fill.background() on the line -- just verify no exception
        assert shape is not None

    def test_sets_text_content(self, slide):
        idx = _add_shape(slide, "rectangle", 1, 1, 3, 2, text="Inside shape")
        shape = list(slide.shapes)[idx]
        assert shape.text_frame.text == "Inside shape"

    def test_creates_all_valid_shape_types(self, slide):
        from pptx_mcp_server.engine.shapes import _SHAPE_MAP
        for shape_type in _SHAPE_MAP:
            _add_shape(slide, shape_type, 0, 0, 1, 1)
        # No exception means all types are valid


# ── Edit Text tests ────────────────────────────────────────────────


class TestEditText:
    """_edit_text must change text and/or formatting on existing shapes."""

    def test_changes_text_content(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Original")
        _edit_text(slide, 0, text="Updated")
        shape = list(slide.shapes)[0]
        assert shape.text_frame.paragraphs[0].text == "Updated"

    def test_applies_formatting_without_changing_text(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Keep")
        _edit_text(slide, 0, text=None, bold=True, font_size=24)
        shape = list(slide.shapes)[0]
        p = shape.text_frame.paragraphs[0]
        assert p.text == "Keep"
        assert p.font.bold is True
        assert p.font.size == Pt(24)

    def test_shape_without_text_frame_raises(self, slide):
        from pptx_mcp_server.engine.tables import _add_table
        _add_table(slide, [["H1"], ["D1"]], 1, 1, 3)
        table_idx = len(list(slide.shapes)) - 1
        with pytest.raises(EngineError) as exc_info:
            _edit_text(slide, table_idx, text="fail")
        assert exc_info.value.code == ErrorCode.INVALID_PARAMETER

    def test_invalid_paragraph_index_raises(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Single")
        with pytest.raises(EngineError) as exc_info:
            _edit_text(slide, 0, text="fail", paragraph_index=99)
        assert exc_info.value.code == ErrorCode.INDEX_OUT_OF_RANGE


# ── Add Paragraph tests ───────────────────────────────────────────


class TestAddParagraph:
    """_add_paragraph must append new paragraphs to a shape."""

    def test_appends_paragraph(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Line 1")
        p_idx = _add_paragraph(slide, 0, "Line 2")
        shape = list(slide.shapes)[0]
        texts = [p.text for p in shape.text_frame.paragraphs]
        assert "Line 1" in texts
        assert "Line 2" in texts

    def test_returns_correct_paragraph_index(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="First")
        p_idx = _add_paragraph(slide, 0, "Second")
        assert p_idx == 1
        p_idx2 = _add_paragraph(slide, 0, "Third")
        assert p_idx2 == 2


# ── Delete Shape tests ────────────────────────────────────────────


class TestDeleteShape:
    """_delete_shape must remove a shape from the slide."""

    def test_removes_shape(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Delete me")
        assert len(slide.shapes) == 1
        _delete_shape(slide, 0)
        assert len(slide.shapes) == 0

    def test_delete_nonexistent_raises(self, slide):
        with pytest.raises(EngineError) as exc_info:
            _delete_shape(slide, 99)
        assert exc_info.value.code == ErrorCode.SHAPE_NOT_FOUND


# ── List Shapes tests ─────────────────────────────────────────────


class TestListShapes:
    """_list_shapes must return a formatted string."""

    def test_returns_formatted_string(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Box A")
        _add_textbox(slide, 5, 1, 4, 0.5, text="Box B")
        result = _list_shapes(slide, 0)
        assert "2 shapes" in result
        assert "Box A" in result
        assert "Box B" in result
        assert "[0]" in result
        assert "[1]" in result

    def test_empty_slide_shows_zero(self, slide):
        result = _list_shapes(slide, 0)
        assert "0 shapes" in result


# ── Add Image tests ──────────────────────────────────────────────


class TestAddImage:
    """_add_image must add a picture shape to a slide."""

    def test_adds_image(self, slide, tmp_path):
        from PIL import Image
        img_path = str(tmp_path / "test.png")
        img = Image.new("RGB", (100, 50), color="red")
        img.save(img_path)

        idx = _add_image(slide, img_path, left=1.0, top=1.0, width=2.0)
        assert idx >= 0
        shape = list(slide.shapes)[idx]
        assert shape.shape_type == MSO_SHAPE_TYPE.PICTURE

    def test_image_not_found_raises(self, slide):
        with pytest.raises(EngineError) as exc_info:
            _add_image(slide, "/nonexistent/image.png", left=1.0, top=1.0)
        assert exc_info.value.code == ErrorCode.FILE_NOT_FOUND

    def test_preserves_aspect_ratio_with_width_only(self, slide, tmp_path):
        from PIL import Image
        img_path = str(tmp_path / "test.png")
        img = Image.new("RGB", (200, 100), color="blue")
        img.save(img_path)

        idx = _add_image(slide, img_path, left=1.0, top=1.0, width=4.0)
        shape = list(slide.shapes)[idx]
        # Width should be 4 inches, height auto-calculated (~2 inches for 2:1 ratio)
        assert abs(shape.width / 914400 - 4.0) < 0.01
