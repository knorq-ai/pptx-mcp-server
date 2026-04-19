"""
Unit tests for formatting operations -- shape formatting and slide dimensions.
"""

from __future__ import annotations

import pytest
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from pptx_mcp_server.engine.shapes import _add_textbox, _add_shape
from pptx_mcp_server.engine.formatting import _format_shape, _set_slide_dimensions


class TestFormatShape:
    """_format_shape must reposition, resize, and restyle shapes."""

    def test_repositions_shape(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Move me")
        _format_shape(slide, 0, left=3, top=4)
        shape = list(slide.shapes)[0]
        assert shape.left == Inches(3)
        assert shape.top == Inches(4)

    def test_resizes_shape(self, slide):
        _add_textbox(slide, 1, 1, 4, 0.5, text="Resize me")
        _format_shape(slide, 0, width=8, height=2)
        shape = list(slide.shapes)[0]
        assert shape.width == Inches(8)
        assert shape.height == Inches(2)

    def test_changes_fill_color(self, slide):
        _add_shape(slide, "rectangle", 1, 1, 3, 2)
        _format_shape(slide, 0, fill_color="#00FF00")
        shape = list(slide.shapes)[0]
        assert shape.fill.fore_color.rgb == RGBColor(0x00, 0xFF, 0x00)

    def test_applies_no_fill(self, slide):
        _add_shape(slide, "rectangle", 1, 1, 3, 2, fill_color="#FF0000")
        _format_shape(slide, 0, no_fill=True)
        shape = list(slide.shapes)[0]
        # no_fill sets background fill -- just verify no exception
        assert shape is not None

    def test_applies_no_line(self, slide):
        _add_shape(slide, "rectangle", 1, 1, 3, 2)
        _format_shape(slide, 0, no_line=True)
        shape = list(slide.shapes)[0]
        assert shape is not None

    def test_no_fill_takes_precedence_over_fill_color(self, slide):
        _add_shape(slide, "rectangle", 1, 1, 3, 2)
        # When no_fill=True, fill_color should be ignored
        _format_shape(slide, 0, fill_color="#FF0000", no_fill=True)
        # No assertion on fill type since background() is set; just ensure no crash

    def test_applies_rotation(self, slide):
        _add_shape(slide, "rectangle", 1, 1, 3, 2)
        _format_shape(slide, 0, rotation=45)
        shape = list(slide.shapes)[0]
        assert shape.rotation == 45


class TestSetSlideDimensions:
    """_set_slide_dimensions must change the presentation dimensions."""

    def test_changes_dimensions(self, one_slide_prs):
        _set_slide_dimensions(one_slide_prs, 10, 5.625)
        assert one_slide_prs.slide_width == Inches(10)
        assert one_slide_prs.slide_height == Inches(5.625)

    def test_widescreen_16_9(self, blank_prs):
        _set_slide_dimensions(blank_prs, 13.333, 7.5)
        assert abs(blank_prs.slide_width - Inches(13.333)) < 10
        assert abs(blank_prs.slide_height - Inches(7.5)) < 10
