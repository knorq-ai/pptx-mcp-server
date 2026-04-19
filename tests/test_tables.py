"""
Unit tests for table operations -- add, edit cell, batch edit, format.
"""

from __future__ import annotations

import pytest
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx_mcp_server.engine.tables import (
    _add_table,
    _edit_table_cell,
    _edit_table_cells,
    _format_table,
)
from pptx_mcp_server.theme import MCKINSEY


class TestAddTable:
    """_add_table must create tables with correct structure and formatting."""

    def test_creates_table_with_correct_dimensions(self, slide):
        rows = [["H1", "H2", "H3"], ["a", "b", "c"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5)
        shape = list(slide.shapes)[idx]
        table = shape.table
        assert len(table.rows) == 2
        assert len(table.columns) == 3

    def test_populates_cell_text(self, slide):
        rows = [["Name", "Value"], ["foo", "42"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5)
        shape = list(slide.shapes)[idx]
        table = shape.table
        assert table.cell(0, 0).text_frame.paragraphs[0].text == "Name"
        assert table.cell(1, 1).text_frame.paragraphs[0].text == "42"

    def test_header_has_bold_font(self, slide):
        rows = [["Header"], ["Data"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5)
        shape = list(slide.shapes)[idx]
        header_p = shape.table.cell(0, 0).text_frame.paragraphs[0]
        assert header_p.font.bold is True

    def test_header_bg_color(self, slide):
        rows = [["H"], ["D"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5, header_bg="051C2C")
        shape = list(slide.shapes)[idx]
        cell = shape.table.cell(0, 0)
        assert cell.fill.fore_color.rgb == RGBColor(0x05, 0x1C, 0x2C)

    def test_header_font_color(self, slide):
        rows = [["H"], ["D"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5, header_fg="FFFFFF")
        shape = list(slide.shapes)[idx]
        p = shape.table.cell(0, 0).text_frame.paragraphs[0]
        assert p.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)

    def test_alternating_row_shading(self, slide):
        rows = [["H"], ["R1"], ["R2"], ["R3"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5, alt_row_bg="F5F5F5")
        shape = list(slide.shapes)[idx]
        # Row 2 (index 2) is even and > 0, so should have alt shading
        cell_r2 = shape.table.cell(2, 0)
        assert cell_r2.fill.fore_color.rgb == RGBColor(0xF5, 0xF5, 0xF5)

    def test_empty_rows_raises_table_error(self, slide):
        with pytest.raises(EngineError) as exc_info:
            _add_table(slide, [], left=1, top=1, width=5)
        assert exc_info.value.code == ErrorCode.TABLE_ERROR

    def test_respects_col_widths(self, slide):
        rows = [["A", "B"], ["1", "2"]]
        idx = _add_table(slide, rows, left=1, top=1, width=10, col_widths=[0.3, 0.7])
        shape = list(slide.shapes)[idx]
        table = shape.table
        c0w = table.columns[0].width
        c1w = table.columns[1].width
        # col 0 should be ~30% of total, col 1 ~70%
        assert c0w < c1w

    def test_mckinsey_borders_no_vertical(self, slide):
        rows = [["H1", "H2"], ["d1", "d2"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5, no_vertical_borders=True)
        shape = list(slide.shapes)[idx]
        # Check that vertical borders have no fill (noFill element present)
        cell = shape.table.cell(1, 0)
        tc = cell._tc
        tcPr = tc.find(qn("a:tcPr"))
        lnL = tcPr.find(qn("a:lnL"))
        assert lnL is not None
        # Should have noFill or width=0
        nf = lnL.find(qn("a:noFill"))
        assert nf is not None or lnL.get("w") == "0"

    def test_mckinsey_borders_thin_horizontal(self, slide):
        rows = [["H"], ["D1"], ["D2"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5)
        shape = list(slide.shapes)[idx]
        # Data row bottom border should have solidFill
        cell = shape.table.cell(1, 0)
        tc = cell._tc
        tcPr = tc.find(qn("a:tcPr"))
        lnB = tcPr.find(qn("a:lnB"))
        assert lnB is not None
        sf = lnB.find(qn("a:solidFill"))
        assert sf is not None

    def test_with_theme_uses_theme_defaults(self, slide, mckinsey_theme):
        rows = [["H"], ["D"]]
        idx = _add_table(slide, rows, left=1, top=1, width=5, theme=mckinsey_theme)
        shape = list(slide.shapes)[idx]
        # Just verify table was created without error
        assert shape.has_table


class TestEditTableCell:
    """_edit_table_cell must modify individual cells."""

    def _make_table(self, slide):
        rows = [["H1", "H2"], ["a", "b"], ["c", "d"]]
        return _add_table(slide, rows, left=1, top=1, width=5)

    def test_changes_text(self, slide):
        idx = self._make_table(slide)
        _edit_table_cell(slide, idx, row=1, col=0, text="updated")
        table = list(slide.shapes)[idx].table
        assert table.cell(1, 0).text_frame.paragraphs[0].text == "updated"

    def test_applies_bold(self, slide):
        idx = self._make_table(slide)
        _edit_table_cell(slide, idx, row=1, col=0, bold=True)
        p = list(slide.shapes)[idx].table.cell(1, 0).text_frame.paragraphs[0]
        assert p.font.bold is True

    def test_applies_bg_color(self, slide):
        idx = self._make_table(slide)
        _edit_table_cell(slide, idx, row=1, col=0, bg_color="#FF0000")
        cell = list(slide.shapes)[idx].table.cell(1, 0)
        assert cell.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_invalid_row_raises(self, slide):
        idx = self._make_table(slide)
        with pytest.raises(EngineError) as exc_info:
            _edit_table_cell(slide, idx, row=99, col=0, text="fail")
        assert exc_info.value.code == ErrorCode.INDEX_OUT_OF_RANGE

    def test_invalid_col_raises(self, slide):
        idx = self._make_table(slide)
        with pytest.raises(EngineError) as exc_info:
            _edit_table_cell(slide, idx, row=0, col=99, text="fail")
        assert exc_info.value.code == ErrorCode.INDEX_OUT_OF_RANGE


class TestEditTableCells:
    """_edit_table_cells must batch-edit multiple cells."""

    def test_batch_edits(self, slide):
        idx = _add_table(slide, [["A", "B"], ["1", "2"]], left=1, top=1, width=5)
        edits = [
            {"row": 1, "col": 0, "text": "X"},
            {"row": 1, "col": 1, "text": "Y"},
        ]
        count = _edit_table_cells(slide, idx, edits)
        assert count == 2
        table = list(slide.shapes)[idx].table
        assert table.cell(1, 0).text_frame.paragraphs[0].text == "X"
        assert table.cell(1, 1).text_frame.paragraphs[0].text == "Y"

    def test_skips_out_of_range(self, slide):
        idx = _add_table(slide, [["A"], ["1"]], left=1, top=1, width=5)
        edits = [
            {"row": 99, "col": 0, "text": "nope"},
            {"row": 0, "col": 0, "text": "yes"},
        ]
        count = _edit_table_cells(slide, idx, edits)
        assert count == 1  # only the valid one counted


class TestFormatTable:
    """_format_table must apply bulk formatting to an entire table."""

    def test_applies_font_name(self, slide):
        idx = _add_table(slide, [["H"], ["D"]], left=1, top=1, width=5)
        _format_table(slide, idx, font_name="Courier")
        table = list(slide.shapes)[idx].table
        p = table.cell(0, 0).text_frame.paragraphs[0]
        assert p.font.name == "Courier"

    def test_applies_header_bg(self, slide):
        idx = _add_table(slide, [["H"], ["D"]], left=1, top=1, width=5)
        _format_table(slide, idx, header_bg="#AA0000")
        cell = list(slide.shapes)[idx].table.cell(0, 0)
        assert cell.fill.fore_color.rgb == RGBColor(0xAA, 0x00, 0x00)

    def test_not_a_table_raises(self, slide):
        from pptx_mcp_server.engine.shapes import _add_textbox
        _add_textbox(slide, 1, 1, 4, 0.5, text="Not a table")
        with pytest.raises(EngineError) as exc_info:
            _format_table(slide, 0)
        assert exc_info.value.code == ErrorCode.TABLE_ERROR
