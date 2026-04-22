"""#131: theme-aware atomic primitives のテスト.

v0.5.0 で ``add_responsive_card_row`` / ``add_data_table`` /
``add_milestone_timeline`` に入った theme string 契約を、v0.6.0 では
atomic primitives (``add_auto_fit_textbox`` + ``_add_shape`` + ``add_shape``
file wrapper) にも展開した。本モジュールは以下を検証する:

- theme token (e.g. ``"primary"``) + theme name (e.g. ``"ir"``) が theme
  registry の hex に解決される。
- raw hex (``"#FF0000"`` / ``"FF0000"``) は passthrough する (``#`` 剥離のみ)。
- ``theme=None`` なら v0.5.1 挙動が byte-for-byte 保たれる (resolver の
  passthrough no-op)。
- MCP tool ``pptx_add_auto_fit_textbox`` / ``pptx_add_shape`` が ``theme``
  kwarg を surface する。
"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from pptx_mcp_server.engine.pptx_io import open_pptx, save_pptx
from pptx_mcp_server.engine.shapes import (
    _add_shape,
    add_auto_fit_textbox,
)
from pptx_mcp_server.server import (
    pptx_add_auto_fit_textbox,
    pptx_add_shape,
)
from pptx_mcp_server.theme import IR


# ---------------------------------------------------------------------------
# add_auto_fit_textbox — theme kwarg
# ---------------------------------------------------------------------------


class TestAddAutoFitTextboxTheme:
    """``add_auto_fit_textbox(color_hex=..., theme=...)`` が color を
    ``resolve_theme_color`` 経由で解決することを確認する (#131)."""

    def test_theme_token_resolves_to_ir_primary_hex(self, slide):
        """``color_hex="primary", theme="ir"`` → IR primary (#0A2540)."""
        shape, _ = add_auto_fit_textbox(
            slide,
            "text",
            left=1.0, top=1.0, width=3.0, height=0.5,
            color_hex="primary",
            theme="ir",
        )
        run = shape.text_frame.paragraphs[0].runs[0]
        expected = IR.colors["primary"].lstrip("#")
        assert run.font.color.rgb == RGBColor.from_string(expected)

    def test_raw_hash_hex_passes_through_with_theme(self, slide):
        """``color_hex="#FF0000", theme="ir"`` → raw hex 経路 (``#`` 剥離)."""
        shape, _ = add_auto_fit_textbox(
            slide,
            "text",
            left=1.0, top=1.0, width=3.0, height=0.5,
            color_hex="#FF0000",
            theme="ir",
        )
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_theme_none_preserves_v0_5_1_behavior(self, slide):
        """``theme=None`` なら resolver は passthrough (raw hex のみ有効)."""
        shape, _ = add_auto_fit_textbox(
            slide,
            "text",
            left=1.0, top=1.0, width=3.0, height=0.5,
            color_hex="333333",
            # theme 省略 (None default)
        )
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(0x33, 0x33, 0x33)

    def test_unknown_token_without_theme_passes_as_noop(self, slide):
        """``theme=None`` + 未登録 token は raw string として下流に流れ、
        ``_parse_color`` 側で ``INVALID_PARAMETER`` が発生する (defensive
        挙動ではなく既存の失敗経路を踏襲する)."""
        from pptx_mcp_server.engine.pptx_io import EngineError
        with pytest.raises(EngineError):
            add_auto_fit_textbox(
                slide,
                "text",
                left=1.0, top=1.0, width=3.0, height=0.5,
                color_hex="primary",
                # theme=None → "primary" は hex として扱われ _parse_color で失敗
            )


# ---------------------------------------------------------------------------
# _add_shape — theme kwarg (fill_color)
# ---------------------------------------------------------------------------


class TestAddShapeTheme:
    """``_add_shape(fill_color=..., theme=...)`` が色を
    ``resolve_theme_color`` 経由で解決することを確認する (#131)."""

    def test_theme_token_resolves_fill_color(self, slide):
        """``fill_color="primary", theme="ir"`` → IR primary hex."""
        idx = _add_shape(
            slide, "rectangle",
            left=1.0, top=1.0, width=2.0, height=1.0,
            fill_color="primary",
            theme="ir",
        )
        shape = list(slide.shapes)[idx]
        expected = IR.colors["primary"].lstrip("#")
        assert shape.fill.fore_color.rgb == RGBColor.from_string(expected)

    def test_raw_hex_passes_through_with_theme(self, slide):
        """``fill_color="#FF0000", theme="ir"`` → raw hex 経路."""
        idx = _add_shape(
            slide, "rectangle",
            left=1.0, top=1.0, width=2.0, height=1.0,
            fill_color="#FF0000",
            theme="ir",
        )
        shape = list(slide.shapes)[idx]
        assert shape.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_theme_none_preserves_v0_5_1_behavior(self, slide):
        """``theme=None`` なら raw hex はそのまま、token は失敗する (既存動作)."""
        idx = _add_shape(
            slide, "rectangle",
            left=1.0, top=1.0, width=2.0, height=1.0,
            fill_color="2251FF",
            # theme=None
        )
        shape = list(slide.shapes)[idx]
        assert shape.fill.fore_color.rgb == RGBColor(0x22, 0x51, 0xFF)

    def test_theme_resolves_line_color(self, slide):
        """``line_color`` も同様に resolve される."""
        idx = _add_shape(
            slide, "rectangle",
            left=1.0, top=1.0, width=2.0, height=1.0,
            fill_color="#FFFFFF",
            line_color="accent",
            line_width=1.0,
            theme="ir",
        )
        shape = list(slide.shapes)[idx]
        expected = IR.colors["accent"].lstrip("#")
        assert shape.line.color.rgb == RGBColor.from_string(expected)


# ---------------------------------------------------------------------------
# MCP tool surface — pptx_add_auto_fit_textbox / pptx_add_shape expose theme
# ---------------------------------------------------------------------------


@pytest.fixture
def blank_pptx(tmp_path):
    """16:9 blank deck を 1 slide で作成して path を返す."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    path = str(tmp_path / "deck.pptx")
    save_pptx(prs, path)
    return path


class TestMcpToolThemeSurface:
    """MCP tool wrappers が ``theme`` kwarg を surface し、engine 側に
    透過することを確認する (#131)."""

    def test_pptx_add_auto_fit_textbox_exposes_theme_kwarg(self, blank_pptx):
        """``pptx_add_auto_fit_textbox(theme="ir", color_hex="primary")`` が
        保存後に IR primary 色を反映する."""
        resp = pptx_add_auto_fit_textbox(
            blank_pptx, 0, "title",
            left=1.0, top=1.0, width=3.0, height=0.5,
            color_hex="primary",
            theme="ir",
        )
        payload = json.loads(resp)
        assert payload["ok"] is True, payload

        prs = open_pptx(blank_pptx)
        slide = prs.slides[0]
        shape = list(slide.shapes)[payload["result"]["shape_index"]]
        run = shape.text_frame.paragraphs[0].runs[0]
        expected = IR.colors["primary"].lstrip("#")
        assert run.font.color.rgb == RGBColor.from_string(expected)

    def test_pptx_add_shape_exposes_theme_kwarg(self, blank_pptx):
        """``pptx_add_shape(theme="ir", fill_color="accent")`` が保存後に
        IR accent 色を反映する."""
        resp = pptx_add_shape(
            blank_pptx, 0, "rectangle",
            1.0, 1.0, 2.0, 1.0,
            fill_color="accent",
            theme="ir",
        )
        payload = json.loads(resp)
        assert payload["ok"] is True, payload

        prs = open_pptx(blank_pptx)
        slide = prs.slides[0]
        # 直近追加 shape は末尾にある
        shape = list(slide.shapes)[-1]
        expected = IR.colors["accent"].lstrip("#")
        assert shape.fill.fore_color.rgb == RGBColor.from_string(expected)
