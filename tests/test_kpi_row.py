"""Tests for the KPIRow block component (Issue #133, v0.6.0).

Covers layout math, auto-fit shrink of long values, card-frame rendering,
theme-token resolution, n==1 / n==6 edge cases, and containment. A
backward-compat smoke test for the renamed legacy MCP tool is included so
the rename from ``pptx_add_kpi_row`` → ``pptx_add_kpi_row_legacy`` cannot
regress silently.
"""

from __future__ import annotations

import json

import pytest

from pptx_mcp_server.engine.components.kpi_row import KPISpec, add_kpi_row
from pptx_mcp_server.engine.validation import check_containment

# Container registry isolation is provided globally by the autouse
# ``_isolated_container_registry`` fixture in ``tests/conftest.py`` (Task G).


def _hex(shape_fill):
    """Return the shape's solid-fill color as a ``#``-less 6-hex string."""
    return str(shape_fill.fore_color.rgb)


# ---------------------------------------------------------------------------
# 1) Even distribution over 12" width with 3 cells
# ---------------------------------------------------------------------------


def test_three_kpis_evenly_distributed_12in(slide):
    kpis = [
        KPISpec(label="Revenue", value="107.8M"),
        KPISpec(label="Margin", value="28.5%"),
        KPISpec(label="Growth", value="+12%"),
    ]
    result = add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=12.0, gap=0.15,
    )

    assert len(result["cells"]) == 3
    # cell_w = (12.0 - 2*0.15) / 3 = 11.7/3 = 3.9
    expected_w = (12.0 - 2 * 0.15) / 3
    for c in result["cells"]:
        assert c["bounds"]["width"] == pytest.approx(expected_w, abs=1e-9)

    # Exact x positions: 1.0, 1.0 + 3.9 + 0.15, 1.0 + 2*(3.9+0.15).
    step = expected_w + 0.15
    expected_x = [1.0, 1.0 + step, 1.0 + 2 * step]
    for exp, c in zip(expected_x, result["cells"]):
        assert c["bounds"]["left"] == pytest.approx(exp, abs=1e-9)


# ---------------------------------------------------------------------------
# 2) Long CJK value auto-shrinks below starting 26pt
# ---------------------------------------------------------------------------


def test_long_cjk_value_auto_fits(slide):
    kpis = [
        KPISpec(label="Brand", value="プレミアムモルツ +156%"),
    ]
    result = add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=3.0,
    )

    assert len(result["cells"]) == 1
    # wrap=False で width に収めるために shrink が走るはず。
    assert result["cells"][0]["value_actual_font_size"] < 26.0


# ---------------------------------------------------------------------------
# 3) card_fill="highlight_row" with theme="ir" resolves through IR palette
# ---------------------------------------------------------------------------


def test_card_fill_highlight_row_ir_theme(slide):
    kpis = [
        KPISpec(label="A", value="1"),
        KPISpec(label="B", value="2"),
    ]
    result = add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=6.0,
        theme="ir",
        card_fill="highlight_row",
    )

    # IR theme: highlight_row = "#F0F0F0" (see theme.py).
    from pptx_mcp_server.theme import resolve_theme_color
    expected = resolve_theme_color("highlight_row", "ir").upper()

    for c in result["cells"]:
        assert c["card_shape"] is not None
        assert _hex(c["card_shape"].fill).upper() == expected


# ---------------------------------------------------------------------------
# 4) Detail line rendered below value when spec.detail is non-empty
# ---------------------------------------------------------------------------


def test_detail_rendered_below_value(slide):
    kpis = [
        KPISpec(label="Revenue", value="107.8M", detail="+12% QoQ"),
    ]
    result = add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=3.0, height=1.0,
    )

    c = result["cells"][0]
    assert c["detail_shape"] is not None
    # detail top should be strictly below value top.
    from pptx.util import Emu
    value_top_emu = c["value_shape"].top
    detail_top_emu = c["detail_shape"].top
    assert detail_top_emu > value_top_emu


# ---------------------------------------------------------------------------
# 5) n == 1: cell spans full width, no gap applied
# ---------------------------------------------------------------------------


def test_single_kpi_full_width_no_gap(slide):
    kpis = [KPISpec(label="Total", value="42")]
    result = add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=4.0, gap=0.15,
    )

    assert len(result["cells"]) == 1
    b = result["cells"][0]["bounds"]
    assert b["width"] == pytest.approx(4.0, abs=1e-9)
    assert b["left"] == pytest.approx(1.0, abs=1e-9)


# ---------------------------------------------------------------------------
# 6) n == 6: many narrow cells
# ---------------------------------------------------------------------------


def test_six_kpis_each_narrow(slide):
    kpis = [KPISpec(label=f"L{i}", value=str(i)) for i in range(6)]
    result = add_kpi_row(
        slide, kpis,
        left=0.5, top=2.0, width=12.0, gap=0.15,
    )

    assert len(result["cells"]) == 6
    # cell_w = (12.0 - 5*0.15) / 6 = 11.25/6 = 1.875
    expected_w = (12.0 - 5 * 0.15) / 6
    for c in result["cells"]:
        assert c["bounds"]["width"] == pytest.approx(expected_w, abs=1e-9)


# ---------------------------------------------------------------------------
# 7) theme="ir" resolves all tokens on value + label + detail + border
# ---------------------------------------------------------------------------


def test_theme_ir_resolves_all_tokens(slide):
    """Passing theme='ir' + mixed tokens should resolve to IR hex everywhere."""
    kpis = [
        KPISpec(label="Rev", value="100M", detail="+5% YoY", value_color="primary"),
    ]
    result = add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=4.0,
        theme="ir",
        card_fill="highlight_row",
        card_border="rule_subtle",
    )

    from pptx_mcp_server.theme import resolve_theme_color

    c = result["cells"][0]
    # Card fill + border resolved.
    assert _hex(c["card_shape"].fill).upper() == resolve_theme_color("highlight_row", "ir").upper()
    # IR primary = #0A2540.
    primary_hex = resolve_theme_color("primary", "ir").upper()
    # Value textbox first run color.
    value_run = c["value_shape"].text_frame.paragraphs[0].runs[0]
    assert str(value_run.font.color.rgb).upper() == primary_hex

    # text_secondary from IR = #6B7280 — label & detail should use it.
    ts_hex = resolve_theme_color("text_secondary", "ir").upper()
    label_run = c["label_shape"].text_frame.paragraphs[0].runs[0]
    detail_run = c["detail_shape"].text_frame.paragraphs[0].runs[0]
    assert str(label_run.font.color.rgb).upper() == ts_hex
    assert str(detail_run.font.color.rgb).upper() == ts_hex


# ---------------------------------------------------------------------------
# 8) Containment: all child shapes stay inside the kpi_row container bbox
# ---------------------------------------------------------------------------


def test_check_containment_zero_findings(one_slide_prs):
    slide = one_slide_prs.slides[0]
    kpis = [
        KPISpec(label="A", value="1"),
        KPISpec(label="B", value="2", detail="note"),
        KPISpec(label="C", value="3"),
    ]
    add_kpi_row(
        slide, kpis,
        left=1.0, top=2.0, width=10.0, height=0.95, gap=0.15,
    )
    findings = check_containment(one_slide_prs)
    # All child shapes must fit inside the declared "kpi_row" container.
    assert findings == [], f"Unexpected containment findings: {findings}"


# ---------------------------------------------------------------------------
# 9) MCP-tool backward compat: legacy rename still callable
# ---------------------------------------------------------------------------


def test_pptx_add_kpi_row_legacy_backward_compat(pptx_file):
    """Verify the renamed legacy MCP tool still works for existing callers."""
    from pptx_mcp_server.server import pptx_add_kpi_row_legacy

    result_json = pptx_add_kpi_row_legacy(
        pptx_file,
        0,
        [{"label": "Rev", "value": "100"}],
        2.0,
    )
    payload = json.loads(result_json)
    assert payload["ok"] is True


# ---------------------------------------------------------------------------
# 10) MCP-tool smoke: new pptx_add_kpi_row writes and validates
# ---------------------------------------------------------------------------


def test_pptx_add_kpi_row_mcp_tool(pptx_file):
    """The new block-component MCP tool should write a file without error."""
    from pptx_mcp_server.server import pptx_add_kpi_row

    result_json = pptx_add_kpi_row(
        file_path=pptx_file,
        slide_index=0,
        kpis=[
            {"label": "Revenue", "value": "107.8M", "detail": "+12% QoQ"},
            {"label": "Margin", "value": "28.5%"},
        ],
        left=1.0,
        top=2.0,
        width=8.0,
    )
    payload = json.loads(result_json)
    assert payload["ok"] is True
    assert len(payload["result"]["cells"]) == 2
    assert payload["result"]["consumed_height"] == pytest.approx(0.95)


# ---------------------------------------------------------------------------
# 11) MCP-tool strict-key: unknown key on a kpi dict is rejected
# ---------------------------------------------------------------------------


def test_pptx_add_kpi_row_rejects_unknown_key(pptx_file):
    from pptx_mcp_server.server import pptx_add_kpi_row

    result_json = pptx_add_kpi_row(
        file_path=pptx_file,
        slide_index=0,
        kpis=[{"label": "A", "value": "1", "typo_field": 1}],
        left=1.0,
        top=2.0,
        width=4.0,
    )
    payload = json.loads(result_json)
    assert payload["ok"] is False
    assert payload["error"]["code"] == "INVALID_PARAMETER"
    assert "typo_field" in payload["error"]["message"]
