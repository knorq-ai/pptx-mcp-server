"""shapes.py の rendering-quality 修正 (#28, #29, #40, #41) のテスト.

1. 実 rPr への font size 書き込み (#28)
2. grapheme-cluster 安全な truncation (#29)
3. 明示的 east-asian typeface (#40)
4. bodyPr anchor を python-pptx API 経由で設定 (#41)
"""

from __future__ import annotations

import zipfile

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches

from pptx_mcp_server.engine.pptx_io import save_pptx
from pptx_mcp_server.engine.shapes import (
    _A_NS,
    _apply_font,
    _strip_last_grapheme,
    _truncate_to_fit,
    _add_textbox,
    add_auto_fit_textbox,
)
from pptx_mcp_server.engine.validation import _dominant_font_size
from pptx_mcp_server.theme import MCKINSEY, DELOITTE, NEUTRAL


# ── Issue #28: rPr に font size が書かれるか ───────────────────────


def test_issue28_auto_fit_writes_run_level_font_size(slide):
    """``add_auto_fit_textbox`` で生成した shape の全 run に
    ``font.size`` が None でない値で設定されていること."""
    shape, actual = add_auto_fit_textbox(
        slide,
        "Hello world",
        left=1.0, top=1.0, width=5.0, height=2.0,
        font_size_pt=10, min_size_pt=7,
    )
    assert actual == 10.0
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    assert runs, "shape には少なくとも 1 つ run が存在すべき"
    for run in runs:
        assert run.font.size is not None
        assert run.font.size.pt == 10.0


def test_issue28_saved_pptx_has_sz_on_every_rpr(tmp_path, slide, one_slide_prs):
    """保存 → unzip → XML 検査で ``<a:r>`` 配下に ``<a:rPr sz="1000"/>`` が
    毎回存在することを確認する."""
    add_auto_fit_textbox(
        slide,
        "run level size test",
        left=1.0, top=1.0, width=5.0, height=2.0,
        font_size_pt=10, min_size_pt=7,
    )
    path = tmp_path / "out.pptx"
    save_pptx(one_slide_prs, str(path))

    import xml.etree.ElementTree as ET
    ns = {"a": _A_NS}
    with zipfile.ZipFile(str(path)) as zf:
        xml_bytes = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(xml_bytes)
    runs = root.findall(".//a:r", ns)
    assert runs, "slide XML に run が 1 つ以上存在するはず"
    for r in runs:
        rPr = r.find("a:rPr", ns)
        assert rPr is not None, "各 run に rPr が必要"
        assert rPr.get("sz") == "1000", f"sz 属性が 1000 (=10pt) であるべき: {rPr.get('sz')}"


def test_issue28_dominant_font_size_from_runs(slide):
    """auto-fit で書いた shape に対し validator の ``_dominant_font_size``
    が run 経由で正しく pt を返す (18 default に落ちない)."""
    shape, actual = add_auto_fit_textbox(
        slide,
        "validator regression",
        left=1.0, top=1.0, width=5.0, height=2.0,
        font_size_pt=9, min_size_pt=7,
    )
    pt = _dominant_font_size(shape.text_frame, default_pt=18.0)
    assert pt == 9.0


# ── Issue #29: grapheme-cluster 安全な truncation ──────────────────


def test_issue29_strip_last_grapheme_zwj_emoji():
    """ZWJ emoji (家族絵文字) を末尾から 1 grapheme 削っても途中の
    ``\\u200D`` が残らない."""
    s = "hello 👨\u200D👩\u200D👧"
    stripped = _strip_last_grapheme(s)
    # 末尾が ZWJ で終わっていないこと。
    assert not stripped.endswith("\u200D")
    # ZWNJ も残らないこと。
    assert not stripped.endswith("\u200C")


def test_issue29_strip_last_grapheme_combining_mark():
    """結合可能な濁点 ``\\u309B`` が単独で残らない."""
    # `か` + 結合マーク U+3099 (濁点) = `が` の decomposed 表現。
    s = "日本語か\u3099"
    stripped = _strip_last_grapheme(s)
    # 削った結果、末尾に combining 文字だけが残ることはない。
    if stripped:
        import unicodedata
        assert unicodedata.combining(stripped[-1]) == 0


def test_issue29_strip_last_grapheme_variation_selector():
    """variation selector (U+FE0F) が単独で残らない."""
    s = "A\u2764\ufe0f"  # ハート + VS-16
    stripped = _strip_last_grapheme(s)
    if stripped:
        cp = ord(stripped[-1])
        assert not (0xFE00 <= cp <= 0xFE0F)


def test_issue29_strip_last_grapheme_plain_ascii():
    """単純な ASCII では末尾 1 文字が削られる."""
    assert _strip_last_grapheme("abcde") == "abcd"
    assert _strip_last_grapheme("a") == ""
    assert _strip_last_grapheme("") == ""


def test_issue29_truncate_does_not_leave_zwj_dangling():
    """``_truncate_to_fit`` に ZWJ emoji sequence を含むテキストを与え、
    切り詰めた結果が dangling ZWJ で終わらない."""
    text = "normal text 👨\u200D👩\u200D👧"
    # 極端に狭い box で確実に truncate させる。
    result = _truncate_to_fit(text, usable_width=0.5, height=0.2,
                              font_name="Arial", size_pt=10)
    # 末尾が ellipsis であればその直前に ZWJ が残っていないこと。
    if result.endswith("\u2026"):
        before_ellipsis = result[:-1]
        assert not before_ellipsis.endswith("\u200D")
        assert not before_ellipsis.endswith("\u200C")


def test_issue29_truncate_plain_ascii_regression(slide):
    """プレーン ASCII の切り詰めが従来どおり動作する (regression)."""
    result = _truncate_to_fit("abcdefghijklmnopqrstuvwxyz" * 4,
                              usable_width=0.5, height=0.2,
                              font_name="Arial", size_pt=10)
    # 1 行以内に丸められ ellipsis で終わる。
    assert result.endswith("\u2026")


# ── Issue #40: 明示的 east-asian typeface ─────────────────────────


def test_issue40_theme_slots_populated():
    """McKinsey / Deloitte / Neutral テーマに east_asian が設定済み."""
    for theme in (MCKINSEY, DELOITTE, NEUTRAL):
        assert "east_asian" in theme.fonts
        assert theme.fonts["east_asian"]


def test_issue40_apply_font_emits_ea(slide):
    """``_apply_font(..., east_asian_font=...)`` が各 run の rPr に
    ``<a:ea typeface="…"/>`` を設定する."""
    from pptx_mcp_server.engine.shapes import _add_textbox
    _add_textbox(slide, 1.0, 1.0, 5.0, 1.0, text="日本語テキスト")
    # 再度 apply_font で east-asian を設定。
    shape = slide.shapes[-1]
    p = shape.text_frame.paragraphs[0]
    _apply_font(p, font_name="Arial", east_asian_font="Yu Gothic")
    for run in p.runs:
        rPr = run._r.find(f"{{{_A_NS}}}rPr")
        assert rPr is not None
        ea = rPr.find(f"{{{_A_NS}}}ea")
        assert ea is not None
        assert ea.get("typeface") == "Yu Gothic"


def test_issue40_auto_fit_with_east_asian_font(slide):
    """auto-fit textbox に east_asian_font を渡すと runs に emit される."""
    shape, _ = add_auto_fit_textbox(
        slide,
        "日本語テスト",
        left=1.0, top=1.0, width=5.0, height=2.0,
        font_size_pt=11, min_size_pt=7,
        east_asian_font="Meiryo",
    )
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    assert runs
    for run in runs:
        rPr = run._r.find(f"{{{_A_NS}}}rPr")
        assert rPr is not None
        ea = rPr.find(f"{{{_A_NS}}}ea")
        assert ea is not None
        assert ea.get("typeface") == "Meiryo"


def test_issue40_no_ea_when_not_requested(slide):
    """``east_asian_font=None`` のデフォルト挙動では ``<a:ea>`` を
    emit しない (後方互換)."""
    shape, _ = add_auto_fit_textbox(
        slide,
        "no ea expected",
        left=1.0, top=1.0, width=5.0, height=2.0,
        font_size_pt=11, min_size_pt=7,
    )
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    for run in runs:
        rPr = run._r.find(f"{{{_A_NS}}}rPr")
        if rPr is None:
            continue
        ea = rPr.find(f"{{{_A_NS}}}ea")
        assert ea is None


def test_issue40_add_textbox_picks_east_asian_from_theme(slide, mckinsey_theme):
    """``_add_textbox`` に theme を渡すと ``theme.fonts['east_asian']``
    が自動で適用される."""
    _add_textbox(
        slide, 1.0, 1.0, 5.0, 1.0,
        text="テーマからの ea",
        theme=mckinsey_theme,
    )
    shape = slide.shapes[-1]
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    assert runs
    for run in runs:
        rPr = run._r.find(f"{{{_A_NS}}}rPr")
        assert rPr is not None
        ea = rPr.find(f"{{{_A_NS}}}ea")
        assert ea is not None
        assert ea.get("typeface") == mckinsey_theme.fonts["east_asian"]


# ── Issue #41: bodyPr anchor を python-pptx API 経由で設定 ───────


@pytest.mark.parametrize("anchor,expected", [
    ("top", MSO_ANCHOR.TOP),
    ("middle", MSO_ANCHOR.MIDDLE),
    ("bottom", MSO_ANCHOR.BOTTOM),
])
def test_issue41_anchor_set_via_python_pptx_api(slide, anchor, expected):
    """3 つの anchor 値が MSO_ANCHOR enum 経由で設定される."""
    _add_textbox(slide, 1.0, 1.0, 5.0, 1.0,
                 text="anchor test", vertical_anchor=anchor)
    shape = slide.shapes[-1]
    assert shape.text_frame.vertical_anchor == expected


def test_issue41_anchor_round_trip(tmp_path, slide, one_slide_prs):
    """save → reopen で anchor 値が保持される."""
    _add_textbox(slide, 1.0, 1.0, 5.0, 1.0,
                 text="roundtrip", vertical_anchor="middle")
    path = tmp_path / "anchor.pptx"
    save_pptx(one_slide_prs, str(path))

    reopened = Presentation(str(path))
    tf = reopened.slides[0].shapes[-1].text_frame
    assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE


def test_issue41_no_duplicate_anchor_attrs(tmp_path, slide, one_slide_prs):
    """保存 XML 内の ``<a:bodyPr>`` に anchor 属性が重複しない."""
    _add_textbox(slide, 1.0, 1.0, 5.0, 1.0,
                 text="dup check", vertical_anchor="middle")
    path = tmp_path / "dup.pptx"
    save_pptx(one_slide_prs, str(path))

    import xml.etree.ElementTree as ET
    ns = {"a": _A_NS}
    with zipfile.ZipFile(str(path)) as zf:
        xml_bytes = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(xml_bytes)
    for bodyPr in root.findall(".//a:bodyPr", ns):
        # attribute key の duplicates はそもそも XML として成立しないが、
        # anchor 属性が 1 つだけ、値が想定通りであることを確認する。
        assert bodyPr.get("anchor") in (None, "t", "ctr", "b")
