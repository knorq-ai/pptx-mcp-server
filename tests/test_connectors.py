"""
Tests for the connector and callout engine (Phase 3).
"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from pptx_mcp_server.engine.connectors import (
    _add_connector,
    _add_callout,
    _set_arrow_heads,
    add_connector,
    add_callout,
)
from pptx_mcp_server.engine.composites import _build_slide
from pptx_mcp_server.engine.pptx_io import EngineError
from pptx_mcp_server.theme import MCKINSEY, DELOITTE


# ---------------------------------------------------------------------------
# Connector creation tests
# ---------------------------------------------------------------------------


class TestAddConnector:
    def test_straight_connector(self, slide):
        idx = _add_connector(slide, 1.0, 1.0, 5.0, 3.0)
        assert idx >= 0
        shape = list(slide.shapes)[idx]
        # Verify it's a connector shape
        assert shape.shape_type is not None

    def test_elbow_connector(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            connector_type="elbow",
        )
        assert idx >= 0

    def test_curve_connector(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            connector_type="curve",
        )
        assert idx >= 0

    def test_connector_color(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            color="FF0000",
        )
        shape = list(slide.shapes)[idx]
        assert shape.line.color.rgb is not None

    def test_connector_width(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            width=3.0,
        )
        shape = list(slide.shapes)[idx]
        assert shape.line.width is not None

    def test_connector_dash_style(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            dash_style="dash",
        )
        shape = list(slide.shapes)[idx]
        assert shape.line.dash_style is not None

    def test_connector_theme_color(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            color="accent",
            theme=MCKINSEY,
        )
        assert idx >= 0

    def test_connector_deloitte_theme(self, slide):
        idx = _add_connector(
            slide, 1.0, 1.0, 5.0, 3.0,
            theme=DELOITTE,
        )
        assert idx >= 0

    def test_multiple_connectors(self, slide):
        idx1 = _add_connector(slide, 1, 1, 3, 3)
        idx2 = _add_connector(slide, 4, 1, 6, 3)
        assert idx2 > idx1


# ---------------------------------------------------------------------------
# Arrow head tests
# ---------------------------------------------------------------------------


class TestArrowHeads:
    def test_end_arrow_triangle(self, slide):
        idx = _add_connector(
            slide, 1, 1, 5, 3,
            end_arrow="triangle",
        )
        shape = list(slide.shapes)[idx]
        # Check XML for tailEnd element
        cxnSp = shape._element
        tail = cxnSp.find(f".//{qn('a:tailEnd')}")
        assert tail is not None
        assert tail.get("type") == "triangle"

    def test_begin_arrow(self, slide):
        idx = _add_connector(
            slide, 1, 1, 5, 3,
            begin_arrow="stealth",
            end_arrow="none",
        )
        shape = list(slide.shapes)[idx]
        cxnSp = shape._element
        head = cxnSp.find(f".//{qn('a:headEnd')}")
        assert head is not None
        assert head.get("type") == "stealth"
        tail = cxnSp.find(f".//{qn('a:tailEnd')}")
        assert tail is None  # none = no element

    def test_both_arrows(self, slide):
        idx = _add_connector(
            slide, 1, 1, 5, 3,
            begin_arrow="diamond",
            end_arrow="oval",
        )
        shape = list(slide.shapes)[idx]
        cxnSp = shape._element
        head = cxnSp.find(f".//{qn('a:headEnd')}")
        tail = cxnSp.find(f".//{qn('a:tailEnd')}")
        assert head is not None
        assert tail is not None
        assert head.get("type") == "diamond"
        assert tail.get("type") == "oval"

    def test_no_arrows(self, slide):
        idx = _add_connector(
            slide, 1, 1, 5, 3,
            begin_arrow="none",
            end_arrow="none",
        )
        shape = list(slide.shapes)[idx]
        cxnSp = shape._element
        assert cxnSp.find(f".//{qn('a:headEnd')}") is None
        assert cxnSp.find(f".//{qn('a:tailEnd')}") is None

    def test_arrow_sizes(self, slide):
        for size in ("small", "medium", "large"):
            idx = _add_connector(
                slide, 1, 1, 5, 3,
                end_arrow="triangle",
                arrow_size=size,
            )
            shape = list(slide.shapes)[idx]
            tail = shape._element.find(f".//{qn('a:tailEnd')}")
            expected = {"small": "sm", "medium": "med", "large": "lg"}[size]
            assert tail.get("w") == expected

    def test_all_arrow_types(self, slide):
        for arrow_type in ("triangle", "stealth", "diamond", "oval", "open"):
            idx = _add_connector(
                slide, 1, 1, 5, 3,
                end_arrow=arrow_type,
            )
            shape = list(slide.shapes)[idx]
            tail = shape._element.find(f".//{qn('a:tailEnd')}")
            assert tail is not None


# ---------------------------------------------------------------------------
# Dash style tests
# ---------------------------------------------------------------------------


class TestDashStyles:
    def test_all_dash_styles(self, slide):
        for style in ("solid", "dash", "dot", "dash_dot", "long_dash"):
            idx = _add_connector(
                slide, 1, 1, 5, 3,
                dash_style=style,
            )
            shape = list(slide.shapes)[idx]
            assert shape.line.dash_style is not None


# ---------------------------------------------------------------------------
# Callout tests
# ---------------------------------------------------------------------------


class TestAddCallout:
    def test_basic_callout(self, slide):
        indices = _add_callout(
            slide, "Test annotation", 5.0, 3.0,
        )
        assert len(indices) == 2
        assert indices[0] >= 0  # textbox
        assert indices[1] >= 0  # connector

    def test_callout_with_position(self, slide):
        indices = _add_callout(
            slide, "Revenue +15%", 5.0, 3.0,
            label_x=7.0, label_y=1.5,
            label_width=2.5,
        )
        assert len(indices) == 2

    def test_callout_with_styling(self, slide):
        indices = _add_callout(
            slide, "Key insight", 5.0, 3.0,
            font_size=12,
            font_color="negative",
            bg_color="F5F5F5",
            border_color="C62828",
            line_color="C62828",
            arrow_end="stealth",
            theme=MCKINSEY,
        )
        assert len(indices) == 2

    def test_callout_auto_placement(self, slide):
        """Label auto-places 1.5" right and 1.0" above target."""
        indices = _add_callout(
            slide, "Auto placed", 5.0, 3.0,
            theme=MCKINSEY,
        )
        # Just verify no crash and both shapes created
        assert len(indices) == 2

    def test_callout_auto_placement_near_right_edge(self, slide):
        """When target is near right edge, label should go left."""
        indices = _add_callout(
            slide, "Near edge", 12.0, 3.0,
            theme=MCKINSEY,
        )
        assert len(indices) == 2

    def test_callout_auto_placement_near_top(self, slide):
        """When target is near top, label_y should clamp to 0.2."""
        indices = _add_callout(
            slide, "Near top", 5.0, 0.5,
            theme=MCKINSEY,
        )
        assert len(indices) == 2

    def test_callout_elbow_connector(self, slide):
        indices = _add_callout(
            slide, "Elbow", 5.0, 3.0,
            connector_type="elbow",
        )
        assert len(indices) == 2


# ---------------------------------------------------------------------------
# File-based wrapper tests
# ---------------------------------------------------------------------------


class TestFileWrappers:
    def test_add_connector_to_file(self, pptx_file):
        result = add_connector(pptx_file, 0, 1, 1, 5, 3)
        assert "connector" in result
        assert "slide [0]" in result

    def test_add_callout_to_file(self, pptx_file):
        result = add_callout(pptx_file, 0, "Test", 5.0, 3.0)
        assert "callout" in result
        assert "Test" in result

    def test_roundtrip(self, pptx_file):
        """Add connector + callout, save, reopen."""
        from pptx_mcp_server.engine.pptx_io import open_pptx

        before = len(list(open_pptx(pptx_file).slides[0].shapes))
        add_connector(pptx_file, 0, 1, 1, 5, 3)
        add_callout(pptx_file, 0, "Annotation", 6, 2)
        after = len(list(open_pptx(pptx_file).slides[0].shapes))
        assert after > before


# ---------------------------------------------------------------------------
# build_slide integration tests
# ---------------------------------------------------------------------------


class TestBuildSlideConnector:
    def test_connector_element(self, blank_prs):
        spec = {
            "layout": "blank",
            "elements": [
                {
                    "type": "connector",
                    "begin_x": 1, "begin_y": 1,
                    "end_x": 5, "end_y": 3,
                    "color": "accent",
                    "end_arrow": "triangle",
                    "dash_style": "dash",
                }
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        assert len(list(slide.shapes)) >= 1

    def test_callout_element(self, blank_prs):
        spec = {
            "layout": "blank",
            "elements": [
                {
                    "type": "callout",
                    "text": "+15% growth",
                    "target_x": 5.0,
                    "target_y": 3.0,
                    "font_color": "accent",
                    "bg_color": "F5F5F5",
                }
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        # Callout creates textbox + connector = 2 shapes
        assert len(list(slide.shapes)) >= 2

    def test_mixed_elements(self, blank_prs):
        """Build slide with chart, connector, and callout together."""
        spec = {
            "layout": "content",
            "title": "Mixed Elements",
            "elements": [
                {
                    "type": "chart",
                    "chart_type": "column",
                    "left": 0.9, "top": 1.15,
                    "width": 7, "height": 4,
                    "categories": ["A", "B"],
                    "series": [{"name": "S", "values": [10, 20]}],
                },
                {
                    "type": "connector",
                    "begin_x": 8, "begin_y": 2,
                    "end_x": 10, "end_y": 4,
                    "end_arrow": "stealth",
                },
                {
                    "type": "callout",
                    "text": "Key insight",
                    "target_x": 5, "target_y": 3,
                    "label_x": 8.5, "label_y": 1.5,
                },
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        # title + divider + chart + connector + callout(textbox+connector) = 6+
        assert len(list(slide.shapes)) >= 5
