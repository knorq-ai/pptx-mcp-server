"""
Shared pytest fixtures for pptx-mcp-server tests.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_mcp_server.theme import MCKINSEY, Theme
from pptx_mcp_server.engine.pptx_io import save_pptx


@pytest.fixture
def blank_prs() -> Presentation:
    """Return a new in-memory Presentation with 16:9 dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


@pytest.fixture
def one_slide_prs(blank_prs: Presentation) -> Presentation:
    """Return a Presentation with exactly one blank slide."""
    layout = blank_prs.slide_layouts[6]
    blank_prs.slides.add_slide(layout)
    return blank_prs


@pytest.fixture
def slide(one_slide_prs: Presentation):
    """Return the first slide from a one-slide presentation."""
    return one_slide_prs.slides[0]


@pytest.fixture
def mckinsey_theme() -> Theme:
    """Return the MCKINSEY theme instance."""
    return MCKINSEY


@pytest.fixture
def pptx_file(tmp_path):
    """Create a blank .pptx file on disk and return its path as a string."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    path = tmp_path / "test.pptx"
    save_pptx(prs, str(path))
    return str(path)
