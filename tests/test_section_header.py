"""Tests for the SectionHeader block component (Issue #136).

Covers:

- Shape layout (two textboxes + one divider rectangle).
- Auto-fit contract: short text keeps base size; long text shrinks.
- ``wrap=False`` is propagated to the resulting textbox.
- Subtitle omission reduces ``consumed_height`` and yields
  ``subtitle_bounds=None``.
- Theme tokens (``primary``, ``text_secondary``) resolve against the IR
  theme's palette.
- ``check_containment`` returns zero findings for a normal header
  (regression safety for ``begin_container`` wrapping).
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_mcp_server.engine.components.section_header import (
    SectionHeaderSpec,
    add_section_header,
    _TITLE_H,
    _SUBTITLE_H,
    _INTRA_GAP_TITLE,
    _INTRA_GAP_SUBTITLE,
)
from pptx_mcp_server.engine.validation import check_containment


def test_title_subtitle_divider_creates_three_shapes(slide):
    """With subtitle → 2 textboxes + 1 filled rectangle divider."""
    top = 0.45
    spec = SectionHeaderSpec(
        title="Market Outlook",
        subtitle="Q4 2026 — Key Takeaways",
    )
    result = add_section_header(
        slide, spec, left=0.9, top=top, width=11.533,
    )

    textboxes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    auto_shapes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]
    assert len(textboxes) == 2, "expected title + subtitle textboxes"
    assert len(auto_shapes) == 1, "expected exactly one divider rectangle"

    # Sanity: result dict exposes all three regions.
    assert result["title_bounds"] is not None
    assert result["subtitle_bounds"] is not None
    assert result["divider_bounds"] is not None

    # Absolute invariant: the divider's bottom edge must equal
    # top + consumed_height, so callers placing body content at
    # ``top + consumed_height`` align flush with the header's footprint.
    divider_bottom = (
        result["divider_bounds"]["top"] + result["divider_bounds"]["height"]
    )
    expected_bottom = top + result["consumed_height"]
    assert abs(divider_bottom - expected_bottom) < 1e-6, (
        f"divider bottom ({divider_bottom}) must equal "
        f"top + consumed_height ({expected_bottom})"
    )


def test_short_title_keeps_base_font_size(slide):
    """Short title easily fits 11.533" width at 32pt → no shrink."""
    spec = SectionHeaderSpec(title="Summary")
    result = add_section_header(
        slide, spec, left=0.9, top=0.45, width=11.533,
    )
    # Default title_size_pt is 32; short text must not shrink.
    assert result["title_actual_font_size"] == 32.0


def test_long_title_auto_fits_single_line(slide):
    """Very long title shrinks below 32pt AND stays single-line (word_wrap=False)."""
    long_title = (
        "Comprehensive strategic review of market dynamics, competitive "
        "positioning, and long-term growth opportunities across all segments"
    )
    spec = SectionHeaderSpec(title=long_title)
    result = add_section_header(
        slide, spec, left=0.9, top=0.45, width=11.533,
    )
    assert result["title_actual_font_size"] < 32.0

    # Find the title textbox (first TEXT_BOX added).
    textboxes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    assert textboxes, "title textbox should exist"
    title_shape = textboxes[0]
    # wrap=False must propagate to the textbox's word_wrap property.
    assert title_shape.text_frame.word_wrap is False


def test_no_subtitle_divider_closer_to_title(slide):
    """subtitle="" → only 1 textbox; consumed_height drops by subtitle strip."""
    top = 0.45
    spec_no_sub = SectionHeaderSpec(title="Title only", subtitle="")
    result = add_section_header(
        slide, spec_no_sub, left=0.9, top=top, width=11.533,
    )

    textboxes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    assert len(textboxes) == 1, "no subtitle → only title textbox"
    assert result["subtitle_bounds"] is None

    # Without subtitle, consumed_height drops by the subtitle strip AND the
    # title→subtitle gap (which is only relevant when a subtitle follows).
    # Layout: title + gap_title + subtitle + gap_subtitle + divider (full)
    #   vs   title              +            gap_subtitle + divider (no-sub)
    # Delta = _INTRA_GAP_TITLE + _SUBTITLE_H.
    # Compute the "with subtitle" height from a matching spec to avoid
    # hard-coding internal constants in the delta check.
    full_spec = SectionHeaderSpec(title="Title only", subtitle="sub")
    full_height = (
        _TITLE_H + _INTRA_GAP_TITLE + _SUBTITLE_H
        + _INTRA_GAP_SUBTITLE + full_spec.divider_thickness
    )
    assert result["consumed_height"] < full_height
    assert abs(
        (full_height - result["consumed_height"])
        - (_INTRA_GAP_TITLE + _SUBTITLE_H)
    ) < 1e-6

    # Absolute invariant: divider bottom edge must equal top + consumed_height.
    # This guards against the "phantom gap" regression where
    # _compute_consumed_height and the rendered divider position disagree.
    divider_bottom = (
        result["divider_bounds"]["top"] + result["divider_bounds"]["height"]
    )
    expected_bottom = top + result["consumed_height"]
    assert abs(divider_bottom - expected_bottom) < 1e-6, (
        f"divider bottom ({divider_bottom}) must equal "
        f"top + consumed_height ({expected_bottom})"
    )


def test_theme_ir_resolves_colors(slide):
    """theme='ir' → title uses IR primary (#0A2540), subtitle uses text_secondary (#6B7280)."""
    spec = SectionHeaderSpec(
        title="IR Headline",
        subtitle="Investor relations subtitle",
    )
    add_section_header(
        slide, spec, left=0.9, top=0.45, width=11.533, theme="ir",
    )

    textboxes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    assert len(textboxes) == 2
    title_shape, subtitle_shape = textboxes[0], textboxes[1]

    title_rgb = title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb
    subtitle_rgb = subtitle_shape.text_frame.paragraphs[0].runs[0].font.color.rgb

    # IR theme: primary = 0A2540, text_secondary = 6B7280.
    assert str(title_rgb) == "0A2540"
    assert str(subtitle_rgb) == "6B7280"


def test_check_containment_zero_findings(one_slide_prs):
    """Rendering SectionHeader into begin_container must not trigger
    `shape_outside_container` findings — nothing should escape bounds."""
    slide = one_slide_prs.slides[0]
    spec = SectionHeaderSpec(
        title="Bounded header",
        subtitle="Subtitle also bounded",
    )
    add_section_header(
        slide, spec, left=0.9, top=0.45, width=11.533,
    )
    findings = check_containment(one_slide_prs)
    assert findings == [], f"expected no containment findings, got {findings}"
