"""
Tests for the icon library (Phase 2).
"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

from pptx_mcp_server.engine.icons import (
    IconRegistry,
    _add_icon,
    _recolor_icon,
    _reassign_shape_ids,
    _resolve_icon_size,
    add_icon,
    list_icons_formatted,
    IconInfo,
)
from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx_mcp_server.engine.composites import _build_slide
from pptx_mcp_server.theme import MCKINSEY

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ---------------------------------------------------------------------------
# Registry tests
# ---------------------------------------------------------------------------


class TestIconRegistry:
    def test_loads_catalog(self):
        reg = IconRegistry.get()
        icons = reg.list_icons()
        assert len(icons) > 500  # We extracted 640

    def test_list_categories(self):
        reg = IconRegistry.get()
        cats = reg.list_categories()
        assert len(cats) >= 5
        cat_ids = {c["id"] for c in cats}
        assert "business" in cat_ids
        assert "people" in cat_ids
        assert "technology" in cat_ids

    def test_get_icon_by_id(self):
        reg = IconRegistry.get()
        icon = reg.get_icon("abacus")
        assert icon.id == "abacus"
        assert icon.name == "abacus"
        assert isinstance(icon.aspect_ratio, float)

    def test_get_icon_xml(self):
        reg = IconRegistry.get()
        elem = reg.get_icon_xml("abacus")
        assert elem is not None
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        assert tag in ("sp", "grpSp")

    def test_get_icon_xml_returns_fresh_copy(self):
        reg = IconRegistry.get()
        e1 = reg.get_icon_xml("abacus")
        e2 = reg.get_icon_xml("abacus")
        assert e1 is not e2  # Different objects

    def test_get_icon_not_found_with_fuzzy_match(self):
        reg = IconRegistry.get()
        with pytest.raises(EngineError) as exc_info:
            reg.get_icon("abacu")
        assert "Did you mean" in str(exc_info.value)
        assert "abacus" in str(exc_info.value)

    def test_get_icon_totally_unknown(self):
        reg = IconRegistry.get()
        with pytest.raises(EngineError) as exc_info:
            reg.get_icon("zzzzzzzzzzz_nonexistent")
        assert "not found" in str(exc_info.value)

    def test_list_icons_by_category(self):
        reg = IconRegistry.get()
        biz = reg.list_icons(category="business")
        assert len(biz) > 0
        assert all(i.category == "business" for i in biz)

    def test_list_icons_by_search(self):
        reg = IconRegistry.get()
        results = reg.list_icons(search="airplane")
        assert len(results) > 0
        assert any("airplane" in i.name.lower() for i in results)

    def test_list_icons_by_category_and_search(self):
        reg = IconRegistry.get()
        results = reg.list_icons(category="transport", search="car")
        assert len(results) > 0


# ---------------------------------------------------------------------------
# Icon size resolution
# ---------------------------------------------------------------------------


class TestIconSize:
    def test_both_dimensions(self):
        info = IconInfo("test", "test", "general", (), "sp", 914400, 914400, 1.0, ())
        w, h = _resolve_icon_size(info, 2.0, 3.0)
        assert w == int(2.0 * 914400)
        assert h == int(3.0 * 914400)

    def test_width_only_preserves_aspect(self):
        info = IconInfo("test", "test", "general", (), "sp", 914400, 457200, 2.0, ())
        w, h = _resolve_icon_size(info, 2.0, None)
        assert w == int(2.0 * 914400)
        assert h == int(w / 2.0)  # aspect ratio 2.0

    def test_height_only_preserves_aspect(self):
        info = IconInfo("test", "test", "general", (), "sp", 914400, 914400, 1.0, ())
        w, h = _resolve_icon_size(info, None, 1.0)
        assert w == h  # aspect ratio 1.0

    def test_default_size(self):
        info = IconInfo("test", "test", "general", (), "sp", 914400, 914400, 1.0, ())
        w, h = _resolve_icon_size(info, None, None)
        assert h == int(0.8 * 914400)  # default 0.8 inches


# ---------------------------------------------------------------------------
# Recoloring tests
# ---------------------------------------------------------------------------


class TestRecolor:
    def _make_elem_with_colors(self, fill="00AAE7", outline="000000"):
        """Create a minimal shape XML with srgbClr elements."""
        xml = f"""
        <p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
          <p:spPr>
            <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
            <a:ln><a:solidFill><a:srgbClr val="{outline}"/></a:solidFill></a:ln>
          </p:spPr>
        </p:sp>
        """
        return etree.fromstring(xml.encode())

    def test_recolor_fill(self):
        elem = self._make_elem_with_colors()
        _recolor_icon(elem, ("00AAE7", "000000"), "FF0000", None, None)
        srgb_vals = [e.get("val") for e in elem.findall(f".//{{{A_NS}}}srgbClr")]
        assert "FF0000" in srgb_vals
        assert "000000" in srgb_vals  # outline unchanged

    def test_recolor_outline(self):
        elem = self._make_elem_with_colors()
        _recolor_icon(elem, ("00AAE7", "000000"), None, "333333", None)
        srgb_vals = [e.get("val") for e in elem.findall(f".//{{{A_NS}}}srgbClr")]
        assert "00AAE7" in srgb_vals  # fill unchanged
        assert "333333" in srgb_vals

    def test_recolor_both(self):
        elem = self._make_elem_with_colors()
        _recolor_icon(elem, ("00AAE7", "000000"), "2251FF", "051C2C", None)
        srgb_vals = [e.get("val") for e in elem.findall(f".//{{{A_NS}}}srgbClr")]
        assert "2251FF" in srgb_vals
        assert "051C2C" in srgb_vals

    def test_recolor_with_theme_token(self):
        elem = self._make_elem_with_colors()
        _recolor_icon(elem, ("00AAE7", "000000"), "accent", "primary", MCKINSEY)
        srgb_vals = [e.get("val") for e in elem.findall(f".//{{{A_NS}}}srgbClr")]
        # accent -> #2251FF -> 2251FF, primary -> #051C2C -> 051C2C
        assert "2251FF" in srgb_vals
        assert "051C2C" in srgb_vals

    def test_recolor_no_op(self):
        elem = self._make_elem_with_colors()
        _recolor_icon(elem, ("00AAE7", "000000"), None, None, None)
        srgb_vals = [e.get("val") for e in elem.findall(f".//{{{A_NS}}}srgbClr")]
        assert "00AAE7" in srgb_vals  # unchanged


# ---------------------------------------------------------------------------
# ID deduplication tests
# ---------------------------------------------------------------------------


class TestIDDedup:
    def test_reassign_ids_unique(self, slide):
        """After injection, all shape IDs on the slide should be unique."""
        # Add two icons
        _add_icon(slide, "abacus", 1, 1, width=0.5)
        _add_icon(slide, "airplane", 3, 1, width=0.5)

        # Collect all cNvPr ids from spTree XML
        spTree = slide.shapes._spTree
        ids = []
        for cNvPr in spTree.findall(f".//{{{P_NS}}}cNvPr"):
            id_val = cNvPr.get("id")
            if id_val and id_val.isdigit():
                ids.append(int(id_val))

        # All IDs should be unique
        assert len(ids) == len(set(ids)), f"Duplicate IDs found: {ids}"


# ---------------------------------------------------------------------------
# In-memory injection tests
# ---------------------------------------------------------------------------


def _count_shapes_in_spTree(slide) -> int:
    """Count shape elements in spTree without using slide.shapes iterator."""
    shape_tags = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}
    spTree = slide.shapes._spTree
    return sum(
        1 for child in spTree
        if (child.tag.split("}")[-1] if "}" in child.tag else child.tag) in shape_tags
    )


class TestAddIcon:
    def test_inject_icon(self, slide):
        before = _count_shapes_in_spTree(slide)
        idx = _add_icon(slide, "abacus", 2.0, 3.0, width=0.8)
        assert idx >= 0
        after = _count_shapes_in_spTree(slide)
        assert after > before

    def test_inject_with_color(self, slide):
        idx = _add_icon(slide, "airplane", 1.0, 1.0, width=0.6, color="FF0000")
        assert idx >= 0

    def test_inject_with_theme(self, slide):
        idx = _add_icon(
            slide, "abacus", 1.0, 1.0,
            width=0.6, color="accent", outline_color="primary",
            theme=MCKINSEY,
        )
        assert idx >= 0

    def test_inject_invalid_icon(self, slide):
        with pytest.raises(EngineError):
            _add_icon(slide, "totally_nonexistent_icon_xyz", 1.0, 1.0)


# ---------------------------------------------------------------------------
# File-based wrapper tests
# ---------------------------------------------------------------------------


class TestAddIconFile:
    def test_add_icon_to_file(self, pptx_file):
        result = add_icon(pptx_file, 0, "abacus", 2.0, 2.0, width=0.8)
        assert "abacus" in result
        assert "slide [0]" in result

    def test_add_icon_roundtrip(self, pptx_file):
        """Add icon, save, reopen — verify shape count increased."""
        from pptx_mcp_server.engine.pptx_io import open_pptx

        # Count shapes before
        prs = open_pptx(pptx_file)
        before = len(list(prs.slides[0].shapes))

        add_icon(pptx_file, 0, "airplane", 1.0, 1.0, width=0.5)

        prs = open_pptx(pptx_file)
        after = len(list(prs.slides[0].shapes))
        assert after > before


# ---------------------------------------------------------------------------
# list_icons_formatted tests
# ---------------------------------------------------------------------------


class TestListIconsFormatted:
    def test_no_filter_shows_categories(self):
        result = list_icons_formatted()
        assert "Categories" in result
        assert "business" in result

    def test_category_filter(self):
        result = list_icons_formatted(category="transport")
        assert "Found" in result
        assert "airplane" in result.lower() or "car" in result.lower()

    def test_search_filter(self):
        result = list_icons_formatted(search="chart")
        assert "Found" in result

    def test_no_results(self):
        result = list_icons_formatted(search="zzzznonexistentzzz")
        assert "No icons found" in result


# ---------------------------------------------------------------------------
# build_slide integration tests
# ---------------------------------------------------------------------------


class TestBuildSlideIcon:
    def test_icon_element_in_build_slide(self, blank_prs):
        spec = {
            "layout": "blank",
            "elements": [
                {
                    "type": "icon",
                    "icon_id": "abacus",
                    "left": 2.0,
                    "top": 2.0,
                    "width": 0.8,
                    "color": "2251FF",
                }
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        assert _count_shapes_in_spTree(slide) >= 1

    def test_multiple_icons(self, blank_prs):
        spec = {
            "layout": "blank",
            "elements": [
                {"type": "icon", "icon_id": "abacus", "left": 1, "top": 1, "width": 0.6},
                {"type": "icon", "icon_id": "airplane", "left": 3, "top": 1, "width": 0.6},
                {"type": "icon", "icon_id": "arrow", "left": 5, "top": 1, "width": 0.6},
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        assert _count_shapes_in_spTree(slide) >= 3
