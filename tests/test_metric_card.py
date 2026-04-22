"""MetricCard block component tests (Issue #132).

Covers:
- Single-card layout (label/title/chart/metrics all inside bounds).
- Empty metrics list → chart area expands.
- Metrics row splits horizontal space evenly.
- ``chart_image_path`` provided → chart shape is a Picture.
- Row variant with N cards (equal widths, correct gap).
- Heterogeneous metrics counts in a row (each card independent).
- Too-small height raises EngineError(INVALID_PARAMETER).
- Theme resolution for ``theme="ir"`` produces IR-palette fills/borders.
- ``check_containment`` finds no overflow for correctly-bounded cards.
"""

from __future__ import annotations

import pytest
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_mcp_server.engine.components.container import clear_container_registry
from pptx_mcp_server.engine.components.metric_card import (
    MetricCardSpec,
    MetricEntry,
    add_metric_card,
    add_metric_card_row,
)
from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx_mcp_server.engine.validation import check_containment


@pytest.fixture(autouse=True)
def _isolated_registry():
    """Each test starts with a clean container registry."""
    clear_container_registry()
    yield
    clear_container_registry()


def _in(emu: int) -> float:
    """EMU → inches."""
    return emu / 914400


def _bbox(shape) -> tuple[float, float, float, float]:
    return _in(shape.left), _in(shape.top), _in(shape.width), _in(shape.height)


def _make_image(tmp_path, name="chart.png", size=(200, 100), color="blue") -> str:
    path = tmp_path / name
    Image.new("RGB", size, color=color).save(path)
    return str(path)


# ---------------------------------------------------------------------------
# 1. Single card: label/title/chart/metrics all inside bounds
# ---------------------------------------------------------------------------


def test_single_card_all_children_inside_bounds(one_slide_prs):
    slide = one_slide_prs.slides[0]
    spec = MetricCardSpec(
        label="KPI",
        title="Revenue 2024",
        metrics=[
            MetricEntry(label="Q1", value="$10M"),
            MetricEntry(label="Q2", value="$12M"),
            MetricEntry(label="Q3", value="$15M"),
            MetricEntry(label="Q4", value="$18M"),
        ],
    )
    result = add_metric_card(
        slide, spec, left=1.0, top=1.0, width=5.0, height=4.5,
    )

    # Bounds reflect the arguments.
    assert result["bounds"] == {
        "left": 1.0, "top": 1.0, "width": 5.0, "height": 4.5,
    }

    # We expect 4 metric cells.
    assert len(result["metric_shapes"]) == 4

    # All child shapes (label, title, chart, 8 metric shapes) should lie
    # inside the outer card bbox.
    outer_l, outer_t, outer_r, outer_b = 1.0, 1.0, 6.0, 5.5
    children = [
        result["label_shape"],
        result["title_shape"],
        result["chart_shape"],
    ]
    for pair in result["metric_shapes"]:
        children.extend(pair)

    for child in children:
        l, t, w, h = _bbox(child)
        assert l >= outer_l - 0.01
        assert t >= outer_t - 0.01
        assert l + w <= outer_r + 0.01
        assert t + h <= outer_b + 0.01


# ---------------------------------------------------------------------------
# 2. Zero metrics → chart area expands (no metrics shapes)
# ---------------------------------------------------------------------------


def test_zero_metrics_expands_chart(one_slide_prs):
    slide = one_slide_prs.slides[0]

    with_metrics_spec = MetricCardSpec(
        label="A", title="B",
        metrics=[MetricEntry(label="x", value="1")],
    )
    no_metrics_spec = MetricCardSpec(label="A", title="B", metrics=[])

    r_with = add_metric_card(
        slide, with_metrics_spec, left=0.5, top=0.5, width=4.0, height=4.5,
    )
    r_without = add_metric_card(
        slide, no_metrics_spec, left=5.0, top=0.5, width=4.0, height=4.5,
    )

    assert r_without["metric_shapes"] == []

    # Chart area for the no-metrics card should be taller than the
    # with-metrics card by (metrics_row_h + chart_metrics_gap) ≈ 1.0".
    _, _, _, chart_h_with = _bbox(r_with["chart_shape"])
    _, _, _, chart_h_without = _bbox(r_without["chart_shape"])
    assert chart_h_without > chart_h_with
    assert chart_h_without - chart_h_with == pytest.approx(1.0, abs=0.01)


# ---------------------------------------------------------------------------
# 3. Four metrics: cells evenly distributed horizontally
# ---------------------------------------------------------------------------


def test_four_metrics_even_split(one_slide_prs):
    slide = one_slide_prs.slides[0]
    spec = MetricCardSpec(
        label="A", title="B",
        metrics=[MetricEntry(label=f"L{i}", value=str(i)) for i in range(4)],
        padding=0.3,
    )
    r = add_metric_card(slide, spec, left=1.0, top=1.0, width=5.0, height=4.5)

    # Inner width = 5.0 - 2*0.3 = 4.4"  →  each cell = 1.1"
    cell_w = 4.4 / 4
    inner_left = 1.0 + 0.3
    for i, (m_label, _m_value) in enumerate(r["metric_shapes"]):
        expected_x = inner_left + i * cell_w
        actual_l, _, actual_w, _ = _bbox(m_label)
        assert actual_l == pytest.approx(expected_x, abs=0.02)
        assert actual_w == pytest.approx(cell_w, abs=0.02)


# ---------------------------------------------------------------------------
# 4. chart_image_path provided → chart_shape is a Picture
# ---------------------------------------------------------------------------


def test_chart_image_path_uses_picture(one_slide_prs, tmp_path):
    slide = one_slide_prs.slides[0]
    img_path = _make_image(tmp_path)
    spec = MetricCardSpec(label="A", title="B", chart_image_path=img_path)

    r = add_metric_card(slide, spec, left=1.0, top=1.0, width=5.0, height=4.5)
    assert r["chart_shape"].shape_type == MSO_SHAPE_TYPE.PICTURE

    # Without image, the chart shape is a plain AUTO_SHAPE rectangle.
    spec_noimg = MetricCardSpec(label="A", title="B")
    r2 = add_metric_card(
        slide, spec_noimg, left=7.0, top=1.0, width=5.0, height=4.5,
    )
    assert r2["chart_shape"].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE


# ---------------------------------------------------------------------------
# 5. Row of 3 cards: equal widths + correct gaps
# ---------------------------------------------------------------------------


def test_row_of_three_cards(one_slide_prs):
    slide = one_slide_prs.slides[0]
    specs = [
        MetricCardSpec(label=f"K{i}", title=f"T{i}",
                       metrics=[MetricEntry(label="x", value="1")])
        for i in range(3)
    ]
    result = add_metric_card_row(
        slide, specs, left=0.5, top=0.5, width=12.0, height=4.5, gap=0.3,
    )

    assert len(result["cards"]) == 3
    assert result["consumed_height"] == 4.5
    assert result["consumed_width"] == 12.0

    # card_w = (12.0 - 2 * 0.3) / 3 = 3.8"
    expected_w = (12.0 - 2 * 0.3) / 3
    xs = [c["left"] for c in result["cards"]]
    widths = [c["width"] for c in result["cards"]]
    assert widths == [pytest.approx(expected_w, abs=1e-6)] * 3
    # Gap invariant: xs[i+1] - (xs[i] + width) == gap
    for i in range(2):
        assert xs[i + 1] - (xs[i] + widths[i]) == pytest.approx(0.3, abs=1e-6)


# ---------------------------------------------------------------------------
# 6. Row of 3 cards with heterogeneous metrics counts [0, 2, 4]
# ---------------------------------------------------------------------------


def test_row_heterogeneous_metrics_counts(one_slide_prs):
    slide = one_slide_prs.slides[0]
    specs = [
        MetricCardSpec(label="A", title="T0", metrics=[]),
        MetricCardSpec(
            label="B", title="T1",
            metrics=[MetricEntry(label="x", value="1"),
                     MetricEntry(label="y", value="2")],
        ),
        MetricCardSpec(
            label="C", title="T2",
            metrics=[MetricEntry(label=f"L{i}", value=str(i))
                     for i in range(4)],
        ),
    ]
    result = add_metric_card_row(
        slide, specs, left=0.5, top=0.5, width=12.0, height=4.5, gap=0.3,
    )
    assert len(result["cards"]) == 3
    # Each card's bounds dict has the same height regardless of metric count.
    assert all(c["height"] == 4.5 for c in result["cards"])


# ---------------------------------------------------------------------------
# 7. Too-small height raises EngineError(INVALID_PARAMETER)
# ---------------------------------------------------------------------------


def test_too_small_height_raises(one_slide_prs):
    slide = one_slide_prs.slides[0]
    spec = MetricCardSpec(label="A", title="B")

    with pytest.raises(EngineError) as exc_info:
        add_metric_card(
            slide, spec, left=1.0, top=1.0, width=5.0, height=0.5,
        )

    err = exc_info.value
    assert err.code == ErrorCode.INVALID_PARAMETER
    assert "MetricCard height too small" in err.args[0]


# ---------------------------------------------------------------------------
# 8. theme="ir" resolves defaults to IR-palette hex values
# ---------------------------------------------------------------------------


def test_theme_ir_resolves_palette(one_slide_prs):
    slide = one_slide_prs.slides[0]
    spec = MetricCardSpec(
        label="KPI",
        title="Revenue",
        fill_color="background",   # IR token → #F8F9F5
        border_color="rule_subtle",  # IR token → #E0E0E0
    )
    r = add_metric_card(
        slide, spec, left=1.0, top=1.0, width=5.0, height=4.5, theme="ir",
    )
    frame = r["frame_shape"]
    assert str(frame.fill.fore_color.rgb) == "F8F9F5"
    assert str(frame.line.color.rgb) == "E0E0E0"


# ---------------------------------------------------------------------------
# 9. check_containment returns no findings for properly-sized card
# ---------------------------------------------------------------------------


def test_check_containment_clean(one_slide_prs):
    slide = one_slide_prs.slides[0]
    specs = [
        MetricCardSpec(
            label=f"K{i}", title=f"Metric {i}",
            metrics=[
                MetricEntry(label="Now", value="10"),
                MetricEntry(label="Target", value="20"),
            ],
        )
        for i in range(3)
    ]
    add_metric_card_row(
        slide, specs, left=0.5, top=1.0, width=12.0, height=5.0, gap=0.3,
    )

    findings = check_containment(one_slide_prs)
    assert findings == [], (
        "MetricCard children should all be inside their container bounds; "
        f"got: {findings}"
    )
