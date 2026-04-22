"""PageMarker / SlideFooter block primitives のテスト (v0.6.0, #135).

配置座標は ``markers.py`` 冒頭の固定オフセット定数で決まる。テストでは
戻り値 ``bounds`` と実 shape の bbox が一致すること、右テキスト空文字
時の skip 挙動、theme 指定時の色解決、MCP ツールからの往復 (tempfile
での save → 再オープン) を検証する。
"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx_mcp_server.engine.components.markers import (
    PageMarkerSpec,
    SlideFooterSpec,
    add_page_marker,
    add_slide_footer,
)
from pptx_mcp_server.theme import IR


# EMU (914400 per inch) → inches
def _in(emu: int) -> float:
    return emu / 914400


def test_page_marker_renders_two_textboxes_top_right(slide):
    """``add_page_marker`` は 2 つの textbox を右上に作り、bounds を返す."""
    slide_w = 13.333
    slide_h = 7.5
    # color="text_secondary" はテーマトークンなので theme 指定必須。
    spec = PageMarkerSpec(section="FINANCIAL SUMMARY", page="P.05 ／ FY Q3")

    before = len(list(slide.shapes))
    result = add_page_marker(
        slide, spec,
        slide_width=slide_w, slide_height=slide_h,
        theme="ir",
    )
    after = len(list(slide.shapes))

    # 2 つの shape が新規に追加される。
    assert after - before == 2
    assert result["section_shape"] is not None
    assert result["page_shape"] is not None

    # bounds は右端が slide_width - 0.5" に達しており、スライド内に収まる。
    b = result["bounds"]
    assert b["left"] + b["width"] == pytest.approx(slide_w - 0.5, abs=1e-6)
    assert b["left"] + b["width"] <= slide_w
    assert b["top"] >= 0
    assert b["top"] + b["height"] <= slide_h

    # section shape の top は page shape の top より上にある (2 行積み)。
    sec_top = _in(result["section_shape"].top)
    page_top = _in(result["page_shape"].top)
    assert sec_top < page_top
    # 両 shape は同じ left / width を持つ (右端揃え)。
    assert _in(result["section_shape"].left) == pytest.approx(
        _in(result["page_shape"].left), abs=1e-6
    )
    assert _in(result["section_shape"].width) == pytest.approx(
        _in(result["page_shape"].width), abs=1e-6
    )


def test_slide_footer_both_texts(slide):
    """left_text と right_text がともに非空なら 2 shape 作る."""
    slide_w = 13.333
    slide_h = 7.5
    spec = SlideFooterSpec(
        left_text="IR Presentation · FY Q3",
        right_text="Confidential",
    )
    before = len(list(slide.shapes))
    result = add_slide_footer(
        slide, spec,
        slide_width=slide_w, slide_height=slide_h,
        theme="ir",
    )
    after = len(list(slide.shapes))

    assert after - before == 2
    assert result["left_shape"] is not None
    assert result["right_shape"] is not None

    # top は slide_height - 0.4".
    assert _in(result["left_shape"].top) == pytest.approx(slide_h - 0.4, abs=1e-6)
    assert _in(result["right_shape"].top) == pytest.approx(slide_h - 0.4, abs=1e-6)

    # 左は left=0.5", 右は right 端が slide_width - 0.5" に揃う。
    assert _in(result["left_shape"].left) == pytest.approx(0.5, abs=1e-6)
    r_left = _in(result["right_shape"].left)
    r_width = _in(result["right_shape"].width)
    assert r_left + r_width == pytest.approx(slide_w - 0.5, abs=1e-6)


def test_slide_footer_right_text_empty_skipped(slide):
    """right_text="" のとき右テキストボックスは作られない."""
    slide_w = 13.333
    slide_h = 7.5
    spec = SlideFooterSpec(left_text="Only left", right_text="")
    before = len(list(slide.shapes))
    result = add_slide_footer(
        slide, spec,
        slide_width=slide_w, slide_height=slide_h,
        theme="ir",
    )
    after = len(list(slide.shapes))

    # 1 shape だけ追加される。
    assert after - before == 1
    assert result["left_shape"] is not None
    assert result["right_shape"] is None


def test_theme_ir_resolves_color(slide):
    """theme="ir" + color="text_secondary" → IR の text_secondary hex が適用される."""
    slide_w = 13.333
    slide_h = 7.5
    spec = PageMarkerSpec(section="SEC", page="P.01")
    result = add_page_marker(
        slide, spec,
        slide_width=slide_w, slide_height=slide_h,
        theme="ir",
    )
    expected = IR.colors["text_secondary"].lstrip("#")
    run = result["section_shape"].text_frame.paragraphs[0].runs[0]
    assert run.font.color.rgb == RGBColor.from_string(expected)


def test_mcp_tool_page_marker_roundtrip(tmp_path):
    """MCP ツール pptx_add_page_marker → save → 再オープンで shape が残る."""
    # server import は module level で重いため、テスト内で遅延 import する。
    from pptx_mcp_server.server import pptx_add_page_marker, pptx_create

    path = str(tmp_path / "marker.pptx")
    create_resp = json.loads(pptx_create(path))
    assert create_resp["ok"] is True

    # blank slide を追加しておく (pptx_create は slide を作らない)。
    from pptx_mcp_server.server import pptx_add_slide
    add_resp = json.loads(pptx_add_slide(path, layout_index=6))
    assert add_resp["ok"] is True

    resp = json.loads(pptx_add_page_marker(
        path, slide_index=0,
        section="FINANCIAL SUMMARY",
        page="P.05 ／ FY Q3",
        theme="ir",
    ))
    assert resp["ok"] is True, resp
    assert "bounds" in resp["result"]

    # 再オープンして shape が残っていることを確認。
    prs = Presentation(path)
    slide = prs.slides[0]
    # 2 shape 追加されている。
    assert len(list(slide.shapes)) >= 2
