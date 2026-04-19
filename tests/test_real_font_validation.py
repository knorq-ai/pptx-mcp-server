"""Issue #91: ``check_text_overflow`` の real-font path 検証.

ヒューリスティックと実フォント advance width の echo chamber を破壊する
ことが目的。本テストスイートは fontTools + 実フォント (Liberation Sans
/ macOS Arial / Windows Arial / Noto CJK 等) が存在する環境でのみ
実行される (``pytest.skip`` でフォント欠損を正直にスキップする)。

テスト対象:
    1. デフォルト ``font_source='heuristic'`` は既存挙動
    2. 同じ入力で heuristic が "fits" と判定する一方、real path が
       overflow を検出することで echo chamber の存在を証明する
    3. ``font_paths`` で解決できないフォントはヒューリスティックに
       フォールバックし、warning finding を発行する
    4. fontTools 未インストール時に friendly message を出す
    5. ``discover_system_fonts`` が macOS/Linux で ≥1 件返す
"""

from __future__ import annotations

import importlib
import os

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


# fontTools 欠損環境では本ファイル全体をスキップ。
pytest.importorskip("fontTools.ttLib")


from pptx_mcp_server.engine.font_metrics import (
    advance_width_inches,
    discover_system_fonts,
    text_width_inches,
)
from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx_mcp_server.engine.validation import (
    check_deck_extended,
    check_text_overflow,
)


_ARIAL_CANDIDATES = [
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    "/Library/Fonts/Arial.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/Windows/Fonts/arial.ttf",
    "C:\\Windows\\Fonts\\arial.ttf",
]


def _first_existing(paths: list[str]) -> str | None:
    for p in paths:
        if os.path.exists(p):
            return p
    return None


@pytest.fixture(scope="module")
def arial_path() -> str:
    p = _first_existing(_ARIAL_CANDIDATES)
    if p is None:
        pytest.skip(
            "Arial/Liberation Sans TTF not found — skip real-font tests. "
            "On Ubuntu CI: apt install fonts-liberation."
        )
    return p


def _make_deck_with_wide_text(
    font_name: str = "Arial",
    frame_width: float = 2.1,
    frame_height: float = 0.35,
    text: str = "W" * 14,
    pt: float = 12.0,
) -> Presentation:
    """heuristic と real で行数が食い違う textbox を 1 つ含む deck を作る.

    ``W`` 14 連は Arial 実測で usable width 2.0" (frame_width=2.1 - padding
    0.10) を越え、real 経路では 2 行に折り返される。ヒューリスティックの
    ASCII WIDE バケット repr (0.01172/pt) は W を ~10% 過小評価するため、
    1 行で収まる判定となる。frame_height=0.35 は 1 行分 (0.20") には
    余裕があり、2 行分 (0.40") は tolerance 1.05 (=0.3675) を超える。
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.0), Inches(frame_width), Inches(frame_height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)
            run.font.name = font_name
    tb.name = "EchoChamberProbe"
    return prs


# ---------------------------------------------------------------------------
# 1. Default (heuristic) path preserved
# ---------------------------------------------------------------------------


def test_default_heuristic_path_unchanged(arial_path: str) -> None:
    """``font_source`` を渡さない既定動作 = 既存の heuristic 挙動."""
    prs = _make_deck_with_wide_text()
    findings = check_text_overflow(prs)
    # ヒューリスティックでは 1 行で収まるので overflow 無し
    overflow = [f for f in findings if f.category == "text_overflow"]
    assert overflow == [], (
        f"heuristic path should find no overflow for 14*'W' in 2.0\" frame; "
        f"got: {overflow}"
    )


def test_explicit_heuristic_path_equivalent(arial_path: str) -> None:
    """``font_source='heuristic'`` は明示指定でも既定と同じ結果を返す."""
    prs = _make_deck_with_wide_text()
    default = check_text_overflow(prs)
    explicit = check_text_overflow(prs, font_source="heuristic")
    assert [f.to_dict() for f in default] == [f.to_dict() for f in explicit]


# ---------------------------------------------------------------------------
# 2. Real path breaks the echo chamber (KEY TEST)
# ---------------------------------------------------------------------------


def test_real_path_catches_overflow_heuristic_misses(arial_path: str) -> None:
    """heuristic は "fits" と言うが、real-font は overflow を正しく検出する.

    この非対称性が echo chamber の核心。heuristic + auto-fit が同じ
    下位モデルを共有する限り、両者は一致して沈黙する。real-font path
    だけが独立した第三者として食い違いを検出できる。
    """
    prs = _make_deck_with_wide_text()

    heur = [
        f
        for f in check_text_overflow(prs, font_source="heuristic")
        if f.category == "text_overflow"
    ]
    real = [
        f
        for f in check_text_overflow(
            prs, font_source="real", font_paths={"Arial": arial_path}
        )
        if f.category == "text_overflow"
    ]
    assert heur == [], f"heuristic must silently accept; got: {heur}"
    assert len(real) == 1, (
        f"real path must flag overflow (found independently of heuristic); "
        f"got: {real}"
    )
    assert real[0].slide_index == 0
    assert real[0].shape_name == "EchoChamberProbe"


# ---------------------------------------------------------------------------
# 3. Missing font fallback
# ---------------------------------------------------------------------------


def test_missing_font_falls_back_to_heuristic_with_warning(
    arial_path: str,
) -> None:
    """``font_paths`` に font name が存在しない場合はヒューリスティック
    にフォールバックし、``font_not_measured`` warning を発行する."""
    prs = _make_deck_with_wide_text(font_name="NonExistentFont")
    findings = check_text_overflow(
        prs,
        font_source="real",
        font_paths={"SomeOtherFont": arial_path},
    )
    warns = [f for f in findings if f.category == "font_not_measured"]
    assert len(warns) == 1, f"expected 1 warning; got: {warns}"
    assert "NonExistentFont" in warns[0].message
    assert warns[0].severity == "warning"
    # overflow は heuristic fallback の結果 = 検出されない
    overflow = [f for f in findings if f.category == "text_overflow"]
    assert overflow == []


# ---------------------------------------------------------------------------
# 4. fontTools absent → friendly error
# ---------------------------------------------------------------------------


def test_fonttools_missing_raises_friendly_error(
    monkeypatch: pytest.MonkeyPatch, arial_path: str
) -> None:
    """``fontTools`` が import できない環境で ``font_source='real'`` を
    呼ぶと ``EngineError(INVALID_PARAMETER)`` + extras ヒントを返す."""
    prs = _make_deck_with_wide_text()

    import pptx_mcp_server.engine.validation as validation_mod

    # validation._check_fonttools_available が import fontTools.ttLib を
    # 試みる。ここでは builtins.__import__ を差し替え、fontTools を
    # import した瞬間に ImportError を送出する。
    import builtins
    real_import = builtins.__import__

    def _blocking_import(name, *args, **kwargs):
        if name == "fontTools" or name.startswith("fontTools"):
            raise ImportError(f"No module named {name!r} (simulated)")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", _blocking_import)

    with pytest.raises(EngineError) as exc_info:
        check_text_overflow(
            prs, font_source="real", font_paths={"Arial": arial_path}
        )
    assert exc_info.value.code == ErrorCode.INVALID_PARAMETER
    msg = str(exc_info.value)
    assert "fontTools" in msg
    assert "[validation]" in msg


# ---------------------------------------------------------------------------
# 5. System font discovery
# ---------------------------------------------------------------------------


def test_discover_system_fonts_returns_at_least_one() -> None:
    """macOS/Linux 標準環境では ≥1 件 返る. CI で全く見つからない場合は
    skip する (test_library_only 等の barebones image への配慮)."""
    fonts = discover_system_fonts()
    if not fonts:
        pytest.skip("no system fonts detected in this environment")
    # 少なくとも Arial 系 or CJK 系のいずれかが見つかる筈
    assert any(
        k in fonts
        for k in ("Arial", "Liberation Sans", "Yu Gothic", "Noto Sans CJK", "Hiragino Sans", "Meiryo")
    ), f"no expected font key found in: {fonts}"


def test_discover_system_fonts_paths_exist() -> None:
    """返り値のパスはすべて実在する。"""
    for name, path in discover_system_fonts().items():
        assert os.path.exists(path), f"{name}: path does not exist: {path}"


# ---------------------------------------------------------------------------
# advance_width_inches / text_width_inches smoke
# ---------------------------------------------------------------------------


def test_advance_width_inches_linear_in_size(arial_path: str) -> None:
    """advance 幅は size_pt に線形比例する."""
    w10 = advance_width_inches(arial_path, "A", 10.0)
    w20 = advance_width_inches(arial_path, "A", 20.0)
    assert w10 > 0
    assert w20 == pytest.approx(w10 * 2, rel=1e-6)


def test_text_width_inches_sums_chars(arial_path: str) -> None:
    """text_width_inches は per-char advance の合計に等しい."""
    text = "Hello"
    summed = sum(advance_width_inches(arial_path, ch, 12.0) for ch in text)
    assert text_width_inches(arial_path, text, 12.0) == pytest.approx(summed, rel=1e-6)


def test_text_width_inches_ignores_newlines(arial_path: str) -> None:
    """``\\n`` は幅 0 として扱う."""
    w_plain = text_width_inches(arial_path, "ab", 12.0)
    w_nl = text_width_inches(arial_path, "a\nb", 12.0)
    assert w_nl == pytest.approx(w_plain, rel=1e-6)


# ---------------------------------------------------------------------------
# check_deck_extended 統合
# ---------------------------------------------------------------------------


def test_check_deck_extended_surfaces_real_path_findings(arial_path: str) -> None:
    """``check_deck_extended`` 経由でも real-font path の結果が流れ込む."""
    prs = _make_deck_with_wide_text()

    # heuristic では text_overflow = 0 のはず
    heur_report = check_deck_extended(prs)
    assert heur_report["summary"]["errors"] == 0

    # real で呼ぶと text_overflow が 1 件検出され errors に加算される
    real_report = check_deck_extended(
        prs, font_source="real", font_paths={"Arial": arial_path}
    )
    slide_0 = real_report["slides"][0]
    assert len(slide_0["text_overflow"]) == 1, slide_0
    assert real_report["summary"]["errors"] >= 1


def test_check_deck_extended_font_not_measured_warning(arial_path: str) -> None:
    """未解決フォントは ``font_not_measured`` warning として slide-level
    に格納され、summary.warnings に加算される."""
    prs = _make_deck_with_wide_text(font_name="GhostFont")
    report = check_deck_extended(
        prs,
        font_source="real",
        font_paths={"OtherFont": arial_path},
    )
    slide_0 = report["slides"][0]
    assert len(slide_0["font_not_measured"]) == 1, slide_0
    assert report["summary"]["warnings"] >= 1


# ---------------------------------------------------------------------------
# calibration_helpers 後方互換
# ---------------------------------------------------------------------------


def test_calibration_helpers_reexport_still_works(arial_path: str) -> None:
    """``tests.calibration_helpers`` は engine モジュールの re-export を返す."""
    from tests.calibration_helpers import advance_width_inches as legacy_adv

    assert legacy_adv is advance_width_inches
    assert legacy_adv(arial_path, "A", 10.0) > 0
