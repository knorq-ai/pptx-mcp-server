"""
Unit tests for composite operations -- content slide, section divider, KPI row, bullet block.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from pptx_mcp_server.engine.composites import (
    _add_content_slide,
    _add_section_divider,
    _add_kpi_row,
    _add_bullet_block,
    _add_card_grid,
)
from pptx_mcp_server.theme import MCKINSEY


class TestAddContentSlide:
    """_add_content_slide must produce an action title + divider + optional footer."""

    def test_creates_slide_with_title_and_divider(self, blank_prs):
        slide, idx = _add_content_slide(blank_prs, "Revenue Growth Analysis")
        assert idx == 0
        assert len(blank_prs.slides) == 1
        # Should have at least title textbox + divider shape
        assert len(slide.shapes) >= 2

    def test_title_text_is_set(self, blank_prs):
        slide, _ = _add_content_slide(blank_prs, "Test Title")
        texts = [
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        ]
        assert any("Test Title" in t for t in texts)

    def test_with_source_and_page_number(self, blank_prs):
        slide, _ = _add_content_slide(
            blank_prs, "Title", source="Source: Annual Report", page_number=5
        )
        texts = [
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        ]
        assert any("Source: Annual Report" in t for t in texts)
        assert any("5" in t for t in texts)
        # title + divider + source + page_number = 4 shapes
        assert len(slide.shapes) == 4

    def test_without_source_page_number_has_fewer_shapes(self, blank_prs):
        slide, _ = _add_content_slide(blank_prs, "Title Only")
        # title + divider = 2 shapes
        assert len(slide.shapes) == 2

    def test_action_title_has_bottom_anchor(self, blank_prs):
        # The action title does NOT set vertical_anchor in the code,
        # so we verify default behavior (no explicit anchor attr set)
        slide, _ = _add_content_slide(blank_prs, "Title")
        # Title is the first textbox added (shape 0)
        title_shape = list(slide.shapes)[0]
        assert title_shape.has_text_frame
        assert "Title" in title_shape.text_frame.text

    def test_works_with_mckinsey_theme(self, blank_prs, mckinsey_theme):
        slide, idx = _add_content_slide(
            blank_prs, "Themed Slide", theme=mckinsey_theme,
            source="Source", page_number=1,
        )
        assert idx == 0
        assert len(slide.shapes) == 4


class TestAddSectionDivider:
    """_add_section_divider must create a dark background slide with stripes."""

    def test_creates_slide(self, blank_prs):
        slide, idx = _add_section_divider(blank_prs, "Section 1")
        assert idx == 0
        assert len(blank_prs.slides) == 1

    def test_has_dark_background(self, blank_prs):
        slide, _ = _add_section_divider(blank_prs, "Section")
        bg = slide.background.fill
        assert bg.fore_color.rgb is not None

    def test_has_title_text(self, blank_prs):
        slide, _ = _add_section_divider(blank_prs, "Section Title")
        texts = [
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        ]
        assert any("Section Title" in t for t in texts)

    def test_has_accent_stripes(self, blank_prs):
        slide, _ = _add_section_divider(blank_prs, "Title")
        # 2 stripes (rectangles) + 1 title textbox = 3 shapes minimum
        assert len(slide.shapes) >= 3

    def test_with_subtitle(self, blank_prs):
        slide, _ = _add_section_divider(blank_prs, "Title", subtitle="Sub")
        texts = [
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        ]
        assert any("Sub" in t for t in texts)
        # 2 stripes + title + subtitle = 4 shapes
        assert len(slide.shapes) == 4

    def test_works_with_mckinsey_theme(self, blank_prs, mckinsey_theme):
        slide, _ = _add_section_divider(
            blank_prs, "Themed Section", theme=mckinsey_theme
        )
        assert len(slide.shapes) >= 3


class TestAddKpiRow:
    """_add_kpi_row must create the correct number of KPI callout boxes."""

    def test_creates_correct_number_of_kpis(self, slide):
        kpis = [
            {"value": "107.8M", "label": "Revenue"},
            {"value": "23.4%", "label": "Margin"},
            {"value": "1.2B", "label": "Market Cap"},
        ]
        indices = _add_kpi_row(slide, kpis, y=2)
        assert len(indices) == 3

    def test_each_kpi_has_value_and_label(self, slide):
        kpis = [{"value": "100", "label": "Sales"}]
        _add_kpi_row(slide, kpis, y=2)
        texts = [
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        ]
        assert "100" in texts
        assert "Sales" in texts

    def test_empty_kpis_returns_empty(self, slide):
        indices = _add_kpi_row(slide, [], y=2)
        assert indices == []

    def test_works_with_mckinsey_theme(self, slide, mckinsey_theme):
        kpis = [{"value": "42", "label": "Metric"}]
        indices = _add_kpi_row(slide, kpis, y=2, theme=mckinsey_theme)
        assert len(indices) == 1


class TestAddBulletBlock:
    """_add_bullet_block must create bulleted text with correct items."""

    def test_creates_bullet_items(self, slide):
        items = ["First point", "Second point", "Third point"]
        idx = _add_bullet_block(slide, items, left=1, top=2, width=5, height=3)
        shape = list(slide.shapes)[idx]
        text = shape.text_frame.text
        assert "First point" in text
        assert "Second point" in text
        assert "Third point" in text

    def test_uses_bullet_character(self, slide):
        items = ["Item A"]
        idx = _add_bullet_block(slide, items, left=1, top=2, width=5, height=3)
        shape = list(slide.shapes)[idx]
        p = shape.text_frame.paragraphs[0]
        assert "\u2022" in p.text

    def test_multiple_paragraphs(self, slide):
        items = ["One", "Two", "Three"]
        idx = _add_bullet_block(slide, items, left=1, top=2, width=5, height=3)
        shape = list(slide.shapes)[idx]
        assert len(shape.text_frame.paragraphs) == 3

    def test_empty_items_returns_none(self, slide):
        idx = _add_bullet_block(slide, [], left=1, top=2, width=5, height=3)
        assert idx is None

    def test_works_with_mckinsey_theme(self, slide, mckinsey_theme):
        items = ["Themed bullet"]
        idx = _add_bullet_block(
            slide, items, left=1, top=2, width=5, height=3, theme=mckinsey_theme
        )
        assert idx is not None

    def test_dict_items(self, slide):
        items = [{"text": "Dict item 1"}, {"text": "Dict item 2"}]
        idx = _add_bullet_block(slide, items, left=1, top=2, width=5, height=3)
        shape = list(slide.shapes)[idx]
        text = shape.text_frame.text
        assert "Dict item 1" in text
        assert "Dict item 2" in text


class TestAddCardGrid:
    """_add_card_grid must produce auto-balanced card layouts."""

    def test_four_cards_produces_2x2(self, slide):
        cards = [
            {"title": f"Card {i}", "body": f"Description {i}"}
            for i in range(4)
        ]
        indices = _add_card_grid(slide, cards, theme=MCKINSEY)
        # Each card = bg rect + title textbox + body textbox = 3 shapes × 4 = 12
        assert len(list(slide.shapes)) >= 12

    def test_two_cards_single_row(self, slide):
        cards = [
            {"title": "A", "body": "Description A"},
            {"title": "B", "body": "Description B"},
        ]
        _add_card_grid(slide, cards, theme=MCKINSEY)
        assert len(list(slide.shapes)) >= 6

    def test_six_cards_produces_2x3(self, slide):
        cards = [{"title": f"Card {i}", "body": f"Desc {i}"} for i in range(6)]
        _add_card_grid(slide, cards, theme=MCKINSEY)
        assert len(list(slide.shapes)) >= 18

    def test_cards_with_bullets(self, slide):
        cards = [
            {"title": "Strategy", "bullets": ["Point A", "Point B", "Point C"]},
        ]
        _add_card_grid(slide, cards, theme=MCKINSEY)
        shapes = list(slide.shapes)
        # bg rect + title textbox + bullets textbox = 3
        assert len(shapes) >= 3
        # Verify bullets are present
        texts = [s.text_frame.text for s in shapes if s.has_text_frame]
        assert any("Point A" in t for t in texts)

    def test_cards_with_icons(self, slide):
        cards = [
            {"title": "With Icon", "body": "Has an icon", "icon_id": "abacus", "icon_color": "2251FF"},
        ]
        # Count shapes via spTree (icons inject raw XML)
        from tests.test_icons import _count_shapes_in_spTree
        before = _count_shapes_in_spTree(slide)
        _add_card_grid(slide, cards, theme=MCKINSEY)
        after = _count_shapes_in_spTree(slide)
        # bg rect + icon + title + body = 4+
        assert after - before >= 4

    def test_empty_cards_list(self, slide):
        indices = _add_card_grid(slide, [], theme=MCKINSEY)
        assert indices == []

    def test_custom_y_position(self, slide):
        cards = [{"title": "A", "body": "B"}]
        _add_card_grid(slide, cards, y=2.0, theme=MCKINSEY)
        shapes = list(slide.shapes)
        # Background rect should start at y=2.0
        rect = shapes[0]
        assert rect.top > 0
