"""
Unit tests for slide operations -- add, delete, duplicate, background, info, read.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode, save_pptx, open_pptx
from pptx_mcp_server.engine.slides import (
    _add_slide,
    _delete_slide,
    _duplicate_slide,
    _set_slide_background,
    _get_presentation_info,
    _read_slide,
    get_presentation_info,
)
from pptx_mcp_server.engine.shapes import _add_textbox


class TestAddSlide:
    """_add_slide must add slides and enforce layout bounds."""

    def test_adds_slide_returns_index(self, blank_prs):
        idx = _add_slide(blank_prs, layout_index=6)
        assert idx == 0
        assert len(blank_prs.slides) == 1

    def test_second_slide_returns_correct_index(self, blank_prs):
        _add_slide(blank_prs, layout_index=6)
        idx2 = _add_slide(blank_prs, layout_index=6)
        assert idx2 == 1
        assert len(blank_prs.slides) == 2

    def test_invalid_layout_raises_index_out_of_range(self, blank_prs):
        with pytest.raises(EngineError) as exc_info:
            _add_slide(blank_prs, layout_index=999)
        assert exc_info.value.code == ErrorCode.INDEX_OUT_OF_RANGE

    def test_negative_layout_raises_index_out_of_range(self, blank_prs):
        with pytest.raises(EngineError) as exc_info:
            _add_slide(blank_prs, layout_index=-1)
        assert exc_info.value.code == ErrorCode.INDEX_OUT_OF_RANGE


class TestDeleteSlide:
    """_delete_slide must remove slides."""

    def test_deletes_slide(self, one_slide_prs):
        """Verify that _delete_slide removes the slide from the presentation."""
        assert len(one_slide_prs.slides) == 1
        _delete_slide(one_slide_prs, 0)
        assert len(one_slide_prs.slides) == 0

    def test_delete_nonexistent_raises(self, blank_prs):
        with pytest.raises(EngineError) as exc_info:
            _delete_slide(blank_prs, 0)
        assert exc_info.value.code == ErrorCode.SLIDE_NOT_FOUND


class TestDuplicateSlide:
    """_duplicate_slide must copy content without placeholder leak."""

    def test_duplicates_slide(self, one_slide_prs):
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1, 1, 4, 0.5, text="Hello")
        original_shape_count = len(slide.shapes)

        new_idx = _duplicate_slide(one_slide_prs, 0)
        assert new_idx == 1
        assert len(one_slide_prs.slides) == 2

        # Verify shape count matches: the duplicate should not have extra
        # placeholder shapes that leaked from the layout
        dup_slide = one_slide_prs.slides[new_idx]
        assert len(dup_slide.shapes) == original_shape_count

    def test_duplicate_preserves_text(self, one_slide_prs):
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1, 1, 4, 0.5, text="Duplicated Text")

        new_idx = _duplicate_slide(one_slide_prs, 0)
        dup_slide = one_slide_prs.slides[new_idx]
        texts = [s.text_frame.text for s in dup_slide.shapes if s.has_text_frame]
        assert "Duplicated Text" in texts

    def test_duplicate_nonexistent_raises(self, blank_prs):
        with pytest.raises(EngineError) as exc_info:
            _duplicate_slide(blank_prs, 0)
        assert exc_info.value.code == ErrorCode.SLIDE_NOT_FOUND


class TestSetSlideBackground:
    """_set_slide_background must set a solid fill color."""

    def test_sets_background_color(self, one_slide_prs):
        _set_slide_background(one_slide_prs, 0, "#FF0000")
        slide = one_slide_prs.slides[0]
        bg = slide.background
        assert bg.fill.fore_color.rgb is not None

    def test_sets_background_with_theme_token(self, one_slide_prs, mckinsey_theme):
        _set_slide_background(one_slide_prs, 0, "primary", theme=mckinsey_theme)
        slide = one_slide_prs.slides[0]
        from pptx.dml.color import RGBColor
        assert slide.background.fill.fore_color.rgb == RGBColor(0x05, 0x1C, 0x2C)


class TestGetPresentationInfo:
    """_get_presentation_info must return formatted overview."""

    def test_returns_slide_count(self, one_slide_prs):
        info = _get_presentation_info(one_slide_prs)
        assert "Slides: 1" in info

    def test_returns_dimensions(self, one_slide_prs):
        info = _get_presentation_info(one_slide_prs)
        assert "13.333" in info
        assert "7.500" in info

    def test_shows_shape_info(self, one_slide_prs):
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 1, 1, 4, 0.5, text="Title text")
        info = _get_presentation_info(one_slide_prs)
        assert "Title text" in info

    def test_empty_presentation(self, blank_prs):
        info = _get_presentation_info(blank_prs)
        assert "Slides: 0" in info


class TestReadSlide:
    """_read_slide must return shape details."""

    def test_returns_shape_details(self, one_slide_prs):
        slide = one_slide_prs.slides[0]
        _add_textbox(slide, 2, 3, 4, 0.5, text="ReadMe")
        result = _read_slide(one_slide_prs, 0)
        assert "ReadMe" in result
        assert "[0]" in result

    def test_shows_table_info(self, one_slide_prs):
        from pptx_mcp_server.engine.tables import _add_table
        slide = one_slide_prs.slides[0]
        _add_table(slide, [["H1", "H2"], ["a", "b"]], 1, 1, 5)
        result = _read_slide(one_slide_prs, 0)
        assert "table" in result.lower()


class TestFileBased:
    """File-based wrappers: get_presentation_info includes file path."""

    def test_get_presentation_info_includes_path(self, pptx_file):
        info = get_presentation_info(pptx_file)
        assert pptx_file in info
        assert "Slides: 1" in info

    def test_file_based_add_slide(self, pptx_file):
        from pptx_mcp_server.engine.slides import add_slide
        result = add_slide(pptx_file, layout_index=6)
        assert "Added slide" in result
        prs = open_pptx(pptx_file)
        assert len(prs.slides) == 2
