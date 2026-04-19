"""
Tests for the chart engine (Phase 1A).
"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_mcp_server.engine.charts import (
    _add_chart,
    _validate_chart_data,
    add_chart,
    _CHART_TYPE_MAP,
)
from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx_mcp_server.engine.composites import _build_slide, build_slide
from pptx_mcp_server.theme import (
    MCKINSEY,
    DELOITTE,
    NEUTRAL,
    get_theme,
    list_themes,
    resolve_color,
    tint_color,
    shade_color,
    get_chart_color,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

SAMPLE_CATEGORIES = ["Q1", "Q2", "Q3", "Q4"]
SAMPLE_SERIES = [
    {"name": "Revenue", "values": [100, 200, 150, 300], "color": "2251FF"},
    {"name": "Cost", "values": [80, 120, 100, 180], "color": "C62828"},
]


# ---------------------------------------------------------------------------
# Theme registry tests
# ---------------------------------------------------------------------------


class TestThemeRegistry:
    def test_list_themes(self):
        names = list_themes()
        assert "mckinsey" in names
        assert "deloitte" in names
        assert "neutral" in names

    def test_get_theme(self):
        theme = get_theme("mckinsey")
        assert theme is MCKINSEY
        assert theme.name == "mckinsey"

    def test_get_theme_deloitte(self):
        theme = get_theme("deloitte")
        assert theme is DELOITTE
        assert theme.colors["primary"] == "#002776"

    def test_get_theme_neutral(self):
        theme = get_theme("neutral")
        assert theme is NEUTRAL

    def test_get_theme_unknown(self):
        assert get_theme("nonexistent") is None

    def test_theme_has_chart_colors(self):
        assert len(MCKINSEY.chart_colors) >= 6
        assert len(DELOITTE.chart_colors) >= 6
        assert len(NEUTRAL.chart_colors) >= 6

    def test_resolve_series_color(self):
        color = resolve_color(MCKINSEY, "series_0")
        assert color == MCKINSEY.chart_colors[0]

    def test_resolve_series_color_wraps(self):
        n = len(MCKINSEY.chart_colors)
        color = resolve_color(MCKINSEY, f"series_{n}")
        assert color == MCKINSEY.chart_colors[0]

    def test_tint_color(self):
        result = tint_color("#000000", 0.5)
        assert result == "#7F7F7F" or result == "#808080"  # rounding

    def test_tint_color_white(self):
        result = tint_color("#000000", 1.0)
        assert result == "#FFFFFF"

    def test_shade_color(self):
        result = shade_color("#FFFFFF", 0.5)
        assert result == "#7F7F7F" or result == "#808080"

    def test_shade_color_black(self):
        result = shade_color("#FFFFFF", 1.0)
        assert result == "#000000"

    def test_get_chart_color(self):
        color = get_chart_color(MCKINSEY, 0)
        assert color == MCKINSEY.chart_colors[0]

    def test_get_chart_color_cycles(self):
        n = len(MCKINSEY.chart_colors)
        assert get_chart_color(MCKINSEY, n) == get_chart_color(MCKINSEY, 0)


# ---------------------------------------------------------------------------
# Validation tests
# ---------------------------------------------------------------------------


class TestChartValidation:
    def test_invalid_chart_type(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("invalid_type", ["A"], [{"name": "S", "values": [1]}])
        assert exc_info.value.code == ErrorCode.CHART_ERROR
        assert "Unknown chart_type" in str(exc_info.value)

    def test_fuzzy_suggestion(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("colum", ["A"], [{"name": "S", "values": [1]}])
        assert "column" in str(exc_info.value)

    def test_empty_categories(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("column", [], [{"name": "S", "values": []}])
        assert "categories must not be empty" in str(exc_info.value)

    def test_empty_series(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("column", ["A"], [])
        assert "series must not be empty" in str(exc_info.value)

    def test_missing_values(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("column", ["A"], [{"name": "S"}])
        assert "missing 'values'" in str(exc_info.value)

    def test_length_mismatch(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("column", ["A", "B"], [{"name": "S", "values": [1]}])
        assert "1 values but 2 categories" in str(exc_info.value)

    def test_non_numeric_value(self):
        with pytest.raises(EngineError) as exc_info:
            _validate_chart_data("column", ["A"], [{"name": "S", "values": ["abc"]}])
        assert "not numeric" in str(exc_info.value)

    def test_none_values_allowed(self):
        # None = gap in chart data, should not raise
        _validate_chart_data("column", ["A", "B"], [{"name": "S", "values": [1, None]}])

    def test_valid_data_passes(self):
        _validate_chart_data("column", SAMPLE_CATEGORIES, SAMPLE_SERIES)


# ---------------------------------------------------------------------------
# In-memory chart creation tests
# ---------------------------------------------------------------------------


class TestAddChart:
    def test_column_chart(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A", "B", "C"],
            series=[{"name": "S1", "values": [10, 20, 30]}],
        )
        assert idx >= 0
        shape = list(slide.shapes)[idx]
        assert shape.has_chart

    def test_bar_chart(self, slide):
        idx = _add_chart(
            slide, "bar", 1, 1, 8, 5,
            categories=["X", "Y"],
            series=[{"name": "S", "values": [5, 10]}],
        )
        shape = list(slide.shapes)[idx]
        assert shape.has_chart

    def test_stacked_column(self, slide):
        idx = _add_chart(
            slide, "stacked_column", 1, 1, 8, 5,
            categories=SAMPLE_CATEGORIES,
            series=SAMPLE_SERIES,
        )
        shape = list(slide.shapes)[idx]
        chart = shape.chart
        assert len(list(chart.plots[0].series)) == 2

    def test_line_chart(self, slide):
        idx = _add_chart(
            slide, "line", 1, 1, 8, 5,
            categories=["Jan", "Feb", "Mar"],
            series=[{"name": "Trend", "values": [10, 15, 12]}],
        )
        shape = list(slide.shapes)[idx]
        assert shape.has_chart

    def test_pie_chart(self, slide):
        idx = _add_chart(
            slide, "pie", 3, 2, 5, 5,
            categories=["A", "B", "C"],
            series=[{"name": "Share", "values": [40, 35, 25]}],
        )
        shape = list(slide.shapes)[idx]
        assert shape.has_chart

    def test_all_chart_types(self, one_slide_prs):
        """Verify every mapped chart type creates successfully."""
        slide = one_slide_prs.slides[0]
        for chart_type in _CHART_TYPE_MAP:
            idx = _add_chart(
                slide, chart_type, 0.5, 0.5, 4, 3,
                categories=["A", "B"],
                series=[{"name": "S", "values": [1, 2]}],
                legend_position=None,
            )
            assert idx >= 0


class TestChartFormatting:
    def test_series_colors(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A", "B"],
            series=[
                {"name": "S1", "values": [10, 20], "color": "FF0000"},
                {"name": "S2", "values": [15, 25], "color": "00FF00"},
            ],
        )
        chart = list(slide.shapes)[idx].chart
        # Verify series exist
        assert len(list(chart.plots[0].series)) == 2

    def test_theme_auto_colors(self, slide):
        """Series without explicit color get theme colors."""
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[
                {"name": "S1", "values": [10]},
                {"name": "S2", "values": [20]},
            ],
            theme=MCKINSEY,
        )
        chart = list(slide.shapes)[idx].chart
        assert len(list(chart.plots[0].series)) == 2

    def test_data_labels(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A", "B"],
            series=[{"name": "S", "values": [10, 20]}],
            data_labels_show=True,
            data_labels_position="outside_end",
            data_labels_number_format="#,##0",
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.plots[0].has_data_labels

    def test_legend_bottom(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
            legend_position="bottom",
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.has_legend

    def test_legend_hidden(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
            legend_position=None,
        )
        chart = list(slide.shapes)[idx].chart
        assert not chart.has_legend

    def test_chart_title(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
            title="My Chart",
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.has_title

    def test_no_title(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
        )
        chart = list(slide.shapes)[idx].chart
        assert not chart.has_title

    def test_axis_value_title(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
            axis_value_title="Revenue (M)",
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.value_axis.has_title

    def test_axis_scale(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
            axis_value_min=0,
            axis_value_max=100,
            axis_value_major_unit=20,
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.value_axis.minimum_scale == 0
        assert chart.value_axis.maximum_scale == 100

    def test_gap_width(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A"],
            series=[{"name": "S", "values": [10]}],
            gap_width=50,
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.plots[0].gap_width == 50

    def test_pie_no_axes(self, slide):
        """Pie charts should not attempt axis formatting."""
        idx = _add_chart(
            slide, "pie", 1, 1, 5, 5,
            categories=["A", "B"],
            series=[{"name": "S", "values": [60, 40]}],
            axis_value_title="Should be ignored",
        )
        # Should not raise

    def test_japanese_text(self, slide):
        """Verify Japanese categories and series names work."""
        idx = _add_chart(
            slide, "stacked_column", 1, 1, 8, 5,
            categories=["2005年", "2015年", "2025年"],
            series=[
                {"name": "0-14歳", "values": [17.6, 15.9, 14.1]},
                {"name": "15-64歳", "values": [84.4, 77.3, 73.7]},
                {"name": "65歳以上", "values": [25.7, 33.9, 39.2]},
            ],
            title="日本の人口構成推移",
            data_labels_show=True,
            data_labels_position="center",
            axis_value_title="人口（百万人）",
            legend_position="bottom",
        )
        chart = list(slide.shapes)[idx].chart
        assert chart.has_title
        assert len(list(chart.plots[0].series)) == 3


# ---------------------------------------------------------------------------
# Deloitte / Neutral theme tests
# ---------------------------------------------------------------------------


class TestThemeIntegration:
    def test_chart_with_deloitte_theme(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A", "B"],
            series=[{"name": "S", "values": [10, 20]}],
            theme=DELOITTE,
        )
        assert idx >= 0

    def test_chart_with_neutral_theme(self, slide):
        idx = _add_chart(
            slide, "column", 1, 1, 8, 5,
            categories=["A", "B"],
            series=[{"name": "S", "values": [10, 20]}],
            theme=NEUTRAL,
        )
        assert idx >= 0


# ---------------------------------------------------------------------------
# File-based wrapper tests
# ---------------------------------------------------------------------------


class TestAddChartFileWrapper:
    def test_add_chart_to_file(self, pptx_file):
        result = add_chart(pptx_file, 0, {
            "chart_type": "column",
            "categories": ["A", "B"],
            "series": [{"name": "S", "values": [10, 20]}],
        })
        assert "column chart" in result
        assert "slide [0]" in result

    def test_add_chart_with_theme(self, pptx_file):
        result = add_chart(pptx_file, 0, {
            "chart_type": "stacked_bar",
            "categories": ["X", "Y"],
            "series": [{"name": "S", "values": [5, 10]}],
            "theme": "deloitte",
        })
        assert "stacked_bar chart" in result


# ---------------------------------------------------------------------------
# build_slide integration tests
# ---------------------------------------------------------------------------


class TestBuildSlideChart:
    def test_chart_element_in_build_slide(self, blank_prs):
        spec = {
            "layout": "content",
            "title": "Chart Test",
            "elements": [
                {
                    "type": "chart",
                    "chart_type": "column",
                    "left": 0.9, "top": 1.15,
                    "width": 11.5, "height": 5.0,
                    "categories": ["A", "B", "C"],
                    "series": [{"name": "S", "values": [10, 20, 30]}],
                    "data_labels_show": True,
                }
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        # Should have title textbox + divider + chart = at least 3 shapes
        assert len(list(slide.shapes)) >= 3

    def test_chart_element_with_theme_in_spec(self, blank_prs):
        spec = {
            "layout": "blank",
            "theme": "deloitte",
            "elements": [
                {
                    "type": "chart",
                    "chart_type": "pie",
                    "left": 3, "top": 1,
                    "width": 5, "height": 5,
                    "categories": ["A", "B"],
                    "series": [{"name": "S", "values": [60, 40]}],
                }
            ],
        }
        slide, idx = _build_slide(blank_prs, spec)
        # Just verify no crash — theme applied via spec key
        assert len(list(slide.shapes)) >= 1

    def test_unknown_element_type_raises(self, blank_prs):
        spec = {
            "layout": "blank",
            "elements": [
                {"type": "nonexistent", "left": 1, "top": 1}
            ],
        }
        with pytest.raises(EngineError) as exc_info:
            _build_slide(blank_prs, spec)
        assert "Unknown element type" in str(exc_info.value)

    def test_build_slide_file_with_chart(self, pptx_file):
        spec = json.dumps({
            "layout": "content",
            "title": "Revenue Breakdown",
            "elements": [
                {
                    "type": "chart",
                    "chart_type": "stacked_column",
                    "left": 0.9, "top": 1.15,
                    "width": 11.5, "height": 5.0,
                    "categories": ["Q1", "Q2", "Q3"],
                    "series": [
                        {"name": "Product A", "values": [10, 15, 20], "color": "2251FF"},
                        {"name": "Product B", "values": [5, 8, 12], "color": "2E7D32"},
                    ],
                    "data_labels_show": True,
                    "data_labels_position": "center",
                    "legend_position": "bottom",
                }
            ],
        })
        result = build_slide(pptx_file, spec)
        assert "Built slide" in result
        assert "1 elements" in result


# ---------------------------------------------------------------------------
# Roundtrip test
# ---------------------------------------------------------------------------


class TestChartRoundtrip:
    def test_save_and_reopen(self, tmp_path):
        """Create a chart, save, reopen — verify chart is intact."""
        from pptx_mcp_server.engine.pptx_io import create_presentation, open_pptx
        from pptx_mcp_server.engine.slides import add_slide

        path = str(tmp_path / "chart_test.pptx")
        create_presentation(path)
        add_slide(path)  # add a blank slide

        result = add_chart(path, 0, {
            "chart_type": "stacked_column",
            "left": 1, "top": 1, "width": 8, "height": 5,
            "categories": ["2020", "2021", "2022"],
            "series": [
                {"name": "Revenue", "values": [100, 150, 200]},
                {"name": "Cost", "values": [80, 100, 130]},
            ],
            "data_labels_show": True,
            "legend_position": "bottom",
            "axis_value_title": "Amount",
        })

        # Reopen
        prs = open_pptx(path)
        slide = prs.slides[0]
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 1

        chart = chart_shapes[0].chart
        assert len(list(chart.plots[0].series)) == 2
        assert chart.has_legend
