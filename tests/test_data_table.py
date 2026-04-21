"""データテーブル (textbox grid) プリミティブのテスト.

``add_data_table`` の shape 数・alt-row / highlight 背景・alignment・
入力検証・空行時の consumed_height を検証する。Issue #110。
"""

from __future__ import annotations

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode
from pptx_mcp_server.engine.tables_grid import (
    TableColumnSpec,
    add_data_table,
)


def _count_auto_shapes(slide) -> int:
    return sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)


def _count_textboxes(slide) -> int:
    return sum(
        1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    )


def test_basic_3x4_grid_shape_count(slide):
    """基本: 4 列 × 3 行 (+ ヘッダー行) のテーブルが期待通りの shape 数を生成する.

    期待内訳:
        - 4 ヘッダー textbox
        - 12 データセル textbox (3 行 × 4 列)
        - 4 罫線 (ヘッダー下 1 + データ行下 3)、alt-row・highlight は無効
    合計 textbox 16、罫線 (auto_shape) 4。
    """
    columns = [
        TableColumnSpec(header="Ticker", align="left", width=1.0),
        TableColumnSpec(header="Name", align="left", width=2.0),
        TableColumnSpec(header="Price", align="right", width=1.0),
        TableColumnSpec(header="YoY", align="right", width=1.0),
    ]
    rows = [
        ["AAPL", "Apple Inc.", 189.2, "+12.3%"],
        ["MSFT", "Microsoft", 412.0, "+8.1%"],
        ["GOOG", "Alphabet", 178.5, "+15.4%"],
    ]

    result = add_data_table(
        slide, rows, columns,
        left=0.5, top=1.0, width=6.0,
        alt_row_color=None,
        highlight_row_index=None,
    )

    assert _count_textboxes(slide) == 16
    # 罫線: ヘッダー下 + 3 データ行下
    assert _count_auto_shapes(slide) == 4
    assert result["shape_count"] == 20
    assert result["header_y_bottom"] == pytest.approx(1.4, abs=1e-6)
    assert result["consumed_height"] == pytest.approx(0.4 + 3 * 0.35, abs=1e-6)


def test_alt_row_color_adds_background_rectangles(slide):
    """alt_row_color 指定で、奇数 index 行の背景シェーディングが追加される.

    4 行データ → index 1, 3 の 2 行がシェーディング。既定罫線 (4 + 1) も
    合わせて auto_shape 数を検証する。
    """
    columns = [
        TableColumnSpec(header="A", width=1.0),
        TableColumnSpec(header="B", width=1.0),
    ]
    rows = [
        ["a0", "b0"],
        ["a1", "b1"],
        ["a2", "b2"],
        ["a3", "b3"],
    ]
    add_data_table(
        slide, rows, columns,
        left=0.5, top=1.0, width=2.0,
        alt_row_color="F8F8F5",
    )

    # alt-row (2) + 罫線 (1 ヘッダー + 4 行下) = 7 auto_shape
    assert _count_auto_shapes(slide) == 7
    # textbox: 2 ヘッダー + 8 データ = 10
    assert _count_textboxes(slide) == 10


def test_highlight_row_adds_one_extra_rectangle(slide):
    """highlight_row_index=2 で 1 行分のハイライト矩形が追加される.

    alt-row 無効時、ハイライトの背景矩形 1 枚 + 罫線のみ。
    """
    columns = [TableColumnSpec(header="X", width=1.0)]
    rows = [["v0"], ["v1"], ["v2"], ["v3"]]
    add_data_table(
        slide, rows, columns,
        left=0.5, top=1.0, width=1.0,
        highlight_row_index=2,
    )

    # highlight (1) + 罫線 (1 ヘッダー + 4 行下) = 6 auto_shape
    assert _count_auto_shapes(slide) == 6


def test_numeric_column_right_aligned(slide):
    """右揃え指定の列はヘッダー・データとも paragraph alignment が RIGHT になる."""
    columns = [
        TableColumnSpec(header="Ticker", align="left", width=1.0),
        TableColumnSpec(header="Price", align="right", width=1.0),
    ]
    rows = [["AAPL", 189.2]]
    add_data_table(
        slide, rows, columns,
        left=0.5, top=1.0, width=2.0,
        alt_row_color=None,
        rule_color=None,  # 罫線を無効化して shape 数を絞る
    )

    # 罫線を無効化したので textbox 4 枚のみ:
    #   [header0-left, header1-right, cell00-left, cell01-right]
    textboxes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    assert len(textboxes) == 4

    # 右揃え: index 1 (header), 3 (data cell)
    assert textboxes[1].text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    assert textboxes[3].text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    # 左揃え: index 0, 2
    assert textboxes[0].text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT
    assert textboxes[2].text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT


def test_mismatched_row_length_raises(slide):
    """行長が列数と一致しない場合 INVALID_PARAMETER を投げる."""
    columns = [
        TableColumnSpec(header="A", width=1.0),
        TableColumnSpec(header="B", width=1.0),
    ]
    rows = [["only-one-value"]]  # 列数 2 に対して 1 要素のみ
    with pytest.raises(EngineError) as ei:
        add_data_table(
            slide, rows, columns,
            left=0.5, top=1.0, width=2.0,
        )
    assert ei.value.code == ErrorCode.INVALID_PARAMETER


def test_column_widths_proportionally_scaled(slide):
    """列幅合計が ``width`` と一致しない場合、比例スケーリングで強制一致する.

    列幅 [2, 2, 2] (合計 6)、``width=3`` を指定した場合、各列幅は 1.0 に
    スケールされる。ヘッダー textbox の left 座標で検証する。
    """
    columns = [
        TableColumnSpec(header="A", width=2.0),
        TableColumnSpec(header="B", width=2.0),
        TableColumnSpec(header="C", width=2.0),
    ]
    rows: list[list[str]] = []
    add_data_table(
        slide, rows, columns,
        left=0.5, top=1.0, width=3.0,
        rule_color=None,
    )

    # ヘッダー textbox 3 枚、left 座標を確認
    textboxes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    assert len(textboxes) == 3
    EMU = 914400
    lefts = [tb.left / EMU for tb in textboxes]
    # 0.5, 1.5, 2.5 と等間隔 (1.0" 幅) に並ぶ
    assert lefts[0] == pytest.approx(0.5, abs=1e-3)
    assert lefts[1] == pytest.approx(1.5, abs=1e-3)
    assert lefts[2] == pytest.approx(2.5, abs=1e-3)


def test_empty_rows_returns_header_height_plus_rule(slide):
    """rows が空のとき consumed_height は ``header_height + rule_thickness`` となる.

    header_rule が有効 (既定) の場合、ヘッダー下罫線の厚みを consumed に
    算入する契約。これにより caller は次要素 y を ``top + consumed_height``
    として安全に計算できる。
    """
    columns = [TableColumnSpec(header="H", width=1.0)]
    result = add_data_table(
        slide, [], columns,
        left=0.5, top=1.0, width=1.0,
        header_height=0.4,
        rule_thickness=0.01,
    )
    # 0.4 (header) + 0.01 (header rule) = 0.41
    assert result["consumed_height"] == pytest.approx(0.41, abs=1e-6)
    assert result["header_y_bottom"] == pytest.approx(1.4, abs=1e-6)


def test_highlight_row_out_of_range_raises(slide):
    """highlight_row_index が rows の範囲外なら INVALID_PARAMETER."""
    columns = [TableColumnSpec(header="X", width=1.0)]
    with pytest.raises(EngineError) as ei:
        add_data_table(
            slide, [["a"], ["b"]], columns,
            left=0.5, top=1.0, width=1.0,
            highlight_row_index=5,
        )
    assert ei.value.code == ErrorCode.INVALID_PARAMETER


# ---------------------------------------------------------------------------
# Issue #121 — rule_thickness validation
# ---------------------------------------------------------------------------


def test_negative_rule_thickness_rejected(one_slide_prs):
    """rule_thickness < 0 must raise INVALID_PARAMETER (#121).

    Pre-fix: produced negative-height rule rectangles silently.
    """
    slide = one_slide_prs.slides[0]
    columns = [TableColumnSpec(header="x")]
    with pytest.raises(EngineError) as ei:
        add_data_table(
            slide, [["a"]], columns,
            left=0.5, top=1.0, width=1.0,
            rule_thickness=-0.01,
        )
    assert ei.value.code == ErrorCode.INVALID_PARAMETER
    assert "rule_thickness" in str(ei.value)


def test_zero_rule_thickness_ok(one_slide_prs):
    """rule_thickness=0 is valid (means no rules rendered)."""
    slide = one_slide_prs.slides[0]
    columns = [TableColumnSpec(header="x")]
    # Should not raise
    result = add_data_table(
        slide, [["a"]], columns,
        left=0.5, top=1.0, width=1.0,
        rule_thickness=0.0,
    )
    assert result["consumed_height"] > 0
