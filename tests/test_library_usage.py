"""
Library-consumer smoke tests for pptx_mcp_server.

These tests exercise the package purely as a Python library — importing engine
and theme modules directly, composing a small deck, and persisting it to disk.
They also include an AST-based guardrail that asserts the ``engine`` subpackage
never imports the ``mcp`` top-level package, which would couple the library
runtime to the MCP SDK.
"""

from __future__ import annotations

import ast
import importlib
import pkgutil
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_mcp_server.engine.pptx_io import create_presentation, open_pptx, save_pptx
from pptx_mcp_server.engine.shapes import add_textbox
from pptx_mcp_server.engine.slides import add_slide
from pptx_mcp_server.theme import MCKINSEY, Theme


def test_theme_mckinsey_is_theme_instance() -> None:
    """MCKINSEY must be importable from the top-level package as a Theme."""
    assert isinstance(MCKINSEY, Theme)
    assert MCKINSEY.name == "mckinsey"


def test_compose_mini_deck_via_library(tmp_path: Path) -> None:
    """Build a tiny deck using only engine + theme, save, reopen, assert state."""
    pptx_path = tmp_path / "library_usage.pptx"

    # Step 1: create a new presentation file on disk.
    create_presentation(str(pptx_path), width_inches=13.333, height_inches=7.5)

    # Step 2: add a blank slide via the file-based engine API.
    add_slide(str(pptx_path), layout_index=6)

    # Step 3: add a textbox on the new slide using engine.shapes.add_textbox.
    sample_text = "Library consumer smoke test"
    add_textbox(
        str(pptx_path),
        slide_index=0,
        left=1.0,
        top=1.0,
        width=6.0,
        height=1.0,
        text=sample_text,
        font_size=18,
        bold=True,
    )

    # Step 4: reopen and verify shape count + text round-trips.
    prs = open_pptx(str(pptx_path))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    shapes = list(slide.shapes)
    assert len(shapes) == 1
    shape = shapes[0]
    assert shape.has_text_frame
    assert shape.text_frame.text == sample_text

    # Step 5: sanity-check that save_pptx works on the reopened Presentation too.
    out_path = tmp_path / "library_usage_resaved.pptx"
    save_pptx(prs, str(out_path))
    assert out_path.exists() and out_path.stat().st_size > 0


def _iter_engine_source_files() -> list[Path]:
    """Return every .py file under the pptx_mcp_server.engine package."""
    engine_pkg = importlib.import_module("pptx_mcp_server.engine")
    files: list[Path] = []
    for module_info in pkgutil.walk_packages(
        engine_pkg.__path__, prefix="pptx_mcp_server.engine."
    ):
        spec = importlib.util.find_spec(module_info.name)
        if spec and spec.origin and spec.origin.endswith(".py"):
            files.append(Path(spec.origin))
    # Include the package __init__ itself.
    init_path = Path(engine_pkg.__file__) if engine_pkg.__file__ else None
    if init_path is not None and init_path not in files:
        files.append(init_path)
    return files


def _top_level_module(name: str) -> str:
    """Return the top-level module name (e.g. 'mcp' from 'mcp.server.fastmcp')."""
    return name.split(".", 1)[0]


def test_engine_has_no_mcp_imports() -> None:
    """AST guardrail: no module in engine/* may import the 'mcp' top-level package.

    This keeps the library usable without the MCP SDK installed. If a future
    change pulls MCP symbols into the engine layer, this test fails loudly so
    the regression is caught at CI time.
    """
    offenders: list[tuple[str, int, str]] = []
    for path in _iter_engine_source_files():
        tree = ast.parse(path.read_text(), filename=str(path))
        for node in ast.walk(tree):
            if isinstance(node, ast.Import):
                for alias in node.names:
                    if _top_level_module(alias.name) == "mcp":
                        offenders.append((str(path), node.lineno, f"import {alias.name}"))
            elif isinstance(node, ast.ImportFrom):
                # Skip relative imports (module=None when level>0).
                if node.level and not node.module:
                    continue
                mod = node.module or ""
                if node.level == 0 and _top_level_module(mod) == "mcp":
                    offenders.append(
                        (str(path), node.lineno, f"from {mod} import ...")
                    )
    assert not offenders, (
        "engine package must not import the 'mcp' top-level package; found: "
        + ", ".join(f"{p}:{ln} ({stmt})" for p, ln, stmt in offenders)
    )


@pytest.mark.parametrize(
    "dotted",
    [
        "pptx_mcp_server.engine.shapes",
        "pptx_mcp_server.engine.pptx_io",
        "pptx_mcp_server.engine.slides",
        "pptx_mcp_server.engine.tables",
        "pptx_mcp_server.engine.composites",
        "pptx_mcp_server.theme",
    ],
)
def test_engine_modules_import_cleanly(dotted: str) -> None:
    """Each engine/theme module must import without any MCP SDK machinery."""
    mod = importlib.import_module(dotted)
    assert mod is not None
