"""End-to-end integration tests for v0.6.0 block components.

Unlike the per-component unit tests (``test_section_header.py``,
``test_kpi_row.py``, ``test_metric_card.py``, ``test_numbered_list.py``,
``test_markers.py``), this module exercises **combinations** of all five
block components on a single slide and validates with
``check_containment``. The intent is to catch interactions that only
surface when the components are composed together:

- Do the components' ``consumed_height`` / ``bounds`` return values
  compose cleanly (no overlap / no phantom gaps)?
- Does every component respect its declared ``begin_container`` bounds
  when sibling components are present on the same slide?
- Does theme resolution stay consistent across all five when a single
  theme name (``"ir"``) is threaded through?
- Does the validator surface real bugs (not just happy-path green)?

All tests build their own ``Presentation`` so they're self-contained and
the module-level container registry stays clean between tests.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_mcp_server.engine.components.container import (
    begin_container,
)
from pptx_mcp_server.engine.components.section_header import (
    SectionHeaderSpec,
    add_section_header,
)
from pptx_mcp_server.engine.components.markers import (
    PageMarkerSpec,
    SlideFooterSpec,
    add_page_marker,
    add_slide_footer,
)
from pptx_mcp_server.engine.components.kpi_row import (
    KPISpec,
    add_kpi_row,
)
from pptx_mcp_server.engine.components.numbered_list import (
    NumberedItem,
    add_numbered_list,
)
from pptx_mcp_server.engine.components.metric_card import (
    MetricCardSpec,
    MetricEntry,
    add_metric_card,
    add_metric_card_row,
)
from pptx_mcp_server.engine.validation import check_containment
from pptx_mcp_server.theme import resolve_theme_color


# Slide dimensions used consistently across these integration tests.
# McKinsey-style 16:9 (13.333 x 7.5). Matches the ``blank_prs`` conftest
# fixture so results compare against the existing unit tests.
_SLIDE_W = 13.333
_SLIDE_H = 7.5


# Container registry isolation is provided globally by the autouse
# ``_isolated_container_registry`` fixture in ``tests/conftest.py`` (Task G).
# ``check_containment`` normally clears registry entries consumed-on-use,
# but tests that set up state without calling it (or that assert on
# findings mid-test) still need explicit isolation to avoid cross-test bleed.


def _new_prs_with_blank_slide():
    """Build a fresh 16:9 Presentation with one blank slide.

    Each test creates its own so the module-level container registry is
    keyed by a fresh ``id(slide)`` — avoiding id reuse across tests.
    """
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    return prs


def _first_run_hex(shape) -> str:
    """Return the first run's font color as an uppercase 6-hex string."""
    run = shape.text_frame.paragraphs[0].runs[0]
    return str(run.font.color.rgb).upper()


# ---------------------------------------------------------------------------
# Test 1: All 5 block components on a single IR-themed slide
# ---------------------------------------------------------------------------


def test_ir_executive_summary_slide_all_components():
    """Compose page_marker + section_header + kpi_row + metric_card_row +
    slide_footer on one slide and verify containment is clean end-to-end.

    This is the canonical "realistic IR exec-summary" layout and the most
    important regression test for v0.6.0: if any two components interact
    badly (e.g. one's declared bounds overlap another's shapes), this
    test surfaces it via a containment finding.
    """
    prs = _new_prs_with_blank_slide()
    slide = prs.slides[0]
    before = len(list(slide.shapes))

    # 1) Page marker (top-right, fixed offsets).
    add_page_marker(
        slide,
        PageMarkerSpec(section="FINANCIAL SUMMARY", page="P.01 / FY Q3"),
        slide_width=_SLIDE_W,
        slide_height=_SLIDE_H,
        theme="ir",
    )

    # 2) Section header near top-left, below page marker row.
    header_left = 0.5
    header_top = 0.9
    header_width = _SLIDE_W - 1.0
    header_result = add_section_header(
        slide,
        SectionHeaderSpec(
            title="FY Q3 Executive Summary",
            subtitle="Revenue, margin and growth highlights",
        ),
        left=header_left,
        top=header_top,
        width=header_width,
        theme="ir",
    )
    header_bottom = header_top + header_result["consumed_height"]

    # 3) KPI row directly below the section header (+ small breathing gap).
    kpi_top = header_bottom + 0.2
    kpi_height = 0.95
    add_kpi_row(
        slide,
        [
            KPISpec(label="Revenue", value="107.8M", detail="+12% QoQ"),
            KPISpec(label="Op Margin", value="28.5%", detail="+2.1pp"),
            KPISpec(label="Net Income", value="22.4M", detail="+8% YoY"),
            KPISpec(label="FCF", value="18.9M", detail="+15% YoY"),
        ],
        left=header_left,
        top=kpi_top,
        width=header_width,
        height=kpi_height,
        gap=0.15,
        theme="ir",
    )
    kpi_bottom = kpi_top + kpi_height

    # 4) Metric card row below the KPI row.
    card_top = kpi_bottom + 0.3
    card_height = 3.2
    add_metric_card_row(
        slide,
        [
            MetricCardSpec(
                label="REVENUE",
                title="Revenue by Segment",
                metrics=[
                    MetricEntry(label="Consumer", value="62M"),
                    MetricEntry(label="Enterprise", value="45M"),
                ],
            ),
            MetricCardSpec(
                label="MARGIN",
                title="Margin Trajectory",
                metrics=[
                    MetricEntry(label="Gross", value="48%"),
                    MetricEntry(label="Operating", value="28%"),
                ],
            ),
        ],
        left=header_left,
        top=card_top,
        width=header_width,
        height=card_height,
        gap=0.3,
        theme="ir",
    )

    # 5) Slide footer at bottom (fixed offsets).
    add_slide_footer(
        slide,
        SlideFooterSpec(
            left_text="IR Presentation - FY Q3",
            right_text="Confidential",
        ),
        slide_width=_SLIDE_W,
        slide_height=_SLIDE_H,
        theme="ir",
    )

    # -------------------------------------------------------------------
    # Assertions
    # -------------------------------------------------------------------

    # All expected shapes exist:
    #   page_marker:       2 (section + page textboxes)
    #   section_header:    3 (title + subtitle textboxes + divider rect)
    #   kpi_row (4 cells): 12 (each cell: label + value + detail = 3)
    #   metric_card_row:   2 cards; each card = frame rect + label + title +
    #                      chart placeholder + 2*(metric label + metric value)
    #                      = 1 + 1 + 1 + 1 + 4 = 8 shapes per card → 16
    #   slide_footer:      2 (left + right textboxes)
    # Total: 2 + 3 + 12 + 16 + 2 = 35.
    after = len(list(slide.shapes))
    total_added = after - before
    assert total_added == 35, (
        f"expected 35 shapes total across all five components, got {total_added}"
    )

    # check_containment must be clean. This is the strongest whole-slide
    # assertion: every child of every declared component container must lie
    # inside its declared bounds.
    findings = check_containment(prs)
    assert findings == [], (
        f"expected no containment findings on integration slide, "
        f"got {findings}"
    )

    # Spot-check IR theme propagation: ONE shape per component must resolve
    # to the expected IR-theme color. We cover the four theme-consuming
    # components (markers provide both page + footer).
    ir_primary = resolve_theme_color("primary", "ir").upper()  # 0A2540
    ir_text_secondary = resolve_theme_color("text_secondary", "ir").upper()  # 6B7280

    textboxes = [s for s in slide.shapes if s.has_text_frame]

    # The section_header title is the first 32pt bold textbox; find by
    # bold+text to avoid fragile index dependencies. Title color = primary.
    header_title = next(
        s for s in textboxes
        if s.text_frame.text == "FY Q3 Executive Summary"
    )
    assert _first_run_hex(header_title) == ir_primary

    # Section header subtitle uses text_secondary.
    header_subtitle = next(
        s for s in textboxes
        if s.text_frame.text == "Revenue, margin and growth highlights"
    )
    assert _first_run_hex(header_subtitle) == ir_text_secondary

    # KPI value "107.8M" must use IR primary.
    kpi_value = next(
        s for s in textboxes if s.text_frame.text == "107.8M"
    )
    assert _first_run_hex(kpi_value) == ir_primary

    # Metric card title "Revenue by Segment" must use IR primary.
    card_title = next(
        s for s in textboxes if s.text_frame.text == "Revenue by Segment"
    )
    assert _first_run_hex(card_title) == ir_primary

    # Page marker section line uses IR text_secondary.
    page_marker = next(
        s for s in textboxes if s.text_frame.text == "FINANCIAL SUMMARY"
    )
    assert _first_run_hex(page_marker) == ir_text_secondary

    # Slide footer left uses IR text_secondary.
    footer_left = next(
        s for s in textboxes if s.text_frame.text == "IR Presentation - FY Q3"
    )
    assert _first_run_hex(footer_left) == ir_text_secondary


# ---------------------------------------------------------------------------
# Test 2: Numbered list placed flush below a section header
# ---------------------------------------------------------------------------


def test_numbered_list_plus_section_header_positioning():
    """Section header's ``consumed_height`` is the correct y-offset for
    placing a numbered list directly below. Both components should share
    their x extent without overlap, and ``check_containment`` should
    produce zero findings.
    """
    prs = _new_prs_with_blank_slide()
    slide = prs.slides[0]

    left = 0.9
    top = 0.5
    width = _SLIDE_W - 1.8  # matches the existing test_section_header.py pattern

    header_result = add_section_header(
        slide,
        SectionHeaderSpec(
            title="Strategic Priorities",
            subtitle="FY26 roadmap anchors",
        ),
        left=left,
        top=top,
        width=width,
        theme="ir",
    )
    consumed = header_result["consumed_height"]

    # Numbered list flush below the header's rendered footprint.
    list_top = top + consumed
    list_height = 4.5  # generous — 3 items at 1.5" each
    list_result = add_numbered_list(
        slide,
        [
            NumberedItem(
                number="01",
                caption="/ Focus",
                title="Double down on enterprise",
                body="Prioritise accounts > $1M ARR.",
            ),
            NumberedItem(
                number="02",
                caption="/ Efficiency",
                title="Reduce unit cost",
                body="Target 15% reduction in COGS.",
            ),
            NumberedItem(
                number="03",
                caption="/ Expansion",
                title="Launch in EU-2 region",
                body="Certification tracks mid-year.",
            ),
        ],
        left=left,
        top=list_top,
        width=width,
        height=list_height,
        theme="ir",
    )

    # Invariant 1: list reports the declared bounds back unchanged.
    assert list_result["consumed_height"] == pytest.approx(list_height)
    assert list_result["consumed_width"] == pytest.approx(width)

    # Invariant 2: header divider bottom edge matches list top (no gap, no
    # overlap). This is the fundamental consumed_height contract.
    divider = header_result["divider_bounds"]
    divider_bottom = divider["top"] + divider["height"]
    assert divider_bottom == pytest.approx(list_top, abs=1e-6)

    # Invariant 3: first numbered item top == list_top (flush).
    first_item_top = list_result["items"][0]["bounds"]["top"]
    assert first_item_top == pytest.approx(list_top, abs=1e-6)

    # Invariant 4: containment is clean.
    findings = check_containment(prs)
    assert findings == [], f"unexpected containment findings: {findings}"


# ---------------------------------------------------------------------------
# Test 3: Metric card row below KPI row, bounds-stacked
# ---------------------------------------------------------------------------


def test_metric_card_with_kpi_row_above_no_overlap():
    """KPI row at top. Metric card row placed directly below using the
    KPI row's bottom (top + height) as the card row's top. No overlap and
    clean containment.
    """
    prs = _new_prs_with_blank_slide()
    slide = prs.slides[0]

    left = 0.5
    top = 0.8
    width = _SLIDE_W - 1.0

    # KPI row.
    kpi_height = 0.95
    kpi_result = add_kpi_row(
        slide,
        [
            KPISpec(label="ARR", value="$120M"),
            KPISpec(label="NRR", value="118%"),
            KPISpec(label="CAC payback", value="14mo"),
        ],
        left=left,
        top=top,
        width=width,
        height=kpi_height,
        gap=0.2,
    )

    # Use first cell's bounds to derive KPI row bottom (should equal
    # top + kpi_height; verifying here ties the two sides of the contract).
    first_cell_top = kpi_result["cells"][0]["bounds"]["top"]
    first_cell_h = kpi_result["cells"][0]["bounds"]["height"]
    kpi_bottom = first_cell_top + first_cell_h
    assert kpi_bottom == pytest.approx(top + kpi_height, abs=1e-6)

    # Metric card row flush below the KPI row (with a small visual gap).
    card_top = kpi_bottom + 0.25
    card_height = 3.5
    card_result = add_metric_card_row(
        slide,
        [
            MetricCardSpec(
                label="GROWTH",
                title="Pipeline health",
                metrics=[
                    MetricEntry(label="New", value="$42M"),
                    MetricEntry(label="Expansion", value="$28M"),
                ],
            ),
            MetricCardSpec(
                label="RETENTION",
                title="Churn profile",
                metrics=[
                    MetricEntry(label="Gross", value="4.2%"),
                    MetricEntry(label="Net", value="-1.8%"),
                ],
            ),
        ],
        left=left,
        top=card_top,
        width=width,
        height=card_height,
        gap=0.3,
    )

    # Card row reports back the declared bounds.
    assert card_result["consumed_height"] == pytest.approx(card_height)

    # Invariant: no KPI-cell bottom exceeds the card row's top — they don't
    # physically overlap.
    for cell in kpi_result["cells"]:
        cell_bottom = cell["bounds"]["top"] + cell["bounds"]["height"]
        assert cell_bottom <= card_top + 1e-6

    # Containment is clean for both components together.
    findings = check_containment(prs)
    assert findings == [], f"unexpected containment findings: {findings}"


# ---------------------------------------------------------------------------
# Test 4: Containment surfaces real overflows in a combined layout
# ---------------------------------------------------------------------------


def test_check_containment_catches_overflow_in_combined_slide():
    """Deliberately under-size a ``kpi_row`` so its internal value/detail
    textboxes escape the declared kpi_row container. Containment must
    surface at least one finding — this validates that the integration
    path surfaces bugs, not just happy paths.

    Why kpi_row: the component unconditionally renders its value textbox
    starting at ``top + _LABEL_H + _LABEL_VALUE_GAP = top + 0.24`` and
    extending ``_VALUE_H = 0.50`` down (so value bottom = top + 0.74).
    Declaring ``height=0.4`` makes the container bottom = top + 0.4, so
    the value textbox escapes by ~0.34" — well beyond the 0.01"
    containment tolerance.
    """
    prs = _new_prs_with_blank_slide()
    slide = prs.slides[0]

    # Add a legitimate header first so the slide has a normal component
    # alongside the intentionally-broken one — mirrors a real-world
    # "I shrank this component and broke it" regression.
    add_section_header(
        slide,
        SectionHeaderSpec(title="Broken layout demo", subtitle="kpi row squeezed"),
        left=0.5,
        top=0.5,
        width=_SLIDE_W - 1.0,
    )

    # Too-small KPI row: height 0.4 < required ~0.94" for detail bottom.
    add_kpi_row(
        slide,
        [
            KPISpec(label="Revenue", value="107.8M", detail="+12% QoQ"),
            KPISpec(label="Margin", value="28.5%", detail="+2.1pp"),
        ],
        left=0.5,
        top=2.0,
        width=_SLIDE_W - 1.0,
        height=0.4,  # deliberately too small
        gap=0.2,
    )

    findings = check_containment(prs)
    assert len(findings) >= 1, (
        "expected at least one containment finding for the squeezed "
        "kpi_row, got none"
    )
    # All findings should be from the kpi_row container (not the header).
    kpi_findings = [
        f for f in findings if "kpi_row" in f.message
    ]
    assert kpi_findings, (
        f"expected at least one finding referencing 'kpi_row', got: "
        f"{[f.message for f in findings]}"
    )
    # Each finding is an error severity + category shape_outside_container.
    for f in kpi_findings:
        assert f.severity == "error"
        assert f.category == "shape_outside_container"


# ---------------------------------------------------------------------------
# Test 5: engine/__init__.py re-exports every v0.6.0 block component
# ---------------------------------------------------------------------------


def test_all_block_components_importable_from_engine_root():
    """The v0.6.0 block components must be importable from the engine root
    package via the names declared in ``engine/__init__.py``'s ``__all__``.

    Regression guard: if a future refactor drops a name from ``__all__``
    (or accidentally renames it), this test fails before downstream code.
    The set below intentionally covers the full v0.6.0 surface the task
    spec calls out.
    """
    from pptx_mcp_server.engine import (
        add_section_header,
        add_page_marker,
        add_slide_footer,
        add_numbered_list,
        add_kpi_row_block,
        add_metric_card,
        add_metric_card_row,
    )

    # Each must be a callable (function). None of them should be ``None``
    # or a leftover attribute that happens to import cleanly.
    for fn in (
        add_section_header,
        add_page_marker,
        add_slide_footer,
        add_numbered_list,
        add_kpi_row_block,
        add_metric_card,
        add_metric_card_row,
    ):
        assert callable(fn), f"{fn!r} is not callable"
