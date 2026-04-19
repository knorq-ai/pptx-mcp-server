"""可変高カード行プリミティブのテスト.

``add_responsive_card_row`` の 3 つの ``height_mode`` (content / max / fill)
と、単一カード・最小高 clamp・アクセントバー描画・消費高返却の各挙動を
網羅する。また #26 (CardSpec mutation + unknown height_mode) と
#27 (padding 二重計上 → text clip) の再発防止テストを含む。
"""

from __future__ import annotations

import json

import pytest

from pptx_mcp_server.engine.cards import (
    CardPlacement,
    CardSpec,
    _content_height,
    _estimate_block_heights,
    add_responsive_card_row,
)
from pptx_mcp_server.engine.pptx_io import EngineError, ErrorCode


# EMU (914400 per inch) → inches 換算ユーティリティ
def _in(emu: int) -> float:
    return emu / 914400


def _shape_bbox(shape) -> tuple[float, float, float, float]:
    """shape の (left, top, width, height) を inches で返す."""
    return _in(shape.left), _in(shape.top), _in(shape.width), _in(shape.height)


# 各カードの先頭 shape (背景矩形) の bbox がそのカード全体の bbox になる。
# アクセントバーや textbox もそこから派生するため、背景矩形を頼りに検証する。
def _card_background_heights(slide, n_cards: int, has_accent: bool) -> list[float]:
    """各カードの背景矩形の高さを順に取り出す.

    ``add_responsive_card_row`` は各カードにつき最初に背景矩形 (AUTO_SHAPE) を
    追加する。その後、任意で accent bar (AUTO_SHAPE)、label/title/body
    (TEXT_BOX) を追加する。本テストでは shape 追加順・種別に基づき「各カード
    ごとに最初に現れる AUTO_SHAPE で、かつ label/title/body の TextBox より幅が
    広いもの」を拾う。簡易に、先頭から走査して AUTO_SHAPE が現れた時点の高さを
    採用し、同一 left の次の AUTO_SHAPE (accent bar) は無視する。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    heights: list[float] = []
    seen_lefts: set[float] = set()
    for s in slide.shapes:
        if s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        key = round(_in(s.left), 3)
        if key in seen_lefts:
            # 同じ left で 2 つ目の AUTO_SHAPE は accent bar のためスキップ
            continue
        seen_lefts.add(key)
        heights.append(_in(s.height))
        if len(heights) == n_cards:
            break
    return heights


def test_max_mode_identical_content(slide):
    """3 カードで同一の title/body を持ち height_mode='max' ならすべて同じ高さになる."""
    cards = [
        CardSpec(title="Strategy", body="Short body text."),
        CardSpec(title="Strategy", body="Short body text."),
        CardSpec(title="Strategy", body="Short body text."),
    ]
    placements, consumed = add_responsive_card_row(
        slide, cards,
        left=0.5, top=1.0, width=12.0, max_height=4.0,
        gap=0.2, height_mode="max", min_card_height=0.5,
    )

    heights = _card_background_heights(slide, n_cards=3, has_accent=False)
    assert len(heights) == 3
    assert heights[0] == pytest.approx(heights[1], abs=1e-4)
    assert heights[1] == pytest.approx(heights[2], abs=1e-4)
    # 消費高 == 各カード高
    assert consumed == pytest.approx(heights[0], abs=1e-4)
    # placements の長さ/高さ一致
    assert len(placements) == 3
    for p, h in zip(placements, heights):
        assert p.height == pytest.approx(h, abs=1e-4)


def test_max_mode_different_content_lengths_align(slide):
    """3 カード + content 長さ差 + height_mode='max' で全カード同一高 (= 最大 content 高)."""
    short = CardSpec(title="T", body="Short.")
    medium = CardSpec(title="T", body="Medium body text that wraps.")
    long_body = "This is a very long body text that will wrap to many lines. " * 6
    long_card = CardSpec(title="T", body=long_body)

    cards = [short, medium, long_card]
    placements, consumed = add_responsive_card_row(
        slide, cards,
        left=0.5, top=1.0, width=12.0, max_height=6.0,
        gap=0.2, height_mode="max",
    )

    heights = _card_background_heights(slide, n_cards=3, has_accent=False)
    # 全カード同一
    assert heights[0] == pytest.approx(heights[1], abs=1e-4)
    assert heights[1] == pytest.approx(heights[2], abs=1e-4)
    # 高さ == 最大 content 高 (placement の width で測る)
    card_w = placements[0].width
    max_content = max(_content_height(c, card_w) for c in cards)
    assert heights[0] == pytest.approx(max_content, abs=1e-3)
    assert consumed == pytest.approx(heights[0], abs=1e-4)


def test_content_mode_heights_differ(slide):
    """height_mode='content' では短い/長いカードで高さが異なる."""
    short = CardSpec(title="T", body="Short.")
    long_body = "Long body text that wraps across multiple lines. " * 8
    long_card = CardSpec(title="T", body=long_body)
    medium = CardSpec(title="T", body="Medium sample body.")

    cards = [short, medium, long_card]
    add_responsive_card_row(
        slide, cards,
        left=0.5, top=1.0, width=12.0, max_height=6.0,
        gap=0.2, height_mode="content", min_card_height=0.2,
    )

    heights = _card_background_heights(slide, n_cards=3, has_accent=False)
    # 短いカードと長いカードの高さは明確に異なる
    assert heights[0] != pytest.approx(heights[2], abs=0.1)
    assert heights[2] > heights[0]


def test_fill_mode_uses_max_height(slide):
    """height_mode='fill' で max_height=5" ならすべてのカード高が 5" になる."""
    cards = [
        CardSpec(title="A", body="Short body."),
        CardSpec(title="B", body="Short body."),
        CardSpec(title="C", body="Short body."),
    ]
    placements, consumed = add_responsive_card_row(
        slide, cards,
        left=0.5, top=0.5, width=12.0, max_height=5.0,
        gap=0.2, height_mode="fill",
    )

    heights = _card_background_heights(slide, n_cards=3, has_accent=False)
    for h in heights:
        assert h == pytest.approx(5.0, abs=1e-4)
    assert consumed == pytest.approx(5.0, abs=1e-4)
    for p in placements:
        assert p.height == pytest.approx(5.0, abs=1e-4)


def test_min_card_height_clamp(slide):
    """min_card_height=2 なら content が ~1" でも全カード 2" に clamp される."""
    # padding+title+body だけの軽いカード (~1" 程度)
    cards = [
        CardSpec(title="X", body="One line."),
        CardSpec(title="Y", body="One line."),
    ]
    add_responsive_card_row(
        slide, cards,
        left=0.5, top=0.5, width=10.0, max_height=6.0,
        gap=0.2, height_mode="max", min_card_height=2.0,
    )

    heights = _card_background_heights(slide, n_cards=2, has_accent=False)
    for h in heights:
        assert h == pytest.approx(2.0, abs=1e-4)


def test_single_card_fills_row_width(slide):
    """単一カード (n == 1): gap を無視し行幅すべてをそのカードが占める."""
    cards = [CardSpec(title="Only", body="Alone in the row.")]
    placements, _ = add_responsive_card_row(
        slide, cards,
        left=1.0, top=1.0, width=8.0, max_height=4.0,
        gap=0.5, height_mode="max",
    )

    # placement の width が 8.0 と一致 (width は layout の出力であり、
    # CardSpec には width フィールド自体が存在しない)
    assert placements[0].width == pytest.approx(8.0, abs=1e-4)

    # 先頭 shape = 背景矩形。width が 8.0 に一致
    first_shape = list(slide.shapes)[0]
    _, _, w, _ = _shape_bbox(first_shape)
    assert w == pytest.approx(8.0, abs=1e-3)


def test_accent_bar_drawn_with_expected_bbox(slide):
    """accent_color 非空のとき、card left と同じ位置に width ≈ 0.08" のバーが生成される."""
    cards = [CardSpec(title="T", body="Body", accent_color="FF0000")]
    add_responsive_card_row(
        slide, cards,
        left=2.0, top=1.0, width=4.0, max_height=3.0,
        gap=0.2, height_mode="max",
    )

    shapes = list(slide.shapes)
    # 背景矩形 → アクセントバー → 任意のテキストの順
    bg = shapes[0]
    accent = shapes[1]
    bg_l, bg_t, bg_w, bg_h = _shape_bbox(bg)
    ac_l, ac_t, ac_w, ac_h = _shape_bbox(accent)

    assert bg_l == pytest.approx(2.0, abs=1e-3)
    assert bg_w == pytest.approx(4.0, abs=1e-3)

    # アクセントバー: 左端一致、幅 0.08"、高さは card と一致
    assert ac_l == pytest.approx(bg_l, abs=1e-3)
    assert ac_t == pytest.approx(bg_t, abs=1e-3)
    assert ac_w == pytest.approx(0.08, abs=1e-3)
    assert ac_h == pytest.approx(bg_h, abs=1e-3)


def test_consumed_height_matches_max_for_max_and_fill(slide, one_slide_prs):
    """consumed_height == max/fill モードの共通カード高."""
    # max モード
    cards_max = [CardSpec(title="T", body="Body A"), CardSpec(title="T", body="Body B")]
    _, consumed_max = add_responsive_card_row(
        slide, cards_max,
        left=0.5, top=0.5, width=10.0, max_height=6.0,
        gap=0.2, height_mode="max",
    )
    heights_max = _card_background_heights(slide, n_cards=2, has_accent=False)
    assert consumed_max == pytest.approx(max(heights_max), abs=1e-4)

    # fill モード (別スライド用意)
    layout = one_slide_prs.slide_layouts[6]
    one_slide_prs.slides.add_slide(layout)
    slide2 = one_slide_prs.slides[1]
    cards_fill = [CardSpec(title="T", body="A"), CardSpec(title="T", body="B")]
    _, consumed_fill = add_responsive_card_row(
        slide2, cards_fill,
        left=0.5, top=0.5, width=10.0, max_height=4.5,
        gap=0.2, height_mode="fill",
    )
    heights_fill = _card_background_heights(slide2, n_cards=2, has_accent=False)
    assert consumed_fill == pytest.approx(4.5, abs=1e-4)
    for h in heights_fill:
        assert h == pytest.approx(4.5, abs=1e-4)


# ── #26: CardSpec mutation + unknown height_mode ───────────────


def test_cardspec_not_mutated_between_calls(slide, one_slide_prs):
    """同じ ``CardSpec`` リストを 2 つの ``add_responsive_card_row`` に渡しても
    2 回目の widths が 1 回目と一致し (幅リーク無し)、``CardSpec`` 自体が
    書き換えられない (#26 Problem A)."""
    cards = [
        CardSpec(title="A", body="Short."),
        CardSpec(title="B", body="Short."),
        CardSpec(title="C", body="Short."),
    ]
    # snapshot: CardSpec の全フィールド値を保存する
    snapshots = [tuple(c.__dict__.items()) for c in cards]

    # 1 回目の呼び出し
    placements1, _ = add_responsive_card_row(
        slide, cards,
        left=0.5, top=0.5, width=12.0, max_height=3.0,
        gap=0.2, height_mode="max",
    )
    widths1 = [p.width for p in placements1]
    # 入力 CardSpec は mutation されない (全フィールドが初期値のまま)
    for c, snap in zip(cards, snapshots):
        assert tuple(c.__dict__.items()) == snap

    # 同じリストで 2 枚目のスライドにもう 1 回レイアウト
    layout = one_slide_prs.slide_layouts[6]
    one_slide_prs.slides.add_slide(layout)
    slide2 = one_slide_prs.slides[1]
    placements2, _ = add_responsive_card_row(
        slide2, cards,
        left=0.5, top=0.5, width=12.0, max_height=3.0,
        gap=0.2, height_mode="max",
    )
    widths2 = [p.width for p in placements2]

    # 幅は一致
    assert widths1 == pytest.approx(widths2, abs=1e-6)
    # 入力 CardSpec はまだ変わっていない
    for c, snap in zip(cards, snapshots):
        assert tuple(c.__dict__.items()) == snap


def test_unknown_height_mode_raises(slide):
    """未知の ``height_mode`` (例: JSON 入力経由の ``'auto'``) は
    ``EngineError(INVALID_PARAMETER)`` を投げる (#26 Problem B)."""
    cards = [CardSpec(title="A", body="a")]
    with pytest.raises(EngineError) as excinfo:
        add_responsive_card_row(
            slide, cards,
            left=0.5, top=0.5, width=10.0, max_height=3.0,
            gap=0.2, height_mode="auto",  # type: ignore[arg-type]
            min_card_height=1.0,
        )
    err = excinfo.value
    assert err.code == ErrorCode.INVALID_PARAMETER
    msg = str(err)
    assert "height_mode" in msg
    assert "auto" in msg


@pytest.mark.parametrize("mode", ["content", "max", "fill"])
def test_valid_modes_still_work(blank_prs, mode):
    """``content`` / ``max`` / ``fill`` は従来通り動作する (#26 回帰防止)."""
    layout = blank_prs.slide_layouts[6]
    blank_prs.slides.add_slide(layout)
    slide = blank_prs.slides[0]

    cards = [CardSpec(title="T", body="B"), CardSpec(title="T", body="B")]
    placements, consumed = add_responsive_card_row(
        slide, cards,
        left=0.5, top=0.5, width=10.0, max_height=3.0,
        gap=0.2, height_mode=mode,  # type: ignore[arg-type]
        min_card_height=1.0,
    )
    assert len(placements) == 2
    assert consumed > 0.0


def test_placements_cover_row_width(slide):
    """``CardPlacement`` の ``left/width`` がギャップを考慮して行全体を覆う."""
    cards = [CardSpec(title=f"T{i}", body="Body") for i in range(3)]
    placements, _ = add_responsive_card_row(
        slide, cards,
        left=1.0, top=1.0, width=12.0, max_height=3.0,
        gap=0.2, height_mode="max",
    )
    assert len(placements) == 3
    # 1 枚目の left と 3 枚目の right が行全体と一致
    assert placements[0].left == pytest.approx(1.0, abs=1e-6)
    last = placements[-1]
    assert last.left + last.width == pytest.approx(1.0 + 12.0, abs=1e-3)
    # top は全カード一致
    for p in placements:
        assert p.top == pytest.approx(1.0, abs=1e-6)


# ── #27: padding 二重計上 → text clip 防止 ────────────────────


def test_no_text_overflow_for_japanese_body(one_slide_prs):
    """3 枚カード × 日本語 80 字 body / 幅 3" で、実描画された body textbox が
    ``estimate_text_height`` ベースで overflow しない (#27 回帰防止).

    ``check_text_overflow`` は font-size を run 単位で読むため、python-pptx
    において paragraph-level の font 設定が run に反映されない既知の挙動
    (本 PR のスコープ外) に引きずられる。そのため本テストでは
    ``check_text_overflow`` の出力にそのまま依存せず、実描画された body
    textbox の ``width × height`` と実際の body font size
    (``CardSpec.body_size_pt``) で ``estimate_text_height`` を再計算し、
    必要高 ≤ frame 高 × tolerance となっていることを直接検証する。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from pptx_mcp_server.engine.text_metrics import estimate_text_height

    slide = one_slide_prs.slides[0]
    # 日本語 80 文字
    body = "あいうえおかきくけこさしすせそたちつてとなにぬねのはひふへほまみむめもやゆよらりるれろわをんがぎぐげござじずぜぞだぢづでどばびぶべぼぱぴぷぺぽ"[:80]
    body_pt = 10.0
    cards = [
        CardSpec(title=f"カード{i}", body=body, body_size_pt=body_pt)
        for i in range(3)
    ]
    # 3 枚 × 3" + gap 0.2" × 2 = 9.4"
    add_responsive_card_row(
        slide, cards,
        left=0.5, top=0.5, width=9.4, max_height=4.0,
        gap=0.2, height_mode="max", min_card_height=1.0,
    )

    # TEXT_BOX は各カードごとに (label なし、) title → body の順に追加される。
    # 3 カード × 2 ブロック = 6 個。奇数 index が body。
    text_shapes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    assert len(text_shapes) == 6
    body_shapes = text_shapes[1::2]

    for bs in body_shapes:
        frame_w = bs.width / 914400
        frame_h = bs.height / 914400
        # validator と同じく左右 0.05" マージン合計 (0.10") を usable_width から引く
        usable = max(0.1, frame_w - 0.10)
        needed = estimate_text_height(body, usable, body_pt)
        # 5% tolerance (validator と同じ)
        assert needed <= frame_h * 1.05 + 1e-6, (
            f"body overflow: needed {needed:.3f}\" > frame {frame_h:.3f}\" × 1.05 "
            f"(width {frame_w:.3f}\", font {body_pt}pt)"
        )


def test_estimate_not_under_observed_body_height(slide):
    """``_estimate_block_heights`` の body 推定高が、実描画後の body textbox
    高さを下回らないこと (allocate >= observed、非過小評価)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 日本語 ~60 字で 3" 幅のカード
    body = "これはカードの本文テキストで、日本語の折り返しが発生するだけの長さを持つサンプル文字列である。" * 1
    card = CardSpec(title="見出し", body=body)

    # 行幅 3" = 単一カード
    placements, _ = add_responsive_card_row(
        slide, [card],
        left=0.5, top=0.5, width=3.0, max_height=5.0,
        gap=0.2, height_mode="content", min_card_height=0.2,
    )
    placement_w = placements[0].width

    # 推定 body 高
    _, _, body_h_est, _ = _estimate_block_heights(card, placement_w)

    # 実描画された body textbox の高さを取得 (最後の TEXT_BOX が body)
    text_shapes = [
        s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    body_shape = text_shapes[-1]
    observed_body_h = _in(body_shape.height)

    # 推定 ≥ 観測 (過小評価しない)。小さなトレランスを許容する。
    assert body_h_est + 1e-3 >= observed_body_h or observed_body_h - body_h_est < 0.02


# ── MCP tool boundary (#26 B) ─────────────────────────────────


def test_mcp_tool_rejects_unknown_height_mode(pptx_file):
    """``pptx_add_responsive_card_row`` MCP ツール側でも ``height_mode='auto'``
    は既存 invalid-param エラー形式で弾かれる (#26 Problem B、境界側)."""
    from pptx_mcp_server.server import pptx_add_responsive_card_row

    cards_json = json.dumps([{"title": "A", "body": "a"}])
    out = pptx_add_responsive_card_row(
        file_path=pptx_file,
        slide_index=0,
        cards_json=cards_json,
        left=0.5,
        top=0.5,
        width=10.0,
        max_height=3.0,
        gap=0.2,
        height_mode="auto",
        min_card_height=1.0,
    )
    # _err フォーマットに INVALID_PARAMETER が含まれる
    assert "INVALID_PARAMETER" in out or "height_mode" in out


def test_card_placement_dataclass_shape():
    """``CardPlacement`` は left/top/width/height の 4 フィールドを持つ."""
    p = CardPlacement(left=1.0, top=2.0, width=3.0, height=4.0)
    assert (p.left, p.top, p.width, p.height) == (1.0, 2.0, 3.0, 4.0)
