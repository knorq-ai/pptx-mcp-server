"""
PPTX file I/O — open, save, create, error types.
"""

from __future__ import annotations

import os
import sys
import tempfile
from enum import Enum
from pptx import Presentation
from pptx.util import Inches, Pt, Emu


class ErrorCode(str, Enum):
    FILE_NOT_FOUND = "FILE_NOT_FOUND"
    INVALID_PPTX = "INVALID_PPTX"
    SLIDE_NOT_FOUND = "SLIDE_NOT_FOUND"
    SHAPE_NOT_FOUND = "SHAPE_NOT_FOUND"
    INDEX_OUT_OF_RANGE = "INDEX_OUT_OF_RANGE"
    INVALID_PARAMETER = "INVALID_PARAMETER"
    TABLE_ERROR = "TABLE_ERROR"
    CHART_ERROR = "CHART_ERROR"


class EngineError(Exception):
    def __init__(self, code: ErrorCode, message: str):
        super().__init__(message)
        self.code = code


def open_pptx(file_path: str) -> Presentation:
    """Open an existing PPTX file."""
    if not os.path.exists(file_path):
        raise EngineError(ErrorCode.FILE_NOT_FOUND, f"File not found: {file_path}")
    try:
        return Presentation(file_path)
    except Exception as e:
        raise EngineError(ErrorCode.INVALID_PPTX, f"Not a valid PPTX file: {file_path} ({e})")


def save_pptx(prs: Presentation, file_path: str, fsync: bool = False) -> None:
    """PPTX を原子的に保存する。

    Atomicity (default): os.replace is syscall-atomic on POSIX and Windows
    within a single filesystem — the old file is preserved if the save raises
    mid-write.

    Durability (fsync=True, opt-in): fsyncs the temp file AND the containing
    directory entry before returning, so the new file survives power loss on
    local ext4/xfs/NTFS. Adds one-to-two disk barriers per call — avoid on
    hot paths. Skipped on Windows for directory fsync (unsupported).

    Caveats:
    - NFS: os.replace atomicity depends on server + client config. Do not
      rely on fsync for cross-mount durability.
    - Windows: os.replace raises PermissionError if the target is held open
      (e.g., the .pptx is open in PowerPoint). Caller should retry or
      surface a clear error.
    """
    abs_path = os.path.abspath(file_path)
    dir_ = os.path.dirname(abs_path) or "."
    # mkstemp で同一ディレクトリ内にユニークな一時パスを確保する。
    # NamedTemporaryFile は open 済み fd を保持するが、python-pptx は
    # パスから自前で open するため fd は不要。ここでは名前だけ欲しい。
    fd, tmp_path = tempfile.mkstemp(
        dir=dir_,
        prefix="." + os.path.basename(abs_path) + ".tmp.",
    )
    os.close(fd)
    try:
        prs.save(tmp_path)
        if fsync:
            # Flush the temp file's contents to disk before swapping it in.
            tmp_fd = os.open(tmp_path, os.O_RDONLY)
            try:
                os.fsync(tmp_fd)
            finally:
                os.close(tmp_fd)
        os.replace(tmp_path, abs_path)
        if fsync and sys.platform != "win32":
            # fsync the containing directory so the rename itself survives
            # power loss. Windows does not support directory fsync.
            dir_fd = os.open(dir_, os.O_DIRECTORY)
            try:
                os.fsync(dir_fd)
            finally:
                os.close(dir_fd)
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


def create_presentation(
    file_path: str,
    width_inches: float = 13.333,
    height_inches: float = 7.5,
) -> str:
    """Create a new blank PPTX file with specified dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)
    save_pptx(prs, file_path)
    return f"Created presentation: {file_path} ({width_inches}\" x {height_inches}\")"


def _get_slide(prs: Presentation, slide_index: int):
    """Get slide by 0-based index with bounds checking."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise EngineError(
            ErrorCode.SLIDE_NOT_FOUND,
            f"Slide index {slide_index} out of range (0-{len(prs.slides) - 1})",
        )
    return prs.slides[slide_index]


def _get_shape(slide, shape_index: int):
    """Get shape by 0-based index with bounds checking."""
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        raise EngineError(
            ErrorCode.SHAPE_NOT_FOUND,
            f"Shape index {shape_index} out of range (0-{len(shapes) - 1})",
        )
    return shapes[shape_index]


def _parse_color(color_str: str):
    """Parse hex color string like '#FF0000' or 'FF0000' to RGBColor."""
    from pptx.dml.color import RGBColor
    color_str = color_str.lstrip("#")
    if len(color_str) != 6:
        raise EngineError(ErrorCode.INVALID_PARAMETER, f"Invalid color: #{color_str}")
    try:
        return RGBColor(
            int(color_str[0:2], 16),
            int(color_str[2:4], 16),
            int(color_str[4:6], 16),
        )
    except ValueError:
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Invalid hex color: '#{color_str}'. Use 6-digit hex like 'FF0000'.",
        )
