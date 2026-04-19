"""
Connector and callout engine — add connector lines and annotation callouts.

Connectors are native PowerPoint connector shapes (straight, elbow, curve)
with optional arrow heads.  Callouts are composites of a textbox + connector
that point from an annotation label to a target coordinate.
"""

from __future__ import annotations

from typing import List, Optional

from lxml import etree
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from .pptx_io import (
    EngineError,
    ErrorCode,
    _get_slide,
    _parse_color,
    open_pptx,
    save_pptx,
)
from .shapes import _add_textbox
from ..theme import Theme, MCKINSEY, resolve_color

# ---------------------------------------------------------------------------
# Mappings
# ---------------------------------------------------------------------------

_CONNECTOR_TYPE_MAP = {
    "straight": MSO_CONNECTOR_TYPE.STRAIGHT,
    "elbow": MSO_CONNECTOR_TYPE.ELBOW,
    "curve": MSO_CONNECTOR_TYPE.CURVE,
}

_DASH_STYLE_MAP = {
    "solid": MSO_LINE_DASH_STYLE.SOLID,
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
    "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
}

# OOXML arrow type values for a:headEnd / a:tailEnd
_ARROW_TYPE_MAP = {
    "none": "none",
    "triangle": "triangle",
    "stealth": "stealth",
    "diamond": "diamond",
    "oval": "oval",
    "open": "arrow",
}

_ARROW_SIZE_MAP = {
    "small": "sm",
    "medium": "med",
    "large": "lg",
}


# ---------------------------------------------------------------------------
# Arrow head XML manipulation
# ---------------------------------------------------------------------------


def _set_arrow_heads(
    connector,
    begin_arrow: str = "none",
    end_arrow: str = "triangle",
    arrow_size: str = "medium",
) -> None:
    """Set arrow heads on a connector via XML manipulation.

    python-pptx does not expose headEnd/tailEnd properties, so we
    manipulate the OOXML directly.
    """
    # Get or create the <a:ln> element on the connector's spPr
    cxnSp = connector._element
    spPr = cxnSp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(cxnSp, qn("p:spPr"))

    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))

    size_val = _ARROW_SIZE_MAP.get(arrow_size, "med")

    # Head (begin) arrow
    head_type = _ARROW_TYPE_MAP.get(begin_arrow, "none")
    # Remove existing
    existing_head = ln.find(qn("a:headEnd"))
    if existing_head is not None:
        ln.remove(existing_head)
    if head_type != "none":
        head_el = etree.SubElement(ln, qn("a:headEnd"))
        head_el.set("type", head_type)
        head_el.set("w", size_val)
        head_el.set("len", size_val)

    # Tail (end) arrow
    tail_type = _ARROW_TYPE_MAP.get(end_arrow, "none")
    existing_tail = ln.find(qn("a:tailEnd"))
    if existing_tail is not None:
        ln.remove(existing_tail)
    if tail_type != "none":
        tail_el = etree.SubElement(ln, qn("a:tailEnd"))
        tail_el.set("type", tail_type)
        tail_el.set("w", size_val)
        tail_el.set("len", size_val)


# ---------------------------------------------------------------------------
# Core connector
# ---------------------------------------------------------------------------


def _add_connector(
    slide,
    begin_x: float,
    begin_y: float,
    end_x: float,
    end_y: float,
    *,
    connector_type: str = "straight",
    color: Optional[str] = None,
    width: Optional[float] = None,
    dash_style: Optional[str] = None,
    begin_arrow: str = "none",
    end_arrow: str = "triangle",
    arrow_size: str = "medium",
    theme: Optional[Theme] = None,
) -> int:
    """Add a connector line between two points.

    All positions in inches.  Returns shape index.

    connector_type: "straight", "elbow", "curve"
    begin_arrow / end_arrow: "none", "triangle", "stealth", "diamond", "oval", "open"
    arrow_size: "small", "medium", "large"
    dash_style: "solid", "dash", "dot", "dash_dot", "long_dash"
    """
    theme = theme or MCKINSEY

    if connector_type not in _CONNECTOR_TYPE_MAP:
        valid = ", ".join(sorted(_CONNECTOR_TYPE_MAP.keys()))
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Unknown connector_type '{connector_type}'. Valid: {valid}",
        )
    ctype = _CONNECTOR_TYPE_MAP[connector_type]

    cxn = slide.shapes.add_connector(
        ctype,
        Inches(begin_x), Inches(begin_y),
        Inches(end_x), Inches(end_y),
    )

    # Color
    color_token = color or theme.connector.get("color", "accent")
    line_color = resolve_color(theme, color_token).lstrip("#")
    cxn.line.color.rgb = _parse_color(line_color)

    # Width
    line_width = width or theme.connector.get("width", 1.5)
    cxn.line.width = Pt(line_width)

    # Dash style
    if dash_style and dash_style in _DASH_STYLE_MAP:
        cxn.line.dash_style = _DASH_STYLE_MAP[dash_style]

    # Arrow heads
    _set_arrow_heads(cxn, begin_arrow, end_arrow, arrow_size)

    return len(list(slide.shapes)) - 1


# ---------------------------------------------------------------------------
# Callout composite (textbox + connector)
# ---------------------------------------------------------------------------


def _add_callout(
    slide,
    text: str,
    target_x: float,
    target_y: float,
    *,
    label_x: Optional[float] = None,
    label_y: Optional[float] = None,
    label_width: float = 2.0,
    label_height: float = 0.4,
    connector_type: str = "straight",
    font_size: float = 10,
    font_color: Optional[str] = None,
    font_bold: bool = True,
    line_color: Optional[str] = None,
    line_width: float = 1.0,
    arrow_end: str = "triangle",
    bg_color: Optional[str] = None,
    border_color: Optional[str] = None,
    theme: Optional[Theme] = None,
) -> List[int]:
    """Add a callout annotation: textbox + connector arrow pointing to target.

    target_x, target_y: the point being annotated (inches).
    label_x, label_y: where to place the label box (auto-calculated if None).

    Returns [textbox_index, connector_index].
    """
    theme = theme or MCKINSEY

    # Auto-place label if not specified
    if label_x is None:
        label_x = target_x + 1.5
        # Clamp to slide right margin
        slide_w = theme.slide.get("width", 13.333)
        if label_x + label_width > slide_w - 0.3:
            label_x = target_x - label_width - 1.5
            label_x = max(0.3, label_x)
    if label_y is None:
        label_y = target_y - 1.0
        label_y = max(0.2, label_y)

    # Resolve colors
    text_color = font_color or "primary"
    connector_color = line_color or theme.connector.get("color", "accent")

    # Add textbox
    tb_idx = _add_textbox(
        slide,
        left=label_x,
        top=label_y,
        width=label_width,
        height=label_height,
        text=text,
        font_size=font_size,
        font_color=text_color,
        bold=font_bold,
        alignment="center",
        vertical_anchor="middle",
        theme=theme,
    )

    # Optionally style the textbox background
    shapes = list(slide.shapes)
    tb_shape = shapes[tb_idx]
    if bg_color:
        bg_hex = resolve_color(theme, bg_color).lstrip("#")
        tb_shape.fill.solid()
        tb_shape.fill.fore_color.rgb = _parse_color(bg_hex)
    if border_color:
        border_hex = resolve_color(theme, border_color).lstrip("#")
        tb_shape.line.color.rgb = _parse_color(border_hex)
        tb_shape.line.width = Pt(0.75)

    # Compute connector endpoints (from label center-edge toward target)
    label_cx = label_x + label_width / 2
    label_cy = label_y + label_height / 2

    # Start from the edge of the label closest to the target
    if abs(target_x - label_cx) > abs(target_y - label_cy):
        # Horizontal dominant — start from left or right edge
        if target_x < label_cx:
            start_x = label_x
        else:
            start_x = label_x + label_width
        start_y = label_cy
    else:
        # Vertical dominant — start from top or bottom edge
        start_x = label_cx
        if target_y < label_cy:
            start_y = label_y
        else:
            start_y = label_y + label_height

    # Add connector
    cxn_idx = _add_connector(
        slide,
        begin_x=start_x,
        begin_y=start_y,
        end_x=target_x,
        end_y=target_y,
        connector_type=connector_type,
        color=connector_color,
        width=line_width,
        begin_arrow="none",
        end_arrow=arrow_end,
        theme=theme,
    )

    return [tb_idx, cxn_idx]


# ---------------------------------------------------------------------------
# File-based wrappers
# ---------------------------------------------------------------------------


def add_connector(
    file_path: str,
    slide_index: int,
    begin_x: float,
    begin_y: float,
    end_x: float,
    end_y: float,
    connector_type: str = "straight",
    color: Optional[str] = None,
    width: Optional[float] = None,
    dash_style: Optional[str] = None,
    begin_arrow: str = "none",
    end_arrow: str = "triangle",
    arrow_size: str = "medium",
) -> str:
    """File-based wrapper: add a connector to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    idx = _add_connector(
        slide, begin_x, begin_y, end_x, end_y,
        connector_type=connector_type, color=color, width=width,
        dash_style=dash_style, begin_arrow=begin_arrow,
        end_arrow=end_arrow, arrow_size=arrow_size,
    )
    save_pptx(prs, file_path)
    return f"Added {connector_type} connector [{idx}] on slide [{slide_index}]"


def add_callout(
    file_path: str,
    slide_index: int,
    text: str,
    target_x: float,
    target_y: float,
    label_x: Optional[float] = None,
    label_y: Optional[float] = None,
    label_width: float = 2.0,
    label_height: float = 0.4,
    connector_type: str = "straight",
    font_size: float = 10,
    font_color: Optional[str] = None,
    font_bold: bool = True,
    line_color: Optional[str] = None,
    line_width: float = 1.0,
    arrow_end: str = "triangle",
    bg_color: Optional[str] = None,
    border_color: Optional[str] = None,
) -> str:
    """File-based wrapper: add a callout annotation to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    indices = _add_callout(
        slide, text, target_x, target_y,
        label_x=label_x, label_y=label_y,
        label_width=label_width, label_height=label_height,
        connector_type=connector_type, font_size=font_size,
        font_color=font_color, font_bold=font_bold,
        line_color=line_color, line_width=line_width,
        arrow_end=arrow_end, bg_color=bg_color,
        border_color=border_color,
    )
    save_pptx(prs, file_path)
    return (
        f"Added callout '{text}' (textbox [{indices[0]}], connector [{indices[1]}]) "
        f"on slide [{slide_index}]"
    )
