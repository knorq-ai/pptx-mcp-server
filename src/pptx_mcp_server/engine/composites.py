"""
Composite operations -- high-level functions that combine multiple primitives.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from .pptx_io import (
    EngineError, ErrorCode,
    open_pptx, save_pptx, _get_slide, _parse_color,
)
from .slides import _add_slide, _set_slide_background
from .shapes import _add_textbox, _add_shape, _add_paragraph, _add_image
from .tables import _add_table
from .charts import _add_chart
from .icons import _add_icon
from .connectors import _add_connector, _add_callout
from ..theme import Theme, MCKINSEY, get_theme, resolve_color


# ── In-memory composite primitives ─────────────────────────────


def _add_content_slide(prs, title, theme=None, source=None, page_number=None):
    """Add a blank slide with action title + divider line + optional source + page number.
    Returns (slide, slide_index)."""
    slide_idx = _add_slide(prs, layout_index=6, theme=theme)
    slide = _get_slide(prs, slide_idx)

    # Defaults from theme or hardcoded
    if theme:
        left_margin = theme.margins.get("left", 0.9)
        right_margin = theme.margins.get("right", 0.9)
        slide_w = theme.slide.get("width", 13.333)
        title_top = theme.layout.get("title_top", 0.45)
        title_height = theme.layout.get("title_height", 0.5)
        divider_top = theme.layout.get("divider_top", 0.95)
        footer_top = theme.layout.get("footer_top", 6.65)
        title_font = theme.fonts.get("title", "Arial")
        title_size = theme.sizes.get("title", 22)
        primary_color = resolve_color(theme, "primary")
        accent_color = resolve_color(theme, "accent")
        footnote_color = resolve_color(theme, "footnote")
        footnote_size = theme.sizes.get("footnote", 8)
        caption_size = theme.sizes.get("caption", 9)
    else:
        left_margin = 0.9
        right_margin = 0.9
        slide_w = 13.333
        title_top = 0.45
        title_height = 0.5
        divider_top = 0.95
        footer_top = 6.65
        title_font = "Arial"
        title_size = 22
        primary_color = "#051C2C"
        accent_color = "#2251FF"
        footnote_color = "#A2AAAD"
        footnote_size = 8
        caption_size = 9

    body_width = slide_w - left_margin - right_margin

    # Warn if title is likely too long for single line (~30 chars for JP at 22pt on 11.5" width)
    title_warning = ""
    if len(title) > 32:
        title_warning = f" [WARNING: Title is {len(title)} chars — may wrap to 2 lines. Consider shortening to ~30 chars.]"

    # Action title — auto-shrink to fit single line, bottom-anchored
    title_idx = _add_textbox(
        slide,
        left=left_margin,
        top=title_top,
        width=body_width,
        height=title_height,
        text=title,
        font_name=title_font,
        font_size=title_size,
        font_color=primary_color,
        bold=True,
        vertical_anchor="bottom",
    )

    # Enable auto-shrink so long titles reduce font size instead of wrapping
    title_shape = list(slide.shapes)[title_idx]
    bodyPr = title_shape.text_frame._txBody.find(qn("a:bodyPr"))
    # Remove any existing autofit elements
    for child in list(bodyPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("noAutofit", "normAutofit", "spAutoFit"):
            bodyPr.remove(child)
    # Add normAutofit — fontScale minimum 60% ensures aggressive shrink for long titles
    norm = etree.SubElement(bodyPr, qn("a:normAutofit"))
    norm.set("fontScale", "60000")  # 60% minimum scale (60000 = 60%)

    # Divider line (thin rectangle)
    _add_shape(
        slide,
        shape_type="rectangle",
        left=left_margin,
        top=divider_top,
        width=body_width,
        height=0.02,
        fill_color=accent_color,
        no_line=True,
    )

    # Source footnote
    if source:
        _add_textbox(
            slide,
            left=left_margin,
            top=footer_top,
            width=body_width * 0.7,
            height=0.3,
            text=source,
            font_name=title_font,
            font_size=caption_size,
            font_color=footnote_color,
            italic=True,
        )

    # Page number
    if page_number is not None:
        _add_textbox(
            slide,
            left=slide_w - right_margin - 0.5,
            top=footer_top,
            width=0.5,
            height=0.3,
            text=str(page_number),
            font_name=title_font,
            font_size=footnote_size,
            font_color=footnote_color,
            alignment="right",
        )

    return slide, slide_idx


def _add_section_divider(prs, title, subtitle="", theme=None):
    """Add a section divider slide (dark background, centered title, accent stripes).
    Returns (slide, slide_index)."""
    slide_idx = _add_slide(prs, layout_index=6, theme=theme)
    slide = _get_slide(prs, slide_idx)

    if theme:
        slide_w = theme.slide.get("width", 13.333)
        slide_h = theme.slide.get("height", 7.5)
        primary_color = resolve_color(theme, "primary")
        accent_color = resolve_color(theme, "accent")
        white_color = resolve_color(theme, "white")
        title_font = theme.fonts.get("title", "Arial")
        title_size = theme.sizes.get("title", 22)
        subtitle_size = theme.sizes.get("subtitle", 16)
    else:
        slide_w = 13.333
        slide_h = 7.5
        primary_color = "#051C2C"
        accent_color = "#2251FF"
        white_color = "#FFFFFF"
        title_font = "Arial"
        title_size = 22
        subtitle_size = 16

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _parse_color(primary_color)

    # Top accent stripe
    _add_shape(
        slide,
        shape_type="rectangle",
        left=0,
        top=0,
        width=slide_w,
        height=0.06,
        fill_color=accent_color,
        no_line=True,
    )

    # Bottom accent stripe
    _add_shape(
        slide,
        shape_type="rectangle",
        left=0,
        top=slide_h - 0.06,
        width=slide_w,
        height=0.06,
        fill_color=accent_color,
        no_line=True,
    )

    # Centered title
    _add_textbox(
        slide,
        left=1.5,
        top=slide_h * 0.35,
        width=slide_w - 3.0,
        height=1.0,
        text=title,
        font_name=title_font,
        font_size=title_size * 1.5,
        font_color=white_color,
        bold=True,
        alignment="center",
        vertical_anchor="middle",
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide,
            left=1.5,
            top=slide_h * 0.35 + 1.2,
            width=slide_w - 3.0,
            height=0.6,
            text=subtitle,
            font_name=title_font,
            font_size=subtitle_size,
            font_color=white_color,
            alignment="center",
            vertical_anchor="top",
        )

    return slide, slide_idx


def _add_kpi_row(slide, kpis, y, theme=None):
    """Add a row of KPI metric cards.

    Design: white card with left accent bar, large value, label below.
    kpis = [{"value": "107.8M", "label": "Revenue"}, ...]
    Returns list of shape indices."""
    if theme:
        left_margin = theme.margins.get("left", 0.9)
        right_margin = theme.margins.get("right", 0.9)
        slide_w = theme.slide.get("width", 13.333)
        primary_color = resolve_color(theme, "primary")
        accent_color = resolve_color(theme, "accent")
        text_secondary = resolve_color(theme, "text_secondary")
        body_font = theme.fonts.get("body", "Arial")
        border_color = resolve_color(theme, "border")
    else:
        left_margin = 0.9
        right_margin = 0.9
        slide_w = 13.333
        primary_color = "#051C2C"
        accent_color = "#2251FF"
        text_secondary = "#666666"
        body_font = "Arial"
        border_color = "#D0D0D0"

    body_width = slide_w - left_margin - right_margin
    n = len(kpis)
    if n == 0:
        return []

    gap = 0.25
    box_w = (body_width - gap * (n - 1)) / n
    box_h = 1.3
    accent_bar_w = 0.06
    indices = []

    for i, kpi in enumerate(kpis):
        x = left_margin + i * (box_w + gap)
        value = kpi.get("value", "")
        label = kpi.get("label", "")

        # White card with subtle border
        _add_shape(
            slide,
            shape_type="rectangle",
            left=x,
            top=y,
            width=box_w,
            height=box_h,
            fill_color="FFFFFF",
            line_color=border_color,
            line_width=0.5,
        )

        # Left accent bar
        _add_shape(
            slide,
            shape_type="rectangle",
            left=x,
            top=y,
            width=accent_bar_w,
            height=box_h,
            fill_color=accent_color,
            no_line=True,
        )

        # Value text — large, accent colored
        content_x = x + accent_bar_w + 0.25
        content_w = box_w - accent_bar_w - 0.4
        _add_textbox(
            slide,
            left=content_x,
            top=y + 0.15,
            width=content_w,
            height=0.6,
            text=value,
            font_name=body_font,
            font_size=28,
            font_color=accent_color,
            bold=True,
            vertical_anchor="bottom",
        )

        # Label text — muted, below value
        idx = _add_textbox(
            slide,
            left=content_x,
            top=y + 0.8,
            width=content_w,
            height=0.35,
            text=label,
            font_name=body_font,
            font_size=11,
            font_color=text_secondary,
        )
        indices.append(idx)

    return indices


def _add_bullet_block(slide, items, left, top, width, height, theme=None):
    """Add a bulleted text block with multiple items.

    Line spacing is auto-computed to distribute bullets across the full height,
    preventing dead whitespace at the bottom.
    Returns shape_index."""
    if theme:
        body_font = theme.fonts.get("body", "Arial")
        body_size = theme.sizes.get("body", 14)
        primary_color = resolve_color(theme, "primary")
    else:
        body_font = "Arial"
        body_size = 14
        primary_color = "#051C2C"

    if not items:
        return None

    # Compute line spacing to fill the height
    from pptx.util import Pt as _Pt

    n_items = len(items)
    min_line_height_pt = body_size * 1.3
    line_height_pt = max(min_line_height_pt, (height / max(n_items, 1)) * 72 * 0.75)
    line_height_pt = min(line_height_pt, body_size * 3.0)  # cap at 3x font size

    # First item
    first_item = items[0] if isinstance(items[0], str) else items[0].get("text", "")
    idx = _add_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        text=f"\u2022  {first_item}",
        font_name=body_font,
        font_size=body_size,
        font_color=primary_color,
        line_spacing=line_height_pt,
    )

    # Additional items as separate paragraphs (inherits line spacing)
    shape = list(slide.shapes)[idx]
    tf = shape.text_frame
    for item in items[1:]:
        text = item if isinstance(item, str) else item.get("text", "")
        p = tf.add_paragraph()
        p.text = f"\u2022  {text}"
        p.font.name = body_font
        p.font.size = _Pt(body_size)
        p.font.color.rgb = _parse_color(primary_color)
        p.line_spacing = _Pt(line_height_pt)

    return idx


def _add_card_grid(slide, cards, y=None, theme=None):
    """Add a grid of content cards that auto-fills the body zone.

    cards = [
        {
            "title": "Strategy 1",
            "body": "Description text...",
            "bullets": ["Point A", "Point B"],   # optional, used instead of body
            "icon_id": "briefcase",               # optional
            "icon_color": "2251FF",               # optional
        },
        ...
    ]

    The grid auto-computes layout:
      1-2 cards  → 1 row
      3-4 cards  → 2×2
      5-6 cards  → 2×3
      7-9 cards  → 3×3

    Returns list of shape indices.
    """
    theme = theme or MCKINSEY

    left_margin = theme.margins.get("left", 0.9)
    right_margin = theme.margins.get("right", 0.9)
    slide_w = theme.slide.get("width", 13.333)
    body_top = y if y is not None else theme.layout.get("body_top", 1.15)
    footer_top = theme.layout.get("footer_top", 6.65)
    primary_color = resolve_color(theme, "primary")
    accent_color = resolve_color(theme, "accent")
    text_secondary = resolve_color(theme, "text_secondary")
    body_font = theme.fonts.get("body", "Arial")
    bg_alt = resolve_color(theme, "bg_alt")

    body_width = slide_w - left_margin - right_margin
    body_height = footer_top - body_top - 0.15  # small bottom padding

    n = len(cards)
    if n == 0:
        return []

    # Determine grid dimensions
    if n <= 2:
        n_cols, n_rows = n, 1
    elif n <= 4:
        n_cols, n_rows = 2, 2
    elif n <= 6:
        n_cols, n_rows = 3, 2
    else:
        n_cols, n_rows = 3, 3

    h_gap = 0.25
    v_gap = 0.3
    card_w = (body_width - h_gap * (n_cols - 1)) / n_cols
    card_h = (body_height - v_gap * (n_rows - 1)) / n_rows

    max_cards = n_cols * n_rows
    if len(cards) > max_cards:
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"card_grid supports max {max_cards} cards ({n_cols}x{n_rows}) "
            f"but {len(cards)} were provided.",
        )

    indices = []
    icon_margin = 0.25
    inner_pad = 0.2

    for idx, card in enumerate(cards):
        col = idx % n_cols
        row = idx // n_cols
        x = left_margin + col * (card_w + h_gap)
        cy = body_top + row * (card_h + v_gap)

        title = card.get("title", "")
        body_text = card.get("body", "")
        bullets = card.get("bullets", [])
        icon_id = card.get("icon_id")
        icon_color = card.get("icon_color", accent_color)

        # Background
        _add_shape(
            slide,
            shape_type="rectangle",
            left=x, top=cy, width=card_w, height=card_h,
            fill_color=bg_alt, no_line=True,
        )

        # Layout: icon + title row, then body fills remaining height
        pad = 0.3
        icon_size = 0.4
        icon_row_h = icon_size + 0.1  # icon height + breathing room
        title_h = 0.35

        # Icon (small, left-aligned)
        title_y = cy + pad
        text_x = x + pad
        if icon_id:
            try:
                _add_icon(
                    slide, icon_id,
                    left=x + pad,
                    top=title_y,
                    height=icon_size,
                    color=icon_color,
                    theme=theme,
                )
                text_x = x + pad + icon_size + 0.15
            except EngineError:
                pass  # icon not found — render card without icon

        # Title (next to icon, vertically centered on icon)
        _add_textbox(
            slide,
            left=text_x, top=title_y + 0.03,
            width=card_w - (text_x - x) - pad,
            height=title_h,
            text=title,
            font_name=body_font,
            font_size=13,
            font_color=primary_color,
            bold=True,
            vertical_anchor="middle",
            theme=theme,
        )

        # Body or bullets — full width, fill from below icon row to card bottom
        content_top = title_y + icon_row_h + 0.1
        content_x = x + pad
        content_w = card_w - 2 * pad
        content_h = card_h - (content_top - cy) - pad

        if bullets:
            # Compute line spacing to distribute bullets evenly
            n_lines = len(bullets)
            line_height_pt = max(14, (content_h / max(n_lines, 1)) * 72 * 0.7)
            line_height_pt = min(line_height_pt, 26)

            text = "\n".join(f"\u2022  {b}" for b in bullets)
            _add_textbox(
                slide,
                left=content_x, top=content_top,
                width=content_w, height=content_h,
                text=text,
                font_name=body_font,
                font_size=10,
                font_color=text_secondary,
                line_spacing=line_height_pt,
                theme=theme,
            )
        elif body_text:
            last_idx = _add_textbox(
                slide,
                left=content_x, top=content_top,
                width=content_w, height=content_h,
                text=body_text,
                font_name=body_font,
                font_size=10,
                font_color=text_secondary,
                theme=theme,
            )
            indices.append(last_idx)

    return indices


def _build_slide(prs, spec, theme=None):
    """Build an entire slide from a JSON spec in one call.

    spec = {
        "layout": "content" | "section_divider" | "blank",
        "title": "Action title text",                    # for content/section_divider
        "subtitle": "...",                               # for section_divider only
        "background": "051C2C",                          # optional solid bg color
        "source": "Source: ...",                          # optional footnote
        "page_number": 1,                                # optional
        "elements": [                                    # list of elements to add
            {
                "type": "textbox",
                "left": 0.9, "top": 1.2, "width": 11.5, "height": 0.3,
                "text": "...",
                "font_size": 14, "font_color": "2251FF", "bold": true,
                "alignment": "left", "vertical_anchor": "top"
            },
            {
                "type": "shape",
                "shape_type": "rectangle",
                "left": 0.9, "top": 1.5, "width": 5.0, "height": 2.0,
                "fill_color": "F5F5F5", "no_line": true
            },
            {
                "type": "table",
                "rows": [["Header1","Header2"],["val1","val2"]],
                "left": 0.9, "top": 3.0, "width": 11.5,
                "col_widths": [0.5, 0.5], "font_size": 12, "row_height": 0.4
            },
            {
                "type": "kpi_row",
                "kpis": [{"value":"100","label":"Metric"}],
                "y": 1.2
            },
            {
                "type": "bullet_block",
                "items": ["Point 1", "Point 2"],
                "left": 0.9, "top": 2.0, "width": 11.5, "height": 2.0
            },
            {
                "type": "image",
                "image_path": "/path/to/image.png",
                "left": 1.0, "top": 1.0, "width": 3.0
            }
        ]
    }

    Returns (slide, slide_index).
    """
    # Theme resolution: explicit param > spec key > MCKINSEY default
    if theme is None:
        theme_name = spec.get("theme")
        if theme_name:
            theme = get_theme(theme_name)
    theme = theme or MCKINSEY

    layout = spec.get("layout", "content")

    if layout == "section_divider":
        slide, slide_idx = _add_section_divider(
            prs,
            spec.get("title", ""),
            spec.get("subtitle", ""),
            theme=theme,
        )
    elif layout == "content":
        slide, slide_idx = _add_content_slide(
            prs,
            spec.get("title", ""),
            theme=theme,
            source=spec.get("source"),
            page_number=spec.get("page_number"),
        )
    else:  # blank
        slide_idx = _add_slide(prs, layout_index=6, theme=theme)
        slide = _get_slide(prs, slide_idx)

    # Apply background if specified
    bg = spec.get("background")
    if bg:
        _set_slide_background(prs, slide_idx, bg, theme=theme)

    # Process elements
    for i, elem in enumerate(spec.get("elements", [])):
        etype = elem.get("type", "")

        try:
            _dispatch_element(slide, elem, etype, theme)
        except KeyError as e:
            raise EngineError(
                ErrorCode.INVALID_PARAMETER,
                f"Element {i} (type='{etype}'): missing required field '{e.args[0]}'",
            ) from None

    return slide, slide_idx


def _dispatch_element(slide, elem, etype, theme):
    """Dispatch a single element spec to the appropriate builder."""
    if etype == "textbox":
        _add_textbox(
            slide,
            left=elem["left"], top=elem["top"],
            width=elem["width"], height=elem["height"],
            text=elem.get("text", ""),
            font_name=elem.get("font_name"),
            font_size=elem.get("font_size"),
            font_color=elem.get("font_color"),
            bold=elem.get("bold"),
            italic=elem.get("italic"),
            alignment=elem.get("alignment"),
            vertical_anchor=elem.get("vertical_anchor"),
            word_wrap=elem.get("word_wrap", True),
            line_spacing=elem.get("line_spacing"),
            underline=elem.get("underline"),
            theme=theme,
        )

    elif etype == "shape":
        _add_shape(
            slide,
            shape_type=elem.get("shape_type", "rectangle"),
            left=elem["left"], top=elem["top"],
            width=elem["width"], height=elem["height"],
            fill_color=elem.get("fill_color"),
            line_color=elem.get("line_color"),
            line_width=elem.get("line_width"),
            no_line=elem.get("no_line", False),
            text=elem.get("text"),
            font_name=elem.get("font_name"),
            font_size=elem.get("font_size"),
            font_color=elem.get("font_color"),
            bold=elem.get("bold"),
            alignment=elem.get("alignment"),
            theme=theme,
        )

    elif etype == "table":
        _add_table(
            slide,
            rows=elem["rows"],
            left=elem["left"], top=elem["top"],
            width=elem["width"],
            col_widths=elem.get("col_widths"),
            row_height=elem.get("row_height", 0.30),
            font_size=elem.get("font_size", 10),
            header_bg=elem.get("header_bg", "051C2C"),
            header_fg=elem.get("header_fg", "FFFFFF"),
            alt_row_bg=elem.get("alt_row_bg", "F5F5F5"),
            border_color=elem.get("border_color", "D0D0D0"),
            border_width=elem.get("border_width", 0.5),
            no_vertical_borders=elem.get("no_vertical_borders", True),
            theme=theme,
        )

    elif etype == "kpi_row":
        _add_kpi_row(
            slide,
            kpis=elem["kpis"],
            y=elem["y"],
            theme=theme,
        )

    elif etype == "bullet_block":
        _add_bullet_block(
            slide,
            items=elem["items"],
            left=elem["left"], top=elem["top"],
            width=elem["width"], height=elem["height"],
            theme=theme,
        )

    elif etype == "image":
        _add_image(
            slide,
            image_path=elem["image_path"],
            left=elem["left"], top=elem["top"],
            width=elem.get("width"),
            height=elem.get("height"),
        )

    elif etype == "chart":
        _add_chart(
            slide,
            chart_type=elem.get("chart_type", "column"),
            left=elem.get("left", 0.9),
            top=elem.get("top", 1.15),
            width=elem.get("width", 11.5),
            height=elem.get("height", 5.0),
            categories=elem.get("categories", []),
            series=elem.get("series", []),
            title=elem.get("title"),
            legend_position=elem.get("legend_position", "bottom"),
            legend_font_size=elem.get("legend_font_size"),
            data_labels_show=elem.get("data_labels_show", False),
            data_labels_position=elem.get("data_labels_position", "outside_end"),
            data_labels_number_format=elem.get("data_labels_number_format"),
            data_labels_font_size=elem.get("data_labels_font_size"),
            data_labels_font_color=elem.get("data_labels_font_color"),
            axis_value_title=elem.get("axis_value_title"),
            axis_value_min=elem.get("axis_value_min"),
            axis_value_max=elem.get("axis_value_max"),
            axis_value_major_unit=elem.get("axis_value_major_unit"),
            axis_value_gridlines=elem.get("axis_value_gridlines", True),
            axis_value_number_format=elem.get("axis_value_number_format"),
            axis_value_visible=elem.get("axis_value_visible", True),
            axis_category_visible=elem.get("axis_category_visible", True),
            gap_width=elem.get("gap_width"),
            overlap=elem.get("overlap"),
            theme=theme,
        )

    elif etype == "icon":
        _add_icon(
            slide,
            icon_id=elem["icon_id"],
            left=elem.get("left", 0),
            top=elem.get("top", 0),
            width=elem.get("width"),
            height=elem.get("height"),
            color=elem.get("color"),
            outline_color=elem.get("outline_color"),
            theme=theme,
        )

    elif etype == "connector":
        _add_connector(
            slide,
            begin_x=elem["begin_x"],
            begin_y=elem["begin_y"],
            end_x=elem["end_x"],
            end_y=elem["end_y"],
            connector_type=elem.get("connector_type", "straight"),
            color=elem.get("color"),
            width=elem.get("width"),
            dash_style=elem.get("dash_style"),
            begin_arrow=elem.get("begin_arrow", "none"),
            end_arrow=elem.get("end_arrow", "triangle"),
            arrow_size=elem.get("arrow_size", "medium"),
            theme=theme,
        )

    elif etype == "callout":
        _add_callout(
            slide,
            text=elem.get("text", ""),
            target_x=elem["target_x"],
            target_y=elem["target_y"],
            label_x=elem.get("label_x"),
            label_y=elem.get("label_y"),
            label_width=elem.get("label_width", 2.0),
            label_height=elem.get("label_height", 0.4),
            connector_type=elem.get("connector_type", "straight"),
            font_size=elem.get("font_size", 10),
            font_color=elem.get("font_color"),
            font_bold=elem.get("font_bold", True),
            line_color=elem.get("line_color"),
            line_width=elem.get("line_width", 1.0),
            arrow_end=elem.get("arrow_end", "triangle"),
            bg_color=elem.get("bg_color"),
            border_color=elem.get("border_color"),
            theme=theme,
        )

    elif etype == "card_grid":
        _add_card_grid(
            slide,
            cards=elem.get("cards", []),
            y=elem.get("y"),
            theme=theme,
        )

    elif etype:
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Unknown element type '{etype}'. Supported: textbox, shape, table, "
            f"kpi_row, bullet_block, image, chart, icon, connector, callout, card_grid",
            )


def build_slide(file_path, spec_json):
    """File-based wrapper: build an entire slide from a JSON spec.
    One file open, one save. Returns slide index and element count."""
    import json
    spec = json.loads(spec_json) if isinstance(spec_json, str) else spec_json
    prs = open_pptx(file_path)
    slide, idx = _build_slide(prs, spec)
    save_pptx(prs, file_path)
    n_elements = len(spec.get("elements", []))
    title = spec.get("title", "(no title)")
    warning = ""
    if len(title) > 32:
        warning = f"\n⚠ Title is {len(title)} chars — consider shortening to ~30 chars."
    return f"Built slide [{idx}] with {n_elements} elements: {title}{warning}"


def build_deck(file_path, slides_json):
    """File-based wrapper: build an entire deck from a list of slide specs.
    Single file open/save for the whole deck."""
    import json
    slides = json.loads(slides_json) if isinstance(slides_json, str) else slides_json
    prs = open_pptx(file_path)
    results = []
    for i, spec in enumerate(slides):
        slide, idx = _build_slide(prs, spec)
        results.append(f"[{idx}] {spec.get('layout', 'content')}: {spec.get('title', '')[:40]}")
    save_pptx(prs, file_path)
    return f"Built {len(slides)} slides:\n" + "\n".join(results)


# ── File-based public wrappers ──────────────────────────────────


def add_content_slide(file_path, title, source=None, page_number=None):
    """File-based wrapper: add a content slide with action title."""
    prs = open_pptx(file_path)
    slide, idx = _add_content_slide(prs, title, source=source, page_number=page_number)
    save_pptx(prs, file_path)
    warning = ""
    if len(title) > 32:
        warning = f"\n⚠ Title is {len(title)} chars (>{32}). May wrap to 2 lines — consider shortening to ~30 chars for single-line display."
    return f"Added content slide [{idx}] with title: {title}{warning}"


def add_section_divider(file_path, title, subtitle=""):
    """File-based wrapper: add a section divider slide."""
    prs = open_pptx(file_path)
    slide, idx = _add_section_divider(prs, title, subtitle)
    save_pptx(prs, file_path)
    return f"Added section divider [{idx}] with title: {title}"


def add_kpi_row(file_path, slide_index, kpis_json, y):
    """File-based wrapper: add a row of KPI callout boxes."""
    import json
    kpis = json.loads(kpis_json) if isinstance(kpis_json, str) else kpis_json
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    indices = _add_kpi_row(slide, kpis, y)
    save_pptx(prs, file_path)
    return f"Added {len(kpis)} KPI boxes on slide [{slide_index}]"


def add_bullet_block(file_path, slide_index, items_json, left, top, width, height):
    """File-based wrapper: add a bulleted text block."""
    import json
    items = json.loads(items_json) if isinstance(items_json, str) else items_json
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    idx = _add_bullet_block(slide, items, left, top, width, height)
    save_pptx(prs, file_path)
    return f"Added bullet block [{idx}] on slide [{slide_index}]"
