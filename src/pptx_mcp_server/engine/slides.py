"""
Slide management -- info, read, add, delete, duplicate, background.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from copy import deepcopy
import lxml.etree as etree

from .pptx_io import (
    EngineError, ErrorCode,
    open_pptx, save_pptx, _get_slide, _parse_color,
)

from ..theme import Theme, resolve_color


# ── In-memory primitives ────────────────────────────────────────


def _get_presentation_info(prs: Presentation) -> str:
    """In-memory: get overview of a presentation."""
    w = prs.slide_width
    h = prs.slide_height
    lines = [
        f"Slides: {len(prs.slides)}",
        f"Dimensions: {w / 914400:.3f}\" x {h / 914400:.3f}\"",
        f"Slide layouts: {len(prs.slide_layouts)}",
        "---",
    ]
    for i, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    title = txt[:80]
                    break
        shape_count = len(slide.shapes)
        lines.append(f"[{i}] {shape_count} shapes | {title or '(no text)'}")
    return "\n".join(lines)


def _read_slide(prs: Presentation, slide_index: int) -> str:
    """In-memory: read detailed content of a single slide."""
    slide = _get_slide(prs, slide_index)
    lines = [f"Slide [{slide_index}] -- {len(slide.shapes)} shapes"]
    lines.append("---")

    for si, shape in enumerate(slide.shapes):
        shape_type = _shape_type_name(shape)
        pos = f"({shape.left / 914400:.2f}\", {shape.top / 914400:.2f}\")"
        size = f"{shape.width / 914400:.2f}\" x {shape.height / 914400:.2f}\""
        lines.append(f"  [{si}] {shape_type} @ {pos} size {size}")

        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                text = para.text
                if text:
                    lines.append(f"       p[{pi}]: {text}")

        if shape.has_table:
            table = shape.table
            lines.append(f"       table: {table.rows.__len__()} rows x {len(table.columns)} cols")
            for ri, row in enumerate(table.rows):
                cells = [row.cells[ci].text for ci in range(len(table.columns))]
                lines.append(f"       row[{ri}]: {' | '.join(cells)}")

    return "\n".join(lines)


def _add_slide(prs: Presentation, layout_index: int = 6, theme: Theme = None) -> int:
    """In-memory: add a blank slide. Returns new slide index."""
    if layout_index < 0 or layout_index >= len(prs.slide_layouts):
        raise EngineError(
            ErrorCode.INDEX_OUT_OF_RANGE,
            f"Layout index {layout_index} out of range (0-{len(prs.slide_layouts) - 1})",
        )
    layout = prs.slide_layouts[layout_index]
    prs.slides.add_slide(layout)
    return len(prs.slides) - 1


def _move_slide(prs: Presentation, from_index: int, to_index: int, theme: Theme = None) -> None:
    """In-memory: move a slide from one position to another."""
    n = len(prs.slides)
    if from_index < 0 or from_index >= n:
        raise EngineError(ErrorCode.SLIDE_NOT_FOUND, f"from_index {from_index} out of range (0-{n-1})")
    if to_index < 0 or to_index >= n:
        raise EngineError(ErrorCode.SLIDE_NOT_FOUND, f"to_index {to_index} out of range (0-{n-1})")
    if from_index == to_index:
        return

    sldIdLst = prs._element.sldIdLst
    sldIds = list(sldIdLst)
    el = sldIds[from_index]
    sldIdLst.remove(el)
    # Re-read after removal
    sldIds = list(sldIdLst)
    if to_index >= len(sldIds):
        sldIdLst.append(el)
    else:
        sldIds[to_index].addprevious(el)


def _delete_slide(prs: Presentation, slide_index: int, theme: Theme = None) -> None:
    """In-memory: delete a slide by index."""
    slide = _get_slide(prs, slide_index)
    rId = None
    for rel in prs.part.rels.values():
        if rel.target_part == slide.part:
            rId = rel.rId
            break
    if rId:
        sldIdLst = prs._element.sldIdLst
        for sldId in list(sldIdLst):
            if sldId.get(qn("r:id")) == rId:
                sldIdLst.remove(sldId)
                break
        prs.part.drop_rel(rId)


def _duplicate_slide(prs: Presentation, slide_index: int, theme: Theme = None) -> int:
    """In-memory: duplicate a slide (appended at end). Returns new slide index."""
    slide = _get_slide(prs, slide_index)
    layout = slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Clear the default placeholder shapes that add_slide creates
    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "sp":
            spTree.remove(child)

    # Copy all elements from source
    for elem in slide.shapes._spTree:
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
            spTree.append(deepcopy(elem))

    return len(prs.slides) - 1


def _set_slide_background(prs: Presentation, slide_index: int, color: str, theme: Theme = None) -> None:
    """In-memory: set solid background color for a slide."""
    if theme:
        color = resolve_color(theme, color)
    slide = _get_slide(prs, slide_index)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _parse_color(color)


# ── File-based public wrappers ──────────────────────────────────


def get_presentation_info(file_path: str) -> str:
    """Get overview of a PPTX file."""
    prs = open_pptx(file_path)
    info = _get_presentation_info(prs)
    return f"File: {file_path}\n{info}"


def read_slide(file_path: str, slide_index: int) -> str:
    """Read detailed content of a single slide."""
    prs = open_pptx(file_path)
    return _read_slide(prs, slide_index)


def add_slide(file_path: str, layout_index: int = 6) -> str:
    """Add a blank slide. Layout 6 is typically 'Blank'."""
    prs = open_pptx(file_path)
    slide_idx = _add_slide(prs, layout_index)
    save_pptx(prs, file_path)
    return f"Added slide [{slide_idx}] with layout {layout_index}"


def move_slide(file_path: str, from_index: int, to_index: int) -> str:
    """Move a slide from one position to another."""
    prs = open_pptx(file_path)
    _move_slide(prs, from_index, to_index)
    save_pptx(prs, file_path)
    return f"Moved slide [{from_index}] to [{to_index}]"


def delete_slide(file_path: str, slide_index: int) -> str:
    """Delete a slide by index."""
    prs = open_pptx(file_path)
    _delete_slide(prs, slide_index)
    save_pptx(prs, file_path)
    return f"Deleted slide [{slide_index}]"


def duplicate_slide(file_path: str, slide_index: int) -> str:
    """Duplicate a slide (appended at end)."""
    prs = open_pptx(file_path)
    new_idx = _duplicate_slide(prs, slide_index)
    save_pptx(prs, file_path)
    return f"Duplicated slide [{slide_index}] -> [{new_idx}]"


def set_slide_background(file_path: str, slide_index: int, color: str) -> str:
    """Set solid background color for a slide."""
    prs = open_pptx(file_path)
    _set_slide_background(prs, slide_index, color)
    save_pptx(prs, file_path)
    return f"Set slide [{slide_index}] background to #{color.lstrip('#')}"


def _shape_type_name(shape) -> str:
    """Human-readable shape type name."""
    if shape.has_table:
        return "TABLE"
    if shape.has_text_frame:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "TEXTBOX"
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "SHAPE"
        return f"TEXT({st})"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "PICTURE"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "GROUP"
    return str(shape.shape_type)
