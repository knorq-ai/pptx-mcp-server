"""
Chart engine — create and format professional charts on slides.

Uses python-pptx's native chart API for bar, column, stacked, line,
pie, doughnut, area, and radar chart types.
"""

from __future__ import annotations

import difflib
from typing import List, Optional

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from .pptx_io import (
    EngineError,
    ErrorCode,
    _get_slide,
    _parse_color,
    open_pptx,
    save_pptx,
)
from ..theme import Theme, get_chart_color, resolve_color, MCKINSEY

# ---------------------------------------------------------------------------
# Chart type mapping
# ---------------------------------------------------------------------------

_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked_bar_100": XL_CHART_TYPE.BAR_STACKED_100,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_column_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "radar": XL_CHART_TYPE.RADAR,
}

_LABEL_POSITION_MAP = {
    "above": XL_LABEL_POSITION.ABOVE,
    "below": XL_LABEL_POSITION.BELOW,
    "best_fit": XL_LABEL_POSITION.BEST_FIT,
    "center": XL_LABEL_POSITION.CENTER,
    "inside_base": XL_LABEL_POSITION.INSIDE_BASE,
    "inside_end": XL_LABEL_POSITION.INSIDE_END,
    "left": XL_LABEL_POSITION.LEFT,
    "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
    "right": XL_LABEL_POSITION.RIGHT,
}

_LEGEND_POSITION_MAP = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "corner": XL_LEGEND_POSITION.CORNER,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
    "top": XL_LEGEND_POSITION.TOP,
}

# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------


def _fuzzy_suggest(value: str, valid: list[str], label: str) -> str:
    """Return a 'did you mean?' suggestion or empty string."""
    matches = difflib.get_close_matches(value, valid, n=1, cutoff=0.5)
    if matches:
        return f" Did you mean '{matches[0]}'?"
    return f" Valid {label}: {', '.join(sorted(valid))}"


def _validate_chart_data(
    chart_type: str,
    categories: list,
    series: list,
) -> None:
    """Validate chart data before creation.  Raises EngineError on problems."""

    # chart_type
    if chart_type not in _CHART_TYPE_MAP:
        hint = _fuzzy_suggest(chart_type, list(_CHART_TYPE_MAP.keys()), "types")
        raise EngineError(
            ErrorCode.CHART_ERROR,
            f"Unknown chart_type '{chart_type}'.{hint}",
        )

    # categories
    if not categories:
        raise EngineError(ErrorCode.CHART_ERROR, "categories must not be empty.")

    # series
    if not series:
        raise EngineError(ErrorCode.CHART_ERROR, "series must not be empty.")

    n_cats = len(categories)
    for i, s in enumerate(series):
        name = s.get("name", f"Series {i}")

        # values presence
        values = s.get("values")
        if values is None:
            raise EngineError(
                ErrorCode.CHART_ERROR,
                f"Series '{name}' is missing 'values'.",
            )

        # length match
        if len(values) != n_cats:
            raise EngineError(
                ErrorCode.CHART_ERROR,
                f"Series '{name}' has {len(values)} values but {n_cats} categories were provided.",
            )

        # numeric check (None allowed for gaps)
        for j, v in enumerate(values):
            if v is not None and not isinstance(v, (int, float)):
                raise EngineError(
                    ErrorCode.CHART_ERROR,
                    f"Series '{name}' value at index {j} is not numeric: {v!r}",
                )


# ---------------------------------------------------------------------------
# Core chart creation
# ---------------------------------------------------------------------------


def _add_chart(
    slide,
    chart_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    categories: list[str],
    series: list[dict],
    *,
    title: Optional[str] = None,
    legend_position: Optional[str] = "bottom",
    legend_font_size: Optional[float] = None,
    data_labels_show: bool = False,
    data_labels_position: str = "outside_end",
    data_labels_number_format: Optional[str] = None,
    data_labels_font_size: Optional[float] = None,
    data_labels_font_color: Optional[str] = None,
    axis_value_title: Optional[str] = None,
    axis_value_min: Optional[float] = None,
    axis_value_max: Optional[float] = None,
    axis_value_major_unit: Optional[float] = None,
    axis_value_gridlines: bool = True,
    axis_value_number_format: Optional[str] = None,
    axis_value_visible: bool = True,
    axis_category_visible: bool = True,
    gap_width: Optional[int] = None,
    overlap: Optional[int] = None,
    theme: Optional[Theme] = None,
) -> int:
    """Add a chart to a slide.  Returns shape index.

    Parameters
    ----------
    chart_type : str
        One of: bar, stacked_bar, stacked_bar_100, column, stacked_column,
        stacked_column_100, line, line_markers, pie, area, area_stacked,
        doughnut, radar.
    categories : list[str]
        Category axis labels.
    series : list[dict]
        Each dict: {"name": str, "values": list[float|None], "color": str (hex, optional)}.
    """
    theme = theme or MCKINSEY

    _validate_chart_data(chart_type, categories, series)

    xl_type = _CHART_TYPE_MAP[chart_type]

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name", ""), tuple(s["values"]))

    # Create chart
    graphic_frame = slide.shapes.add_chart(
        xl_type,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data,
    )
    chart = graphic_frame.chart

    # --- Series colors ---
    _apply_series_colors(chart, series, theme)

    # --- Chart title ---
    if title:
        chart.has_title = True
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = title
        tf = chart.chart_title.text_frame
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.name = theme.fonts.get("body", "Arial")
                run.font.color.rgb = _parse_color(
                    resolve_color(theme, "primary")
                )
    else:
        chart.has_title = False

    # --- Legend ---
    _apply_legend(chart, legend_position, legend_font_size, theme)

    # --- Data labels ---
    if data_labels_show:
        _apply_data_labels(
            chart, data_labels_position, data_labels_number_format,
            data_labels_font_size, data_labels_font_color, theme,
        )

    # --- Axes ---
    _apply_axes(
        chart, chart_type,
        axis_value_title, axis_value_min, axis_value_max,
        axis_value_major_unit, axis_value_gridlines,
        axis_value_number_format, axis_value_visible,
        axis_category_visible, theme,
    )

    # --- Gap / Overlap ---
    if chart.plots:
        plot = chart.plots[0]
        if gap_width is not None and hasattr(plot, "gap_width"):
            plot.gap_width = gap_width
        if overlap is not None and hasattr(plot, "overlap"):
            plot.overlap = overlap

    return len(list(slide.shapes)) - 1


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------


def _apply_series_colors(chart, series: list[dict], theme: Theme) -> None:
    """Apply fill colors to each chart series."""
    if not chart.plots:
        return
    plot = chart.plots[0]
    for i, series_obj in enumerate(plot.series):
        # Precedence: explicit series color > theme chart_colors > fallback
        color_hex = None
        if i < len(series) and series[i].get("color"):
            color_hex = resolve_color(theme, series[i]["color"])
        else:
            color_hex = get_chart_color(theme, i)

        if color_hex:
            color_hex = color_hex.lstrip("#")
            fill = series_obj.format.fill
            fill.solid()
            fill.fore_color.rgb = _parse_color(color_hex)

            # Also set line color for line/radar/area charts
            line = series_obj.format.line
            line.color.rgb = _parse_color(color_hex)


def _apply_legend(
    chart,
    position: Optional[str],
    font_size: Optional[float],
    theme: Theme,
) -> None:
    """Apply legend settings."""
    if position is None:
        chart.has_legend = False
        return

    chart.has_legend = True
    pos_key = position.lower()
    if pos_key in _LEGEND_POSITION_MAP:
        chart.legend.position = _LEGEND_POSITION_MAP[pos_key]
    chart.legend.include_in_layout = False

    size = font_size or 9
    chart.legend.font.size = Pt(size)
    chart.legend.font.name = theme.fonts.get("body", "Arial")


def _apply_data_labels(
    chart,
    position: str,
    number_format: Optional[str],
    font_size: Optional[float],
    font_color: Optional[str],
    theme: Theme,
) -> None:
    """Apply data labels to all series."""
    if not chart.plots:
        return
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.show_category_name = False
    data_labels.show_series_name = False

    # Position
    pos_key = position.lower() if position else "outside_end"
    if pos_key in _LABEL_POSITION_MAP:
        data_labels.position = _LABEL_POSITION_MAP[pos_key]

    # Number format
    if number_format:
        data_labels.number_format = number_format
        data_labels.number_format_is_linked = False

    # Font
    size = font_size or 8
    data_labels.font.size = Pt(size)
    data_labels.font.name = theme.fonts.get("body", "Arial")

    if font_color:
        data_labels.font.color.rgb = _parse_color(
            resolve_color(theme, font_color).lstrip("#")
        )


def _apply_axes(
    chart,
    chart_type: str,
    value_title: Optional[str],
    value_min: Optional[float],
    value_max: Optional[float],
    value_major_unit: Optional[float],
    value_gridlines: bool,
    value_number_format: Optional[str],
    value_visible: bool,
    category_visible: bool,
    theme: Theme,
) -> None:
    """Apply axis formatting.  Pie/doughnut charts have no axes."""
    if chart_type in ("pie", "doughnut"):
        return

    font_name = theme.fonts.get("body", "Arial")
    axis_color = resolve_color(theme, "text_secondary").lstrip("#")

    # Value axis
    try:
        value_axis = chart.value_axis
    except (ValueError, AttributeError):
        return

    value_axis.visible = value_visible

    if value_title:
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = value_title
        for p in value_axis.axis_title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)
                run.font.name = font_name

    if value_min is not None:
        value_axis.minimum_scale = value_min
    if value_max is not None:
        value_axis.maximum_scale = value_max
    if value_major_unit is not None:
        value_axis.major_unit = value_major_unit

    value_axis.has_major_gridlines = value_gridlines
    if value_gridlines and value_axis.major_gridlines:
        gl_format = value_axis.major_gridlines.format
        gl_format.line.color.rgb = _parse_color(
            resolve_color(theme, "border").lstrip("#")
        )
        gl_format.line.width = Pt(0.3)

    if value_number_format:
        value_axis.tick_labels.number_format = value_number_format
        value_axis.tick_labels.number_format_is_linked = False

    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.font.name = font_name
    value_axis.tick_labels.font.color.rgb = _parse_color(axis_color)

    # Category axis
    try:
        cat_axis = chart.category_axis
    except (ValueError, AttributeError):
        return

    cat_axis.visible = category_visible
    cat_axis.has_major_gridlines = False

    cat_axis.tick_labels.font.size = Pt(8)
    cat_axis.tick_labels.font.name = font_name
    cat_axis.tick_labels.font.color.rgb = _parse_color(axis_color)


# ---------------------------------------------------------------------------
# File-based wrapper
# ---------------------------------------------------------------------------


def add_chart(file_path: str, slide_index: int, chart_spec: dict) -> str:
    """File-based wrapper: add a chart to a slide from a spec dict.

    Returns a summary string.
    """
    from ..theme import get_theme

    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)

    # Resolve theme from spec
    theme = None
    theme_name = chart_spec.get("theme")
    if theme_name:
        theme = get_theme(theme_name)
    theme = theme or MCKINSEY

    idx = _add_chart(
        slide,
        chart_type=chart_spec.get("chart_type", "column"),
        left=chart_spec.get("left", 0.9),
        top=chart_spec.get("top", 1.15),
        width=chart_spec.get("width", 11.5),
        height=chart_spec.get("height", 5.0),
        categories=chart_spec.get("categories", []),
        series=chart_spec.get("series", []),
        title=chart_spec.get("title"),
        legend_position=chart_spec.get("legend_position", "bottom"),
        legend_font_size=chart_spec.get("legend_font_size"),
        data_labels_show=chart_spec.get("data_labels_show", False),
        data_labels_position=chart_spec.get("data_labels_position", "outside_end"),
        data_labels_number_format=chart_spec.get("data_labels_number_format"),
        data_labels_font_size=chart_spec.get("data_labels_font_size"),
        data_labels_font_color=chart_spec.get("data_labels_font_color"),
        axis_value_title=chart_spec.get("axis_value_title"),
        axis_value_min=chart_spec.get("axis_value_min"),
        axis_value_max=chart_spec.get("axis_value_max"),
        axis_value_major_unit=chart_spec.get("axis_value_major_unit"),
        axis_value_gridlines=chart_spec.get("axis_value_gridlines", True),
        axis_value_number_format=chart_spec.get("axis_value_number_format"),
        axis_value_visible=chart_spec.get("axis_value_visible", True),
        axis_category_visible=chart_spec.get("axis_category_visible", True),
        gap_width=chart_spec.get("gap_width"),
        overlap=chart_spec.get("overlap"),
        theme=theme,
    )

    n_series = len(chart_spec.get("series", []))
    n_cats = len(chart_spec.get("categories", []))
    save_pptx(prs, file_path)
    return (
        f"Added {chart_spec.get('chart_type', 'column')} chart [{idx}] on slide [{slide_index}] "
        f"({n_series} series, {n_cats} categories)"
    )
