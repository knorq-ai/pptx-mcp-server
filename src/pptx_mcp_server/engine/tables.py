"""
Table operations -- add table, edit cells, format.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from .pptx_io import (
    EngineError, ErrorCode,
    open_pptx, save_pptx, _get_slide, _get_shape, _parse_color,
)

from ..theme import Theme, resolve_color


# ── In-memory primitives ────────────────────────────────────────


def _add_table(
    slide,
    rows,
    left,
    top,
    width,
    col_widths=None,
    row_height=0.36,
    font_size=12,
    header_bg="051C2C",
    header_fg="FFFFFF",
    alt_row_bg="F5F5F5",
    border_color="D0D0D0",
    border_width=0.5,
    no_vertical_borders=True,
    theme=None,
):
    """In-memory: add a professionally formatted table. Returns shape_index."""
    if theme:
        header_bg = header_bg or resolve_color(theme, theme.table.get("header_bg", "primary"))
        header_fg = header_fg or resolve_color(theme, theme.table.get("header_fg", "white"))
        alt_row_bg = alt_row_bg or resolve_color(theme, theme.table.get("alt_row_bg", "bg_alt"))
        border_color = border_color or resolve_color(theme, theme.table.get("border_color", "border"))
        font_size = font_size or theme.sizes.get("table", 12)
        row_height = row_height or theme.table.get("row_height", 0.36)
        border_width = border_width or theme.table.get("border_width", 0.5)
        no_vertical_borders = theme.table.get("no_vertical_borders", True)
        font_name = theme.fonts.get("body", "Arial")
    else:
        font_name = "Arial"

    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    if n_cols == 0:
        raise EngineError(ErrorCode.TABLE_ERROR, "Table must have at least one column")

    table_height = Inches(row_height * n_rows)
    table_width = Inches(width)

    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), table_width, table_height)
    table = shape.table

    # Column widths
    if col_widths:
        for i, pct in enumerate(col_widths):
            if i < n_cols:
                table.columns[i].width = int(table_width * pct)

    # Populate cells
    for r_idx, row_data in enumerate(rows):
        is_header = header_bg and r_idx == 0
        table.rows[r_idx].height = Inches(row_height)

        for c_idx, cell_text in enumerate(row_data):
            if c_idx >= n_cols:
                break
            cell = table.cell(r_idx, c_idx)
            cell.text = ""

            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.name = font_name
            p.font.size = Pt(font_size)

            if is_header:
                p.font.bold = True
                p.font.color.rgb = _parse_color(header_fg) if header_fg else RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _parse_color(header_bg) if header_bg else RGBColor(0x05, 0x1C, 0x2C)
            else:
                p.font.color.rgb = RGBColor(0x05, 0x1C, 0x2C)
                # Alternating row shading
                if alt_row_bg and r_idx % 2 == 0 and r_idx > 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _parse_color(alt_row_bg)
                else:
                    cell.fill.background()

                # Right-align numeric values
                stripped = str(cell_text).strip()
                if stripped and (stripped[0].isdigit() or stripped[0] in "¥$€△▲-+~≈∼%"):
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)

            # Borders
            _set_cell_borders(cell, is_header, r_idx == n_rows - 1,
                              border_color or "D0D0D0", border_width,
                              no_vertical_borders, header_bg or "051C2C")

    return len(list(slide.shapes)) - 1


def _edit_table_cell(
    slide,
    shape_index,
    row,
    col,
    text=None,
    font_size=None,
    font_color=None,
    bold=None,
    bg_color=None,
    alignment=None,
    theme=None,
):
    """In-memory: edit a single table cell."""
    shape = _get_shape(slide, shape_index)

    if not shape.has_table:
        raise EngineError(ErrorCode.TABLE_ERROR, f"Shape [{shape_index}] is not a table")

    if theme:
        if font_color:
            font_color = resolve_color(theme, font_color)
        if bg_color:
            bg_color = resolve_color(theme, bg_color)

    table = shape.table
    if row < 0 or row >= len(table.rows):
        raise EngineError(ErrorCode.INDEX_OUT_OF_RANGE, f"Row {row} out of range")
    if col < 0 or col >= len(table.columns):
        raise EngineError(ErrorCode.INDEX_OUT_OF_RANGE, f"Col {col} out of range")

    cell = table.cell(row, col)
    p = cell.text_frame.paragraphs[0]

    if text is not None:
        p.text = text
    if font_size is not None:
        p.font.size = Pt(font_size)
    if font_color is not None:
        p.font.color.rgb = _parse_color(font_color)
    if bold is not None:
        p.font.bold = bold
    if bg_color is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _parse_color(bg_color)
    if alignment is not None:
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        if alignment in align_map:
            p.alignment = align_map[alignment]


def _edit_table_cells(
    slide,
    shape_index,
    edits,
    theme=None,
):
    """In-memory: batch edit multiple table cells."""
    shape = _get_shape(slide, shape_index)

    if not shape.has_table:
        raise EngineError(ErrorCode.TABLE_ERROR, f"Shape [{shape_index}] is not a table")

    table = shape.table
    count = 0
    for edit in edits:
        r = edit.get("row", 0)
        c = edit.get("col", 0)
        if r < 0 or r >= len(table.rows) or c < 0 or c >= len(table.columns):
            continue
        cell = table.cell(r, c)
        p = cell.text_frame.paragraphs[0]
        if "text" in edit:
            p.text = str(edit["text"])
        if "font_size" in edit:
            p.font.size = Pt(edit["font_size"])
        if "font_color" in edit:
            color = edit["font_color"]
            if theme:
                color = resolve_color(theme, color)
            p.font.color.rgb = _parse_color(color)
        if "bold" in edit:
            p.font.bold = edit["bold"]
        if "bg_color" in edit:
            color = edit["bg_color"]
            if theme:
                color = resolve_color(theme, color)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _parse_color(color)
        count += 1

    return count


def _format_table(
    slide,
    shape_index,
    font_name=None,
    font_size=None,
    header_bg=None,
    header_fg=None,
    alt_row_bg=None,
    theme=None,
):
    """In-memory: apply bulk formatting to an entire table."""
    shape = _get_shape(slide, shape_index)

    if not shape.has_table:
        raise EngineError(ErrorCode.TABLE_ERROR, f"Shape [{shape_index}] is not a table")

    if theme:
        if header_bg:
            header_bg = resolve_color(theme, header_bg)
        if header_fg:
            header_fg = resolve_color(theme, header_fg)
        if alt_row_bg:
            alt_row_bg = resolve_color(theme, alt_row_bg)

    table = shape.table
    for r_idx, row in enumerate(table.rows):
        for c_idx in range(len(table.columns)):
            cell = table.cell(r_idx, c_idx)
            for p in cell.text_frame.paragraphs:
                if font_name:
                    p.font.name = font_name
                if font_size:
                    p.font.size = Pt(font_size)
                if r_idx == 0:
                    if header_bg:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _parse_color(header_bg)
                    if header_fg:
                        p.font.color.rgb = _parse_color(header_fg)
                else:
                    if alt_row_bg and r_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _parse_color(alt_row_bg)


# ── File-based public wrappers ──────────────────────────────────


def add_table(
    file_path,
    slide_index,
    rows,
    left,
    top,
    width,
    col_widths=None,
    row_height=0.36,
    font_size=12,
    header_bg="051C2C",
    header_fg="FFFFFF",
    alt_row_bg="F5F5F5",
    border_color="D0D0D0",
    border_width=0.5,
    no_vertical_borders=True,
):
    """File-based wrapper: add a table to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    shape_idx = _add_table(
        slide, rows, left, top, width,
        col_widths, row_height, font_size,
        header_bg, header_fg, alt_row_bg, border_color,
        border_width, no_vertical_borders,
    )
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    save_pptx(prs, file_path)
    return f"Added table [{shape_idx}] ({n_rows}x{n_cols}) on slide [{slide_index}]"


def edit_table_cell(
    file_path,
    slide_index,
    shape_index,
    row,
    col,
    text=None,
    font_size=None,
    font_color=None,
    bold=None,
    bg_color=None,
    alignment=None,
):
    """File-based wrapper: edit a single table cell."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    _edit_table_cell(slide, shape_index, row, col, text, font_size, font_color, bold, bg_color, alignment)
    save_pptx(prs, file_path)
    return f"Edited cell ({row},{col}) in table [{shape_index}] on slide [{slide_index}]"


def edit_table_cells(
    file_path,
    slide_index,
    shape_index,
    edits,
):
    """File-based wrapper: batch edit multiple table cells."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    count = _edit_table_cells(slide, shape_index, edits)
    save_pptx(prs, file_path)
    return f"Edited {count} cells in table [{shape_index}] on slide [{slide_index}]"


def format_table(
    file_path,
    slide_index,
    shape_index,
    font_name=None,
    font_size=None,
    header_bg=None,
    header_fg=None,
    alt_row_bg=None,
):
    """File-based wrapper: apply bulk formatting to a table."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    _format_table(slide, shape_index, font_name, font_size, header_bg, header_fg, alt_row_bg)
    save_pptx(prs, file_path)
    return f"Formatted table [{shape_index}] on slide [{slide_index}]"


# ── Border helpers ──────────────────────────────────────────────


def _set_cell_borders(
    cell, is_header, is_last_row,
    border_color, border_width,
    no_vertical, header_border_color,
):
    """Set McKinsey-style cell borders: horizontal only, no vertical."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    def set_border(name, color_hex, width_pt):
        border_el = tcPr.find(qn(f"a:{name}"))
        if border_el is None:
            border_el = etree.SubElement(tcPr, qn(f"a:{name}"))
        border_el.set("w", str(int(width_pt * 12700)))
        border_el.set("cmpd", "sng")
        sf = border_el.find(qn("a:solidFill"))
        if sf is None:
            sf = etree.SubElement(border_el, qn("a:solidFill"))
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
        srgb.set("val", color_hex.lstrip("#"))

    def no_border(name):
        border_el = tcPr.find(qn(f"a:{name}"))
        if border_el is None:
            border_el = etree.SubElement(tcPr, qn(f"a:{name}"))
        border_el.set("w", "0")
        border_el.set("cmpd", "sng")
        sf = border_el.find(qn("a:solidFill"))
        if sf is not None:
            border_el.remove(sf)
        nf = border_el.find(qn("a:noFill"))
        if nf is None:
            etree.SubElement(border_el, qn("a:noFill"))

    if no_vertical:
        no_border("lnL")
        no_border("lnR")
    else:
        set_border("lnL", border_color, border_width)
        set_border("lnR", border_color, border_width)

    if is_header:
        no_border("lnT")
        set_border("lnB", header_border_color, 2.0)
    else:
        no_border("lnT")
        set_border("lnB", border_color, border_width)
