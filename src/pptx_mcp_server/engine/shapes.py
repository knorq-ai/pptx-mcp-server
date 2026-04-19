"""
Shape operations -- add textbox, add shape, edit text, add paragraph, delete, list.
"""

from __future__ import annotations

import unicodedata

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from .pptx_io import (
    EngineError, ErrorCode,
    open_pptx, save_pptx, _get_slide, _get_shape, _parse_color,
)
from .text_metrics import estimate_text_height, estimate_text_width, wrap_text

from ..theme import Theme, resolve_color

# DrawingML 名前空間 (OOXML 直接操作用)。
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Mapping of alignment strings to enums
_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

_ANCHOR_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}

# Common auto shapes
_SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "chevron": MSO_SHAPE.CHEVRON,
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "arrow_left": MSO_SHAPE.LEFT_ARROW,
    "arrow_up": MSO_SHAPE.UP_ARROW,
    "arrow_down": MSO_SHAPE.DOWN_ARROW,
    "callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "star_5": MSO_SHAPE.STAR_5_POINT,
}


# ── Internal helpers ────────────────────────────────────────────


def _set_east_asian_font(run, typeface: str) -> None:
    """run の ``<a:rPr>`` に ``<a:ea typeface="..."/>`` を設定する.

    python-pptx は east-asian typeface を直接公開しないため、lxml で
    ``rPr`` 要素を取得し ``<a:ea>`` 子要素を差し替える。既存の ``<a:ea>``
    は削除してから新しく追加することで重複を避ける。
    """
    rPr = run._r.get_or_add_rPr()
    for existing in rPr.findall(qn("a:ea")):
        rPr.remove(existing)
    ea = rPr.makeelement(qn("a:ea"), {"typeface": typeface})
    rPr.append(ea)


def _apply_font(
    paragraph,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    underline=None,
    theme=None,
    east_asian_font=None,
):
    """paragraph とその全 run に font 書式を適用する.

    挙動:
        - paragraph レベル (``pPr/defRPr`` 経由) にも値を書き込むが、
          authoritative な値は各 run の ``<a:rPr>`` に直接書き込む。これに
          より Keynote / Google Slides / PowerPoint Online のような
          ``defRPr`` を異なる優先順位で解決するクライアントでも意図した
          サイズ・フォントで描画される (issue #28)。
        - ``east_asian_font`` を指定した場合、各 run の ``rPr`` に
          ``<a:ea typeface="…"/>`` を追加する (issue #40)。
        - paragraph に run が存在しない場合は run レベル書き込みを skip する
          (caller 側で p.text 等により run を追加している想定)。
    """
    font = paragraph.font
    if font_name is not None:
        font.name = font_name
    if font_size is not None:
        font.size = Pt(font_size)
    if font_color is not None:
        color_hex = resolve_color(theme, font_color) if theme else font_color
        font.color.rgb = _parse_color(color_hex)
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic
    if underline is not None:
        font.underline = underline

    # run レベルにも同じ属性を書き込む。これにより ``<a:r><a:rPr sz="…"/>``
    # が確実に emit され、クライアント依存の resolve 経路に左右されない。
    for run in paragraph.runs:
        rfont = run.font
        if font_name is not None:
            rfont.name = font_name
        if font_size is not None:
            rfont.size = Pt(font_size)
        if font_color is not None:
            color_hex = resolve_color(theme, font_color) if theme else font_color
            rfont.color.rgb = _parse_color(color_hex)
        if bold is not None:
            rfont.bold = bold
        if italic is not None:
            rfont.italic = italic
        if underline is not None:
            rfont.underline = underline
        if east_asian_font is not None:
            _set_east_asian_font(run, east_asian_font)


def _apply_paragraph(
    paragraph,
    alignment=None,
    line_spacing=None,
):
    """Apply paragraph-level formatting."""
    if alignment and alignment in _ALIGN_MAP:
        paragraph.alignment = _ALIGN_MAP[alignment]
    if line_spacing is not None:
        paragraph.line_spacing = Pt(line_spacing)


# ── In-memory primitives ────────────────────────────────────────


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    alignment=None,
    vertical_anchor=None,
    word_wrap=True,
    line_spacing=None,
    underline=None,
    theme=None,
    east_asian_font=None,
):
    """In-memory: add textbox to an existing slide object. Returns shape_index.

    ``vertical_anchor`` は python-pptx の ``tf.vertical_anchor`` 経由で設定
    する (issue #41)。直接 ``bodyPr.set("anchor", ...)`` を呼ぶと round-trip
    で属性が重複する危険があるため、lxml 直接操作はしない。
    """
    if theme:
        font_name = font_name or theme.fonts.get("body")
        east_asian_font = east_asian_font or theme.fonts.get("east_asian")
        if font_color:
            font_color = resolve_color(theme, font_color)

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    if vertical_anchor and vertical_anchor in _ANCHOR_MAP:
        tf.vertical_anchor = _ANCHOR_MAP[vertical_anchor]

    p = tf.paragraphs[0]
    p.text = text

    _apply_font(
        p, font_name, font_size, font_color, bold, italic, underline,
        theme=theme, east_asian_font=east_asian_font,
    )
    _apply_paragraph(p, alignment, line_spacing)

    return len(list(slide.shapes)) - 1


def _add_image(slide, image_path, left, top, width=None, height=None):
    """In-memory: add an image to a slide. Returns shape_index.

    If only width is given, height is auto-calculated to maintain aspect ratio (and vice versa).
    If both are given, image is stretched to fit.
    If neither is given, image is placed at its native size.
    """
    import os
    if not os.path.exists(image_path):
        raise EngineError(ErrorCode.FILE_NOT_FOUND, f"Image not found: {image_path}")

    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)

    slide.shapes.add_picture(image_path, **kwargs)
    return len(list(slide.shapes)) - 1


def _add_shape(
    slide,
    shape_type,
    left,
    top,
    width,
    height,
    fill_color=None,
    line_color=None,
    line_width=None,
    no_line=False,
    text=None,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    alignment=None,
    theme=None,
):
    """In-memory: add an auto shape to a slide. Returns shape_index."""
    if theme:
        font_name = font_name or theme.fonts.get("body")
        if fill_color:
            fill_color = resolve_color(theme, fill_color)
        if line_color:
            line_color = resolve_color(theme, line_color)
        if font_color:
            font_color = resolve_color(theme, font_color)

    shape_enum = _SHAPE_MAP.get(shape_type.lower())
    if shape_enum is None:
        available = ", ".join(sorted(_SHAPE_MAP.keys()))
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Unknown shape type '{shape_type}'. Available: {available}",
        )

    shape = slide.shapes.add_shape(
        shape_enum, Inches(left), Inches(top), Inches(width), Inches(height),
    )

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(fill_color)

    if no_line:
        shape.line.fill.background()
    elif line_color:
        shape.line.color.rgb = _parse_color(line_color)
    if line_width is not None:
        shape.line.width = Pt(line_width)

    if text is not None:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        _apply_font(p, font_name, font_size, font_color, bold, None, None, theme=None)
        _apply_paragraph(p, alignment, None)

    return len(list(slide.shapes)) - 1


def _edit_text(
    slide,
    shape_index,
    text=None,
    paragraph_index=0,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    underline=None,
    alignment=None,
    line_spacing=None,
    theme=None,
):
    """In-memory: edit text content and formatting in a shape's paragraph."""
    shape = _get_shape(slide, shape_index)

    if not shape.has_text_frame:
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Shape [{shape_index}] does not have a text frame",
        )

    tf = shape.text_frame
    if paragraph_index < 0 or paragraph_index >= len(tf.paragraphs):
        raise EngineError(
            ErrorCode.INDEX_OUT_OF_RANGE,
            f"Paragraph index {paragraph_index} out of range (0-{len(tf.paragraphs) - 1})",
        )

    if theme:
        font_name = font_name or None  # don't override with theme default for edit
        if font_color:
            font_color = resolve_color(theme, font_color)

    p = tf.paragraphs[paragraph_index]
    if text is not None:
        p.text = text
    _apply_font(p, font_name, font_size, font_color, bold, italic, underline, theme=None)
    _apply_paragraph(p, alignment, line_spacing)


def _add_paragraph(
    slide,
    shape_index,
    text,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    underline=None,
    alignment=None,
    line_spacing=None,
    theme=None,
):
    """In-memory: append a new paragraph to an existing shape's text frame.
    Returns paragraph index."""
    shape = _get_shape(slide, shape_index)

    if not shape.has_text_frame:
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Shape [{shape_index}] does not have a text frame",
        )

    if theme:
        if font_color:
            font_color = resolve_color(theme, font_color)

    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = text
    _apply_font(p, font_name, font_size, font_color, bold, italic, underline, theme=None)
    _apply_paragraph(p, alignment, line_spacing)

    return len(tf.paragraphs) - 1


def _delete_shape(slide, shape_index, theme=None):
    """In-memory: delete a shape from a slide by index."""
    shape = _get_shape(slide, shape_index)
    sp = shape._element
    sp.getparent().remove(sp)


def _list_shapes(slide, slide_index, theme=None):
    """In-memory: list all shapes on a slide."""
    lines = [f"Slide [{slide_index}] -- {len(slide.shapes)} shapes"]
    for si, shape in enumerate(slide.shapes):
        stype = "TABLE" if shape.has_table else ("TEXTBOX" if shape.has_text_frame else str(shape.shape_type))
        pos = f"({shape.left / 914400:.2f}\", {shape.top / 914400:.2f}\")"
        size = f"{shape.width / 914400:.2f}\" x {shape.height / 914400:.2f}\""
        text_preview = ""
        if shape.has_text_frame:
            txt = shape.text_frame.text[:50]
            if txt:
                text_preview = f" \"{txt}\""
        lines.append(f"  [{si}] {stype} @ {pos} {size}{text_preview}")
    return "\n".join(lines)


# ── File-based public wrappers ──────────────────────────────────


def add_textbox(
    file_path,
    slide_index,
    left,
    top,
    width,
    height,
    text="",
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    alignment=None,
    vertical_anchor=None,
    word_wrap=True,
    line_spacing=None,
    underline=None,
):
    """File-based wrapper: add a text box to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    idx = _add_textbox(
        slide, left, top, width, height, text,
        font_name, font_size, font_color, bold, italic,
        alignment, vertical_anchor, word_wrap, line_spacing, underline,
    )
    save_pptx(prs, file_path)
    return f"Added textbox [{idx}] on slide [{slide_index}]"


def add_image(file_path, slide_index, image_path, left, top, width=None, height=None):
    """File-based wrapper: add an image to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    idx = _add_image(slide, image_path, left, top, width, height)
    save_pptx(prs, file_path)
    return f"Added image [{idx}] on slide [{slide_index}]"


def add_shape(
    file_path,
    slide_index,
    shape_type,
    left,
    top,
    width,
    height,
    fill_color=None,
    line_color=None,
    line_width=None,
    no_line=False,
    text=None,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    alignment=None,
):
    """File-based wrapper: add an auto shape to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    idx = _add_shape(
        slide, shape_type, left, top, width, height,
        fill_color, line_color, line_width, no_line,
        text, font_name, font_size, font_color, bold, alignment,
    )
    save_pptx(prs, file_path)
    return f"Added {shape_type} [{idx}] on slide [{slide_index}]"


def edit_text(
    file_path,
    slide_index,
    shape_index,
    text=None,
    paragraph_index=0,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    underline=None,
    alignment=None,
    line_spacing=None,
):
    """File-based wrapper: edit text content and formatting."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    _edit_text(
        slide, shape_index, text, paragraph_index,
        font_name, font_size, font_color, bold, italic, underline,
        alignment, line_spacing,
    )
    save_pptx(prs, file_path)
    return f"Edited shape [{shape_index}] paragraph [{paragraph_index}] on slide [{slide_index}]"


def add_paragraph(
    file_path,
    slide_index,
    shape_index,
    text,
    font_name=None,
    font_size=None,
    font_color=None,
    bold=None,
    italic=None,
    underline=None,
    alignment=None,
    line_spacing=None,
):
    """File-based wrapper: append a paragraph to a shape's text frame."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    p_idx = _add_paragraph(
        slide, shape_index, text,
        font_name, font_size, font_color, bold, italic, underline,
        alignment, line_spacing,
    )
    save_pptx(prs, file_path)
    return f"Added paragraph [{p_idx}] to shape [{shape_index}] on slide [{slide_index}]"


def delete_shape(file_path, slide_index, shape_index):
    """File-based wrapper: delete a shape from a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    _delete_shape(slide, shape_index)
    save_pptx(prs, file_path)
    return f"Deleted shape [{shape_index}] from slide [{slide_index}]"


def list_shapes(file_path, slide_index):
    """File-based wrapper: list all shapes on a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    return _list_shapes(slide, slide_index)


# ── Auto-fit textbox ────────────────────────────────────────────

# 内側 padding (inches、左右各側)。左右に同量の padding を見込むため、
# usable width は width - 2 * _AUTO_FIT_PADDING_PER_SIDE となる。
# cards.py / validation.py と単一定義を共有するため ``layout_constants`` から
# import する。alias 名は既存 import 互換のため維持する。
from .layout_constants import (
    TEXTBOX_INNER_PADDING_PER_SIDE as _AUTO_FIT_PADDING_PER_SIDE,
)

# 後方互換のための別名 (旧名)。新規コードは _AUTO_FIT_PADDING_PER_SIDE を使うこと。
_AUTO_FIT_PADDING: float = _AUTO_FIT_PADDING_PER_SIDE

# font size 縮小ステップ (pt)
_AUTO_FIT_STEP_PT: float = 0.5

# 高さ推定に使う行高倍率
_AUTO_FIT_LINE_HEIGHT: float = 1.2

# 省略記号
_ELLIPSIS: str = "\u2026"


def _fit_font_size(
    text: str,
    usable_width: float,
    height: float,
    font_name: str,
    font_size_pt: float,
    min_size_pt: float,
) -> tuple[float, bool]:
    """指定 box に収まる font size を二分探索ではなく 0.5pt ステップで決定する.

    Returns:
        (size, fits): ``size`` は採用する font size。``fits`` は最終的にその
        size で text が box 内に収まるかどうか (min に達してもなおオーバー
        フローする場合は False)。
    """
    size = float(font_size_pt)
    min_size = float(min_size_pt)
    while size > min_size:
        height_est = estimate_text_height(
            text, usable_width, size, font_name,
            line_height_factor=_AUTO_FIT_LINE_HEIGHT,
        )
        if height_est <= height:
            return size, True
        size = round(size - _AUTO_FIT_STEP_PT, 2)
    # min_size でチェック
    height_est = estimate_text_height(
        text, usable_width, min_size, font_name,
        line_height_factor=_AUTO_FIT_LINE_HEIGHT,
    )
    return min_size, height_est <= height


def _fit_font_size_single_line(
    text: str,
    usable_width: float,
    font_name: str,
    font_size_pt: float,
    min_size_pt: float,
) -> tuple[float, bool]:
    """単一行描画を前提に width に収まる最大 font size を 0.5pt 刻みで決定する.

    ``wrap=False`` モードで使う。高さは参照せず、``estimate_text_width`` が
    ``usable_width`` 以下となる最大 size を返す。min_size に達してもなお
    幅を超える場合は ``(min_size, False)`` を返す。

    Returns:
        (size, fits): ``size`` は採用する font size、``fits`` は決定 size で
        実際に単一行 width に収まるかどうか。
    """
    size = float(font_size_pt)
    min_size = float(min_size_pt)
    while size > min_size:
        width_est = estimate_text_width(text, size, font_name)
        if width_est <= usable_width:
            return size, True
        size = round(size - _AUTO_FIT_STEP_PT, 2)
    # min_size でチェック
    width_est = estimate_text_width(text, min_size, font_name)
    return min_size, width_est <= usable_width


def _truncate_to_fit_single_line(
    text: str,
    usable_width: float,
    font_name: str,
    size_pt: float,
) -> str:
    """``size_pt`` の size で ``usable_width`` に収まるよう末尾を省略記号で切り
    詰める (単一行版).

    ``_truncate_to_fit`` は wrap を前提としているため、``wrap=False`` モード用に
    行頭から grapheme cluster を 1 つずつ削りながら ``text + 省略記号`` の幅が
    ``usable_width`` 以下になるまで縮める専用実装を用意する。省略記号単体でも
    収まらない場合は省略記号のみを返す。
    """
    if not text:
        return text
    candidate = text
    while candidate and estimate_text_width(candidate + _ELLIPSIS, size_pt, font_name) > usable_width:
        new_candidate = _strip_last_grapheme(candidate)
        if new_candidate == candidate:
            # 保険: 進まない場合は code unit で 1 文字削る。
            new_candidate = candidate[:-1]
        candidate = new_candidate
    return (candidate + _ELLIPSIS) if candidate else _ELLIPSIS


def _strip_last_grapheme(s: str) -> str:
    """``s`` の末尾 grapheme cluster (近似) を 1 つ除去した文字列を返す.

    Unicode の正式な grapheme cluster は ``regex`` の ``\\X`` や ICU が必要
    だが、stdlib のみで ``_truncate_to_fit`` が壊しやすいケース (ZWJ emoji
    sequence / 結合文字 / variation selector) を安全に扱えるよう近似実装
    する。具体的には:

        1. 末尾から結合マーク (``unicodedata.combining(ch) != 0``)、
           ZWJ/ZWNJ (``\u200C`` / ``\u200D``) を飛ばす。
        2. その直前が base 文字なら 1 文字削る。さらに再度 (1) を行い、
           削った base に付随する combining 系列も丸ごと落とす。
        3. 削った結果に末尾 ZWJ が残る (ZWJ emoji sequence の途中) 場合
           はそれも除去する。

    Args:
        s: 対象文字列。

    Returns:
        末尾 grapheme を除去した文字列。空文字列は空文字列を返す。

    Issue:
        #29 — code-unit ベースの切り詰めが ZWJ emoji や濁点などを半端に
        残すのを防ぐ。
    """
    if not s:
        return s

    # (1) 末尾の combining / ZW-joiner / variation selector を飛ばして base を探す。
    i = len(s)
    while i > 0:
        ch = s[i - 1]
        cp = ord(ch)
        if (
            unicodedata.combining(ch) != 0
            or cp in (0x200C, 0x200D)
            or 0xFE00 <= cp <= 0xFE0F
        ):
            i -= 1
            continue
        break

    # (2) base 文字を 1 つ落とす。全部 combiner だった場合は全削除となる。
    if i > 0:
        i -= 1

    # (3) 残った末尾に ZWJ 等の「次の cluster を期待する」文字が残って
    # いる場合はさらに剥がす (例: `A<ZWJ>B` を右から削るとき `A<ZWJ>`
    # が中途半端に残らないようにする)。
    while i > 0:
        ch = s[i - 1]
        cp = ord(ch)
        if (
            unicodedata.combining(ch) != 0
            or cp in (0x200C, 0x200D)
            or 0xFE00 <= cp <= 0xFE0F
        ):
            i -= 1
            continue
        break

    return s[:i]


def _truncate_to_fit(
    text: str,
    usable_width: float,
    height: float,
    font_name: str,
    size_pt: float,
) -> str:
    """``size_pt`` の size で ``(usable_width, height)`` に収まるよう末尾を
    省略記号で切り詰める.

    方針:
    - まず wrap_text で折り返し結果を得る。
    - 高さが許す行数 ``max_lines`` を算出する。
    - 最終行は末尾から 1 grapheme cluster ずつ削りながら ``line + 省略記号``
      の幅が ``usable_width`` 以下になるまで縮める (issue #29)。
    - ``max_lines == 0`` の場合は空文字列を返す。
    """
    line_h = size_pt * 0.0139 * _AUTO_FIT_LINE_HEIGHT
    if line_h <= 0:
        return text
    max_lines = int(height // line_h)
    if max_lines <= 0:
        return ""

    lines = wrap_text(text, usable_width, size_pt, font_name)
    if not lines:
        return ""

    if len(lines) <= max_lines:
        # 高さには収まるが、幅の再確認は wrap で保証済み。そのまま返す。
        return "\n".join(lines)

    retained = lines[:max_lines]
    last = retained[-1]
    # last + ellipsis が usable_width を超える間、末尾から grapheme
    # cluster 単位で削る (code-unit 単位だと ZWJ emoji / 結合文字を壊す)。
    while last and estimate_text_width(last + _ELLIPSIS, size_pt, font_name) > usable_width:
        new_last = _strip_last_grapheme(last)
        if new_last == last:
            # 保険: 進まない場合は code unit で 1 文字削る。
            new_last = last[:-1]
        last = new_last
    retained[-1] = (last + _ELLIPSIS) if last else _ELLIPSIS
    return "\n".join(retained)


def add_auto_fit_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_name: str = "Arial",
    font_size_pt: float = 11,
    min_size_pt: float = 7,
    bold: bool = False,
    color_hex: str = "333333",
    align: str = "left",
    vertical_anchor: str = "top",
    truncate_with_ellipsis: bool = True,
    east_asian_font: str | None = None,
    wrap: bool = True,
) -> tuple[object, float]:
    """指定 box に収まる最大 font size でテキストを描画する.

    動作 (``wrap=True``、既定):
        (a) 内側 padding 0.05" を左右に見込んで usable width を計算する。
        (b) ``font_size_pt`` から開始し、``estimate_text_height(wrapped)`` が
            ``height`` 以下になるまで 0.5pt 刻みで縮小する。
        (c) ``min_size_pt`` に達しても収まらない場合、
            ``truncate_with_ellipsis=True`` なら末尾を省略記号で切り詰めて
            描画する。False なら full text をそのまま描画しオーバーフロー
            を許容する。
        (d) 決定した font size で ``_add_textbox`` 経由で textframe を生成し、
            ``vertical_anchor`` に応じて ``MSO_ANCHOR`` を設定する。

    動作 (``wrap=False``):
        単一行描画モード。アクションタイトル・KPI 値のように折り返しを
        許容しないユースケース向け。高さではなく幅で font size を縮小する。
        具体的には (a) で得た usable width に対し、``estimate_text_width``
        が usable width 以下となる最大 size を 0.5pt 刻みで探す。``min_size_pt``
        に達しても幅を超える場合の振る舞いは ``truncate_with_ellipsis`` に
        従う (True なら末尾を省略記号で切り詰め、False なら clip を許容して
        そのまま描画)。生成 textbox には ``word_wrap=False`` を設定し、
        PowerPoint 側でも自動折り返しが発生しないようにする。

    Args:
        slide: 対象スライドオブジェクト。
        text: 描画するテキスト。
        left, top, width, height: 位置と寸法 (inches)。
        font_name: フォント名。デフォルトは ``"Arial"``。
        font_size_pt: 開始 font size (pt)。
        min_size_pt: 縮小の下限 (pt)。
        bold: 太字にするかどうか。
        color_hex: 文字色 (hex、``#`` なし)。
        align: 水平方向の揃え ``"left" | "center" | "right"``。
        vertical_anchor: 垂直方向の揃え ``"top" | "middle" | "bottom"``。
        truncate_with_ellipsis: min size でも収まらない場合に末尾を
            省略記号で切り詰めるかどうか。
        east_asian_font: 東アジア用フォント名 (省略時は theme の既定値等)。
        wrap: True なら高さベースで auto-fit し折り返しを許容する。False なら
            幅ベースで auto-fit し単一行を維持する (``word_wrap=False`` を
            textbox に適用)。既定は後方互換のため True。

    Returns:
        ``(shape, actual_font_size)`` のタプル。テスト・デバッグ用途。
    """
    usable_width = max(width - 2 * _AUTO_FIT_PADDING, 0.01)

    if wrap:
        actual_size, fits = _fit_font_size(
            text, usable_width, height, font_name, font_size_pt, min_size_pt,
        )
        rendered_text = text
        if not fits and truncate_with_ellipsis:
            rendered_text = _truncate_to_fit(
                text, usable_width, height, font_name, actual_size,
            )
    else:
        actual_size, fits = _fit_font_size_single_line(
            text, usable_width, font_name, font_size_pt, min_size_pt,
        )
        rendered_text = text
        if not fits and truncate_with_ellipsis:
            rendered_text = _truncate_to_fit_single_line(
                text, usable_width, font_name, actual_size,
            )

    idx = _add_textbox(
        slide,
        left, top, width, height,
        text=rendered_text,
        font_name=font_name,
        font_size=actual_size,
        font_color=color_hex,
        bold=bold if bold else None,
        italic=None,
        alignment=align,
        vertical_anchor=vertical_anchor,
        word_wrap=wrap,
        line_spacing=None,
        underline=None,
        east_asian_font=east_asian_font,
    )

    shape = slide.shapes[idx]
    tf = shape.text_frame
    tf.word_wrap = wrap
    # auto_size は手動でサイズ決定済みなので明示的に None にする。
    tf.auto_size = None
    if vertical_anchor in _ANCHOR_MAP:
        tf.vertical_anchor = _ANCHOR_MAP[vertical_anchor]

    return shape, actual_size


def add_auto_fit_textbox_file(
    file_path: str,
    slide_index: int,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Arial",
    font_size_pt: float = 11,
    min_size_pt: float = 7,
    bold: bool = False,
    color_hex: str = "333333",
    align: str = "left",
    vertical_anchor: str = "top",
    truncate_with_ellipsis: bool = True,
    wrap: bool = True,
) -> dict:
    """File-based wrapper: 指定 box に収まるよう自動縮小する textbox を追加する.

    Returns:
        ``{"shape_index": int, "slide_index": int, "actual_font_size": float}``
        を含む dict。MCP ツール経由で呼び出される想定。
    """
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    shape, actual_size = add_auto_fit_textbox(
        slide, text, left, top, width, height,
        font_name=font_name,
        font_size_pt=font_size_pt,
        min_size_pt=min_size_pt,
        bold=bold,
        color_hex=color_hex,
        align=align,
        vertical_anchor=vertical_anchor,
        truncate_with_ellipsis=truncate_with_ellipsis,
        wrap=wrap,
    )
    # shape_index を決定 (slide 内で shape を線形検索)
    shape_index = -1
    for i, s in enumerate(slide.shapes):
        if s is shape:
            shape_index = i
            break
    save_pptx(prs, file_path)
    return {
        "slide_index": slide_index,
        "shape_index": shape_index,
        "shape_name": shape.name,
        "actual_font_size": actual_size,
    }
