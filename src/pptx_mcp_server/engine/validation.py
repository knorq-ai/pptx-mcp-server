"""
Layout validation — detect overlapping shapes, out-of-bounds elements,
text overflow, unreadable font, title/divider collision, inconsistent gaps.

Run after build_slide / build_deck to catch layout issues before delivery.
"""

from __future__ import annotations

import re
import statistics
from dataclasses import dataclass, field, asdict
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu

from .layout_constants import (
    TEXTBOX_INNER_PADDING_PER_SIDE,
    TEXTBOX_INNER_PADDING_TOTAL,
)
from .text_metrics import estimate_text_height, estimate_text_width

# ---------------------------------------------------------------------------
# 既存互換ヘルパ
# ---------------------------------------------------------------------------


def _get_shape_bounds(shape) -> Tuple[float, float, float, float]:
    """shape の (left, top, right, bottom) を inches で返す."""
    emu_to_in = 1 / 914400
    left = shape.left * emu_to_in
    top = shape.top * emu_to_in
    right = left + shape.width * emu_to_in
    bottom = top + shape.height * emu_to_in
    return (left, top, right, bottom)


def _boxes_overlap(a: Tuple, b: Tuple, margin: float = 0.0) -> bool:
    """2 つの (left, top, right, bottom) box が重なるかを返す.

    margin: shape 間に要求する最小間隔 (inches)。
    負値を与えると指定量以上重なった場合のみ True となる。
    """
    a_left, a_top, a_right, a_bottom = a
    b_left, b_top, b_right, b_bottom = b

    if a_right + margin <= b_left:
        return False
    if b_right + margin <= a_left:
        return False
    if a_bottom + margin <= b_top:
        return False
    if b_bottom + margin <= a_top:
        return False
    return True


def _is_small_decorative(shape) -> bool:
    """小さな装飾 shape (divider、accent bar 等) かを判定する."""
    emu_to_in = 1 / 914400
    w = shape.width * emu_to_in
    h = shape.height * emu_to_in
    if h < 0.08 or w < 0.08:
        return True
    if w * h < 0.05:
        return True
    return False


def _is_container(bounds_a: Tuple, bounds_b: Tuple) -> bool:
    """一方が他方を完全包含するか (背景矩形内のテキスト配置など)."""
    a_left, a_top, a_right, a_bottom = bounds_a
    b_left, b_top, b_right, b_bottom = bounds_b

    if a_left <= b_left and a_top <= b_top and a_right >= b_right and a_bottom >= b_bottom:
        return True
    if b_left <= a_left and b_top <= a_top and b_right >= a_right and b_bottom >= a_bottom:
        return True
    return False


# ---------------------------------------------------------------------------
# 既存チェッカ (後方互換のためシグネチャ・挙動を維持)
# ---------------------------------------------------------------------------


def check_slide_overlaps(
    slide,
    margin: float = 0.03,
    slide_width: float = 13.333,
    slide_height: float = 7.5,
) -> List[str]:
    """1 slide 内の重なり・はみ出しを検出する.

    margin: shape 間に要求する最小間隔 (inches)。

    戻り値は警告文字列のリスト。空リストなら問題なし。
    """
    warnings: List[str] = []
    shapes = list(slide.shapes)

    content_shapes = [s for s in shapes if not _is_small_decorative(s)]

    for i in range(len(content_shapes)):
        for j in range(i + 1, len(content_shapes)):
            si = content_shapes[i]
            sj = content_shapes[j]

            bi = _get_shape_bounds(si)
            bj = _get_shape_bounds(sj)

            if _boxes_overlap(bi, bj, margin):
                if _is_container(bi, bj):
                    continue

                overlap_x = max(0, min(bi[2], bj[2]) - max(bi[0], bj[0]))
                overlap_y = max(0, min(bi[3], bj[3]) - max(bi[1], bj[1]))
                overlap_area = overlap_x * overlap_y

                area_i = (bi[2] - bi[0]) * (bi[3] - bi[1])
                area_j = (bj[2] - bj[0]) * (bj[3] - bj[1])
                # 2026-04 #77: 旧閾値 1.0 sq in は通常の単行 text label (2.3"×0.35"=0.805 sq in)
                # を取り逃がしていた。本当に "tiny" な shape は既に `_is_small_decorative`
                # (h<0.08 or w<0.08 or area<0.05) で除外されているため、0.15 に緩める。
                if min(area_i, area_j) < 0.15:
                    continue

                name_i = si.name or f"Shape {i}"
                name_j = sj.name or f"Shape {j}"

                if overlap_area >= 0.15:
                    warnings.append(
                        f"OVERLAP ({overlap_area:.1f} sq in): "
                        f"'{name_i}' at ({bi[0]:.1f},{bi[1]:.1f},{bi[2]:.1f},{bi[3]:.1f}) "
                        f"× '{name_j}' at ({bj[0]:.1f},{bj[1]:.1f},{bj[2]:.1f},{bj[3]:.1f})"
                    )
                else:
                    h_gap = max(0, max(bi[0], bj[0]) - min(bi[2], bj[2]))
                    v_gap = max(0, max(bi[1], bj[1]) - min(bi[3], bj[3]))
                    min_gap = max(h_gap, v_gap)
                    # 同 threshold に揃える (#77)。
                    if area_i > 0.15 and area_j > 0.15:
                        warnings.append(
                            f"TOO CLOSE (gap {min_gap:.2f}\"): "
                            f"'{name_i}' at ({bi[0]:.1f},{bi[1]:.1f},{bi[2]:.1f},{bi[3]:.1f}) "
                            f"× '{name_j}' at ({bj[0]:.1f},{bj[1]:.1f},{bj[2]:.1f},{bj[3]:.1f})"
                        )

    for s in shapes:
        if _is_small_decorative(s):
            continue
        b = _get_shape_bounds(s)
        if b[2] > slide_width + 0.1:
            warnings.append(
                f"OUT OF BOUNDS (right): '{s.name}' extends to {b[2]:.1f}\" "
                f"(slide width={slide_width}\")"
            )
        if b[3] > slide_height + 0.1:
            warnings.append(
                f"OUT OF BOUNDS (bottom): '{s.name}' extends to {b[3]:.1f}\" "
                f"(slide height={slide_height}\")"
            )

    return warnings


def check_deck_overlaps(pptx_path: str) -> str:
    """deck 全体の重なり・はみ出しを報告する.

    戻り値はレポート文字列。空のレポートは `All slides clean…` を返す。
    """
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400

    all_warnings: List[str] = []
    for i, slide in enumerate(prs.slides):
        warnings = check_slide_overlaps(
            slide, margin=-0.05,
            slide_width=slide_w, slide_height=slide_h,
        )
        for w in warnings:
            all_warnings.append(f"Slide {i + 1}: {w}")

    if not all_warnings:
        return "All slides clean — no overlaps or out-of-bounds detected."

    return f"Found {len(all_warnings)} layout issues:\n" + "\n".join(all_warnings)


# ---------------------------------------------------------------------------
# 拡張バリデーション (Issue #5)
# ---------------------------------------------------------------------------


@dataclass
class ValidationFinding:
    """バリデーション検出結果の 1 件分を表す.

    severity: ``"error"`` | ``"warning"`` | ``"info"``
    category: ``"text_overflow"`` | ``"unreadable_text"``
        | ``"divider_collision"`` | ``"inconsistent_gap"``
    """

    severity: str
    slide_index: int
    shape_name: str
    category: str
    message: str
    suggested_fix: str = ""

    def to_dict(self) -> Dict[str, Any]:
        """dict 形式で返す (JSON シリアライズ用)."""
        return asdict(self)


# フッタ等の小文字ホワイトリストパターン (unreadable チェックをスキップ)
_UNREADABLE_WHITELIST_RE = re.compile(r"footer|page_number|source|footnote", re.IGNORECASE)


def _emu_to_in(emu: Optional[int]) -> float:
    """EMU → inches。None は 0 を返す."""
    if emu is None:
        return 0.0
    return emu / 914400


def _shape_name(shape, fallback: str = "") -> str:
    """shape の name を取得する。欠落時は fallback を返す."""
    try:
        name = shape.name
    except Exception:
        name = None
    if not name:
        return fallback
    return name


def _iter_runs(text_frame) -> List[Any]:
    """text_frame の全 run をフラットに返す."""
    runs: List[Any] = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            runs.append(run)
    return runs


# OOXML DrawingML 名前空間。pPr / defRPr の lxml 検索で使う。
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _pptx_defrpr_size(paragraph) -> Optional[float]:
    """``<a:pPr><a:defRPr sz="2400"/></a:pPr>`` から pt 値を解決する.

    OOXML では ``sz`` は 1/100 pt 単位。未設定の場合は ``None`` を返す。
    """
    try:
        p_pr = paragraph._pPr  # type: ignore[attr-defined]
    except Exception:
        p_pr = None
    if p_pr is None:
        return None
    defrpr = p_pr.find(f"{{{_A_NS}}}defRPr")
    if defrpr is None:
        return None
    sz = defrpr.get("sz")
    if sz is None:
        return None
    try:
        return int(sz) / 100.0
    except (TypeError, ValueError):
        return None


def _effective_paragraph_size(paragraph, default_pt: float = 18.0) -> float:
    """paragraph の実効フォントサイズ (pt) を解決する.

    優先順:
        1. paragraph 内 run の最大 ``font.size`` (run に明示されたもの)
        2. ``paragraph.font.size`` (paragraph の ``rPr`` 由来)
        3. ``pPr/defRPr/@sz`` (段落既定値)
        4. ``default_pt`` (フォールバック)

    ``add_auto_fit_textbox`` のように ``defRPr`` にのみサイズを書き込む
    経路でも正しく pt を抽出できるよう (3) を明示的にサポートする。
    """
    # (1) run レベル
    max_run_pt: Optional[float] = None
    for run in paragraph.runs:
        size = run.font.size
        if size is None:
            continue
        pt = size.pt
        if max_run_pt is None or pt > max_run_pt:
            max_run_pt = pt
    if max_run_pt is not None:
        return max_run_pt

    # (2) paragraph.font.size
    try:
        pf_size = paragraph.font.size
    except Exception:
        pf_size = None
    if pf_size is not None:
        return pf_size.pt

    # (3) pPr/defRPr/@sz
    defrpr_pt = _pptx_defrpr_size(paragraph)
    if defrpr_pt is not None:
        return defrpr_pt

    # (4) デフォルト
    return default_pt


def _effective_run_size(run, paragraph, default_pt: float = 18.0) -> float:
    """run の実効フォントサイズ (pt) を解決する.

    優先順:
        1. ``run.font.size`` (run に明示された値)
        2. ``paragraph.font.size`` (paragraph の ``rPr`` 由来)
        3. ``pPr/defRPr/@sz`` (段落既定値)
        4. ``default_pt`` (フォールバック)

    ``_effective_paragraph_size`` が段落全体の代表サイズ (run の最大値) を
    返すのに対し、こちらは個別 run の継承サイズを解決する。同一段落内で
    明示 size を持つ sibling run がいても、``size=None`` の run が defRPr
    を継承する場合にその pt を正しく得る目的で使う (Issue #60)。
    """
    # (1) run レベル
    size = run.font.size
    if size is not None:
        return size.pt

    # (2) paragraph.font.size
    try:
        pf_size = paragraph.font.size
    except Exception:
        pf_size = None
    if pf_size is not None:
        return pf_size.pt

    # (3) pPr/defRPr/@sz
    defrpr_pt = _pptx_defrpr_size(paragraph)
    if defrpr_pt is not None:
        return defrpr_pt

    # (4) デフォルト
    return default_pt


def _dominant_font_size(text_frame, default_pt: float = 18.0) -> float:
    """text_frame 内の段落で最大の実効フォントサイズ (pt) を返す.

    段落単位の解決には ``_effective_paragraph_size`` を利用する。メッセージ
    文字列の表示用 (代表サイズ) に残してあるユーティリティである。
    """
    max_pt: Optional[float] = None
    for para in text_frame.paragraphs:
        pt = _effective_paragraph_size(para, default_pt=default_pt)
        if max_pt is None or pt > max_pt:
            max_pt = pt
    if max_pt is None:
        return default_pt
    return max_pt


def _full_text(text_frame) -> str:
    """text_frame の全段落テキストを改行結合して返す."""
    lines: List[str] = []
    for para in text_frame.paragraphs:
        lines.append("".join(run.text for run in para.runs))
    return "\n".join(lines)


def _paragraph_text(paragraph) -> str:
    """paragraph 内の全 run テキストを連結して返す."""
    return "".join(run.text for run in paragraph.runs)


def _estimate_frame_needed_height(
    text_frame,
    usable_width: float,
    default_pt: float = 18.0,
) -> Tuple[float, float]:
    """段落単位で推定高さを合算した ``(total_in, max_pt)`` を返す.

    段落ごとに ``_effective_paragraph_size`` で実効 pt を解決し、その pt
    で ``estimate_text_height`` を呼び合算する。空段落も 1 行分の高さ
    (pt × 0.0139 × 1.2) を確保する。代表サイズとして最大 pt も返す。
    """
    total = 0.0
    max_pt: Optional[float] = None
    paragraphs = list(text_frame.paragraphs)
    for para in paragraphs:
        pt = _effective_paragraph_size(para, default_pt=default_pt)
        if max_pt is None or pt > max_pt:
            max_pt = pt
        text = _paragraph_text(para)
        if text:
            total += estimate_text_height(text, usable_width, pt)
        else:
            total += pt * 0.0139 * 1.2
    if max_pt is None:
        max_pt = default_pt
    return total, max_pt


def check_text_overflow(
    presentation: Presentation,
    *,
    min_readable_pt: float = 8.0,
    overflow_tolerance_pct: float = 5.0,
) -> List[ValidationFinding]:
    """テキストフレームの高さ溢れを検出する.

    各 text frame の幅から左右の ``TEXTBOX_INNER_PADDING_TOTAL`` (= 0.10") を
    差し引いた usable width に基づき、段落単位で有効フォントサイズを解決し、
    段落ごとの推定高さを合算する。合計が frame_height × (1 + tolerance/100)
    を超える場合に ``error`` finding を返す。
    """
    findings: List[ValidationFinding] = []
    tolerance_factor = 1 + overflow_tolerance_pct / 100.0

    for slide_index, slide in enumerate(presentation.slides):
        for shape_i, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.has_text_frame:
                continue
            frame_width = _emu_to_in(shape.width)
            frame_height = _emu_to_in(shape.height)
            if frame_width <= 0.1 or frame_height <= 0.05:
                continue

            tf = shape.text_frame
            text = _full_text(tf)
            if not text.strip():
                continue

            usable_width = max(0.1, frame_width - TEXTBOX_INNER_PADDING_TOTAL)
            needed_height, font_size = _estimate_frame_needed_height(
                tf, usable_width
            )
            max_allowed = frame_height * tolerance_factor

            if needed_height > max_allowed:
                # 収まる最大代表 pt を 0.5pt 刻みで探索する。段落間の比率を
                # 保って一様スケールする想定で needed_height に線形近似する。
                fit_pt: Optional[float] = None
                candidate = font_size - 0.5
                while candidate >= min_readable_pt:
                    scale = candidate / font_size if font_size > 0 else 1.0
                    if needed_height * scale <= frame_height:
                        fit_pt = candidate
                        break
                    candidate -= 0.5

                if fit_pt is not None:
                    suggested = (
                        f"reduce font {font_size:.1f}pt → {fit_pt:.1f}pt"
                    )
                else:
                    delta = max(0.05, needed_height - frame_height)
                    suggested = f'increase box height by {delta:.2f}"'

                name = _shape_name(shape, f"Shape {shape_i}")
                findings.append(
                    ValidationFinding(
                        severity="error",
                        slide_index=slide_index,
                        shape_name=name,
                        category="text_overflow",
                        message=(
                            f"text overflow: needed {needed_height:.2f}\" "
                            f"> frame {frame_height:.2f}\" "
                            f"(max font {font_size:.1f}pt, width {frame_width:.2f}\")"
                        ),
                        suggested_fix=suggested,
                    )
                )

    return findings


def _matches_whitelist_names(name: str, whitelist_names: Optional[List[str]]) -> bool:
    """``whitelist_names`` の任意要素が ``name`` に部分一致するかを返す.

    大文字小文字を区別しない。``whitelist_names`` が None/空の場合は False。
    """
    if not whitelist_names:
        return False
    lowered = (name or "").lower()
    for needle in whitelist_names:
        if not needle:
            continue
        if needle.lower() in lowered:
            return True
    return False


def _is_footer_zone_shape(
    shape,
    text_frame,
    slide_height: float,
    min_readable_pt: float,
) -> bool:
    """フッタゾーン (スライド下端 0.6" 以内) にある小型テキストかをヒューリスティックで判定する.

    条件 (すべて満たす場合に True):
        1. shape の bounding box 上端が ``slide_height - 0.6`` 以上
           (すなわち底辺から 0.6" 以内に完全に収まる)
        2. 推定折返し後の行数 ≤ 3
        3. 代表フォントサイズが ``min_readable_pt - 1`` 以上
           (極小フォントはフッタ扱いせず別途 unreadable で検出する)

    名前ベースのホワイトリストを補完するために用いる。
    """
    try:
        s_top = _emu_to_in(shape.top)
        s_height = _emu_to_in(shape.height)
        s_width = _emu_to_in(shape.width)
    except Exception:
        return False
    # bbox が完全に底辺 0.6" 以内に収まるか
    if s_top < slide_height - 0.6:
        return False

    # 代表フォント pt
    max_pt: Optional[float] = None
    for para in text_frame.paragraphs:
        pt = _effective_paragraph_size(para, default_pt=min_readable_pt)
        if max_pt is None or pt > max_pt:
            max_pt = pt
    if max_pt is None:
        max_pt = min_readable_pt
    if max_pt < min_readable_pt - 1:
        return False

    # 行数 (段落ごとに折返し推定)
    usable_width = max(0.1, s_width - TEXTBOX_INNER_PADDING_TOTAL)
    total_lines = 0
    for para in text_frame.paragraphs:
        text = _paragraph_text(para)
        if not text:
            total_lines += 1
            continue
        pt = _effective_paragraph_size(para, default_pt=max_pt)
        # estimate_text_height は pt × 0.0139 × 1.2 × n_lines 相当
        line_h = max(pt * 0.0139 * 1.2, 0.01)
        h = estimate_text_height(text, usable_width, pt)
        n_lines = max(1, int(round(h / line_h)))
        total_lines += n_lines
        if total_lines > 3:
            return False
    if total_lines > 3:
        return False
    return True


def check_unreadable_text(
    presentation: Presentation,
    *,
    min_readable_pt: float = 8.0,
    whitelist_names: Optional[List[str]] = None,
) -> List[ValidationFinding]:
    """``min_readable_pt`` 未満の font サイズを警告する.

    以下のいずれかに該当する shape は検査対象から除外する:

    1. ``whitelist_names`` (呼び出し側が明示した部分一致リスト) に shape 名が
       部分一致する (大文字小文字を区別しない)。日本語名のフッタ等にも対応する。
    2. フッタゾーン ヒューリスティック (``_is_footer_zone_shape``) に合致する
       (底辺 0.6" 以内に収まり、行数 ≤ 3、フォントが ``min_readable_pt - 1`` 以上)。
    3. shape 名に ``footer`` / ``page_number`` / ``source`` / ``footnote`` を
       含む (英語既定の正規表現フォールバック、後方互換)。
    """
    findings: List[ValidationFinding] = []
    slide_height_in = _emu_to_in(presentation.slide_height)

    for slide_index, slide in enumerate(presentation.slides):
        for shape_i, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.has_text_frame:
                continue
            name = _shape_name(shape, f"Shape {shape_i}")
            # 1. 明示 whitelist (substring, case-insensitive)
            if _matches_whitelist_names(name, whitelist_names):
                continue
            # 3. 既存英語正規表現フォールバック
            if _UNREADABLE_WHITELIST_RE.search(name or ""):
                continue

            tf = shape.text_frame

            # 2. フッタゾーン ヒューリスティック
            if _is_footer_zone_shape(shape, tf, slide_height_in, min_readable_pt):
                continue
            smallest: Optional[float] = None
            # run 単位で実効サイズを解決する。明示 size のある sibling run が
            # いても、``size=None`` の run が defRPr を継承するケースを取り
            # こぼさない (Issue #60)。
            for para in tf.paragraphs:
                runs = list(para.runs)
                if runs:
                    for run in runs:
                        # 明示されない場合は min_readable_pt をフォールバックに
                        # することで、defRPr 等で明示されない限り警告しない
                        # ("size 不明 = 読みにくい" と誤警告しない) 挙動を保つ。
                        eff_pt = _effective_run_size(
                            run, para, default_pt=min_readable_pt
                        )
                        if eff_pt < min_readable_pt:
                            if smallest is None or eff_pt < smallest:
                                smallest = eff_pt
                else:
                    # run を持たない段落 (空段落) は paragraph / defRPr のみを見る
                    eff_pt = _effective_paragraph_size(
                        para, default_pt=min_readable_pt
                    )
                    if eff_pt < min_readable_pt:
                        if smallest is None or eff_pt < smallest:
                            smallest = eff_pt

            if smallest is not None:
                findings.append(
                    ValidationFinding(
                        severity="warning",
                        slide_index=slide_index,
                        shape_name=name,
                        category="unreadable_text",
                        message=(
                            f"font size {smallest:.1f}pt is below readable "
                            f"threshold {min_readable_pt:.1f}pt"
                        ),
                        suggested_fix=(
                            f"increase font size to at least {min_readable_pt:.1f}pt"
                        ),
                    )
                )

    return findings


def _get_vertical_anchor(text_frame) -> Any:
    """text_frame の ``vertical_anchor`` を安全に取得する.

    未設定または例外時は ``None`` を返す (呼び出し側で TOP 相当に扱う)。
    """
    try:
        return text_frame.vertical_anchor
    except Exception:
        return None


def _projected_text_range(
    shape, text_frame, needed_height: float
) -> Tuple[float, float]:
    """text_frame の vertical_anchor に基づいた実描画垂直範囲 (top, bottom) を返す.

    - ``TOP`` (または未設定/None): ``[s_top, s_top + needed_height]``
    - ``BOTTOM``: ``[s_bottom - needed_height, s_bottom]``
    - ``MIDDLE``: ``[s_center - needed_height/2, s_center + needed_height/2]``

    戻り値は inches の (top, bottom) タプル。
    """
    s_top = _emu_to_in(shape.top)
    s_height = _emu_to_in(shape.height)
    s_bottom = s_top + s_height
    s_center = s_top + s_height / 2.0

    anchor = _get_vertical_anchor(text_frame)
    if anchor == MSO_ANCHOR.BOTTOM:
        return (s_bottom - needed_height, s_bottom)
    if anchor == MSO_ANCHOR.MIDDLE:
        half = needed_height / 2.0
        return (s_center - half, s_center + half)
    # TOP / None / unset → TOP 相当
    return (s_top, s_top + needed_height)


def check_divider_collision(
    presentation: Presentation,
) -> List[ValidationFinding]:
    """タイトル直下の divider line にタイトルテキストが食い込まないかを検証する.

    divider 条件: 高さ < 0.05" かつ 幅 > 2"。
    その上 0.5" 以内に存在する text frame について、vertical_anchor を考慮した
    投影範囲 (``_projected_text_range``) が divider 帯と交差したら
    ``error`` finding を返す。

    bottom-anchor の場合 (McKinsey 風タイトル: y=0.45..0.95 を anchor=bottom で
    テキストが上方向に伸びる) は、``s_bottom <= divider_top`` の限り
    collision とみなさない。
    """
    findings: List[ValidationFinding] = []

    for slide_index, slide in enumerate(presentation.slides):
        shapes = list(slide.shapes)
        dividers: List[Tuple[Any, Tuple[float, float, float, float]]] = []
        for shape in shapes:
            w = _emu_to_in(shape.width)
            h = _emu_to_in(shape.height)
            if h < 0.05 and w > 2.0:
                dividers.append((shape, _get_shape_bounds(shape)))

        if not dividers:
            continue

        for shape_i, shape in enumerate(shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.has_text_frame:
                continue

            bounds = _get_shape_bounds(shape)
            s_left, s_top, s_right, s_bottom = bounds

            for div_shape, div_bounds in dividers:
                if div_shape is shape:
                    continue
                d_left, d_top, d_right, d_bottom = div_bounds
                # divider は text frame より下にあり、かつ text frame 下端から 0.5" 以内
                if d_top < s_top:
                    continue
                if d_top - s_bottom > 0.5:
                    continue

                tf = shape.text_frame
                text = _full_text(tf)
                if not text.strip():
                    continue
                frame_width = max(
                    0.1, (s_right - s_left) - TEXTBOX_INNER_PADDING_TOTAL
                )
                needed_height, font_size = _estimate_frame_needed_height(
                    tf, frame_width
                )
                proj_top, proj_bottom = _projected_text_range(
                    shape, tf, needed_height
                )
                limit = d_top - 0.02

                # divider 帯との交差判定 (projected range が divider 上端を
                # 越えて食い込むか)。divider 自体の厚みは 0.05" 未満なので
                # ``limit`` (= d_top - 0.02) を越えたら error とする。
                anchor = _get_vertical_anchor(tf)
                collides = False
                if anchor == MSO_ANCHOR.BOTTOM:
                    # bottom-anchor: テキストは上方向に伸びる。
                    # s_bottom (= proj_bottom) が divider top を越えて、
                    # かつ proj_top が divider top より上にあるときのみ衝突。
                    # s_bottom <= d_top の場合は絶対に衝突しない。
                    if proj_bottom > d_top and proj_top < d_top:
                        collides = True
                elif anchor == MSO_ANCHOR.MIDDLE:
                    # center 基準の投影範囲が divider 上端を跨ぐ場合に衝突
                    if proj_bottom > limit:
                        collides = True
                else:
                    # TOP (既定): 従来どおり bottom が limit を越えたら衝突
                    if proj_bottom > limit:
                        collides = True

                if collides:
                    name = _shape_name(shape, f"Shape {shape_i}")
                    div_name = _shape_name(div_shape, "Divider")
                    findings.append(
                        ValidationFinding(
                            severity="error",
                            slide_index=slide_index,
                            shape_name=name,
                            category="divider_collision",
                            message=(
                                f"text extends to {proj_bottom:.2f}\" "
                                f"overlapping divider '{div_name}' at "
                                f"{d_top:.2f}\" (limit {limit:.2f}\")"
                            ),
                            suggested_fix=(
                                "reduce title length, shrink font, or "
                                "move the divider line lower"
                            ),
                        )
                    )
                    # divider ごとに最大 1 件で十分
                    break

    return findings


def _axis_groups(
    values: List[Tuple[int, float]],
    tolerance: float,
) -> List[List[int]]:
    """(index, axis_value) を tolerance 以内で clustering して返す.

    返り値は各 group の index リスト。3 要素以上のクラスタのみ返す。

    clustering は group 先頭 (anchor) との差分で single-linkage 判定を
    行う。running-mean 更新はドリフトを生む (例: tolerance=0.1 で
    ``[0.0, 0.05, 0.10, 0.15, 0.20]`` が全て同一 group に吸収される)
    ため採用しない。
    """
    if not values:
        return []
    sorted_vals = sorted(values, key=lambda x: x[1])
    groups: List[List[Tuple[int, float]]] = []
    current: List[Tuple[int, float]] = [sorted_vals[0]]
    for item in sorted_vals[1:]:
        anchor = current[0][1]
        if abs(item[1] - anchor) <= tolerance:
            current.append(item)
        else:
            groups.append(current)
            current = [item]
    groups.append(current)
    return [[idx for idx, _ in g] for g in groups if len(g) >= 3]


def check_inconsistent_gaps(
    presentation: Presentation,
    *,
    axis_tolerance: float = 0.1,
    gap_tolerance: float = 0.05,
) -> List[ValidationFinding]:
    """同じ行・列に並ぶ shape 群の gap のばらつきを検出する.

    行: top が ``axis_tolerance`` 以内で揃う 3 shape 以上。隣接 x gap の
    stddev または max-min が ``gap_tolerance`` を超えたら ``info`` finding。
    列: 同様に left 基準で y gap を比較する。
    """
    findings: List[ValidationFinding] = []

    for slide_index, slide in enumerate(presentation.slides):
        shapes = list(slide.shapes)
        content_shapes: List[Tuple[int, Any, Tuple[float, float, float, float]]] = []
        for idx, shape in enumerate(shapes):
            if _is_small_decorative(shape):
                continue
            bounds = _get_shape_bounds(shape)
            content_shapes.append((idx, shape, bounds))

        if len(content_shapes) < 3:
            continue

        # 行方向 (top 同一)
        top_values = [(i, b[1]) for i, (_, _, b) in enumerate(content_shapes)]
        row_groups = _axis_groups(top_values, axis_tolerance)
        for group in row_groups:
            # left 順
            sorted_group = sorted(group, key=lambda ii: content_shapes[ii][2][0])
            gaps: List[float] = []
            for a, b in zip(sorted_group, sorted_group[1:]):
                ba = content_shapes[a][2]
                bb = content_shapes[b][2]
                gap = bb[0] - ba[2]
                gaps.append(gap)
            if not gaps:
                continue
            gap_range = max(gaps) - min(gaps)
            std = statistics.pstdev(gaps) if len(gaps) > 1 else 0.0
            if std > gap_tolerance or gap_range > gap_tolerance:
                names = [
                    _shape_name(content_shapes[i][1], f"Shape {content_shapes[i][0]}")
                    for i in sorted_group
                ]
                gaps_str = ", ".join(f"{g:.2f}\"" for g in gaps)
                findings.append(
                    ValidationFinding(
                        severity="info",
                        slide_index=slide_index,
                        shape_name=names[0],
                        category="inconsistent_gap",
                        message=(
                            f"row of {len(sorted_group)} shapes "
                            f"({', '.join(names)}) has inconsistent x gaps: "
                            f"[{gaps_str}]"
                        ),
                        suggested_fix=(
                            "equalise horizontal spacing between shapes"
                        ),
                    )
                )

        # 列方向 (left 同一)
        left_values = [(i, b[0]) for i, (_, _, b) in enumerate(content_shapes)]
        col_groups = _axis_groups(left_values, axis_tolerance)
        for group in col_groups:
            sorted_group = sorted(group, key=lambda ii: content_shapes[ii][2][1])
            gaps = []
            for a, b in zip(sorted_group, sorted_group[1:]):
                ba = content_shapes[a][2]
                bb = content_shapes[b][2]
                gap = bb[1] - ba[3]
                gaps.append(gap)
            if not gaps:
                continue
            gap_range = max(gaps) - min(gaps)
            std = statistics.pstdev(gaps) if len(gaps) > 1 else 0.0
            if std > gap_tolerance or gap_range > gap_tolerance:
                names = [
                    _shape_name(content_shapes[i][1], f"Shape {content_shapes[i][0]}")
                    for i in sorted_group
                ]
                gaps_str = ", ".join(f"{g:.2f}\"" for g in gaps)
                findings.append(
                    ValidationFinding(
                        severity="info",
                        slide_index=slide_index,
                        shape_name=names[0],
                        category="inconsistent_gap",
                        message=(
                            f"column of {len(sorted_group)} shapes "
                            f"({', '.join(names)}) has inconsistent y gaps: "
                            f"[{gaps_str}]"
                        ),
                        suggested_fix=(
                            "equalise vertical spacing between shapes"
                        ),
                    )
                )

    return findings


def check_deck_extended(
    presentation: Presentation,
    *,
    min_readable_pt: float = 8.0,
    overflow_tolerance_pct: float = 5.0,
    axis_tolerance: float = 0.1,
    gap_tolerance: float = 0.05,
    whitelist_names: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """deck 全体を走査し、既存 + 拡張チェックをまとめて返す.

    戻り値構造::

        {
            "slides": [
                {
                    "index": int,
                    "overlaps": [...],
                    "out_of_bounds": [...],
                    "text_overflow": [...],
                    "unreadable_text": [...],
                    "divider_collision": [...],
                    "inconsistent_gaps": [...],
                },
                ...
            ],
            "summary": {"errors": int, "warnings": int, "infos": int},
        }

    既存 ``check_slide_overlaps`` の警告文字列を ``overlaps`` と
    ``out_of_bounds`` に分割して格納する (既存キーは文字列リストのまま)。

    Severity mapping — ``summary`` の各カウンタに寄与するカテゴリは以下のとおり:

    - ``errors``:
        - ``overlaps`` (``check_slide_overlaps`` 由来の overlap 警告)
        - ``out_of_bounds`` (slide 外に出ている shape)
        - ``text_overflow`` (``check_text_overflow``)
        - ``divider_collision`` (``check_divider_collision``)
    - ``warnings``:
        - ``unreadable_text`` (``check_unreadable_text``)
    - ``infos``:
        - ``inconsistent_gaps`` (``check_inconsistent_gaps``)

    CI ゲートで ``summary.errors == 0`` を要求する場合、overlap / out_of_bounds
    を含めて error と扱う必要があるため、本関数は legacy の文字列リスト長も
    errors に加算する。
    """
    slide_w = presentation.slide_width / 914400
    slide_h = presentation.slide_height / 914400

    text_overflow = check_text_overflow(
        presentation,
        min_readable_pt=min_readable_pt,
        overflow_tolerance_pct=overflow_tolerance_pct,
    )
    unreadable = check_unreadable_text(
        presentation,
        min_readable_pt=min_readable_pt,
        whitelist_names=whitelist_names,
    )
    divider = check_divider_collision(presentation)
    gaps = check_inconsistent_gaps(
        presentation,
        axis_tolerance=axis_tolerance,
        gap_tolerance=gap_tolerance,
    )

    # slide index ごとに集約
    by_slide: Dict[int, Dict[str, List[Any]]] = {}
    for i, slide in enumerate(presentation.slides):
        warnings = check_slide_overlaps(
            slide, margin=-0.05,
            slide_width=slide_w, slide_height=slide_h,
        )
        overlap_warnings = [w for w in warnings if not w.startswith("OUT OF BOUNDS")]
        oob_warnings = [w for w in warnings if w.startswith("OUT OF BOUNDS")]
        by_slide[i] = {
            "index": i,
            "overlaps": overlap_warnings,
            "out_of_bounds": oob_warnings,
            "text_overflow": [],
            "unreadable_text": [],
            "divider_collision": [],
            "inconsistent_gaps": [],
        }

    def _place(findings: List[ValidationFinding], key: str) -> None:
        for f in findings:
            # slide_index がプレゼンテーション範囲外のものはスキップ (防御)
            if f.slide_index not in by_slide:
                continue
            by_slide[f.slide_index][key].append(f.to_dict())

    _place(text_overflow, "text_overflow")
    _place(unreadable, "unreadable_text")
    _place(divider, "divider_collision")
    _place(gaps, "inconsistent_gaps")

    slides_out = [by_slide[i] for i in sorted(by_slide.keys())]

    # summary: ValidationFinding 由来 + legacy (overlaps / out_of_bounds) を集計する。
    # 詳細な severity mapping は docstring 参照。
    errors = sum(1 for f in text_overflow + divider if f.severity == "error")
    for slide_data in slides_out:
        errors += len(slide_data["overlaps"])
        errors += len(slide_data["out_of_bounds"])
    warnings_count = sum(1 for f in unreadable if f.severity == "warning")
    infos = sum(1 for f in gaps if f.severity == "info")

    return {
        "slides": slides_out,
        "summary": {
            "errors": errors,
            "warnings": warnings_count,
            "infos": infos,
        },
    }
