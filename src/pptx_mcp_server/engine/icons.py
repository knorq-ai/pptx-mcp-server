"""
Icon library — load, search, and inject vector icons from the built-in catalog.

Icons are pre-extracted from iconography_NEW.pptx at build time and stored
as XML strings in assets/icons.json.  The IconRegistry lazy-loads the catalog
on first access and caches parsed XML elements for fast injection.
"""

from __future__ import annotations

import difflib
import json
from dataclasses import dataclass
from importlib import resources
from typing import ClassVar, Dict, List, Optional

from lxml import etree
from pptx.util import Inches

from .pptx_io import (
    EngineError,
    ErrorCode,
    _get_slide,
    _parse_color,
    open_pptx,
    save_pptx,
)
from ..theme import Theme, MCKINSEY, resolve_color

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class IconInfo:
    """Metadata for a single icon in the library."""
    id: str
    name: str
    category: str
    keywords: tuple
    type: str              # "sp" or "grpSp"
    orig_width_emu: int
    orig_height_emu: int
    aspect_ratio: float
    colors: tuple


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------

class IconRegistry:
    """Singleton registry for the built-in icon library.

    Lazy-loads the catalog JSON from package assets on first access.
    Caches parsed lxml elements for fast repeated injection.

    Note: Not thread-safe. The MCP server uses stdio transport (single-threaded).
    If migrating to SSE/HTTP transport with concurrent requests, add a
    threading.Lock around _ensure_loaded().
    """
    _instance: ClassVar[Optional["IconRegistry"]] = None

    def __init__(self):
        self._catalog: Optional[dict] = None
        self._icons: Dict[str, IconInfo] = {}
        self._by_category: Dict[str, List[str]] = {}
        self._xml_strings: Dict[str, str] = {}  # id -> raw XML string

    @classmethod
    def get(cls) -> "IconRegistry":
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    @classmethod
    def _reset(cls) -> None:
        """Reset the singleton (for testing)."""
        cls._instance = None

    def _ensure_loaded(self) -> None:
        if self._catalog is not None:
            return
        self._load_catalog()

    def _load_catalog(self) -> None:
        """Load icons.json from package assets."""
        assets = resources.files("pptx_mcp_server") / "assets" / "icons.json"
        raw = assets.read_text(encoding="utf-8")
        self._catalog = json.loads(raw)

        for entry in self._catalog.get("icons", []):
            info = IconInfo(
                id=entry["id"],
                name=entry["name"],
                category=entry["category"],
                keywords=tuple(entry.get("keywords", [])),
                type=entry.get("type", "sp"),
                orig_width_emu=entry.get("orig_width_emu", 914400),
                orig_height_emu=entry.get("orig_height_emu", 914400),
                aspect_ratio=entry.get("aspect_ratio", 1.0),
                colors=tuple(entry.get("colors", [])),
            )
            self._icons[info.id] = info
            self._by_category.setdefault(info.category, []).append(info.id)
            self._xml_strings[info.id] = entry["xml"]

    def list_icons(
        self,
        category: Optional[str] = None,
        search: Optional[str] = None,
    ) -> List[IconInfo]:
        """Return icons, optionally filtered by category and/or keyword search."""
        self._ensure_loaded()
        results = list(self._icons.values())

        if category:
            cat_ids = set(self._by_category.get(category, []))
            results = [i for i in results if i.id in cat_ids]

        if search:
            terms = search.lower().split()
            results = [
                i for i in results
                if all(
                    t in i.name.lower()
                    or t in i.category.lower()
                    or any(t in kw for kw in i.keywords)
                    for t in terms
                )
            ]

        return results

    def list_categories(self) -> List[dict]:
        """Return all categories with counts."""
        self._ensure_loaded()
        return self._catalog.get("categories", [])

    def get_icon(self, icon_id: str) -> IconInfo:
        """Get icon metadata by ID.  Fuzzy match on miss."""
        self._ensure_loaded()
        if icon_id in self._icons:
            return self._icons[icon_id]

        # Fuzzy match
        matches = difflib.get_close_matches(icon_id, list(self._icons.keys()), n=3, cutoff=0.5)
        if matches:
            hint = ", ".join(matches)
            raise EngineError(
                ErrorCode.INVALID_PARAMETER,
                f"Icon '{icon_id}' not found. Did you mean: {hint}?",
            )
        raise EngineError(
            ErrorCode.INVALID_PARAMETER,
            f"Icon '{icon_id}' not found. Use pptx_list_icons to browse available icons.",
        )

    def get_icon_xml(self, icon_id: str) -> etree._Element:
        """Get a fresh parsed copy of the icon's XML element.

        Uses tostring/fromstring round-trip instead of deepcopy for performance.
        """
        self._ensure_loaded()
        if icon_id not in self._xml_strings:
            self.get_icon(icon_id)  # will raise with fuzzy match
        xml_str = self._xml_strings[icon_id]
        return etree.fromstring(xml_str.encode("utf-8"))


# ---------------------------------------------------------------------------
# Icon injection helpers
# ---------------------------------------------------------------------------


def _resolve_icon_size(
    info: IconInfo,
    width_in: Optional[float],
    height_in: Optional[float],
) -> tuple:
    """Calculate target EMU dimensions preserving aspect ratio."""
    EMU_PER_INCH = 914400
    if width_in is not None and height_in is not None:
        return int(width_in * EMU_PER_INCH), int(height_in * EMU_PER_INCH)
    if width_in is not None:
        w = int(width_in * EMU_PER_INCH)
        h = int(w / info.aspect_ratio) if info.aspect_ratio else w
        return w, h
    if height_in is not None:
        h = int(height_in * EMU_PER_INCH)
        w = int(h * info.aspect_ratio) if info.aspect_ratio else h
        return w, h
    # Default: 0.8 inch height
    h = int(0.8 * EMU_PER_INCH)
    w = int(h * info.aspect_ratio) if info.aspect_ratio else h
    return w, h


def _rewrite_xfrm(elem: etree._Element, left_emu: int, top_emu: int, cx: int, cy: int) -> None:
    """Rewrite the shape's position and size.

    For <p:sp>:     targets p:spPr/a:xfrm
    For <p:grpSp>:  targets p:grpSpPr/a:xfrm (the outer group transform)
    """
    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

    # Find the correct xfrm: group shapes use grpSpPr, individual shapes use spPr
    if tag == "grpSp":
        grpSpPr = elem.find(f"{{{P_NS}}}grpSpPr")
        if grpSpPr is None:
            return
        xfrm = grpSpPr.find(f"{{{A_NS}}}xfrm")
    else:
        spPr = elem.find(f"{{{P_NS}}}spPr")
        if spPr is None:
            return
        xfrm = spPr.find(f"{{{A_NS}}}xfrm")

    if xfrm is None:
        return

    off = xfrm.find(f"{{{A_NS}}}off")
    if off is not None:
        off.set("x", str(left_emu))
        off.set("y", str(top_emu))

    ext = xfrm.find(f"{{{A_NS}}}ext")
    if ext is not None:
        ext.set("cx", str(cx))
        ext.set("cy", str(cy))


def _recolor_icon(
    elem: etree._Element,
    original_colors: tuple,
    new_fill_color: Optional[str],
    new_outline_color: Optional[str],
    theme: Optional[Theme],
) -> None:
    """Recolor icon by swapping fill/outline srgbClr values in XML."""
    if theme:
        if new_fill_color:
            new_fill_color = resolve_color(theme, new_fill_color)
        if new_outline_color:
            new_outline_color = resolve_color(theme, new_outline_color)

    if new_fill_color:
        new_fill_color = new_fill_color.lstrip("#").upper()
    if new_outline_color:
        new_outline_color = new_outline_color.lstrip("#").upper()

    # Identify original fill (non-black) and outline (black) colors
    fill_originals = set()
    outline_original = None
    for c in original_colors:
        cu = c.upper()
        if cu == "000000":
            outline_original = cu
        else:
            fill_originals.add(cu)

    # Build color mapping
    color_map: Dict[str, str] = {}
    if new_fill_color and fill_originals:
        for fo in fill_originals:
            color_map[fo] = new_fill_color
    if new_outline_color and outline_original:
        color_map[outline_original] = new_outline_color

    if not color_map:
        return

    # Walk all srgbClr elements and replace
    for srgb in elem.findall(f".//{{{A_NS}}}srgbClr"):
        val = srgb.get("val", "").upper()
        if val in color_map:
            srgb.set("val", color_map[val])


def _reassign_shape_ids(slide, new_elem: etree._Element) -> None:
    """Assign unique shape IDs to avoid collisions with existing shapes."""
    spTree = slide.shapes._spTree

    # Find max existing ID
    max_id = 0
    for cNvPr in spTree.findall(f".//{{{P_NS}}}cNvPr"):
        id_val = cNvPr.get("id", "0")
        if id_val.isdigit():
            max_id = max(max_id, int(id_val))

    # Reassign IDs in the new element
    next_id = max_id + 1
    for cNvPr in new_elem.findall(f".//{{{P_NS}}}cNvPr"):
        cNvPr.set("id", str(next_id))
        next_id += 1


# ---------------------------------------------------------------------------
# Core icon injection
# ---------------------------------------------------------------------------


def _add_icon(
    slide,
    icon_id: str,
    left: float,
    top: float,
    width: Optional[float] = None,
    height: Optional[float] = None,
    color: Optional[str] = None,
    outline_color: Optional[str] = None,
    theme: Optional[Theme] = None,
) -> int:
    """In-memory: inject an icon from the library onto a slide.

    Returns shape index.
    """
    theme = theme or MCKINSEY
    registry = IconRegistry.get()
    info = registry.get_icon(icon_id)
    icon_elem = registry.get_icon_xml(icon_id)

    # Calculate target size
    EMU_PER_INCH = 914400
    target_w, target_h = _resolve_icon_size(info, width, height)

    # Rewrite position and size
    _rewrite_xfrm(icon_elem, int(left * EMU_PER_INCH), int(top * EMU_PER_INCH), target_w, target_h)

    # Recolor if requested
    if color or outline_color:
        _recolor_icon(icon_elem, info.colors, color, outline_color, theme)

    # Reassign shape IDs
    _reassign_shape_ids(slide, icon_elem)

    # Inject into slide's spTree
    spTree = slide.shapes._spTree
    spTree.append(icon_elem)

    # Count shapes via spTree children (not slide.shapes which requires oxml proxies)
    shape_tags = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}
    count = sum(
        1 for child in spTree
        if (child.tag.split("}")[-1] if "}" in child.tag else child.tag) in shape_tags
    )
    return count - 1


# ---------------------------------------------------------------------------
# File-based wrapper
# ---------------------------------------------------------------------------


def add_icon(
    file_path: str,
    slide_index: int,
    icon_id: str,
    left: float,
    top: float,
    width: Optional[float] = None,
    height: Optional[float] = None,
    color: Optional[str] = None,
    outline_color: Optional[str] = None,
) -> str:
    """File-based wrapper: add an icon to a slide."""
    prs = open_pptx(file_path)
    slide = _get_slide(prs, slide_index)
    idx = _add_icon(slide, icon_id, left, top, width, height, color, outline_color)
    save_pptx(prs, file_path)
    return f"Added icon '{icon_id}' [{idx}] on slide [{slide_index}]"


def list_icons_formatted(
    category: Optional[str] = None,
    search: Optional[str] = None,
) -> str:
    """Return a formatted string listing available icons."""
    registry = IconRegistry.get()

    if not category and not search:
        # Show categories overview
        cats = registry.list_categories()
        lines = ["Icon Library — Categories:"]
        for cat in cats:
            lines.append(f"  {cat['id']:15s} ({cat['count']} icons)")
        lines.append(f"\nTotal: {sum(c['count'] for c in cats)} icons")
        lines.append("Use category or search parameter to filter.")
        return "\n".join(lines)

    icons = registry.list_icons(category=category, search=search)
    if not icons:
        return f"No icons found for category='{category}' search='{search}'"

    lines = [f"Found {len(icons)} icons:"]
    for icon in icons[:50]:  # cap at 50 results
        lines.append(f"  {icon.id:30s} | {icon.category:12s} | {icon.name}")
    if len(icons) > 50:
        lines.append(f"  ... and {len(icons) - 50} more")
    return "\n".join(lines)
