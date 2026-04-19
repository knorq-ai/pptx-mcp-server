#!/usr/bin/env python3
"""
pptx-mcp-server -- MCP server for PPTX presentation editing.

Parameter conventions (new tools):
- ``*_json``: JSON-stringified array/object input (例: ``rows_json``, ``kpis_json``,
  ``cards_json``, ``items_json``). 生の Python list/dict ではなく必ず JSON 文字列で渡す。
- ``*_pt``: フォントサイズなど「ポイント単位」を明示する (例: ``font_size_pt``,
  ``min_size_pt``)。旧 tool の素の ``font_size`` は後方互換のため温存する。
- ``colors``: ``"#"`` を含まない 6 桁 hex (例: ``"2251FF"``)。
- coordinates: inches (float)。

Response shape (v0.2.0; BREAKING change — see issue #88):

All tool calls return a JSON string. Success payloads are wrapped as::

    {"ok": true, "result": <legacy return value>}

Error payloads are structured as::

    {"ok": false, "error": {
        "code": "INVALID_PARAMETER",
        "parameter": "items_json",   // optional
        "message": "...",
        "hint": "...",                // optional
        "issue": 35                    // optional GitHub issue reference
    }}

Consumers should ``json.loads`` the response and branch on the ``ok`` field.
The ``error.code`` field mirrors ``EngineError.code`` enum values
(``INVALID_PARAMETER``, ``FILE_NOT_FOUND``, ``SLIDE_NOT_FOUND``, etc.).

Auto-render (v0.2.0; BREAKING change — see issue #86):

Composite / batch-build tools previously invoked LibreOffice for a PNG preview
after every successful edit. This is now **opt-in** via the
``PPTX_MCP_AUTO_RENDER=1`` environment variable, with a hard timeout controlled
by ``PPTX_MCP_RENDER_TIMEOUT`` (default 10 seconds). If rendering times out
or fails, the primary tool still succeeds — the render outcome is surfaced as
a ``render_warning`` field in the result payload.
"""

import concurrent.futures
import json
import os
from dataclasses import fields
from typing import Any, Dict, Optional

try:
    from mcp.server.fastmcp import FastMCP
except ImportError as e:
    raise ImportError(
        "pptx_mcp_server.server requires the 'mcp' package. "
        "Install with: pip install 'pptx-mcp-server[mcp]'"
    ) from e

from .engine import (
    EngineError,
    ErrorCode,
    create_presentation,
    get_presentation_info,
    read_slide,
    add_slide,
    move_slide,
    delete_slide,
    duplicate_slide,
    set_slide_background,
    add_textbox,
    add_auto_fit_textbox_file,
    add_flex_container_file,
    add_shape,
    add_image,
    edit_text,
    add_paragraph,
    delete_shape,
    list_shapes,
    add_table,
    edit_table_cell,
    edit_table_cells,
    format_table,
    format_shape,
    set_slide_dimensions,
    add_content_slide,
    add_section_divider,
    add_kpi_row,
    add_bullet_block,
    build_slide,
    build_deck,
    render_slide,
    add_chart,
    add_icon,
    list_icons_formatted,
    add_connector,
    add_callout,
    check_deck_overlaps,
    check_deck_extended,
    CardSpec,
    add_responsive_card_row,
)
from .engine.pptx_io import open_pptx, save_pptx, _get_slide

INSTRUCTIONS = """
# pptx-editor — Professional Presentation Builder

This MCP server is a neutral capability provider. It exposes tools for
building PowerPoint decks. It does not prescribe user-facing UX (whether
to prompt the caller for a theme, when to confirm with a user, etc.) —
that belongs in the agent's system prompt or the calling application.

## Quick Start
1. `pptx_create` — create a new 16:9 PPTX file
2. `pptx_build_deck` — build an ENTIRE deck from a JSON spec (most efficient)
3. `pptx_render_slide` — render to PNG for visual verification (optional,
   requires LibreOffice)

## Recommended Workflow
- Use `pptx_build_deck` for new decks (1 call = all slides, 1 file I/O).
- Use `pptx_build_slide` to add individual slides.
- Use `pptx_render_slide` explicitly to verify visually (returns PNG path
  — read with Read tool). Composite tools do NOT auto-render by default.
- Use `pptx_check_layout` after building to detect overlaps before delivery.
- Use primitive tools (`pptx_edit_text`, `pptx_format_shape`) for fine-grained
  edits.

## Response Shape (v0.2.0+)
All tools return a JSON string. Parse with `json.loads`:

- Success: `{"ok": true, "result": ...}`
- Error:   `{"ok": false, "error": {"code": "INVALID_PARAMETER",
            "parameter": "items_json", "message": "...", "hint": "..."}}`

Error codes mirror `EngineError.code`: `INVALID_PARAMETER`, `FILE_NOT_FOUND`,
`SLIDE_NOT_FOUND`, `SHAPE_NOT_FOUND`, `INDEX_OUT_OF_RANGE`, `INVALID_PPTX`,
`TABLE_ERROR`, `CHART_ERROR`, `INTERNAL_ERROR`.

## Auto-Render (opt-in; off by default)
Composite tools (`pptx_add_content_slide`, `pptx_build_slide`, etc.) can
auto-render a PNG preview after each successful edit. This forks LibreOffice
and adds ~1.5s per call, so it is **OFF** by default. Enable via the
`PPTX_MCP_AUTO_RENDER=1` environment variable. Timeout is controlled by
`PPTX_MCP_RENDER_TIMEOUT` (default 10 seconds). If rendering times out or
fails, the primary tool still succeeds — the failure is reported in the
result's `render_warning` field. For explicit rendering, use
`pptx_render_slide` directly.

## Available Themes
Pass `"theme": "<name>"` in slide specs for `build_slide` / `build_deck`.
Available: `mckinsey` (default), `deloitte`, `neutral`. Custom palettes are
supported by passing explicit `primary_color` / `accent_color` hex values on
individual elements (e.g., `font_color`, `fill_color`) instead of a named
theme.

## McKinsey-Style Layout Rules
- **Slide dimensions**: 16:9 widescreen (13.333" x 7.5")
- **Margins**: 0.9" left/right, content width = 11.533"
- **Action title**: 22pt Arial Bold, max ~30 chars for single line (auto-shrinks if longer)
- **Title zone**: 0.45" to 0.95" (bottom-anchored, divider line at 0.95")
- **Body zone**: 1.15" to 6.5" — USE THE FULL 5.35" HEIGHT. Distribute content evenly.
- **Footer zone**: 6.65" (source line + page number)
- **Font**: Arial for everything (sans-serif works best with Japanese)

## Color Palette (McKinsey-inspired, for reference)
- Primary text: `051C2C` (dark navy)
- Accent/highlight: `2251FF` (bright blue)
- Secondary text: `666666`
- Footnote: `A2AAAD`
- Background alt: `F5F5F5`
- Positive: `2E7D32` (green)
- Negative: `C62828` (red)
- Table border: `D0D0D0`

## Table Formatting (McKinsey-style, automatic)
- Dark navy header (`051C2C`) with white text
- No vertical borders, thin horizontal borders only
- Alternating row shading (`F5F5F5`)
- Numbers right-aligned, text left-aligned
- Use `pptx_add_table` or `"type": "table"` in build_slide spec

## Data Density Guidelines (Consulting Quality)
- **Fill every slide**: Body zone (1.15" to 6.5") should be 90%+ utilized. No large blank areas.
- **Use small fonts for detail**: Tables at 9-10pt, bullet text at 10-12pt. Only titles at 22pt.
- **Pack information**: A chart+sidebar slide should have 4+ text bullets AND a table in the sidebar. A table slide should have 8+ rows.
- **Multi-zone layouts**: Combine chart (left 60%) + sidebar (right 40%) with title + bullets + table in the sidebar.
- **Bullet blocks should describe, not list**: Each bullet should be a full sentence with specifics, not a 3-word fragment.
- **For CAGR/growth annotations**: Use a floating `rounded_rectangle` shape badge, NOT a callout with arrow.

## Common Pitfalls to Avoid
1. **Text behind shapes**: Don't put text on a shape if you'll overlay other shapes on top. Use `pptx_add_textbox` for labels over background shapes.
2. **Sparse slides**: Don't leave the bottom half empty. If you have 4 bullets, use the full body height — line spacing auto-distributes.
3. **Long titles**: Keep action titles under 30 chars. The tool auto-shrinks and warns if too long.
4. **Too many tool calls**: Use `pptx_build_deck` (1 call for whole deck) instead of individual `pptx_add_*` calls.
5. **Arrows to nowhere**: Callout arrows must point at a specific data point. For general annotations (CAGR, labels), use a shaped textbox instead.

## Slide Spec Format (for build_slide / build_deck)
```json
{
  "layout": "content",        // "content" | "section_divider" | "blank"
  "title": "Action title",    // for content/section_divider
  "subtitle": "...",           // section_divider only
  "background": "051C2C",     // optional hex color
  "source": "Source: ...",     // optional footnote
  "page_number": 1,           // optional
  "elements": [
    {"type": "textbox", "left": 0.9, "top": 1.2, "width": 11.5, "height": 0.3,
     "text": "...", "font_size": 14, "font_color": "2251FF", "bold": true,
     "alignment": "left", "vertical_anchor": "top"},
    {"type": "shape", "shape_type": "rectangle", ...},
    {"type": "table", "rows": [["H1","H2"],["v1","v2"]], "left": 0.9, "top": 3.0,
     "width": 11.5, "col_widths": [0.5, 0.5]},
    {"type": "kpi_row", "kpis": [{"value":"100","label":"Metric"}], "y": 1.2},
    {"type": "bullet_block", "items": ["Point 1","Point 2"],
     "left": 0.9, "top": 2.0, "width": 11.5, "height": 2.0},
    {"type": "image", "image_path": "/path/img.png", "left": 1.0, "top": 1.0, "width": 3.0},
    {"type": "chart", "chart_type": "stacked_column",
     "left": 0.9, "top": 1.15, "width": 7.0, "height": 5.0,
     "categories": ["2020","2021","2022"],
     "series": [{"name":"Revenue","values":[10,20,30],"color":"2251FF"}],
     "data_labels_show": true, "legend_position": "bottom"}
  ]
}
```

## Chart Element (for build_slide / pptx_add_chart)
chart_type: bar, stacked_bar, stacked_bar_100, column, stacked_column, stacked_column_100, line, line_markers, pie, area, area_stacked, doughnut, radar.

Key fields (all flat — no nesting):
- `categories`: ["A","B","C"] — category labels
- `series`: [{"name":"S1","values":[1,2,3],"color":"2251FF"}] — data series (color optional, auto-assigned from theme)
- `data_labels_show`: true/false — show values on chart
- `data_labels_position`: center, outside_end, inside_end, inside_base, above, below
- `data_labels_number_format`: "#,##0", "0.0%", etc.
- `legend_position`: bottom, top, right, left, null (hidden)
- `axis_value_title`: "Revenue (M)" — Y-axis label
- `axis_value_min/max/major_unit`: scale control
- `axis_value_gridlines`: true/false
- `gap_width`: 150 (bar gap), `overlap`: 100 (stacked)

## Icon Element (for build_slide / pptx_add_icon)
640 built-in vector icons. Use `pptx_list_icons` to browse, or use directly:
```json
{"type": "icon", "icon_id": "briefcase", "left": 2.0, "top": 3.0, "width": 0.8,
 "color": "2251FF", "outline_color": "051C2C"}
```
Common icons: briefcase, chart, person, globe, airplane, laptop, phone, car, building, arrow,
calendar, clock, document, email, gear, handshake, key, lock, money, star, target, trophy.

## Card Grid Element (for build_slide) — auto-balanced layout
Use `card_grid` for 2×2 frameworks, feature grids, strategy cards, etc. Cards auto-size to fill the body zone.
```json
{"type": "card_grid", "cards": [
  {"title": "Strategy 1", "bullets": ["Point A", "Point B"], "icon_id": "target", "icon_color": "2251FF"},
  {"title": "Strategy 2", "body": "Description text", "icon_id": "globe"}
]}
```
Each card: title (required), body or bullets, optional icon_id + icon_color.
Grid auto-computes: 1-2 cards → 1 row, 3-4 → 2×2, 5-6 → 2×3, 7-9 → 3×3.

## Connector Element (for build_slide / pptx_add_connector)
```json
{"type": "connector", "begin_x": 2.0, "begin_y": 3.5, "end_x": 6.0, "end_y": 5.0,
 "color": "accent", "end_arrow": "triangle", "dash_style": "dash"}
```
connector_type: straight (default), elbow, curve.
Arrows: none, triangle, stealth, diamond, oval, open.
dash_style: solid, dash, dot, dash_dot, long_dash.

## Callout Element (for build_slide / pptx_add_callout)
Annotation textbox + arrow pointing to target. Auto-places label if label_x/label_y omitted.
```json
{"type": "callout", "text": "+15% YoY", "target_x": 5.5, "target_y": 3.0,
 "font_color": "negative", "bg_color": "F5F5F5", "arrow_end": "stealth"}
```

## Rendering (Optional)
`pptx_render_slide` requires LibreOffice. Install:
- macOS: `brew install --cask libreoffice`
- Ubuntu/Debian: `sudo apt install libreoffice`
- Windows: Download from https://www.libreoffice.org/download/
"""

mcp = FastMCP("pptx-editor", instructions=INSTRUCTIONS)


# ── Structured response helpers (issue #88) ────────────────────────────

def _success(result: Any) -> str:
    """Wrap a successful tool result in ``{"ok": true, "result": ...}``.

    ``result`` は legacy tool の戻り値 (通常は human-readable string) を
    そのまま格納する。JSON で表現できない object は呼び出し側で事前に
    serialize すること。
    """
    return json.dumps({"ok": True, "result": result}, ensure_ascii=False)


def _error(
    code: str,
    message: str,
    *,
    parameter: Optional[str] = None,
    hint: Optional[str] = None,
    issue: Optional[int] = None,
) -> str:
    """Build a structured error payload and return JSON-string.

    Shape::

        {"ok": false, "error": {"code": <str>, "message": <str>,
         "parameter": <optional str>, "hint": <optional str>,
         "issue": <optional int>}}
    """
    err: Dict[str, Any] = {"code": code, "message": message}
    if parameter is not None:
        err["parameter"] = parameter
    if hint is not None:
        err["hint"] = hint
    if issue is not None:
        err["issue"] = issue
    return json.dumps({"ok": False, "error": err}, ensure_ascii=False)


def _err(e: Exception) -> str:
    """Translate an exception into a structured error JSON string.

    ``EngineError`` は ``code`` enum をそのまま error.code として流用する。
    それ以外は ``INTERNAL_ERROR`` に fall back する。
    """
    if isinstance(e, EngineError):
        return _error(e.code.value, str(e))
    return _error("INTERNAL_ERROR", f"{type(e).__name__}: {e}")


# ── Auto-render gate (issue #86) ───────────────────────────────────────

_DEFAULT_RENDER_TIMEOUT_S = 10.0


def _auto_render_enabled() -> bool:
    """``PPTX_MCP_AUTO_RENDER`` が truthy なら auto-render を実行する."""
    v = os.environ.get("PPTX_MCP_AUTO_RENDER", "").strip().lower()
    return v in {"1", "true", "yes", "on"}


def _auto_render_timeout() -> float:
    """``PPTX_MCP_RENDER_TIMEOUT`` (秒) を float で返す. 既定 10 秒."""
    raw = os.environ.get("PPTX_MCP_RENDER_TIMEOUT", "").strip()
    if not raw:
        return _DEFAULT_RENDER_TIMEOUT_S
    try:
        v = float(raw)
        if v <= 0:
            return _DEFAULT_RENDER_TIMEOUT_S
        return v
    except ValueError:
        return _DEFAULT_RENDER_TIMEOUT_S


def _auto_render(file_path: str, slide_index: int) -> Dict[str, Any]:
    """Render a slide preview if enabled; else return a neutral "skipped" payload.

    Always returns a dict — never raises, never fails the caller. Shape::

        {"rendered": false, "reason": "disabled"}                       # off
        {"rendered": true, "preview_path": "/.../slide-01.png"}         # ok
        {"rendered": false, "reason": "timeout", "timeout_s": 10.0}     # slow
        {"rendered": false, "reason": "failed", "error": "<msg>"}       # crash

    Opt-in via ``PPTX_MCP_AUTO_RENDER=1``. Timeout via
    ``PPTX_MCP_RENDER_TIMEOUT`` (default 10s). The caller should only invoke
    this AFTER the primary action has succeeded.
    """
    if not _auto_render_enabled():
        return {"rendered": False, "reason": "disabled"}

    timeout = _auto_render_timeout()

    def _do_render() -> str:
        return render_slide(file_path, slide_index=slide_index, dpi=100)

    # ThreadPoolExecutor で走らせ future.result(timeout) で上限を掛ける。
    # `with` block を使うと __exit__ で shutdown(wait=True) が呼ばれ、
    # 裏の slow スレッドが終わるまでブロックするため timeout が機能しない。
    # 代わりに明示的に shutdown(wait=False) を呼ぶ。
    # timeout 到達時はスレッドが裏で生きたままだが、subprocess 側にも
    # 120 秒 / 60 秒の独自 timeout があるため無限ハングはしない。
    ex = concurrent.futures.ThreadPoolExecutor(max_workers=1)
    try:
        future = ex.submit(_do_render)
        try:
            out = future.result(timeout=timeout)
        except concurrent.futures.TimeoutError:
            return {
                "rendered": False,
                "reason": "timeout",
                "timeout_s": timeout,
            }
        except Exception as e:  # renderer itself raised
            return {"rendered": False, "reason": "failed", "error": f"{type(e).__name__}: {e}"}
    finally:
        ex.shutdown(wait=False)

    try:
        # render_slide may return multiple lines; take the last one (target slide)
        lines = out.strip().split("\n")
        return {"rendered": True, "preview_path": lines[-1]}
    except Exception as e:
        return {"rendered": False, "reason": "failed", "error": f"{type(e).__name__}: {e}"}


def _success_with_render(primary: Any, render_info: Dict[str, Any]) -> str:
    """Compose a success payload plus the auto-render outcome.

    - Render disabled → plain ``{"ok": true, "result": <primary>}``.
    - Render succeeded → result wraps ``{"value": primary, "preview_path": ...}``.
    - Render failed/timed out → result wraps ``{"value": primary,
      "render_warning": {...}}``.
    """
    if not render_info.get("rendered") and render_info.get("reason") == "disabled":
        return _success(primary)
    if render_info.get("rendered"):
        return _success(
            {"value": primary, "preview_path": render_info.get("preview_path")}
        )
    # Failed / timeout — primary still succeeded.
    return _success({"value": primary, "render_warning": render_info})


# --- Presentation ------------------------------------------------

@mcp.tool()
def pptx_create(
    file_path: str,
    width_inches: float = 13.333,
    height_inches: float = 7.5,
) -> str:
    """Create a new blank PPTX file. Default is 16:9 widescreen."""
    try:
        return _success(create_presentation(file_path, width_inches, height_inches))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_get_info(file_path: str) -> str:
    """Get presentation overview: slide count, dimensions, shape summaries."""
    try:
        return _success(get_presentation_info(file_path))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_read_slide(file_path: str, slide_index: int) -> str:
    """Read detailed content of a slide -- all shapes, text, tables."""
    try:
        return _success(read_slide(file_path, slide_index))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_list_shapes(file_path: str, slide_index: int) -> str:
    """List all shapes on a slide with indices, types, positions, text preview."""
    try:
        return _success(list_shapes(file_path, slide_index))
    except Exception as e:
        return _err(e)


# --- Slides -------------------------------------------------------

@mcp.tool()
def pptx_add_slide(file_path: str, layout_index: int = 6) -> str:
    """Add a new slide. Layout 6 = Blank (most common)."""
    try:
        return _success(add_slide(file_path, layout_index))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_move_slide(file_path: str, from_index: int, to_index: int) -> str:
    """Move a slide from one position to another. 0-based indices."""
    try:
        return _success(move_slide(file_path, from_index, to_index))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_delete_slide(file_path: str, slide_index: int) -> str:
    """Delete a slide by 0-based index."""
    try:
        return _success(delete_slide(file_path, slide_index))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_duplicate_slide(file_path: str, slide_index: int) -> str:
    """Duplicate a slide (appended at end)."""
    try:
        return _success(duplicate_slide(file_path, slide_index))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_set_slide_background(file_path: str, slide_index: int, color: str) -> str:
    """Set solid background color for a slide. Color as hex e.g. '051C2C' (without #)."""
    try:
        return _success(set_slide_background(file_path, slide_index, color))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_set_dimensions(file_path: str, width: float, height: float) -> str:
    """Set presentation slide dimensions in inches (applies to all slides)."""
    try:
        return _success(set_slide_dimensions(file_path, width, height))
    except Exception as e:
        return _err(e)


# --- Textboxes ----------------------------------------------------

@mcp.tool()
def pptx_add_textbox(
    file_path: str,
    slide_index: int,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    font_name: str = None,
    font_size: float = None,
    font_color: str = None,
    bold: bool = None,
    italic: bool = None,
    alignment: str = None,
    vertical_anchor: str = None,
    word_wrap: bool = True,
    line_spacing: float = None,
    underline: bool = None,
) -> str:
    """Add a text box to a slide. Position and size in inches. Alignment: left/center/right. Vertical anchor: top/middle/bottom."""
    try:
        return _success(add_textbox(
            file_path, slide_index, left, top, width, height, text,
            font_name, font_size, font_color, bold, italic,
            alignment, vertical_anchor, word_wrap, line_spacing, underline,
        ))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_auto_fit_textbox(
    file_path: str,
    slide_index: int,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Arial",
    font_size_pt: float = 11,
    min_size_pt: float = 7,
    bold: bool = False,
    color_hex: str = "333333",
    align: str = "left",
    vertical_anchor: str = "top",
    truncate_with_ellipsis: bool = True,
) -> str:
    """Add a textbox that auto-shrinks font size to fit a fixed box. Starts from font_size_pt and steps down 0.5pt until text fits height, or reaches min_size_pt. If still overflowing at min and truncate_with_ellipsis=True, trailing chars are replaced with an ellipsis. Returns a JSON object with shape_index, shape_name, and actual_font_size."""
    try:
        result = add_auto_fit_textbox_file(
            file_path, slide_index, text, left, top, width, height,
            font_name=font_name,
            font_size_pt=font_size_pt,
            min_size_pt=min_size_pt,
            bold=bold,
            color_hex=color_hex,
            align=align,
            vertical_anchor=vertical_anchor,
            truncate_with_ellipsis=truncate_with_ellipsis,
        )
        return _success(result)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_flex_container(
    file_path: str,
    slide_index: int,
    items_json: str,
    left: float,
    top: float,
    width: float,
    height: float,
    direction: str = "row",
    gap: float = 0.15,
    padding: float = 0.0,
    align: str = "stretch",
) -> str:
    """Add a CSS-flexbox-style container that lays out child items along a main axis.

    ``items_json`` は JSON-stringified array (``*_json`` 命名規約に沿う)。各要素 dict
    のキーは以下:
      - `sizing`: "fixed" | "grow" | "content"
      - `type`: "text" | "rectangle"
      - `size` (for fixed), `grow` (for grow, default 1), `content_size` (for content)
      - optional `min_size`, `max_size`
      - for type=text: `text`, `font_size_pt`, `bold`, `color_hex`, `align`, `vertical_anchor`, `truncate_with_ellipsis`
      - for type=rectangle: `fill_color`, `line_color`, `line_width`, `no_line`

    direction: "row" (horizontal) | "column" (vertical). gap and padding in inches.
    align cross-axis: "stretch" のみ現状サポート。"start" / "center" / "end" は
    `INVALID_PARAMETER` を返す (将来対応予定; #24 参照)。

    例: ``items_json='[{"sizing":"fixed","size":2,"type":"rectangle"}]'``

    Returns JSON with allocations (per-item [left, top, width, height]) and shape identifiers created.
    """
    try:
        if not isinstance(items_json, str):
            return _error(
                "INVALID_PARAMETER",
                "items_json must be a JSON string, not a raw Python list.",
                parameter="items_json",
                hint=(
                    "Pass a JSON-stringified array, e.g., "
                    "'[{\"sizing\":\"fixed\",\"size\":2,\"type\":\"rectangle\"}]'."
                ),
                issue=35,
            )
        try:
            items = json.loads(items_json)
        except json.JSONDecodeError as e:
            return _error(
                "INVALID_PARAMETER",
                f"Invalid JSON in items_json: {e}",
                parameter="items_json",
                hint="items_json must be a JSON-stringified array.",
                issue=35,
            )
        if not isinstance(items, list):
            return _error(
                "INVALID_PARAMETER",
                "items_json must decode to a JSON array.",
                parameter="items_json",
                issue=35,
            )
        result = add_flex_container_file(
            file_path, slide_index, items,
            left=left, top=top, width=width, height=height,
            direction=direction, gap=gap, padding=padding, align=align,
        )
        return _success(result)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_edit_text(
    file_path: str,
    slide_index: int,
    shape_index: int,
    text: str = None,
    paragraph_index: int = 0,
    font_name: str = None,
    font_size: float = None,
    font_color: str = None,
    bold: bool = None,
    italic: bool = None,
    underline: bool = None,
    alignment: str = None,
    line_spacing: float = None,
) -> str:
    """Edit text content and formatting in an existing shape's paragraph. Supports all formatting: font, color, bold, italic, underline, alignment, line spacing."""
    try:
        return _success(edit_text(
            file_path, slide_index, shape_index, text, paragraph_index,
            font_name, font_size, font_color, bold, italic, underline,
            alignment, line_spacing,
        ))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_paragraph(
    file_path: str,
    slide_index: int,
    shape_index: int,
    text: str,
    font_name: str = None,
    font_size: float = None,
    font_color: str = None,
    bold: bool = None,
    italic: bool = None,
    underline: bool = None,
    alignment: str = None,
    line_spacing: float = None,
) -> str:
    """Append a new paragraph to an existing shape's text frame. Useful for multi-line text."""
    try:
        return _success(add_paragraph(
            file_path, slide_index, shape_index, text,
            font_name, font_size, font_color, bold, italic, underline,
            alignment, line_spacing,
        ))
    except Exception as e:
        return _err(e)


# --- Shapes -------------------------------------------------------

@mcp.tool()
def pptx_add_shape(
    file_path: str,
    slide_index: int,
    shape_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = None,
    line_color: str = None,
    line_width: float = None,
    no_line: bool = False,
    text: str = None,
    font_name: str = None,
    font_size: float = None,
    font_color: str = None,
    bold: bool = None,
    alignment: str = None,
) -> str:
    """Add an auto shape. Types: rectangle, rounded_rectangle, oval, triangle, diamond, chevron, arrow_right, arrow_left, arrow_up, arrow_down, callout, star_5, hexagon, pentagon. Position/size in inches. Colors as hex. WARNING: text inside shapes renders BEHIND any shapes placed on top. For labels over background shapes, use pptx_add_textbox instead."""
    try:
        return _success(add_shape(
            file_path, slide_index, shape_type, left, top, width, height,
            fill_color, line_color, line_width, no_line,
            text, font_name, font_size, font_color, bold, alignment,
        ))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_image(
    file_path: str,
    slide_index: int,
    image_path: str,
    left: float,
    top: float,
    width: float = None,
    height: float = None,
) -> str:
    """Add an image (PNG, JPG, SVG) to a slide. Position in inches. If only width or height is given, aspect ratio is preserved. If both given, image stretches to fit."""
    try:
        return _success(add_image(file_path, slide_index, image_path, left, top, width, height))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_delete_shape(file_path: str, slide_index: int, shape_index: int) -> str:
    """Delete a shape from a slide by its 0-based index."""
    try:
        return _success(delete_shape(file_path, slide_index, shape_index))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_format_shape(
    file_path: str,
    slide_index: int,
    shape_index: int,
    left: float = None,
    top: float = None,
    width: float = None,
    height: float = None,
    fill_color: str = None,
    no_fill: bool = False,
    line_color: str = None,
    line_width: float = None,
    no_line: bool = False,
    rotation: float = None,
) -> str:
    """Reposition, resize, or restyle an existing shape. Dimensions in inches."""
    try:
        return _success(format_shape(
            file_path, slide_index, shape_index,
            left, top, width, height,
            fill_color, no_fill, line_color, line_width, no_line, rotation,
        ))
    except Exception as e:
        return _err(e)


# --- Tables -------------------------------------------------------

@mcp.tool()
def pptx_add_table(
    file_path: str,
    slide_index: int,
    rows_json: str,
    left: float,
    top: float,
    width: float,
    col_widths_json: str = "",
    row_height: float = 0.36,
    font_size: float = 12,
    header_bg: str = "051C2C",
    header_fg: str = "FFFFFF",
    alt_row_bg: str = "F5F5F5",
    border_color: str = "D0D0D0",
    no_vertical_borders: bool = True,
) -> str:
    """Add a professionally formatted table. rows_json: JSON 2D array e.g. '[["A","B"],["1","2"]]'. First row = header. col_widths_json: JSON array of fractions e.g. '[0.5, 0.5]'."""
    try:
        rows = json.loads(rows_json)
        col_widths = json.loads(col_widths_json) if col_widths_json else None
        return _success(add_table(
            file_path, slide_index, rows, left, top, width,
            col_widths, row_height, font_size,
            header_bg, header_fg, alt_row_bg, border_color,
            0.5, no_vertical_borders,
        ))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_edit_table_cell(
    file_path: str,
    slide_index: int,
    shape_index: int,
    row: int,
    col: int,
    text: str = None,
    font_size: float = None,
    font_color: str = None,
    bold: bool = None,
    bg_color: str = None,
    alignment: str = None,
) -> str:
    """Edit a single table cell's text and formatting."""
    try:
        return _success(edit_table_cell(
            file_path, slide_index, shape_index, row, col,
            text, font_size, font_color, bold, bg_color, alignment,
        ))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_edit_table_cells(
    file_path: str,
    slide_index: int,
    shape_index: int,
    edits_json: str,
) -> str:
    """Batch edit multiple table cells. edits_json: JSON array of objects e.g. '[{"row":0,"col":1,"text":"new"}]'. Each: {row, col, text?, font_size?, font_color?, bold?, bg_color?}."""
    try:
        edits = json.loads(edits_json)
        return _success(edit_table_cells(file_path, slide_index, shape_index, edits))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_format_table(
    file_path: str,
    slide_index: int,
    shape_index: int,
    font_name: str = None,
    font_size: float = None,
    header_bg: str = None,
    header_fg: str = None,
    alt_row_bg: str = None,
) -> str:
    """Apply bulk formatting to an entire table (font, header colors, alternating rows)."""
    try:
        return _success(format_table(
            file_path, slide_index, shape_index,
            font_name, font_size, header_bg, header_fg, alt_row_bg,
        ))
    except Exception as e:
        return _err(e)


# --- Composites ---------------------------------------------------

@mcp.tool()
def pptx_add_content_slide(
    file_path: str,
    title: str,
    source: str = None,
    page_number: int = None,
) -> str:
    """Add a content slide with action title (auto-shrink to fit), divider line, optional source footnote and page number. McKinsey-style layout. Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1. LAYOUT GUIDE: Body area is 1.15\" to 6.65\" (5.5\" usable height). Distribute content evenly across this range — avoid clustering content in the top half with empty bottom space."""
    try:
        result = add_content_slide(file_path, title, source, page_number)
        idx = int(result.split("[")[1].split("]")[0])
        render_info = _auto_render(file_path, idx)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_section_divider(
    file_path: str,
    title: str,
    subtitle: str = "",
) -> str:
    """Add a section divider slide with dark background, centered title, and accent stripes. Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1."""
    try:
        result = add_section_divider(file_path, title, subtitle)
        idx = int(result.split("[")[1].split("]")[0])
        render_info = _auto_render(file_path, idx)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_kpi_row(
    file_path: str,
    slide_index: int,
    kpis_json: str,
    y: float,
) -> str:
    """Add a row of KPI callout boxes. kpis_json: JSON array e.g. '[{"value":"107.8M","label":"Revenue"}]'. y = vertical position in inches. Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1."""
    try:
        result = add_kpi_row(file_path, slide_index, kpis_json, y)
        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_bullet_block(
    file_path: str,
    slide_index: int,
    items_json: str,
    left: float,
    top: float,
    width: float,
    height: float,
) -> str:
    """Add a bulleted text block with multiple items. items_json: JSON array of strings e.g. '["Item 1","Item 2"]'. Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1."""
    try:
        result = add_bullet_block(file_path, slide_index, items_json, left, top, width, height)
        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_responsive_card_row(
    file_path: str,
    slide_index: int,
    cards_json: str,
    left: float,
    top: float,
    width: float,
    max_height: float,
    gap: float = 0.2,
    height_mode: str = "max",
    min_card_height: float = 1.0,
) -> str:
    """Add a row of variable-height cards that auto-size to content and align to the max content height.

    cards_json: JSON array of card dicts. Each card supports keys: title, body, label,
    accent_color (hex, "" disables bar), fill_color, title_size_pt, body_size_pt,
    title_color, body_color, label_size_pt, label_color, padding.

    height_mode:
      - "content": each card uses its own content height (heights may differ).
      - "max":     all cards take the max individual content height (bottoms align).
      - "fill":    all cards fill max_height (short content is centered vertically).

    Returns a JSON object: {"cards": [{"left","top","width","height"}, ...], "consumed_height": float}.
    Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1.
    """
    try:
        card_dicts = json.loads(cards_json) if isinstance(cards_json, str) else cards_json
        # #43: CardSpec dataclass は未知キーを TypeError として弾くが、
        # MCP ツール層で明示的に ``INVALID_PARAMETER`` として報告し、
        # どのカード・どのキーが原因かを LLM エージェントが再試行時に
        # 解釈できる形式で返す。
        _card_known_keys = {f.name for f in fields(CardSpec)}
        for i, spec in enumerate(card_dicts):
            unknown = set(spec.keys()) - _card_known_keys
            if unknown:
                raise EngineError(
                    ErrorCode.INVALID_PARAMETER,
                    (
                        f"card[{i}]: unknown keys {sorted(unknown)}; "
                        f"known keys: {sorted(_card_known_keys)}."
                    ),
                )
        cards = [CardSpec(**d) for d in card_dicts]
        prs = open_pptx(file_path)
        slide = _get_slide(prs, slide_index)

        placements, consumed = add_responsive_card_row(
            slide,
            cards,
            left=left, top=top, width=width, max_height=max_height,
            gap=gap,
            height_mode=height_mode,  # type: ignore[arg-type]
            min_card_height=min_card_height,
        )

        # CardPlacement を JSON 化可能な dict に変換する (save 前に行う)。
        # ここで serialize に失敗しても disk 上のファイルは変更されない (#34)。
        result = {
            "cards": [
                {
                    "left": p.left,
                    "top": p.top,
                    "width": p.width,
                    "height": p.height,
                }
                for p in placements
            ],
            "consumed_height": consumed,
        }

        # すべての in-memory 処理と return 値構築が成功した最後に保存する。
        # これにより中途半端な save による破損状態を防ぐ (#34)。
        save_pptx(prs, file_path)

        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


# --- Batch Build --------------------------------------------------

@mcp.tool()
def pptx_build_slide(
    file_path: str,
    spec_json: str,
) -> str:
    """Build an entire slide in ONE call from a JSON spec. Single file open/save. Much faster than individual tool calls. Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1.

    spec_json format:
    {
        "layout": "content" | "section_divider" | "blank",
        "title": "Action title",
        "background": "051C2C",  (optional)
        "source": "Source: ...", (optional)
        "page_number": 1,       (optional)
        "elements": [
            {"type": "textbox", "left": 0.9, "top": 1.2, "width": 11.5, "height": 0.3,
             "text": "...", "font_size": 14, "font_color": "2251FF", "bold": true},
            {"type": "shape", "shape_type": "rectangle", "left": 0.9, "top": 2.0,
             "width": 5.0, "height": 1.0, "fill_color": "F5F5F5", "no_line": true},
            {"type": "table", "rows": [["H1","H2"],["v1","v2"]], "left": 0.9,
             "top": 3.0, "width": 11.5, "col_widths": [0.5, 0.5]},
            {"type": "kpi_row", "kpis": [{"value":"100","label":"Metric"}], "y": 1.2},
            {"type": "bullet_block", "items": ["Point 1","Point 2"],
             "left": 0.9, "top": 2.0, "width": 11.5, "height": 2.0},
            {"type": "image", "image_path": "/path/img.png", "left": 1.0, "top": 1.0, "width": 3.0}
        ]
    }
    LAYOUT GUIDE: Body area is 1.15" to 6.65" (5.5" usable). Distribute elements evenly."""
    try:
        result = build_slide(file_path, spec_json)
        idx_str = result.split("[")[1].split("]")[0]
        render_info = _auto_render(file_path, int(idx_str))
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_build_deck(
    file_path: str,
    slides_json: str,
) -> str:
    """Build an ENTIRE DECK in ONE call from a JSON array of slide specs. Single file open/save for all slides. Use this for generating complete presentations efficiently. Auto-renders a preview PNG of the last slide ONLY when PPTX_MCP_AUTO_RENDER=1.

    slides_json: JSON array where each element is a slide spec (same format as pptx_build_slide).
    Example: '[{"layout":"content","title":"Slide 1","elements":[...]},{"layout":"section_divider","title":"Section"}]'"""
    try:
        result = build_deck(file_path, slides_json)
        render_info = _auto_render(file_path, -1)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


# --- Connectors & Callouts ----------------------------------------

@mcp.tool()
def pptx_add_connector(
    file_path: str,
    slide_index: int,
    begin_x: float,
    begin_y: float,
    end_x: float,
    end_y: float,
    connector_type: str = "straight",
    color: str = None,
    width: float = None,
    dash_style: str = None,
    begin_arrow: str = "none",
    end_arrow: str = "triangle",
    arrow_size: str = "medium",
) -> str:
    """Add a connector line between two points. Position in inches.
    connector_type: straight/elbow/curve.
    Arrows: none/triangle/stealth/diamond/oval/open.
    dash_style: solid/dash/dot/dash_dot/long_dash.
    Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1."""
    try:
        result = add_connector(
            file_path, slide_index, begin_x, begin_y, end_x, end_y,
            connector_type, color, width, dash_style,
            begin_arrow, end_arrow, arrow_size,
        )
        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_callout(
    file_path: str,
    slide_index: int,
    text: str,
    target_x: float,
    target_y: float,
    label_x: float = None,
    label_y: float = None,
    label_width: float = 2.0,
    label_height: float = 0.4,
    connector_type: str = "straight",
    font_size: float = 10,
    font_color: str = None,
    font_bold: bool = True,
    line_color: str = None,
    line_width: float = 1.0,
    arrow_end: str = "triangle",
    bg_color: str = None,
    border_color: str = None,
) -> str:
    """Add a callout annotation: textbox + connector arrow pointing to target.
    Auto-places label if label_x/label_y omitted. Position in inches.
    Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1."""
    try:
        result = add_callout(
            file_path, slide_index, text, target_x, target_y,
            label_x, label_y, label_width, label_height,
            connector_type, font_size, font_color, font_bold,
            line_color, line_width, arrow_end, bg_color, border_color,
        )
        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


# --- Icons --------------------------------------------------------

@mcp.tool()
def pptx_list_icons(
    category: str = "",
    search: str = "",
) -> str:
    """List available icons from the built-in icon library (640 vector icons).
    Filter by category and/or keyword search.

    Categories: business, people, technology, transport, medical, education, nature, objects, general.
    Example: pptx_list_icons(category="business") or pptx_list_icons(search="chart")

    Common icons: briefcase, chart, person, globe, airplane, laptop, phone, car, building, arrow,
    calendar, clock, document, email, gear, handshake, key, lock, money, star, target, trophy, user."""
    try:
        return _success(list_icons_formatted(
            category=category or None,
            search=search or None,
        ))
    except Exception as e:
        return _err(e)


@mcp.tool()
def pptx_add_icon(
    file_path: str,
    slide_index: int,
    icon_id: str,
    left: float,
    top: float,
    width: float = None,
    height: float = None,
    color: str = None,
    outline_color: str = None,
) -> str:
    """Add a vector icon from the built-in library to a slide.
    Position in inches. If only width or height given, aspect ratio is preserved.
    Colors as hex (e.g. '2251FF') or theme token ('accent', 'primary').
    Use pptx_list_icons to browse available icons.
    Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1."""
    try:
        result = add_icon(file_path, slide_index, icon_id, left, top, width, height, color, outline_color)
        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


# --- Charts -------------------------------------------------------

@mcp.tool()
def pptx_add_chart(
    file_path: str,
    slide_index: int,
    chart_json: str,
) -> str:
    """Add a professional chart to a slide. chart_json is a JSON object with:

    Required: chart_type (column/stacked_column/bar/stacked_bar/line/pie/area/doughnut/radar),
    categories (array of labels), series (array of {name, values, color?}).

    Optional: title, legend_position (bottom/top/right/left/null), data_labels_show (bool),
    data_labels_position, data_labels_number_format, axis_value_title, axis_value_min/max,
    axis_value_gridlines, gap_width, overlap, theme (mckinsey/deloitte/neutral).

    Auto-renders a preview PNG ONLY when PPTX_MCP_AUTO_RENDER=1.

    Example: '{"chart_type":"stacked_column","categories":["Q1","Q2"],"series":[{"name":"Rev","values":[10,20],"color":"2251FF"}],"data_labels_show":true}'"""
    try:
        try:
            spec = json.loads(chart_json) if isinstance(chart_json, str) else chart_json
        except json.JSONDecodeError as e:
            return _error(
                "INVALID_PARAMETER",
                f"Invalid JSON in chart_json: {e}",
                parameter="chart_json",
                hint="chart_json must be a JSON-stringified object.",
            )
        result = add_chart(file_path, slide_index, spec)
        render_info = _auto_render(file_path, slide_index)
        return _success_with_render(result, render_info)
    except Exception as e:
        return _err(e)


# --- Rendering ----------------------------------------------------

@mcp.tool()
def pptx_render_slide(
    file_path: str,
    slide_index: int = -1,
    dpi: int = 150,
) -> str:
    """Render PPTX slide(s) to PNG image(s) for visual verification. Returns path(s) to PNG files that can be viewed with the Read tool. slide_index: 0-based (-1 = all slides). dpi: 150 for review, 300 for print."""
    try:
        return _success(render_slide(file_path, slide_index, dpi=dpi))
    except Exception as e:
        return _err(e)


# --- Entry Point --------------------------------------------------

def _format_check_layout_summary(result: Dict[str, Any]) -> str:
    """``check_deck_extended`` の dict を legacy 形式の人間可読 string に整形する.

    Clean の場合:
        ``"All slides clean — no overlaps, out-of-bounds, text overflow, or
        readability issues detected."``

    問題がある場合:
        ``"Found N layout issues:\\n- Slide <i> [severity] <category>: <msg>"``
    """
    lines: list[str] = []
    for slide in result.get("slides", []):
        idx = slide.get("index", 0)
        # overlaps / out_of_bounds は legacy 文字列リスト (severity = error 固定)。
        for msg in slide.get("overlaps", []) or []:
            lines.append(f"- Slide {idx} [error] overlap: {msg}")
        for msg in slide.get("out_of_bounds", []) or []:
            lines.append(f"- Slide {idx} [error] out_of_bounds: {msg}")
        # ValidationFinding 由来カテゴリ (dict)
        for key in (
            "text_overflow",
            "unreadable_text",
            "divider_collision",
            "inconsistent_gaps",
        ):
            for f in slide.get(key, []) or []:
                sev = f.get("severity", "info") if isinstance(f, dict) else "info"
                msg = f.get("message", "") if isinstance(f, dict) else str(f)
                lines.append(f"- Slide {idx} [{sev}] {key}: {msg}")

    if not lines:
        return (
            "All slides clean — no overlaps, out-of-bounds, text overflow, "
            "or readability issues detected."
        )
    return f"Found {len(lines)} layout issues:\n" + "\n".join(lines)


@mcp.tool()
def pptx_check_layout(
    file_path: str,
    detailed: bool = False,
    min_readable_pt: float = 8.0,
    overflow_tolerance_pct: float = 5.0,
) -> str:
    """Validate slide layouts: overlaps, out-of-bounds, text overflow,
    unreadable font, title/divider collision, inconsistent gaps.

    v0.2.0+: tool 戻り値は ``{"ok": true, "result": <payload>}`` で包まれる。
    ``result`` 内は従来どおり:

    - ``detailed=False`` (既定): legacy human-readable string
      (``"All slides clean …"`` または ``"Found N layout issues:\\n…"``)。
      #33 で導入された文字列フォーマットはそのまま維持される。
    - ``detailed=True``: JSON 文字列 (中身は以下のスキーマ)。

    ``detailed=True`` schema::

        {
            "slides": [
                {"index": int, "overlaps": [...], "out_of_bounds": [...],
                 "text_overflow": [...], "unreadable_text": [...],
                 "divider_collision": [...], "inconsistent_gaps": [...]},
                ...
            ],
            "summary": {"errors": int, "warnings": int, "infos": int}
        }

    Run after building a deck to catch layout issues before delivery."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        result = check_deck_extended(
            prs,
            min_readable_pt=min_readable_pt,
            overflow_tolerance_pct=overflow_tolerance_pct,
        )
        if detailed:
            return _success(json.dumps(result, ensure_ascii=False, indent=2))
        return _success(_format_check_layout_summary(result))
    except Exception as e:
        return _err(e)


def main():
    mcp.run(transport="stdio")


if __name__ == "__main__":
    main()
